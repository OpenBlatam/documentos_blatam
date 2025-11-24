from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_pitch_deck(filename):
    prs = Presentation()

    # --- Design Constants ---
    NAVY_BLUE = RGBColor(32, 55, 100)
    LIGHT_BLUE = RGBColor(47, 117, 181)
    WHITE = RGBColor(255, 255, 255)
    GREY = RGBColor(128, 128, 128)

    def add_slide(prs, layout_index, title_text, content_text=None):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
        title.text_frame.paragraphs[0].font.bold = True

        # Content
        if content_text:
            # Handling different layouts might require finding the body placeholder
            if len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.text = content_text
                for p in tf.paragraphs:
                    p.font.size = Pt(20)
        
        return slide

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide
    title = slide.shapes.title
    title.text = "VENTURA CAPITAL"
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = "SaaS IA Marketing & Copywriting\nInvestment Opportunity - Series A"
    subtitle.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE

    # --- Slide 2: Problem ---
    s2 = add_slide(prs, 1, "El Problema")
    content = s2.placeholders[1].text_frame
    content.text = "El marketing digital en LATAM es ineficiente y costoso."
    p = content.add_paragraph()
    p.text = "• Agencias tradicionales son lentas y caras."
    p = content.add_paragraph()
    p.text = "• Herramientas de IA actuales (Jasper/Copy.ai) no entienden el contexto cultural LATAM."
    p = content.add_paragraph()
    p.text = "• Las empresas pierden 40% de presupuesto en contenido que no convierte."

    # --- Slide 3: Solution ---
    s3 = add_slide(prs, 1, "La Solución: Ventura AI")
    content = s3.placeholders[1].text_frame
    content.text = "Plataforma de IA Generativa nativa para LATAM."
    p = content.add_paragraph()
    p.text = "• Entrenada con dialectos y contextos culturales específicos de la región."
    p = content.add_paragraph()
    p.text = "• Generación de campañas completas en segundos, no días."
    p = content.add_paragraph()
    p.text = "• 30% más económica que competidores globales."

    # --- Slide 4: Market Size ---
    s4 = add_slide(prs, 1, "Oportunidad de Mercado")
    content = s4.placeholders[1].text_frame
    content.text = "Mercado Total Direccionable (TAM) en expansión."
    p = content.add_paragraph()
    p.text = "TAM: $2.8 Billones (Mercado MarTech LATAM)"
    p.font.bold = True
    p = content.add_paragraph()
    p.text = "SAM: $280 Millones (Empresas Digitales)"
    p = content.add_paragraph()
    p.text = "SOM: $28 Millones (Objetivo a 3 años)"

    # --- Slide 5: Traction ---
    s5 = add_slide(prs, 1, "Tracción Actual")
    content = s5.placeholders[1].text_frame
    content.text = "Crecimiento validado y eficiente."
    p = content.add_paragraph()
    p.text = "• ARR: $139K (Creciendo 20% MoM)"
    p = content.add_paragraph()
    p.text = "• Usuarios Activos: 500+"
    p = content.add_paragraph()
    p.text = "• LTV/CAC: 9:1 (Unit Economics Saludables)"
    p = content.add_paragraph()
    p.text = "• Churn: 5% (vs 8% promedio industria)"

    # --- Slide 6: Business Model ---
    s6 = add_slide(prs, 1, "Modelo de Negocio")
    content = s6.placeholders[1].text_frame
    content.text = "SaaS B2B Recurrente"
    p = content.add_paragraph()
    p.text = "• Freemium: Entrada gratuita limitada."
    p = content.add_paragraph()
    p.text = "• Pro Plan: $29/mes (Solopreneurs)."
    p = content.add_paragraph()
    p.text = "• Business Plan: $99/mes (Agencias/Pymes)."
    p = content.add_paragraph()
    p.text = "• Enterprise: Custom pricing."

    # --- Slide 7: Competition ---
    s7 = add_slide(prs, 1, "Ventaja Competitiva")
    content = s7.placeholders[1].text_frame
    content.text = "Por qué ganamos en LATAM"
    p = content.add_paragraph()
    p.text = "• Global Players (Jasper, Copy.ai): Costosos, inglés-centricos, genéricos."
    p = content.add_paragraph()
    p.text = "• Ventura AI: Localización profunda, soporte en español nativo, integraciones locales (MercadoLibre, TiendaNube)."

    # --- Slide 8: Financial Projections ---
    s8 = add_slide(prs, 1, "Proyecciones Financieras")
    content = s8.placeholders[1].text_frame
    content.text = "Camino a $3.6M ARR en 2026"
    p = content.add_paragraph()
    p.text = "• 2024 (Actual): $139K ARR"
    p = content.add_paragraph()
    p.text = "• 2025 (Est): $1.46M ARR"
    p = content.add_paragraph()
    p.text = "• 2026 (Est): $3.60M ARR"
    p = content.add_paragraph()
    p.text = "• EBITDA Positivo proyectado para Q4 2026"

    # --- Slide 9: The Ask ---
    s9 = add_slide(prs, 1, "La Propuesta de Inversión")
    content = s9.placeholders[1].text_frame
    content.text = "Buscamos $2.0M Series A"
    p = content.add_paragraph()
    p.text = "Uso de Fondos:"
    p = content.add_paragraph()
    p.text = "• 40% Desarrollo de Producto (Modelos propios)"
    p.level = 1
    p = content.add_paragraph()
    p.text = "• 35% Ventas y Marketing (Expansión MX y BR)"
    p.level = 1
    p = content.add_paragraph()
    p.text = "• 25% Operaciones y Equipo"
    p.level = 1

    # --- Slide 10: Contact ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Gracias"
    title.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
    
    subtitle = slide.placeholders[1]
    subtitle.text = "ceo@venturacapital.ai\nwww.venturacapital.ai\n+52 55 1234 5678"

    prs.save(filename)
    print(f"Pitch Deck created: {filename}")

if __name__ == "__main__":
    create_pitch_deck("Ventura_Capital_Pitch_Deck.pptx")


