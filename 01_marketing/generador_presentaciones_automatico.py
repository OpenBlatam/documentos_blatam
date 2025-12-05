#!/usr/bin/env python3
"""
Generador de Presentaciones Automáticas - Crea presentaciones PowerPoint
automáticamente desde datos estructurados.
"""

import os
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import json

class GeneradorPresentaciones:
    """Genera presentaciones automáticas"""
    
    def __init__(self):
        self.directorio = '/Users/adan/Documents/documentos_blatam/01_marketing'
        self.directorio_presentaciones = os.path.join(self.directorio, 'Presentaciones_Automaticas')
        os.makedirs(self.directorio_presentaciones, exist_ok=True)
    
    def crear_presentacion_ejecutiva(self, titulo="Presentación Ejecutiva", datos=None):
        """Crea presentación ejecutiva"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Portada
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        left = top = Inches(1)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide1.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = titulo
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)  # #1F4E78
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtítulo
        subtitle_box = slide1.shapes.add_textbox(left, Inches(3), width, Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Generado el {datetime.now().strftime('%d/%m/%Y')}"
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Slide 2: Resumen Ejecutivo
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide2.shapes.title.text = "Resumen Ejecutivo"
        
        content = slide2.placeholders[1]
        tf = content.text_frame
        tf.text = "• Análisis completo del sistema"
        p = tf.add_paragraph()
        p.text = "• Métricas clave de rendimiento"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Conclusiones y recomendaciones"
        p.level = 0
        
        # Slide 3: Métricas Clave
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = "Métricas Clave"
        
        if datos:
            content = slide3.placeholders[1]
            tf = content.text_frame
            tf.text = f"Total de Documentos: {datos.get('total_documentos', 'N/A')}"
            p = tf.add_paragraph()
            p.text = f"Scripts Activos: {datos.get('scripts_activos', 'N/A')}"
            p.level = 0
            p = tf.add_paragraph()
            p.text = f"Formatos Soportados: {datos.get('formatos', 'N/A')}"
            p.level = 0
        
        # Slide 4: Análisis
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        slide4.shapes.title.text = "Análisis Detallado"
        
        content = slide4.placeholders[1]
        tf = content.text_frame
        tf.text = "• Análisis estadístico completo"
        p = tf.add_paragraph()
        p.text = "• Predicciones con Machine Learning"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Análisis de sentimiento"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Métricas y KPIs en tiempo real"
        p.level = 0
        
        # Slide 5: Conclusiones
        slide5 = prs.slides.add_slide(prs.slide_layouts[1])
        slide5.shapes.title.text = "Conclusiones"
        
        content = slide5.placeholders[1]
        tf = content.text_frame
        tf.text = "• Sistema completamente funcional"
        p = tf.add_paragraph()
        p.text = "• Alta calidad en todos los formatos"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Automatización completa implementada"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Listo para producción"
        p.level = 0
        
        # Slide 6: Próximos Pasos
        slide6 = prs.slides.add_slide(prs.slide_layouts[1])
        slide6.shapes.title.text = "Próximos Pasos"
        
        content = slide6.placeholders[1]
        tf = content.text_frame
        tf.text = "• Continuar mejoras incrementales"
        p = tf.add_paragraph()
        p.text = "• Monitoreo continuo de KPIs"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Expansión de funcionalidades"
        p.level = 0
        
        return prs
    
    def crear_presentacion_desde_datos(self, datos_json):
        """Crea presentación desde datos JSON"""
        if isinstance(datos_json, str):
            with open(datos_json, 'r', encoding='utf-8') as f:
                datos = json.load(f)
        else:
            datos = datos_json
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Portada
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = top = Inches(1)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = datos.get('titulo', 'Presentación de Datos')
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Crear slides para cada sección
        for seccion in datos.get('secciones', []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = seccion.get('titulo', 'Sección')
            
            content = slide.placeholders[1]
            tf = content.text_frame
            
            for item in seccion.get('items', []):
                if tf.text == "":
                    tf.text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
        
        return prs
    
    def guardar_presentacion(self, prs, nombre=None):
        """Guarda presentación"""
        if nombre is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            nombre = f"PRESENTACION_{timestamp}.pptx"
        
        archivo = os.path.join(self.directorio_presentaciones, nombre)
        prs.save(archivo)
        print(f"✓ Presentación guardada: {archivo}")
        return archivo

def main():
    """Función principal"""
    generador = GeneradorPresentaciones()
    
    print("📊 Generando presentación ejecutiva...\n")
    
    # Datos de ejemplo
    datos = {
        'total_documentos': 50,
        'scripts_activos': 33,
        'formatos': 10
    }
    
    # Crear presentación
    prs = generador.crear_presentacion_ejecutiva(
        "Sistema Premium 2.0 - Resumen Ejecutivo",
        datos
    )
    
    # Guardar
    archivo = generador.guardar_presentacion(prs)
    
    print(f"\n✅ Presentación generada exitosamente!")
    print(f"📁 Archivo: {archivo}")

if __name__ == "__main__":
    main()








