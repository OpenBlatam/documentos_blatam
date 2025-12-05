#!/usr/bin/env python3
"""
Script para convertir documentos importantes a PDF, Word y Excel
con gráficas y formato profesional de alta calidad
"""

import os
import re
import json
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Tuple, Any, Optional
from collections import Counter, defaultdict
from difflib import SequenceMatcher
import hashlib
import markdown
from markdown.extensions import tables, fenced_code, codehilite

# PDF
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("⚠️  reportlab no disponible. Instala con: pip install reportlab")

# Word
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    WORD_AVAILABLE = True
except ImportError:
    WORD_AVAILABLE = False
    print("⚠️  python-docx no disponible. Instala con: pip install python-docx")

# Excel
try:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, PieChart, ScatterChart, Reference
    from openpyxl.chart.series import DataPoint
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("⚠️  openpyxl no disponible. Instala con: pip install openpyxl")

# PowerPoint
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx no disponible. Instala con: pip install python-pptx")

# Gráficas
try:
    import matplotlib
    matplotlib.use('Agg')  # Backend sin GUI
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from matplotlib.backends.backend_pdf import PdfPages
    import numpy as np
    import seaborn as sns
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False
    print("⚠️  matplotlib no disponible. Instala con: pip install matplotlib numpy seaborn")

# Configuración
OUTPUT_DIR = Path("documentos_exportados")
OUTPUT_DIR.mkdir(exist_ok=True)

# Documentos importantes a convertir
IMPORTANT_DOCS = [
    {
        "path": "airflow_automation_prompt.md",
        "title": "Guía de Automatización con Airflow",
        "category": "Automatización"
    },
    {
        "path": "ARCHITECTURE.md",
        "title": "Arquitectura del Proyecto",
        "category": "Arquitectura"
    },
    {
        "path": "README.md",
        "title": "Documentación Principal - Documentos BLATAM",
        "category": "Documentación"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/ARCHITECTURE_IMPROVEMENTS.md",
        "title": "Mejoras Arquitectónicas",
        "category": "Arquitectura"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/REFACTORING_PLAN.md",
        "title": "Plan de Refactorización",
        "category": "Desarrollo"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/README.md",
        "title": "Código de Producción - README",
        "category": "Desarrollo"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/RESUMEN_FINAL_MEJORAS.md",
        "title": "Resumen Final - Mejoras Arquitectónicas",
        "category": "Resumen"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/MEJORAS_ARQUITECTURA_COMPLETAS.md",
        "title": "Mejoras Arquitectónicas Completas",
        "category": "Arquitectura"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/MEJORAS_FINALES_CONSOLIDADO.md",
        "title": "Mejoras Finales - Resumen Consolidado",
        "category": "Resumen"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/MEJORAS_FINALES_COMPLETAS.md",
        "title": "Mejoras Finales Completas",
        "category": "Mejoras"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/MEJORAS_ADICIONALES_RECOMENDADAS.md",
        "title": "Mejoras Adicionales Recomendadas",
        "category": "Mejoras"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/INDICE_DOCUMENTACION.md",
        "title": "Índice de Documentación - Production Code",
        "category": "Documentación"
    }
]


class DocumentConverter:
    """Clase principal para convertir documentos"""
    
    def __init__(self, base_path: Path):
        self.base_path = base_path
        self.graphs_dir = OUTPUT_DIR / "graficas"
        self.graphs_dir.mkdir(exist_ok=True)
        self.graphs_created = []
        
    def read_markdown(self, file_path: Path) -> str:
        """Lee un archivo markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"❌ Error leyendo {file_path}: {e}")
            return ""
    
    def parse_markdown(self, content: str) -> Tuple[Dict[str, Any], str]:
        """Parsea markdown y extrae información estructurada"""
        # Convertir markdown a HTML
        md = markdown.Markdown(extensions=['tables', 'fenced_code', 'codehilite', 'nl2br'])
        html_content = md.convert(content)
        
        # Extraer datos para gráficas
        data = {
            'sections': [],
            'code_blocks': [],
            'tables': [],
            'metrics': {},
            'lists': []
        }
        
        # Extraer secciones (headers)
        section_pattern = r'^#+\s+(.+)$'
        for line in content.split('\n'):
            match = re.match(section_pattern, line)
            if match:
                level = len(line) - len(line.lstrip('#'))
                title = match.group(1).strip()
                # Limitar longitud del título
                if len(title) > 100:
                    title = title[:97] + "..."
                data['sections'].append({
                    'title': title,
                    'level': level
                })
        
        # Extraer bloques de código
        code_block_pattern = r'```[\w]*\n(.*?)```'
        code_blocks = re.findall(code_block_pattern, content, re.DOTALL)
        data['code_blocks'] = [cb[:200] for cb in code_blocks[:50]]  # Limitar cantidad y longitud
        
        # Extraer tablas markdown estructuradas
        table_lines = []
        in_table = False
        for line in content.split('\n'):
            if '|' in line and line.count('|') >= 2:
                if not in_table:
                    in_table = True
                    table_lines = []
                table_lines.append(line)
            elif in_table and line.strip() and '|' not in line:
                # Fin de tabla
                if len(table_lines) >= 2:  # Al menos encabezado y separador
                    data['tables'].append(table_lines)
                in_table = False
                table_lines = []
        
        # Agregar última tabla si existe
        if in_table and len(table_lines) >= 2:
            data['tables'].append(table_lines)
        
        # Análisis de palabras clave (palabras más frecuentes)
        words = re.findall(r'\b[a-zA-Z]{4,}\b', content.lower())
        word_freq = {}
        stop_words = {'this', 'that', 'with', 'from', 'have', 'been', 'will', 'were', 'what', 
                     'when', 'where', 'which', 'there', 'their', 'these', 'those', 'them', 'then'}
        for word in words:
            if word not in stop_words and len(word) > 3:
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # Top 20 palabras más frecuentes
        top_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:20]
        data['top_keywords'] = [{'word': w, 'count': c} for w, c in top_words]
        
        # Calcular tiempo estimado de lectura (250 palabras por minuto)
        total_words = len(content.split())
        reading_time_minutes = max(1, total_words // 250)
        data['metrics']['reading_time_minutes'] = reading_time_minutes
        
        # Extraer enlaces
        link_pattern = r'\[([^\]]+)\]\(([^\)]+)\)'
        links = re.findall(link_pattern, content)
        data['links'] = links[:50]  # Limitar a 50
        data['metrics']['total_links'] = len(links)
        
        # Extraer imágenes
        img_pattern = r'!\[([^\]]*)\]\(([^\)]+)\)'
        images = re.findall(img_pattern, content)
        data['images'] = images[:30]  # Limitar a 30
        data['metrics']['total_images'] = len(images)
        
        # Análisis de lenguajes de código
        code_languages = {}
        for code_block in data['code_blocks']:
            # Detectar lenguaje común por palabras clave
            code_lower = code_block.lower()
            if 'def ' in code_lower or 'import ' in code_lower or 'print(' in code_lower:
                code_languages['Python'] = code_languages.get('Python', 0) + 1
            elif 'function ' in code_lower or 'const ' in code_lower or 'let ' in code_lower:
                code_languages['JavaScript'] = code_languages.get('JavaScript', 0) + 1
            elif 'public class' in code_lower or 'private ' in code_lower:
                code_languages['Java'] = code_languages.get('Java', 0) + 1
            elif '#include' in code_lower or 'int main' in code_lower:
                code_languages['C/C++'] = code_languages.get('C/C++', 0) + 1
            elif 'SELECT' in code_block or 'FROM' in code_block:
                code_languages['SQL'] = code_languages.get('SQL', 0) + 1
            else:
                code_languages['Otro'] = code_languages.get('Otro', 0) + 1
        
        data['code_languages'] = code_languages
        
        # Análisis de complejidad
        # Complejidad basada en: secciones, código, tablas, longitud
        complexity_score = 0
        complexity_score += len(data['sections']) * 2
        complexity_score += len(data['code_blocks']) * 5
        complexity_score += len(data['tables']) * 3
        complexity_score += min(total_words // 1000, 50)  # Máximo 50 puntos por palabras
        
        if complexity_score < 20:
            complexity_level = "Baja"
        elif complexity_score < 50:
            complexity_level = "Media"
        elif complexity_score < 100:
            complexity_level = "Alta"
        else:
            complexity_level = "Muy Alta"
        
        data['metrics']['complexity_score'] = complexity_score
        data['metrics']['complexity_level'] = complexity_level
        
        # Generar resumen ejecutivo automático (primeras 3 secciones principales)
        main_sections = [s for s in data['sections'] if s['level'] == 1][:3]
        data['executive_summary'] = [s['title'] for s in main_sections]
        
        # ANÁLISIS DE CALIDAD DEL DOCUMENTO
        quality_analysis = {
            'scores': {},
            'issues': [],
            'recommendations': [],
            'overall_score': 0
        }
        
        # 1. Análisis de estructura (0-25 puntos)
        structure_score = 0
        if len(data['sections']) >= 3:
            structure_score += 10
        if len(data['sections']) >= 10:
            structure_score += 5
        # Verificar jerarquía correcta (H1 -> H2 -> H3)
        levels = [s['level'] for s in data['sections']]
        if 1 in levels and 2 in levels:
            structure_score += 5
        if 1 in levels and 2 in levels and 3 in levels:
            structure_score += 5
        quality_analysis['scores']['structure'] = min(25, structure_score)
        
        # 2. Análisis de contenido (0-25 puntos)
        content_score = 0
        if total_words >= 500:
            content_score += 10
        if total_words >= 2000:
            content_score += 10
        if len(data['code_blocks']) > 0:
            content_score += 5
        quality_analysis['scores']['content'] = min(25, content_score)
        
        # 3. Análisis de enlaces y referencias (0-20 puntos)
        links_score = 0
        if len(links) >= 5:
            links_score += 10
        if len(links) >= 15:
            links_score += 10
        quality_analysis['scores']['links'] = min(20, links_score)
        
        # 4. Análisis de tablas y datos (0-15 puntos)
        tables_score = 0
        if len(data['tables']) > 0:
            tables_score += 10
        if len(data['tables']) >= 3:
            tables_score += 5
        quality_analysis['scores']['tables'] = min(15, tables_score)
        
        # 5. Análisis de código (0-15 puntos)
        code_score = 0
        if len(data['code_blocks']) > 0:
            code_score += 10
        if len(data['code_blocks']) >= 5:
            code_score += 5
        quality_analysis['scores']['code'] = min(15, code_score)
        
        # Calcular score total
        quality_analysis['overall_score'] = sum(quality_analysis['scores'].values())
        
        # Detectar issues
        if len(data['sections']) < 3:
            quality_analysis['issues'].append("Pocas secciones (menos de 3)")
            quality_analysis['recommendations'].append("Agregar más secciones para mejorar la estructura")
        
        if total_words < 500:
            quality_analysis['issues'].append("Contenido muy corto (menos de 500 palabras)")
            quality_analysis['recommendations'].append("Expandir el contenido con más detalles")
        
        if len(links) < 5:
            quality_analysis['issues'].append("Pocos enlaces (menos de 5)")
            quality_analysis['recommendations'].append("Agregar más enlaces a recursos relacionados")
        
        if len(data['code_blocks']) == 0 and 'código' in content.lower():
            quality_analysis['issues'].append("Se menciona código pero no hay bloques de código")
            quality_analysis['recommendations'].append("Agregar ejemplos de código cuando se mencionen")
        
        # Detectar errores comunes
        errors = []
        # Enlaces rotos (sin http/https)
        for link_text, link_url in links[:20]:
            if not link_url.startswith(('http://', 'https://', '/', '#')) and '.' in link_url:
                errors.append(f"Enlace posiblemente roto: [{link_text}]({link_url})")
        
        # Headers sin contenido después
        lines = content.split('\n')
        for i, line in enumerate(lines):
            if line.strip().startswith('#'):
                # Verificar si hay contenido después
                next_lines = lines[i+1:i+4]
                if all(not l.strip() or l.strip().startswith('#') for l in next_lines):
                    errors.append(f"Header sin contenido: {line.strip()}")
        
        quality_analysis['errors'] = errors[:10]  # Limitar a 10 errores
        
        # Determinar nivel de calidad
        if quality_analysis['overall_score'] >= 80:
            quality_level = "Excelente"
        elif quality_analysis['overall_score'] >= 60:
            quality_level = "Buena"
        elif quality_analysis['overall_score'] >= 40:
            quality_level = "Regular"
        else:
            quality_level = "Necesita Mejora"
        
        quality_analysis['quality_level'] = quality_level
        data['quality_analysis'] = quality_analysis
        data['metrics']['quality_score'] = quality_analysis['overall_score']
        data['metrics']['quality_level'] = quality_level
        
        # Análisis de longitud de secciones
        section_lengths = []
        current_section = None
        current_length = 0
        
        for line in lines:
            if line.strip().startswith('#'):
                if current_section is not None:
                    section_lengths.append({
                        'title': current_section,
                        'length': current_length
                    })
                current_section = line.strip()
                current_length = 0
            else:
                current_length += len(line)
        
        if current_section is not None:
            section_lengths.append({
                'title': current_section,
                'length': current_length
            })
        
        data['section_lengths'] = section_lengths[:20]  # Top 20 secciones más largas
        
        # Extraer métricas numéricas avanzadas
        numbers = re.findall(r'\b(\d+)\b', content)
        if numbers:
            data['metrics']['total_numbers'] = len(numbers)
            numeric_values = [int(n) for n in numbers if n.isdigit() and int(n) < 1000000]
            if numeric_values:
                data['metrics']['max_number'] = max(numeric_values)
                data['metrics']['avg_number'] = sum(numeric_values) // len(numeric_values)
        
        # Extraer fases y porcentajes
        phase_pattern = r'(?:Phase|Fase)\s*(\d+)[:\s]*(?:✅|Completa|Complete|Completada)'
        phases = re.findall(phase_pattern, content, re.IGNORECASE)
        if phases:
            data['metrics']['phases_completed'] = len(set(phases))
            data['metrics']['phases_list'] = sorted([int(p) for p in set(phases)])
        
        # Extraer porcentajes
        percent_pattern = r'(\d+)%'
        percents = re.findall(percent_pattern, content)
        if percents:
            data['metrics']['percentages'] = [int(p) for p in percents]
            data['metrics']['max_percentage'] = max([int(p) for p in percents])
            data['metrics']['avg_percentage'] = sum([int(p) for p in percents]) // len(percents) if percents else 0
        
        # Contar elementos
        data['metrics']['total_lines'] = len(content.split('\n'))
        data['metrics']['total_words'] = len(content.split())
        data['metrics']['total_sections'] = len(data['sections'])
        data['metrics']['code_blocks'] = len(data['code_blocks'])
        
        # ANÁLISIS DE LEGIBILIDAD (Flesch Reading Ease)
        try:
            readability = self.calculate_readability(content)
            data['readability'] = readability
            data['metrics']['flesch_score'] = readability.get('flesch_score', 0)
            data['metrics']['reading_level'] = readability.get('reading_level', 'N/A')
        except Exception:
            data['readability'] = {}
            data['metrics']['flesch_score'] = 0
            data['metrics']['reading_level'] = 'N/A'
        
        # ANÁLISIS DE SENTIMIENTO/TONO BÁSICO
        try:
            sentiment = self.analyze_sentiment_basic(content)
            data['sentiment'] = sentiment
            data['metrics']['sentiment_tone'] = sentiment.get('primary_tone', 'Neutral')
        except Exception:
            data['sentiment'] = {}
            data['metrics']['sentiment_tone'] = 'Neutral'
        
        # ANÁLISIS DE COHERENCIA
        try:
            coherence = self.analyze_coherence(content, data['sections'])
            data['coherence'] = coherence
            data['metrics']['coherence_score'] = coherence.get('score', 0)
        except Exception:
            data['coherence'] = {}
            data['metrics']['coherence_score'] = 0
        
        return data, html_content
    
    def calculate_readability(self, content: str) -> Dict[str, Any]:
        """Calcula métricas de legibilidad Flesch Reading Ease"""
        # Remover código y enlaces para análisis de texto
        text_only = re.sub(r'```.*?```', '', content, flags=re.DOTALL)
        text_only = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text_only)
        text_only = re.sub(r'`[^`]+`', '', text_only)
        
        sentences = re.split(r'[.!?]+', text_only)
        sentences = [s.strip() for s in sentences if s.strip() and len(s) > 10]
        
        words = re.findall(r'\b[a-zA-ZáéíóúÁÉÍÓÚñÑ]{2,}\b', text_only.lower())
        
        if len(sentences) == 0 or len(words) == 0:
            return {
                'flesch_score': 0,
                'reading_level': 'N/A',
                'avg_sentence_length': 0,
                'avg_syllables_per_word': 0
            }
        
        # Contar sílabas (aproximado para español)
        total_syllables = 0
        for word in words[:1000]:  # Limitar para rendimiento
            syllables = max(1, len(re.findall(r'[aeiouáéíóúAEIOUÁÉÍÓÚ]', word)))
            total_syllables += syllables
        
        avg_syllables = total_syllables / len(words) if words else 0
        avg_sentence_length = len(words) / len(sentences)
        
        # Fórmula Flesch adaptada para español
        flesch_score = 206.84 - (1.015 * avg_sentence_length) - (84.6 * avg_syllables)
        flesch_score = max(0, min(100, flesch_score))
        
        # Determinar nivel
        if flesch_score >= 80:
            level = "Muy Fácil"
        elif flesch_score >= 70:
            level = "Fácil"
        elif flesch_score >= 60:
            level = "Bastante Fácil"
        elif flesch_score >= 50:
            level = "Estándar"
        elif flesch_score >= 30:
            level = "Bastante Difícil"
        else:
            level = "Difícil"
        
        return {
            'flesch_score': round(flesch_score, 2),
            'reading_level': level,
            'avg_sentence_length': round(avg_sentence_length, 2),
            'avg_syllables_per_word': round(avg_syllables, 2),
            'sentence_count': len(sentences),
            'word_count': len(words)
        }
    
    def analyze_sentiment_basic(self, content: str) -> Dict[str, Any]:
        """Análisis básico de sentimiento/tono del texto"""
        text_lower = content.lower()
        
        # Palabras positivas
        positive_words = ['excelente', 'bueno', 'mejor', 'éxito', 'completado', 'logrado', 
                         'perfecto', 'genial', 'fantástico', 'increíble', 'satisfactorio']
        # Palabras negativas
        negative_words = ['error', 'problema', 'fallo', 'malo', 'difícil', 'complejo',
                         'frustrante', 'lento', 'incompleto', 'roto']
        # Palabras técnicas
        technical_words = ['implementar', 'configurar', 'desplegar', 'optimizar', 'arquitectura',
                          'sistema', 'módulo', 'función', 'clase', 'método']
        
        positive_count = sum(1 for word in positive_words if word in text_lower)
        negative_count = sum(1 for word in negative_words if word in text_lower)
        technical_count = sum(1 for word in technical_words if word in text_lower)
        
        # Determinar tono principal
        if positive_count > negative_count and positive_count > 0:
            primary_tone = "Positivo"
        elif negative_count > positive_count and negative_count > 0:
            primary_tone = "Neutro-Técnico" if technical_count > 5 else "Negativo"
        elif technical_count > 10:
            primary_tone = "Técnico"
        else:
            primary_tone = "Neutral"
        
        return {
            'primary_tone': primary_tone,
            'positive_words': positive_count,
            'negative_words': negative_count,
            'technical_words': technical_count,
            'sentiment_score': positive_count - negative_count
        }
    
    def analyze_coherence(self, content: str, sections: List[Dict]) -> Dict[str, Any]:
        """Analiza la coherencia y estructura del documento"""
        score = 0
        issues = []
        
        # Verificar jerarquía de headers
        if len(sections) > 0:
            levels = [s['level'] for s in sections]
            if 1 in levels:
                score += 20
            else:
                issues.append("Falta header de nivel 1 (H1)")
            
            if 1 in levels and 2 in levels:
                score += 20
            else:
                if 2 in levels and 1 not in levels:
                    issues.append("Headers H2 sin H1 principal")
            
            # Verificar distribución de niveles
            level_dist = {}
            for level in levels:
                level_dist[level] = level_dist.get(level, 0) + 1
            
            if len(level_dist) >= 2:
                score += 20
            else:
                issues.append("Poca variación en niveles de headers")
        
        # Verificar longitud de secciones
        if len(sections) >= 3:
            score += 20
        else:
            issues.append("Muy pocas secciones (menos de 3)")
        
        # Verificar transiciones (presencia de conectores)
        connectors = ['además', 'por otro lado', 'sin embargo', 'por lo tanto', 
                     'en consecuencia', 'finalmente', 'en resumen', 'por ejemplo']
        connector_count = sum(1 for conn in connectors if conn in content.lower())
        if connector_count >= 3:
            score += 20
        else:
            issues.append("Pocos conectores de transición")
        
        return {
            'score': min(100, score),
            'level': 'Alta' if score >= 80 else 'Media' if score >= 60 else 'Baja',
            'issues': issues[:5]
        }
    
    def create_summary_graphs(self, doc_data: Dict[str, Any], doc_title: str) -> List[str]:
        """Crea gráficas de resumen del documento con mejor calidad"""
        graphs = []
        
        if not MATPLOTLIB_AVAILABLE:
            return graphs
        
        try:
            # Configurar estilo de gráficas
            plt.style.use('seaborn-v0_8-darkgrid')
            sns.set_palette("husl")
            
            # Gráfica 1: Distribución de secciones por nivel
            if doc_data.get('sections'):
                levels = {}
                for section in doc_data['sections']:
                    level = section['level']
                    levels[level] = levels.get(level, 0) + 1
                
                if levels:
                    fig, ax = plt.subplots(figsize=(12, 7))
                    levels_sorted = sorted(levels.items())
                    colors_list = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E', '#06A77D']
                    bars = ax.bar([f'Nivel {l}' for l, _ in levels_sorted], 
                                 [c for _, c in levels_sorted],
                                 color=colors_list[:len(levels_sorted)])
                    
                    # Agregar valores en las barras
                    for bar in bars:
                        height = bar.get_height()
                        ax.text(bar.get_x() + bar.get_width()/2., height,
                               f'{int(height)}',
                               ha='center', va='bottom', fontweight='bold', fontsize=11)
                    
                    ax.set_title(f'Distribución de Secciones por Nivel\n{doc_title}', 
                               fontsize=16, fontweight='bold', pad=20)
                    ax.set_xlabel('Nivel de Sección', fontsize=13, fontweight='bold')
                    ax.set_ylabel('Cantidad de Secciones', fontsize=13, fontweight='bold')
                    ax.grid(axis='y', alpha=0.3, linestyle='--')
                    
                    graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_')}_secciones.png"
                    plt.tight_layout()
                    plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                    plt.close()
                    graphs.append(str(graph_path))
            
            # Gráfica 2: Métricas del documento (mejorada)
            if doc_data.get('metrics'):
                metrics = doc_data['metrics']
                fig, ax = plt.subplots(figsize=(12, 7))
                
                metric_names = ['Líneas', 'Palabras', 'Secciones', 'Bloques Código']
                metric_values = [
                    metrics.get('total_lines', 0),
                    metrics.get('total_words', 0),
                    metrics.get('total_sections', 0),
                    metrics.get('code_blocks', 0)
                ]
                
                # Filtrar métricas con valor > 0
                filtered_data = [(name, val) for name, val in zip(metric_names, metric_values) if val > 0]
                if filtered_data:
                    metric_names_filtered = [d[0] for d in filtered_data]
                    metric_values_filtered = [d[1] for d in filtered_data]
                    
                    # Normalizar para visualización
                    max_val = max(metric_values_filtered) if metric_values_filtered else 1
                    normalized = [v / max_val * 100 for v in metric_values_filtered]
                    
                    bars = ax.barh(metric_names_filtered, normalized, 
                                 color=['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E'])
                    
                    # Agregar valores reales
                    for i, (bar, val) in enumerate(zip(bars, metric_values_filtered)):
                        width = bar.get_width()
                        ax.text(width + 1, bar.get_y() + bar.get_height()/2,
                               f'{val:,}', ha='left', va='center', 
                               fontweight='bold', fontsize=11)
                    
                    ax.set_title(f'Métricas del Documento\n{doc_title}', 
                               fontsize=16, fontweight='bold', pad=20)
                    ax.set_xlabel('Valor Normalizado (%)', fontsize=13, fontweight='bold')
                    ax.grid(axis='x', alpha=0.3, linestyle='--')
                    
                    graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_')}_metricas.png"
                    plt.tight_layout()
                    plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                    plt.close()
                    graphs.append(str(graph_path))
            
            # Gráfica 3: Progreso de fases (si hay fases)
            if doc_data.get('metrics', {}).get('phases_list'):
                phases = doc_data['metrics']['phases_list']
                fig, ax = plt.subplots(figsize=(12, 7))
                
                phase_labels = [f'Fase {p}' for p in phases]
                completion = [100] * len(phases)  # Asumiendo completas
                
                bars = ax.bar(phase_labels, completion, 
                            color=['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E', '#06A77D'])
                
                # Agregar valores
                for bar in bars:
                    height = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2., height,
                           f'{int(height)}%',
                           ha='center', va='bottom', fontweight='bold', fontsize=11)
                
                ax.set_ylim(0, 110)
                ax.set_ylabel('Porcentaje de Completitud (%)', fontsize=13, fontweight='bold')
                ax.set_xlabel('Fases', fontsize=13, fontweight='bold')
                ax.set_title(f'Progreso de Fases Completadas\n{doc_title}', 
                           fontsize=16, fontweight='bold', pad=20)
                ax.grid(axis='y', alpha=0.3, linestyle='--')
                
                graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_')}_progreso.png"
                plt.tight_layout()
                plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                plt.close()
                graphs.append(str(graph_path))
            
            # Gráfica 4: Timeline de fechas (si hay fechas)
            if doc_data.get('metrics', {}).get('dates'):
                dates = doc_data['metrics']['dates']
                if len(dates) >= 2:
                    fig, ax = plt.subplots(figsize=(12, 7))
                    
                    # Simular progreso basado en fechas
                    progress = np.linspace(0, 100, len(dates))
                    
                    ax.plot(dates, progress, marker='o', linewidth=3, markersize=10, 
                           color='#2E86AB', markerfacecolor='#06A77D', 
                           markeredgewidth=2, markeredgecolor='#2E86AB')
                    
                    ax.fill_between(dates, progress, alpha=0.3, color='#2E86AB')
                    
                    # Agregar valores en puntos
                    for date, prog in zip(dates, progress):
                        ax.text(date, prog + 5, f'{int(prog)}%', 
                               ha='center', fontweight='bold', fontsize=10)
                    
                    ax.set_xlabel('Fecha', fontsize=13, fontweight='bold')
                    ax.set_ylabel('Progreso (%)', fontsize=13, fontweight='bold')
                    ax.set_title(f'Línea de Tiempo de Progreso\n{doc_title}', 
                               fontsize=16, fontweight='bold', pad=20)
                    ax.grid(True, alpha=0.3, linestyle='--')
                    
                    plt.xticks(rotation=45, ha='right')
                    graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_')}_timeline.png"
                    plt.tight_layout()
                    plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                    plt.close()
                    graphs.append(str(graph_path))
            
            # Gráfica 5: Distribución de contenido (pie chart)
            if doc_data.get('metrics'):
                metrics = doc_data['metrics']
                fig, ax = plt.subplots(figsize=(12, 7))
                
                labels = []
                sizes = []
                colors_pie = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E', '#06A77D']
                
                if metrics.get('total_sections', 0) > 0:
                    labels.append('Secciones')
                    sizes.append(metrics.get('total_sections', 0))
                if metrics.get('code_blocks', 0) > 0:
                    labels.append('Bloques Código')
                    sizes.append(metrics.get('code_blocks', 0))
                if metrics.get('total_tables', 0) > 0:
                    labels.append('Tablas')
                    sizes.append(metrics.get('total_tables', 0))
                if metrics.get('total_numbers', 0) > 0:
                    labels.append('Números')
                    sizes.append(min(metrics.get('total_numbers', 0), 1000))  # Limitar para visualización
                
                if labels and sizes:
                    wedges, texts, autotexts = ax.pie(sizes, labels=labels, autopct='%1.1f%%',
                                                      colors=colors_pie[:len(labels)],
                                                      startangle=90, textprops={'fontsize': 12, 'fontweight': 'bold'})
                    
                    # Mejorar texto de porcentajes
                    for autotext in autotexts:
                        autotext.set_color('white')
                        autotext.set_fontweight('bold')
                        autotext.set_fontsize(11)
                    
                    ax.set_title(f'Distribución de Contenido\n{doc_title}', 
                               fontsize=16, fontweight='bold', pad=20)
                    
                    graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_')}_distribucion.png"
                    plt.tight_layout()
                    plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                    plt.close()
                    graphs.append(str(graph_path))
            
            # Gráfica 6: Comparación de métricas (si hay múltiples valores)
            if doc_data.get('metrics', {}).get('percentages'):
                percents = doc_data['metrics']['percentages'][:20]  # Limitar a 20
                if len(percents) >= 3:
                    fig, ax = plt.subplots(figsize=(12, 7))
                    
                    # Crear histograma de porcentajes
                    bins = np.linspace(0, 100, 11)
                    n, bins, patches = ax.hist(percents, bins=bins, color='#2E86AB', 
                                              edgecolor='black', linewidth=1.5, alpha=0.7)
                    
                    # Colorear barras según rango
                    for i, (patch, val) in enumerate(zip(patches, n)):
                        if val > 0:
                            if bins[i] < 25:
                                patch.set_facecolor('#C73E1D')
                            elif bins[i] < 50:
                                patch.set_facecolor('#F18F01')
                            elif bins[i] < 75:
                                patch.set_facecolor('#06A77D')
                            else:
                                patch.set_facecolor('#2E86AB')
                    
                    ax.set_xlabel('Rango de Porcentaje (%)', fontsize=13, fontweight='bold')
                    ax.set_ylabel('Frecuencia', fontsize=13, fontweight='bold')
                    ax.set_title(f'Distribución de Porcentajes\n{doc_title}', 
                               fontsize=16, fontweight='bold', pad=20)
                    ax.grid(axis='y', alpha=0.3, linestyle='--')
                    
                    graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_')}_porcentajes.png"
                    plt.tight_layout()
                    plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                    plt.close()
                    graphs.append(str(graph_path))
            
            # Gráfica 7: Análisis de Calidad
            if doc_data.get('quality_analysis'):
                quality = doc_data['quality_analysis']
                fig, ax = plt.subplots(figsize=(12, 7))
                
                factors = list(quality['scores'].keys())
                scores = list(quality['scores'].values())
                max_scores = [25, 25, 20, 15, 15]  # Máximos por factor
                
                x = np.arange(len(factors))
                width = 0.35
                
                bars1 = ax.bar(x - width/2, scores, width, label='Score Actual', color='#2E86AB', alpha=0.8)
                bars2 = ax.bar(x + width/2, max_scores, width, label='Score Máximo', color='#E0E0E0', alpha=0.5)
                
                # Agregar valores
                for i, (bar, score) in enumerate(zip(bars1, scores)):
                    height = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2., height,
                           f'{int(score)}',
                           ha='center', va='bottom', fontweight='bold', fontsize=10)
                
                ax.set_xlabel('Factores de Calidad', fontsize=13, fontweight='bold')
                ax.set_ylabel('Score', fontsize=13, fontweight='bold')
                ax.set_title(f'Análisis de Calidad del Documento\n{doc_title}', 
                           fontsize=16, fontweight='bold', pad=20)
                ax.set_xticks(x)
                ax.set_xticklabels([f.title() for f in factors], rotation=45, ha='right')
                ax.legend()
                ax.grid(axis='y', alpha=0.3, linestyle='--')
                
                graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_')}_calidad.png"
                plt.tight_layout()
                plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                plt.close()
                graphs.append(str(graph_path))
            
            # Gráfica 8: Análisis de Legibilidad
            if doc_data.get('readability'):
                readability = doc_data['readability']
                fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6))
                
                # Gráfica 1: Score de legibilidad
                flesch_score = readability.get('flesch_score', 0)
                levels = ['Muy\nDifícil', 'Difícil', 'Bastante\nDifícil', 
                         'Estándar', 'Bastante\nFácil', 'Fácil', 'Muy\nFácil']
                level_scores = [0, 30, 50, 60, 70, 80, 100]
                
                # Encontrar nivel actual
                current_level_idx = 0
                for i, score in enumerate(level_scores):
                    if flesch_score >= score:
                        current_level_idx = i
                
                colors_list = ['#C73E1D', '#F18F01', '#FFE699', '#06A77D', '#2E86AB', '#1a5490', '#0d3d61']
                bars = ax1.barh(levels, [100] * len(levels), color=colors_list, alpha=0.3)
                bars[current_level_idx].set_alpha(1.0)
                bars[current_level_idx].set_color(colors_list[current_level_idx])
                
                ax1.axvline(flesch_score, color='red', linestyle='--', linewidth=2, label=f'Score Actual: {flesch_score:.1f}')
                ax1.set_xlabel('Flesch Reading Ease Score', fontsize=12, fontweight='bold')
                ax1.set_title('Nivel de Legibilidad', fontsize=14, fontweight='bold')
                ax1.legend()
                ax1.grid(axis='x', alpha=0.3)
                
                # Gráfica 2: Métricas de legibilidad
                metrics_data = {
                    'Oraciones': readability.get('sentence_count', 0),
                    'Palabras': readability.get('word_count', 0),
                    'Long. Prom. Oración': readability.get('avg_sentence_length', 0),
                }
                
                bars2 = ax2.bar(metrics_data.keys(), metrics_data.values(), color=['#2E86AB', '#A23B72', '#F18F01'])
                ax2.set_ylabel('Valor', fontsize=12, fontweight='bold')
                ax2.set_title('Métricas de Legibilidad', fontsize=14, fontweight='bold')
                ax2.grid(axis='y', alpha=0.3)
                
                # Agregar valores
                for bar in bars2:
                    height = bar.get_height()
                    ax2.text(bar.get_x() + bar.get_width()/2., height,
                           f'{int(height)}',
                           ha='center', va='bottom', fontweight='bold')
                
                plt.suptitle(f'Análisis de Legibilidad - {doc_title}', 
                           fontsize=16, fontweight='bold', y=1.02)
                
                graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_')}_legibilidad.png"
                plt.tight_layout()
                plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                plt.close()
                graphs.append(str(graph_path))
            
            # Gráfica 9: Palabras clave más frecuentes
            if doc_data.get('top_keywords'):
                keywords = doc_data['top_keywords'][:15]  # Top 15
                if keywords:
                    fig, ax = plt.subplots(figsize=(12, 7))
                    
                    words = [k['word'].title() for k in keywords]
                    counts = [k['count'] for k in keywords]
                    
                    bars = ax.barh(words, counts, color='#2E86AB', alpha=0.8)
                    
                    # Agregar valores
                    for i, (bar, count) in enumerate(zip(bars, counts)):
                        width = bar.get_width()
                        ax.text(width + max(counts) * 0.01, bar.get_y() + bar.get_height()/2,
                               f'{count}', ha='left', va='center', 
                               fontweight='bold', fontsize=10)
                    
                    ax.set_xlabel('Frecuencia', fontsize=13, fontweight='bold')
                    ax.set_ylabel('Palabras Clave', fontsize=13, fontweight='bold')
                    ax.set_title(f'Palabras Clave Más Frecuentes\n{doc_title}', 
                               fontsize=16, fontweight='bold', pad=20)
                    ax.grid(axis='x', alpha=0.3, linestyle='--')
                    
                    graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_')}_keywords.png"
                    plt.tight_layout()
                    plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                    plt.close()
                    graphs.append(str(graph_path))
            
        except Exception as e:
            print(f"⚠️  Error creando gráficas: {e}")
            import traceback
            traceback.print_exc()
        
        return graphs
    
    def create_pdf(self, content: str, doc_data: Dict[str, Any], 
                   doc_title: str, output_path: Path, graphs: List[str]):
        """Crea un PDF profesional del documento"""
        if not PDF_AVAILABLE:
            print("⚠️  PDF no disponible. Instala reportlab")
            return
        
        try:
            # Clase para números de página y encabezados
            class NumberedCanvas(canvas.Canvas):
                def __init__(self, *args, **kwargs):
                    canvas.Canvas.__init__(self, *args, **kwargs)
                    self._saved_page_states = []
                    self.doc_title = doc_title
                
                def showPage(self):
                    self._saved_page_states.append(dict(self.__dict__))
                    self._startPage()
                
                def save(self):
                    num_pages = len(self._saved_page_states)
                    for state in self._saved_page_states:
                        self.__dict__.update(state)
                        self.draw_header_footer(num_pages)
                        canvas.Canvas.showPage(self)
                    canvas.Canvas.save(self)
                
                def draw_header_footer(self, total_pages):
                    self.saveState()
                    # Encabezado
                    self.setFont("Helvetica", 9)
                    self.setFillColor(colors.HexColor('#2E86AB'))
                    header_text = self.doc_title[:50] + "..." if len(self.doc_title) > 50 else self.doc_title
                    self.drawString(72, A4[1] - 50, header_text)
                    
                    # Línea decorativa
                    self.setStrokeColor(colors.HexColor('#2E86AB'))
                    self.setLineWidth(0.5)
                    self.line(72, A4[1] - 55, A4[0] - 72, A4[1] - 55)
                    
                    # Pie de página - Número de página
                    self.setFont("Helvetica", 9)
                    self.setFillColor(colors.HexColor('#666666'))
                    page_text = f"Página {self._pageNumber} de {total_pages}"
                    self.drawRightString(A4[0] - 72, 30, page_text)
                    
                    # Fecha en pie de página
                    date_str = datetime.now().strftime("%d/%m/%Y")
                    self.drawString(72, 30, f"Generado: {date_str}")
                    
                    self.restoreState()
            
            doc = SimpleDocTemplate(str(output_path), pagesize=A4,
                                  rightMargin=72, leftMargin=72,
                                  topMargin=90, bottomMargin=50)
            
            # Estilos
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#2E86AB'),
                spaceAfter=30,
                alignment=TA_CENTER,
                fontName='Helvetica-Bold'
            )
            
            heading1_style = ParagraphStyle(
                'CustomHeading1',
                parent=styles['Heading1'],
                fontSize=18,
                textColor=colors.HexColor('#2E86AB'),
                spaceAfter=12,
                spaceBefore=12,
                fontName='Helvetica-Bold'
            )
            
            heading2_style = ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=colors.HexColor('#A23B72'),
                spaceAfter=10,
                spaceBefore=10,
                fontName='Helvetica-Bold'
            )
            
            normal_style = ParagraphStyle(
                'CustomNormal',
                parent=styles['Normal'],
                fontSize=11,
                leading=14,
                spaceAfter=6,
                alignment=TA_JUSTIFY
            )
            
            # Contenido
            story = []
            
            # PORTADA PROFESIONAL
            story.append(Spacer(1, 2*inch))
            
            # Título principal en portada
            cover_title_style = ParagraphStyle(
                'CoverTitle',
                parent=styles['Heading1'],
                fontSize=28,
                textColor=colors.HexColor('#2E86AB'),
                spaceAfter=30,
                alignment=TA_CENTER,
                fontName='Helvetica-Bold',
                leading=34
            )
            story.append(Paragraph(doc_title, cover_title_style))
            story.append(Spacer(1, 0.5*inch))
            
            # Línea decorativa
            story.append(Paragraph("─" * 60, styles['Normal']))
            story.append(Spacer(1, 0.5*inch))
            
            # Información de portada
            metrics = doc_data.get('metrics', {})
            cover_info = []
            cover_info.append(f"Líneas: {metrics.get('total_lines', 0):,}")
            cover_info.append(f"Palabras: {metrics.get('total_words', 0):,}")
            cover_info.append(f"Secciones: {metrics.get('total_sections', 0)}")
            cover_info.append(f"Tiempo de lectura: ~{metrics.get('reading_time_minutes', 0)} minutos")
            cover_info.append(f"Complejidad: {metrics.get('complexity_level', 'N/A')}")
            
            for info in cover_info:
                story.append(Paragraph(f"<i>{info}</i>", styles['Normal']))
                story.append(Spacer(1, 0.1*inch))
            
            story.append(Spacer(1, 0.5*inch))
            date_str = datetime.now().strftime("%d de %B de %Y")
            story.append(Paragraph(f"<i>Generado el {date_str}</i>", styles['Normal']))
            
            story.append(PageBreak())
            
            # ÍNDICE AUTOMÁTICO
            story.append(Paragraph("📑 Índice de Contenido", heading1_style))
            story.append(Spacer(1, 0.3*inch))
            
            # Generar índice desde secciones
            sections = doc_data.get('sections', [])
            toc_items = []
            current_level = 0
            item_number = 1
            
            for section in sections[:50]:  # Limitar a 50 secciones en índice
                level = section['level']
                title = section['title'][:80]  # Limitar longitud
                
                if level == 1:
                    toc_items.append(f"{item_number}. {title}")
                    item_number += 1
                    current_level = 1
                elif level == 2 and current_level <= 2:
                    toc_items.append(f"   {item_number-1}.{len([s for s in sections[:sections.index(section)] if s['level'] == 2 and sections.index(s) < sections.index(section)]) + 1} {title}")
                elif level == 3 and current_level <= 3:
                    toc_items.append(f"      • {title}")
            
            for item in toc_items[:30]:  # Mostrar máximo 30 items
                story.append(Paragraph(item, normal_style))
            
            story.append(PageBreak())
            
            # RESUMEN EJECUTIVO
            if doc_data.get('executive_summary'):
                story.append(Paragraph("📊 Resumen Ejecutivo", heading1_style))
                story.append(Spacer(1, 0.2*inch))
                story.append(Paragraph("Este documento contiene las siguientes secciones principales:", normal_style))
                story.append(Spacer(1, 0.1*inch))
                for i, section_title in enumerate(doc_data['executive_summary'], 1):
                    story.append(Paragraph(f"{i}. {section_title}", normal_style))
                story.append(Spacer(1, 0.2*inch))
                
                # Métricas clave
                metrics = doc_data.get('metrics', {})
                quality = doc_data.get('quality_analysis', {})
                metrics_text = f"""
                <b>Métricas Clave:</b><br/>
                • Complejidad: {metrics.get('complexity_level', 'N/A')} (Score: {metrics.get('complexity_score', 0)})<br/>
                • Calidad: {metrics.get('quality_level', 'N/A')} (Score: {metrics.get('quality_score', 0)}/100)<br/>
                • Enlaces: {metrics.get('total_links', 0)}<br/>
                • Tablas: {metrics.get('total_tables', 0)}<br/>
                • Bloques de código: {metrics.get('code_blocks', 0)}
                """
                story.append(Paragraph(metrics_text, normal_style))
                
                # Issues y recomendaciones
                if quality.get('issues'):
                    story.append(Spacer(1, 0.2*inch))
                    story.append(Paragraph("<b>⚠️ Issues Detectados:</b>", heading2_style))
                    for issue in quality['issues'][:5]:
                        story.append(Paragraph(f"• {issue}", normal_style))
                
                if quality.get('recommendations'):
                    story.append(Spacer(1, 0.2*inch))
                    story.append(Paragraph("<b>💡 Recomendaciones:</b>", heading2_style))
                    for rec in quality['recommendations'][:5]:
                        story.append(Paragraph(f"→ {rec}", normal_style))
                
                story.append(PageBreak())
            
            # Título en contenido (más pequeño)
            story.append(Paragraph(doc_title, heading1_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Ya no necesitamos esto aquí porque está en la portada
            
            # Gráficas
            if graphs:
                story.append(Paragraph("📊 Resumen Visual", heading1_style))
                for graph_path in graphs:
                    if os.path.exists(graph_path):
                        try:
                            img = Image(graph_path, width=6*inch, height=3.6*inch)
                            story.append(img)
                            story.append(Spacer(1, 0.2*inch))
                        except Exception as e:
                            print(f"⚠️  Error agregando gráfica {graph_path}: {e}")
                story.append(PageBreak())
            
            # Contenido del documento
            story.append(Paragraph("📄 Contenido del Documento", heading1_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Procesar contenido markdown línea por línea
            lines = content.split('\n')
            in_code_block = False
            code_lines = []
            in_table = False
            table_lines = []
            
            for i, line in enumerate(lines):
                # Headers
                if line.startswith('# '):
                    story.append(Paragraph(line[2:], heading1_style))
                elif line.startswith('## '):
                    story.append(Paragraph(line[3:], heading2_style))
                elif line.startswith('### '):
                    story.append(Paragraph(line[4:], heading2_style))
                # Code blocks
                elif line.startswith('```'):
                    if in_code_block:
                        # Fin del bloque de código
                        if code_lines:
                            code_text = '\n'.join(code_lines)
                            story.append(Paragraph(f"<font face='Courier' size=9>{code_text}</font>", normal_style))
                            code_lines = []
                        in_code_block = False
                    else:
                        in_code_block = True
                elif in_code_block:
                    code_lines.append(line)
                # Tablas markdown
                elif '|' in line and line.count('|') >= 2:
                    if not in_table:
                        in_table = True
                        table_lines = []
                    table_lines.append(line)
                elif in_table:
                    # Procesar tabla completa
                    if len(table_lines) >= 2:
                        table_data = []
                        for tbl_line in table_lines:
                            if '---' not in tbl_line:  # Saltar separador
                                cells = [cell.strip() for cell in tbl_line.split('|')[1:-1]]
                                if cells:
                                    table_data.append(cells)
                        
                        if table_data:
                            # Crear tabla en PDF
                            table = Table(table_data)
                            table.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2E86AB')),
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                ('FONTSIZE', (0, 0), (-1, 0), 11),
                                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F5F5F5')])
                            ]))
                            story.append(table)
                            story.append(Spacer(1, 0.2*inch))
                    in_table = False
                    table_lines = []
                    # Procesar línea actual si no es parte de la tabla
                    if '|' not in line:
                        if line.strip().startswith('# '):
                            story.append(Paragraph(line[2:], heading1_style))
                        elif line.strip().startswith('## '):
                            story.append(Paragraph(line[3:], heading2_style))
                        elif line.strip():
                            clean_line = line
                            clean_line = clean_line.replace('&', '&amp;')
                            clean_line = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', clean_line)
                            clean_line = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r'<i>\1</i>', clean_line)
                            parts = re.split(r'(<[bi]>|</[bi]>)', clean_line)
                            result_parts = []
                            for part in parts:
                                if re.match(r'<[bi]>|</[bi]>', part):
                                    result_parts.append(part)
                                else:
                                    result_parts.append(part.replace('<', '&lt;').replace('>', '&gt;'))
                            clean_line = ''.join(result_parts)
                            try:
                                story.append(Paragraph(clean_line, normal_style))
                            except Exception:
                                plain_text = re.sub(r'[*_`]', '', line)
                                story.append(Paragraph(plain_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'), normal_style))
                # Listas
                elif line.strip().startswith('- ') or line.strip().startswith('* '):
                    text = line.strip()[2:].strip()
                    story.append(Paragraph(f"• {text}", normal_style))
                # Líneas vacías
                elif not line.strip():
                    story.append(Spacer(1, 0.1*inch))
                # Texto normal
                elif line.strip() and not in_table:
                    # Convertir markdown a HTML de forma segura
                    clean_line = line
                    # Escapar primero todos los caracteres especiales
                    clean_line = clean_line.replace('&', '&amp;')
                    # Convertir markdown a HTML (después de escapar)
                    clean_line = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', clean_line)
                    clean_line = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r'<i>\1</i>', clean_line)
                    # Escapar < y > que no sean parte de nuestros tags
                    # Dividir por tags HTML conocidos
                    parts = re.split(r'(<[bi]>|</[bi]>)', clean_line)
                    result_parts = []
                    for part in parts:
                        if re.match(r'<[bi]>|</[bi]>', part):
                            result_parts.append(part)  # Tag HTML, mantenerlo
                        else:
                            result_parts.append(part.replace('<', '&lt;').replace('>', '&gt;'))
                    clean_line = ''.join(result_parts)
                    
                    # Intentar agregar el párrafo
                    try:
                        story.append(Paragraph(clean_line, normal_style))
                    except Exception:
                        # Si falla, usar texto plano
                        plain_text = re.sub(r'[*_`]', '', line)
                        story.append(Paragraph(plain_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'), normal_style))
            
            # Construir PDF con canvas personalizado
            doc.build(story, canvasmaker=NumberedCanvas)
            print(f"✅ PDF creado: {output_path}")
            
        except Exception as e:
            print(f"❌ Error creando PDF: {e}")
            import traceback
            traceback.print_exc()
    
    def create_word(self, content: str, doc_data: Dict[str, Any],
                    doc_title: str, output_path: Path, graphs: List[str]):
        """Crea un documento Word profesional"""
        if not WORD_AVAILABLE:
            print("⚠️  Word no disponible. Instala python-docx")
            return
        
        try:
            doc = Document()
            
            # Estilos personalizados
            styles = doc.styles
            
            # PORTADA PROFESIONAL
            # Título principal
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(doc_title)
            title_run.font.name = 'Calibri'
            title_run.font.size = Pt(28)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(46, 134, 171)  # #2E86AB
            
            doc.add_paragraph()  # Espacio
            doc.add_paragraph()  # Espacio
            
            # Línea decorativa
            line_para = doc.add_paragraph("─" * 60)
            line_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph()  # Espacio
            
            # Información de portada
            metrics = doc_data.get('metrics', {})
            cover_info = [
                f"Líneas: {metrics.get('total_lines', 0):,}",
                f"Palabras: {metrics.get('total_words', 0):,}",
                f"Secciones: {metrics.get('total_sections', 0)}",
                f"Tiempo de lectura: ~{metrics.get('reading_time_minutes', 0)} minutos",
                f"Complejidad: {metrics.get('complexity_level', 'N/A')}"
            ]
            
            for info in cover_info:
                info_para = doc.add_paragraph()
                info_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                info_run = info_para.add_run(info)
                info_run.italic = True
                info_run.font.size = Pt(11)
            
            doc.add_paragraph()  # Espacio
            
            # Fecha
            date_para = doc.add_paragraph()
            date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            date_run = date_para.add_run(f"Generado el {datetime.now().strftime('%d de %B de %Y')}")
            date_run.italic = True
            date_run.font.size = Pt(10)
            
            doc.add_page_break()
            
            # ÍNDICE AUTOMÁTICO
            doc.add_heading('📑 Índice de Contenido', 1)
            doc.add_paragraph()
            
            sections = doc_data.get('sections', [])
            toc_items = []
            item_number = 1
            
            for section in sections[:50]:
                level = section['level']
                title = section['title'][:80]
                
                if level == 1:
                    toc_items.append(f"{item_number}. {title}")
                    item_number += 1
                elif level == 2:
                    toc_items.append(f"   {item_number-1}.{len([s for s in sections[:sections.index(section)] if s['level'] == 2]) + 1} {title}")
                elif level == 3:
                    toc_items.append(f"      • {title}")
            
            for item in toc_items[:30]:
                doc.add_paragraph(item, style='List Bullet')
            
            doc.add_page_break()
            
            # RESUMEN EJECUTIVO
            if doc_data.get('executive_summary'):
                doc.add_heading('📊 Resumen Ejecutivo', 1)
                doc.add_paragraph("Este documento contiene las siguientes secciones principales:")
                doc.add_paragraph()
                
                for i, section_title in enumerate(doc_data['executive_summary'], 1):
                    doc.add_paragraph(f"{i}. {section_title}", style='List Number')
                
                doc.add_paragraph()
                metrics = doc_data.get('metrics', {})
                metrics_para = doc.add_paragraph()
                metrics_para.add_run("Métricas Clave: ").bold = True
                quality = doc_data.get('quality_analysis', {})
                metrics_para.add_run(
                    f"Complejidad: {metrics.get('complexity_level', 'N/A')} "
                    f"(Score: {metrics.get('complexity_score', 0)}), "
                    f"Calidad: {metrics.get('quality_level', 'N/A')} "
                    f"(Score: {metrics.get('quality_score', 0)}/100), "
                    f"Enlaces: {metrics.get('total_links', 0)}, "
                    f"Tablas: {metrics.get('total_tables', 0)}, "
                    f"Bloques de código: {metrics.get('code_blocks', 0)}"
                )
                
                # Issues y recomendaciones
                if quality.get('issues'):
                    doc.add_paragraph()
                    doc.add_heading('⚠️ Issues Detectados', 3)
                    for issue in quality['issues'][:5]:
                        doc.add_paragraph(f"• {issue}", style='List Bullet')
                
                if quality.get('recommendations'):
                    doc.add_paragraph()
                    doc.add_heading('💡 Recomendaciones', 3)
                    for rec in quality['recommendations'][:5]:
                        doc.add_paragraph(f"→ {rec}", style='List Bullet')
                
                doc.add_page_break()
            
            # Título en contenido
            doc.add_heading(doc_title, 1)
            doc.add_paragraph()  # Espacio
            
            # Gráficas
            if graphs:
                doc.add_heading('📊 Resumen Visual', 1)
                for graph_path in graphs:
                    if os.path.exists(graph_path):
                        try:
                            doc.add_picture(graph_path, width=Inches(6))
                            doc.add_paragraph()  # Espacio después de la imagen
                        except Exception as e:
                            print(f"⚠️  Error agregando gráfica {graph_path}: {e}")
                doc.add_page_break()
            
            # Contenido
            doc.add_heading('📄 Contenido del Documento', 1)
            
            # Procesar contenido
            lines = content.split('\n')
            in_code_block = False
            code_lines = []
            in_table = False
            table_lines = []
            
            for i, line in enumerate(lines):
                if line.startswith('# '):
                    doc.add_heading(line[2:], 1)
                elif line.startswith('## '):
                    doc.add_heading(line[3:], 2)
                elif line.startswith('### '):
                    doc.add_heading(line[4:], 3)
                elif line.startswith('```'):
                    if in_code_block:
                        if code_lines:
                            code_para = doc.add_paragraph('\n'.join(code_lines))
                            code_para.style = 'No Spacing'
                            for run in code_para.runs:
                                run.font.name = 'Courier New'
                                run.font.size = Pt(9)
                            code_lines = []
                        in_code_block = False
                    else:
                        in_code_block = True
                elif in_code_block:
                    code_lines.append(line)
                # Tablas markdown
                elif '|' in line and line.count('|') >= 2:
                    if not in_table:
                        in_table = True
                        table_lines = []
                    table_lines.append(line)
                elif in_table:
                    # Procesar tabla completa
                    if len(table_lines) >= 2:
                        table_data = []
                        for tbl_line in table_lines:
                            if '---' not in tbl_line:  # Saltar separador
                                cells = [cell.strip() for cell in tbl_line.split('|')[1:-1]]
                                if cells:
                                    table_data.append(cells)
                        
                        if table_data and len(table_data) > 0:
                            # Crear tabla en Word
                            table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                            table.style = 'Light Grid Accent 1'
                            
                            for row_idx, row_data in enumerate(table_data):
                                for col_idx, cell_data in enumerate(row_data):
                                    if col_idx < len(table.rows[row_idx].cells):
                                        cell = table.rows[row_idx].cells[col_idx]
                                        cell.text = cell_data
                                        # Encabezado en negrita
                                        if row_idx == 0:
                                            for paragraph in cell.paragraphs:
                                                for run in paragraph.runs:
                                                    run.bold = True
                                                    run.font.color.rgb = RGBColor(255, 255, 255)
                                            # Usar shading correctamente
                                            try:
                                                from docx.oxml import parse_xml
                                                from docx.oxml.ns import nsdecls, qn
                                                shading = parse_xml(
                                                    r'<w:shd {} w:fill="2E86AB"/>'.format(nsdecls('w'))
                                                )
                                                cell._element.get_or_add_tcPr().append(shading)
                                            except Exception:
                                                pass  # Si falla, continuar sin shading
                            
                            doc.add_paragraph()  # Espacio después de tabla
                    in_table = False
                    table_lines = []
                    # Procesar línea actual si no es parte de la tabla
                    if '|' not in line and line.strip():
                        if line.strip().startswith('# '):
                            doc.add_heading(line[2:], 1)
                        elif line.strip().startswith('## '):
                            doc.add_heading(line[3:], 2)
                        elif line.strip().startswith('### '):
                            doc.add_heading(line[4:], 3)
                        elif line.strip().startswith('- ') or line.strip().startswith('* '):
                            doc.add_paragraph(line.strip()[2:].strip(), style='List Bullet')
                        else:
                            para = doc.add_paragraph()
                            text = line
                            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
                            text = re.sub(r'\*(.+?)\*', r'\1', text)
                            para.add_run(text)
                elif line.strip().startswith('- ') or line.strip().startswith('* '):
                    doc.add_paragraph(line.strip()[2:].strip(), style='List Bullet')
                elif line.strip() and not in_table:
                    para = doc.add_paragraph()
                    # Procesar formato markdown básico
                    text = line
                    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
                    text = re.sub(r'\*(.+?)\*', r'\1', text)
                    para.add_run(text)
            
            # Guardar
            doc.save(str(output_path))
            print(f"✅ Word creado: {output_path}")
            
        except Exception as e:
            print(f"❌ Error creando Word: {e}")
            import traceback
            traceback.print_exc()
    
    def create_excel(self, doc_data: Dict[str, Any], doc_title: str, 
                     output_path: Path, graphs: List[str]):
        """Crea un archivo Excel con datos y gráficas"""
        if not EXCEL_AVAILABLE:
            print("⚠️  Excel no disponible. Instala openpyxl")
            return
        
        try:
            wb = Workbook()
            
            # Hoja 1: Resumen
            ws1 = wb.active
            ws1.title = "Resumen"
            
            # Estilos
            title_fill = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
            header_fill = PatternFill(start_color="A23B72", end_color="A23B72", fill_type="solid")
            title_font = Font(bold=True, size=16, color="FFFFFF")
            header_font = Font(bold=True, size=12, color="FFFFFF")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Título
            ws1['A1'] = doc_title
            ws1['A1'].font = title_font
            ws1['A1'].fill = title_fill
            ws1['A1'].alignment = Alignment(horizontal='center', vertical='center')
            ws1.merge_cells('A1:D1')
            ws1.row_dimensions[1].height = 30
            
            # Fecha
            ws1['A2'] = f"Generado el {datetime.now().strftime('%d/%m/%Y %H:%M')}"
            ws1['A2'].font = Font(italic=True, size=10)
            ws1.merge_cells('A2:D2')
            
            # Métricas
            row = 4
            ws1[f'A{row}'] = "Métrica"
            ws1[f'B{row}'] = "Valor"
            ws1[f'A{row}'].font = header_font
            ws1[f'B{row}'].font = header_font
            ws1[f'A{row}'].fill = header_fill
            ws1[f'B{row}'].fill = header_fill
            ws1[f'A{row}'].border = border
            ws1[f'B{row}'].border = border
            
            row += 1
            metrics = doc_data.get('metrics', {})
            metric_data = [
                ("Total de Líneas", metrics.get('total_lines', 0)),
                ("Total de Palabras", metrics.get('total_words', 0)),
                ("Total de Secciones", metrics.get('total_sections', 0)),
                ("Bloques de Código", metrics.get('code_blocks', 0)),
                ("Tablas Encontradas", metrics.get('total_tables', 0)),
                ("Fases Completadas", metrics.get('phases_completed', 0)),
                ("Porcentaje Máximo", f"{metrics.get('max_percentage', 0)}%"),
                ("Números Encontrados", metrics.get('total_numbers', 0)),
            ]
            
            for metric_name, metric_value in metric_data:
                ws1[f'A{row}'] = metric_name
                ws1[f'B{row}'] = metric_value
                ws1[f'A{row}'].border = border
                ws1[f'B{row}'].border = border
                ws1[f'B{row}'].alignment = Alignment(horizontal='right')
                row += 1
            
            # Ajustar ancho de columnas
            ws1.column_dimensions['A'].width = 30
            ws1.column_dimensions['B'].width = 20
            
            # Formato condicional para valores altos
            for r in range(5, row):
                cell = ws1[f'B{r}']
                if isinstance(cell.value, (int, float)) and cell.value > 1000:
                    cell.number_format = '#,##0'
            
            # Gráfica de barras
            if metric_data and len(metric_data) > 0:
                try:
                    chart1 = BarChart()
                    chart1.type = "col"
                    chart1.style = 10
                    chart1.title = "Métricas del Documento"
                    chart1.y_axis.title = 'Valor'
                    chart1.x_axis.title = 'Métrica'
                    
                    # Crear referencias correctas
                    data_ref = Reference(ws1, min_col=2, min_row=5, max_row=row-1)
                    cats_ref = Reference(ws1, min_col=1, min_row=5, max_row=row-1)
                    
                    chart1.add_data(data_ref, titles_from_data=False)
                    chart1.set_categories(cats_ref)
                    chart1.height = 10
                    chart1.width = 15
                    
                    ws1.add_chart(chart1, "D4")
                except Exception as e:
                    print(f"⚠️  Error creando gráfica en Excel: {e}")
                    import traceback
                    traceback.print_exc()
            
            # Hoja 2: Secciones
            if doc_data.get('sections'):
                ws2 = wb.create_sheet("Secciones")
                
                ws2['A1'] = "Nivel"
                ws2['B1'] = "Título"
                ws2['C1'] = "Tipo"
                ws2['A1'].font = header_font
                ws2['B1'].font = header_font
                ws2['C1'].font = header_font
                ws2['A1'].fill = header_fill
                ws2['B1'].fill = header_fill
                ws2['C1'].fill = header_fill
                ws2['A1'].border = border
                ws2['B1'].border = border
                ws2['C1'].border = border
                
                row = 2
                for section in doc_data['sections'][:200]:  # Limitar a 200 secciones
                    ws2[f'A{row}'] = section['level']
                    ws2[f'B{row}'] = section['title'][:100]  # Limitar longitud
                    ws2[f'C{row}'] = f"H{section['level']}"
                    ws2[f'A{row}'].border = border
                    ws2[f'B{row}'].border = border
                    ws2[f'C{row}'].border = border
                    ws2[f'A{row}'].alignment = Alignment(horizontal='center')
                    # Color alternado para filas
                    if row % 2 == 0:
                        fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                        ws2[f'A{row}'].fill = fill
                        ws2[f'B{row}'].fill = fill
                        ws2[f'C{row}'].fill = fill
                    row += 1
                
                ws2.column_dimensions['A'].width = 10
                ws2.column_dimensions['B'].width = 70
                ws2.column_dimensions['C'].width = 10
                
                # Gráfica de secciones por nivel
                levels = {}
                for section in doc_data['sections']:
                    level = section['level']
                    levels[level] = levels.get(level, 0) + 1
                
                if levels:
                    # Crear datos en la hoja para el gráfico
                    chart_row = 2
                    for level in sorted(levels.keys()):
                        ws2[f'C{chart_row}'] = f"Nivel {level}"
                        ws2[f'D{chart_row}'] = levels[level]
                        chart_row += 1
                    
                    chart2 = PieChart()
                    chart2.title = "Distribución de Secciones por Nivel"
                    
                    data_end_row = 2 + len(levels) - 1
                    data_ref = Reference(ws2, min_col=4, min_row=2, max_row=data_end_row)
                    cats_ref = Reference(ws2, min_col=3, min_row=2, max_row=data_end_row)
                    
                    chart2.add_data(data_ref, titles_from_data=False)
                    chart2.set_categories(cats_ref)
                    chart2.height = 10
                    chart2.width = 15
                    
                    ws2.add_chart(chart2, "F2")
            
            # Hoja 3: Métricas Detalladas
            if doc_data.get('metrics'):
                ws3 = wb.create_sheet("Métricas Detalladas")
                metrics = doc_data['metrics']
                
                ws3['A1'] = "Categoría"
                ws3['B1'] = "Métrica"
                ws3['C1'] = "Valor"
                ws3['A1'].font = header_font
                ws3['B1'].font = header_font
                ws3['C1'].font = header_font
                ws3['A1'].fill = header_fill
                ws3['B1'].fill = header_fill
                ws3['C1'].fill = header_fill
                ws3['A1'].border = border
                ws3['B1'].border = border
                ws3['C1'].border = border
                
                row = 2
                detailed_metrics = [
                    ("Contenido", "Total de Líneas", metrics.get('total_lines', 0)),
                    ("Contenido", "Total de Palabras", metrics.get('total_words', 0)),
                    ("Estructura", "Total de Secciones", metrics.get('total_sections', 0)),
                    ("Código", "Bloques de Código", metrics.get('code_blocks', 0)),
                    ("Tablas", "Tablas Encontradas", metrics.get('total_tables', 0)),
                    ("Progreso", "Fases Completadas", metrics.get('phases_completed', 0)),
                    ("Datos", "Números Encontrados", metrics.get('total_numbers', 0)),
                    ("Datos", "Número Máximo", metrics.get('max_number', 0)),
                    ("Datos", "Número Promedio", metrics.get('avg_number', 0)),
                    ("Progreso", "Porcentaje Máximo", f"{metrics.get('max_percentage', 0)}%"),
                    ("Progreso", "Porcentaje Promedio", f"{metrics.get('avg_percentage', 0)}%"),
                    ("Fechas", "Fechas Encontradas", metrics.get('dates_found', 0)),
                    ("Lectura", "Tiempo Estimado (min)", metrics.get('reading_time_minutes', 0)),
                    ("Complejidad", "Nivel", metrics.get('complexity_level', 'N/A')),
                    ("Complejidad", "Score", metrics.get('complexity_score', 0)),
                    ("Enlaces", "Total de Enlaces", metrics.get('total_links', 0)),
                    ("Imágenes", "Total de Imágenes", metrics.get('total_images', 0)),
                    ("Calidad", "Score General", f"{metrics.get('quality_score', 0)}/100"),
                    ("Calidad", "Nivel", metrics.get('quality_level', 'N/A')),
                    ("Legibilidad", "Flesch Score", metrics.get('flesch_score', 0)),
                    ("Legibilidad", "Nivel Lectura", metrics.get('reading_level', 'N/A')),
                    ("Coherencia", "Score", metrics.get('coherence_score', 0)),
                    ("Sentimiento", "Tono Principal", metrics.get('sentiment_tone', 'N/A')),
                ]
                
                for cat, metric, value in detailed_metrics:
                    if value and value != 0 and str(value) != "0%":
                        ws3[f'A{row}'] = cat
                        ws3[f'B{row}'] = metric
                        ws3[f'C{row}'] = value
                        ws3[f'A{row}'].border = border
                        ws3[f'B{row}'].border = border
                        ws3[f'C{row}'].border = border
                        if row % 2 == 0:
                            fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                            ws3[f'A{row}'].fill = fill
                            ws3[f'B{row}'].fill = fill
                            ws3[f'C{row}'].fill = fill
                        row += 1
                
                ws3.column_dimensions['A'].width = 20
                ws3.column_dimensions['B'].width = 30
                ws3.column_dimensions['C'].width = 20
            
            # Hoja 4: Análisis de Fases (si hay fases)
            if doc_data.get('metrics', {}).get('phases_list'):
                ws4 = wb.create_sheet("Análisis de Fases")
                phases = doc_data['metrics']['phases_list']
                
                ws4['A1'] = "Fase"
                ws4['B1'] = "Estado"
                ws4['C1'] = "Completitud"
                ws4['A1'].font = header_font
                ws4['B1'].font = header_font
                ws4['C1'].font = header_font
                ws4['A1'].fill = header_fill
                ws4['B1'].fill = header_fill
                ws4['C1'].fill = header_fill
                ws4['A1'].border = border
                ws4['B1'].border = border
                ws4['C1'].border = border
                
                row = 2
                for phase in phases:
                    ws4[f'A{row}'] = f"Fase {phase}"
                    ws4[f'B{row}'] = "✅ Completa"
                    ws4[f'C{row}'] = "100%"
                    ws4[f'A{row}'].border = border
                    ws4[f'B{row}'].border = border
                    ws4[f'C{row}'].border = border
                    ws4[f'C{row}'].fill = PatternFill(start_color="C6E0B4", end_color="C6E0B4", fill_type="solid")
                    if row % 2 == 0:
                        fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                        ws4[f'A{row}'].fill = fill
                        ws4[f'B{row}'].fill = fill
                    row += 1
                
                ws4.column_dimensions['A'].width = 15
                ws4.column_dimensions['B'].width = 20
                ws4.column_dimensions['C'].width = 15
            
            # Hoja 5: Palabras Clave (si hay)
            if doc_data.get('top_keywords'):
                ws5 = wb.create_sheet("Palabras Clave")
                keywords = doc_data['top_keywords'][:30]  # Top 30
                
                ws5['A1'] = "Palabra"
                ws5['B1'] = "Frecuencia"
                ws5['C1'] = "Porcentaje"
                ws5['A1'].font = header_font
                ws5['B1'].font = header_font
                ws5['C1'].font = header_font
                ws5['A1'].fill = header_fill
                ws5['B1'].fill = header_fill
                ws5['C1'].fill = header_fill
                ws5['A1'].border = border
                ws5['B1'].border = border
                ws5['C1'].border = border
                
                total_keyword_count = sum(k['count'] for k in keywords)
                
                row = 2
                for keyword in keywords:
                    ws5[f'A{row}'] = keyword['word'].title()
                    ws5[f'B{row}'] = keyword['count']
                    percentage = (keyword['count'] / total_keyword_count * 100) if total_keyword_count > 0 else 0
                    ws5[f'C{row}'] = f"{percentage:.2f}%"
                    ws5[f'A{row}'].border = border
                    ws5[f'B{row}'].border = border
                    ws5[f'C{row}'].border = border
                    ws5[f'B{row}'].alignment = Alignment(horizontal='right')
                    ws5[f'C{row}'].alignment = Alignment(horizontal='right')
                    if row % 2 == 0:
                        fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                        ws5[f'A{row}'].fill = fill
                        ws5[f'B{row}'].fill = fill
                        ws5[f'C{row}'].fill = fill
                    row += 1
                
                ws5.column_dimensions['A'].width = 25
                ws5.column_dimensions['B'].width = 15
                ws5.column_dimensions['C'].width = 15
            
            # Hoja 6: Enlaces y Referencias (si hay)
            if doc_data.get('links'):
                ws6 = wb.create_sheet("Enlaces y Referencias")
                links = doc_data['links'][:100]  # Top 100
                
                ws6['A1'] = "Texto del Enlace"
                ws6['B1'] = "URL"
                ws6['A1'].font = header_font
                ws6['B1'].font = header_font
                ws6['A1'].fill = header_fill
                ws6['B1'].fill = header_fill
                ws6['A1'].border = border
                ws6['B1'].border = border
                
                row = 2
                for link_text, link_url in links:
                    ws6[f'A{row}'] = link_text[:100]  # Limitar longitud
                    ws6[f'B{row}'] = link_url[:200]  # Limitar longitud
                    ws6[f'A{row}'].border = border
                    ws6[f'B{row}'].border = border
                    if row % 2 == 0:
                        fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                        ws6[f'A{row}'].fill = fill
                        ws6[f'B{row}'].fill = fill
                    row += 1
                
                ws6.column_dimensions['A'].width = 40
                ws6.column_dimensions['B'].width = 60
            
            # Hoja 7: Análisis de Código (si hay)
            if doc_data.get('code_languages'):
                ws7 = wb.create_sheet("Análisis de Código")
                code_langs = doc_data['code_languages']
                
                ws7['A1'] = "Lenguaje"
                ws7['B1'] = "Bloques"
                ws7['C1'] = "Porcentaje"
                ws7['A1'].font = header_font
                ws7['B1'].font = header_font
                ws7['C1'].font = header_font
                ws7['A1'].fill = header_fill
                ws7['B1'].fill = header_fill
                ws7['C1'].fill = header_fill
                ws7['A1'].border = border
                ws7['B1'].border = border
                ws7['C1'].border = border
                
                total_blocks = sum(code_langs.values())
                
                row = 2
                for lang, count in sorted(code_langs.items(), key=lambda x: x[1], reverse=True):
                    ws7[f'A{row}'] = lang
                    ws7[f'B{row}'] = count
                    percentage = (count / total_blocks * 100) if total_blocks > 0 else 0
                    ws7[f'C{row}'] = f"{percentage:.2f}%"
                    ws7[f'A{row}'].border = border
                    ws7[f'B{row}'].border = border
                    ws7[f'C{row}'].border = border
                    ws7[f'B{row}'].alignment = Alignment(horizontal='right')
                    ws7[f'C{row}'].alignment = Alignment(horizontal='right')
                    if row % 2 == 0:
                        fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                        ws7[f'A{row}'].fill = fill
                        ws7[f'B{row}'].fill = fill
                        ws7[f'C{row}'].fill = fill
                    row += 1
                
                ws7.column_dimensions['A'].width = 20
                ws7.column_dimensions['B'].width = 15
                ws7.column_dimensions['C'].width = 15
                
                # Gráfica de lenguajes
                if len(code_langs) > 0:
                    chart3 = PieChart()
                    chart3.title = "Distribución de Lenguajes de Código"
                    
                    chart_row = 2
                    data_end_row = 2 + len(code_langs) - 1
                    data_ref = Reference(ws7, min_col=2, min_row=2, max_row=data_end_row)
                    cats_ref = Reference(ws7, min_col=1, min_row=2, max_row=data_end_row)
                    
                    chart3.add_data(data_ref, titles_from_data=False)
                    chart3.set_categories(cats_ref)
                    chart3.height = 10
                    chart3.width = 15
                    
                    ws7.add_chart(chart3, "E2")
            
            # Hoja 8: Análisis de Calidad
            if doc_data.get('quality_analysis'):
                ws8 = wb.create_sheet("Análisis de Calidad")
                quality = doc_data['quality_analysis']
                
                ws8['A1'] = "Factor de Calidad"
                ws8['B1'] = "Score"
                ws8['C1'] = "Máximo"
                ws8['A1'].font = header_font
                ws8['B1'].font = header_font
                ws8['C1'].font = header_font
                ws8['A1'].fill = header_fill
                ws8['B1'].fill = header_fill
                ws8['C1'].fill = header_fill
                ws8['A1'].border = border
                ws8['B1'].border = border
                ws8['C1'].border = border
                
                row = 2
                quality_factors = [
                    ("Estructura", quality['scores'].get('structure', 0), 25),
                    ("Contenido", quality['scores'].get('content', 0), 25),
                    ("Enlaces", quality['scores'].get('links', 0), 20),
                    ("Tablas", quality['scores'].get('tables', 0), 15),
                    ("Código", quality['scores'].get('code', 0), 15),
                ]
                
                for factor, score, max_score in quality_factors:
                    ws8[f'A{row}'] = factor
                    ws8[f'B{row}'] = score
                    ws8[f'C{row}'] = max_score
                    ws8[f'A{row}'].border = border
                    ws8[f'B{row}'].border = border
                    ws8[f'C{row}'].border = border
                    ws8[f'B{row}'].alignment = Alignment(horizontal='right')
                    ws8[f'C{row}'].alignment = Alignment(horizontal='right')
                    
                    # Color según score
                    percentage = (score / max_score * 100) if max_score > 0 else 0
                    if percentage >= 80:
                        fill_color = "C6E0B4"  # Verde
                    elif percentage >= 60:
                        fill_color = "FFE699"  # Amarillo
                    else:
                        fill_color = "F8CBAD"  # Naranja
                    
                    ws8[f'B{row}'].fill = PatternFill(start_color=fill_color, end_color=fill_color, fill_type="solid")
                    
                    if row % 2 == 0:
                        fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                        ws8[f'A{row}'].fill = fill
                        ws8[f'C{row}'].fill = fill
                    row += 1
                
                # Score total
                ws8[f'A{row}'] = "TOTAL"
                ws8[f'B{row}'] = quality.get('overall_score', 0)
                ws8[f'C{row}'] = 100
                ws8[f'A{row}'].font = Font(bold=True, size=12)
                ws8[f'B{row}'].font = Font(bold=True, size=12)
                ws8[f'C{row}'].font = Font(bold=True, size=12)
                ws8[f'A{row}'].border = border
                ws8[f'B{row}'].border = border
                ws8[f'C{row}'].border = border
                
                # Nivel de calidad
                ws8[f'A{row+2}'] = "Nivel de Calidad:"
                ws8[f'B{row+2}'] = quality.get('quality_level', 'N/A')
                ws8[f'A{row+2}'].font = Font(bold=True)
                ws8[f'B{row+2}'].font = Font(bold=True, size=14)
                
                # Issues y recomendaciones
                if quality.get('issues'):
                    ws8[f'A{row+4}'] = "Issues Detectados:"
                    ws8[f'A{row+4}'].font = Font(bold=True)
                    issue_row = row + 5
                    for issue in quality['issues'][:10]:
                        ws8[f'A{issue_row}'] = f"• {issue}"
                        issue_row += 1
                
                if quality.get('recommendations'):
                    start_row = issue_row if quality.get('issues') else row + 4
                    ws8[f'A{start_row}'] = "Recomendaciones:"
                    ws8[f'A{start_row}'].font = Font(bold=True)
                    rec_row = start_row + 1
                    for rec in quality['recommendations'][:10]:
                        ws8[f'A{rec_row}'] = f"→ {rec}"
                        rec_row += 1
                
                ws8.column_dimensions['A'].width = 40
                ws8.column_dimensions['B'].width = 15
                ws8.column_dimensions['C'].width = 15
                
                # Gráfica de calidad
                if len(quality_factors) > 0:
                    chart4 = BarChart()
                    chart4.type = "col"
                    chart4.style = 10
                    chart4.title = "Análisis de Calidad por Factor"
                    chart4.y_axis.title = 'Score'
                    
                    data_ref = Reference(ws8, min_col=2, min_row=2, max_row=row-1)
                    cats_ref = Reference(ws8, min_col=1, min_row=2, max_row=row-1)
                    
                    chart4.add_data(data_ref, titles_from_data=False)
                    chart4.set_categories(cats_ref)
                    chart4.height = 10
                    chart4.width = 15
                    
                    ws8.add_chart(chart4, "E2")
            
            # Hoja 9: Análisis de Legibilidad
            if doc_data.get('readability'):
                ws9 = wb.create_sheet("Legibilidad")
                readability = doc_data['readability']
                
                ws9['A1'] = "Métrica"
                ws9['B1'] = "Valor"
                ws9['A1'].font = header_font
                ws9['B1'].font = header_font
                ws9['A1'].fill = header_fill
                ws9['B1'].fill = header_fill
                ws9['A1'].border = border
                ws9['B1'].border = border
                
                row = 2
                readability_metrics = [
                    ("Flesch Reading Ease Score", readability.get('flesch_score', 0)),
                    ("Nivel de Lectura", readability.get('reading_level', 'N/A')),
                    ("Longitud Promedio de Oraciones", readability.get('avg_sentence_length', 0)),
                    ("Sílabas Promedio por Palabra", readability.get('avg_syllables_per_word', 0)),
                    ("Total de Oraciones", readability.get('sentence_count', 0)),
                    ("Total de Palabras Analizadas", readability.get('word_count', 0)),
                ]
                
                for metric, value in readability_metrics:
                    ws9[f'A{row}'] = metric
                    ws9[f'B{row}'] = value
                    ws9[f'A{row}'].border = border
                    ws9[f'B{row}'].border = border
                    if row % 2 == 0:
                        fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                        ws9[f'A{row}'].fill = fill
                        ws9[f'B{row}'].fill = fill
                    row += 1
                
                ws9.column_dimensions['A'].width = 35
                ws9.column_dimensions['B'].width = 25
            
            # Guardar
            wb.save(str(output_path))
            print(f"✅ Excel creado: {output_path}")
            
        except Exception as e:
            print(f"❌ Error creando Excel: {e}")
            import traceback
            traceback.print_exc()
    
    def export_to_json(self, doc_data: Dict[str, Any], doc_title: str, output_path: Path):
        """Exporta datos estructurados a JSON"""
        try:
            export_data = {
                'title': doc_title,
                'generated_at': datetime.now().isoformat(),
                'metrics': doc_data.get('metrics', {}),
                'sections': doc_data.get('sections', [])[:50],
                'quality_analysis': doc_data.get('quality_analysis', {}),
                'readability': doc_data.get('readability', {}),
                'sentiment': doc_data.get('sentiment', {}),
                'coherence': doc_data.get('coherence', {}),
                'top_keywords': doc_data.get('top_keywords', [])[:20],
                'links_count': doc_data.get('metrics', {}).get('total_links', 0),
                'images_count': doc_data.get('metrics', {}).get('total_images', 0)
            }
            
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(export_data, f, indent=2, ensure_ascii=False)
            
            print(f"✅ JSON exportado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error exportando JSON: {e}")
    
    def export_to_html(self, content: str, doc_data: Dict[str, Any], 
                      doc_title: str, output_path: Path, graphs: List[str]):
        """Exporta documento a HTML interactivo"""
        try:
            html_content = f"""<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{doc_title}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            line-height: 1.6;
            color: #333;
            background: #f5f5f5;
        }}
        .container {{
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            padding: 20px;
            box-shadow: 0 0 10px rgba(0,0,0,0.1);
        }}
        .header {{
            background: linear-gradient(135deg, #2E86AB 0%, #A23B72 100%);
            color: white;
            padding: 40px;
            text-align: center;
            margin: -20px -20px 30px -20px;
        }}
        .header h1 {{
            font-size: 2.5em;
            margin-bottom: 10px;
        }}
        .header .meta {{
            font-size: 0.9em;
            opacity: 0.9;
        }}
        .metrics-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
            gap: 15px;
            margin: 30px 0;
        }}
        .metric-card {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            border-left: 4px solid #2E86AB;
        }}
        .metric-card h3 {{
            color: #2E86AB;
            font-size: 0.9em;
            margin-bottom: 5px;
        }}
        .metric-card .value {{
            font-size: 2em;
            font-weight: bold;
            color: #333;
        }}
        .graphs-section {{
            margin: 40px 0;
        }}
        .graphs-section h2 {{
            color: #2E86AB;
            margin-bottom: 20px;
            padding-bottom: 10px;
            border-bottom: 2px solid #2E86AB;
        }}
        .graph-container {{
            margin: 20px 0;
            text-align: center;
        }}
        .graph-container img {{
            max-width: 100%;
            height: auto;
            border-radius: 8px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
        }}
        .content-section {{
            margin: 40px 0;
        }}
        .content-section h2 {{
            color: #A23B72;
            margin: 30px 0 15px 0;
        }}
        .content-section h3 {{
            color: #F18F01;
            margin: 20px 0 10px 0;
        }}
        .content-section p {{
            margin: 10px 0;
            text-align: justify;
        }}
        .quality-badge {{
            display: inline-block;
            padding: 5px 15px;
            border-radius: 20px;
            font-weight: bold;
            margin: 5px;
        }}
        .quality-excellent {{ background: #C6E0B4; color: #2d5016; }}
        .quality-good {{ background: #FFE699; color: #856404; }}
        .quality-regular {{ background: #F8CBAD; color: #854d0e; }}
        .quality-poor {{ background: #F4B084; color: #843c0c; }}
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 20px 0;
        }}
        table th, table td {{
            padding: 12px;
            text-align: left;
            border-bottom: 1px solid #ddd;
        }}
        table th {{
            background: #2E86AB;
            color: white;
        }}
        table tr:hover {{
            background: #f5f5f5;
        }}
        .footer {{
            text-align: center;
            padding: 20px;
            color: #666;
            font-size: 0.9em;
            margin-top: 40px;
            border-top: 1px solid #ddd;
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>{doc_title}</h1>
            <div class="meta">
                Generado el {datetime.now().strftime('%d de %B de %Y')}
            </div>
        </div>
        
        <div class="metrics-grid">
"""
            
            metrics = doc_data.get('metrics', {})
            quality = doc_data.get('quality_analysis', {})
            readability = doc_data.get('readability', {})
            
            metrics_html = [
                ("Líneas", f"{metrics.get('total_lines', 0):,}"),
                ("Palabras", f"{metrics.get('total_words', 0):,}"),
                ("Secciones", metrics.get('total_sections', 0)),
                ("Calidad", f"{metrics.get('quality_score', 0)}/100"),
                ("Legibilidad", readability.get('reading_level', 'N/A')),
                ("Tiempo Lectura", f"~{metrics.get('reading_time_minutes', 0)} min"),
            ]
            
            for metric_name, metric_value in metrics_html:
                html_content += f"""
            <div class="metric-card">
                <h3>{metric_name}</h3>
                <div class="value">{metric_value}</div>
            </div>
"""
            
            html_content += """
        </div>
        
        <div class="graphs-section">
            <h2>📊 Visualizaciones</h2>
"""
            
            for graph_path in graphs[:6]:  # Mostrar hasta 6 gráficas
                if os.path.exists(graph_path):
                    graph_name = Path(graph_path).name
                    html_content += f"""
            <div class="graph-container">
                <img src="{graph_path}" alt="{graph_name}">
            </div>
"""
            
            html_content += """
        </div>
        
        <div class="content-section">
            <h2>📄 Contenido del Documento</h2>
"""
            
            # Convertir markdown a HTML básico
            html_content += markdown.markdown(content[:5000], extensions=['tables', 'fenced_code', 'codehilite'])
            
            html_content += """
        </div>
        
        <div class="footer">
            <p>Documento generado automáticamente con análisis avanzado</p>
        </div>
    </div>
</body>
</html>
"""
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            print(f"✅ HTML exportado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error exportando HTML: {e}")
            import traceback
            traceback.print_exc()
    
    def analyze_document_similarity(self, all_docs_data: List[Dict]) -> Dict[str, Dict[str, float]]:
        """Analiza similitud entre documentos usando múltiples métodos"""
        similarity_matrix = {}
        
        for i, doc1 in enumerate(all_docs_data):
            doc1_title = doc1['title']
            similarity_matrix[doc1_title] = {}
            
            for j, doc2 in enumerate(all_docs_data):
                if i == j:
                    similarity_matrix[doc1_title][doc2['title']] = 1.0
                    continue
                
                # Método 1: Similitud de texto (SequenceMatcher)
                text1 = doc1.get('content', '')
                text2 = doc2.get('content', '')
                text_similarity = SequenceMatcher(None, text1[:2000], text2[:2000]).ratio()
                
                # Método 2: Similitud de keywords
                keywords1 = set(doc1.get('top_keywords', []))
                keywords2 = set(doc2.get('top_keywords', []))
                if keywords1 or keywords2:
                    keyword_similarity = len(keywords1 & keywords2) / max(len(keywords1 | keywords2), 1)
                else:
                    keyword_similarity = 0
                
                # Método 3: Similitud de secciones
                sections1 = set(doc1.get('sections', []))
                sections2 = set(doc2.get('sections', []))
                if sections1 or sections2:
                    section_similarity = len(sections1 & sections2) / max(len(sections1 | sections2), 1)
                else:
                    section_similarity = 0
                
                # Score combinado
                combined_score = (
                    text_similarity * 0.5 +
                    keyword_similarity * 0.3 +
                    section_similarity * 0.2
                )
                
                similarity_matrix[doc1_title][doc2['title']] = round(combined_score, 3)
        
        return similarity_matrix
    
    def analyze_document_dependencies(self, all_docs_data: List[Dict]) -> List[Dict]:
        """Analiza dependencias y referencias entre documentos"""
        dependencies = []
        
        for doc in all_docs_data:
            content = doc.get('content', '')
            doc_title = doc['title']
            
            # Buscar referencias a otros documentos
            for other_doc in all_docs_data:
                if other_doc['title'] == doc_title:
                    continue
                
                other_title = other_doc['title']
                # Buscar menciones del título (simplificado)
                title_words = other_title.split()
                if len(title_words) > 0:
                    # Buscar al menos 2 palabras del título
                    matches = sum(1 for word in title_words[:3] if word.lower() in content.lower())
                    if matches >= 2:
                        dependencies.append({
                            'source': doc_title,
                            'target': other_title,
                            'strength': matches / len(title_words[:3]),
                            'type': 'reference'
                        })
            
            # Buscar enlaces a otros documentos
            links = re.findall(r'\[([^\]]+)\]\(([^\)]+)\)', content)
            for link_text, link_url in links:
                # Verificar si el enlace apunta a otro documento
                for other_doc in all_docs_data:
                    if other_doc['title'] == doc_title:
                        continue
                    
                    other_path = other_doc.get('path', '')
                    if other_path in link_url or link_url.endswith('.md'):
                        dependencies.append({
                            'source': doc_title,
                            'target': other_doc['title'],
                            'strength': 1.0,
                            'type': 'link',
                            'link_text': link_text
                        })
        
        return dependencies
    
    def create_similarity_report(self, similarity_matrix: Dict[str, Dict[str, float]], 
                                 output_path: Path):
        """Crea reporte Excel de similitud entre documentos"""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            
            wb = Workbook()
            ws = wb.active
            ws.title = "Matriz de Similitud"
            
            # Estilos
            header_font = Font(bold=True, size=12, color="FFFFFF")
            header_fill = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Título
            ws['A1'] = "Matriz de Similitud entre Documentos"
            ws['A1'].font = Font(bold=True, size=14)
            ws.merge_cells('A1:Z1')
            
            # Headers
            doc_titles = list(similarity_matrix.keys())
            ws['A2'] = "Documento"
            ws['A2'].font = header_font
            ws['A2'].fill = header_fill
            ws['A2'].border = border
            
            for col_idx, title in enumerate(doc_titles, 2):
                cell = ws.cell(row=2, column=col_idx)
                cell.value = title[:30]  # Limitar longitud
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
                cell.alignment = Alignment(horizontal='center', vertical='center', text_rotation=90)
            
            # Datos
            for row_idx, doc1_title in enumerate(doc_titles, 3):
                ws.cell(row=row_idx, column=1).value = doc1_title[:50]
                ws.cell(row=row_idx, column=1).border = border
                
                for col_idx, doc2_title in enumerate(doc_titles, 2):
                    similarity = similarity_matrix[doc1_title].get(doc2_title, 0)
                    cell = ws.cell(row=row_idx, column=col_idx)
                    cell.value = similarity
                    cell.border = border
                    cell.number_format = '0.000'
                    
                    # Color según similitud
                    if similarity >= 0.7:
                        cell.fill = PatternFill(start_color="C6E0B4", end_color="C6E0B4", fill_type="solid")
                    elif similarity >= 0.4:
                        cell.fill = PatternFill(start_color="FFE699", end_color="FFE699", fill_type="solid")
                    elif similarity >= 0.1:
                        cell.fill = PatternFill(start_color="F8CBAD", end_color="F8CBAD", fill_type="solid")
            
            # Ajustar anchos
            ws.column_dimensions['A'].width = 50
            for col_idx in range(2, len(doc_titles) + 2):
                ws.column_dimensions[get_column_letter(col_idx)].width = 15
            
            wb.save(str(output_path))
            print(f"✅ Reporte de similitud guardado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando reporte de similitud: {e}")
    
    def create_dependencies_report(self, dependencies: List[Dict], output_path: Path):
        """Crea reporte Excel de dependencias entre documentos"""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            
            wb = Workbook()
            ws = wb.active
            ws.title = "Dependencias"
            
            # Estilos
            header_font = Font(bold=True, size=12, color="FFFFFF")
            header_fill = PatternFill(start_color="A23B72", end_color="A23B72", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Título
            ws['A1'] = "Análisis de Dependencias entre Documentos"
            ws['A1'].font = Font(bold=True, size=14)
            ws.merge_cells('A1:E1')
            
            # Headers
            headers = ['Documento Origen', 'Documento Destino', 'Tipo', 'Fuerza', 'Detalles']
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=2, column=col_idx)
                cell.value = header
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
            
            # Datos
            for row_idx, dep in enumerate(dependencies, 3):
                ws.cell(row=row_idx, column=1).value = dep['source'][:50]
                ws.cell(row=row_idx, column=2).value = dep['target'][:50]
                ws.cell(row=row_idx, column=3).value = dep['type']
                ws.cell(row=row_idx, column=4).value = dep['strength']
                ws.cell(row=row_idx, column=5).value = dep.get('link_text', '')
                
                for col in range(1, 6):
                    cell = ws.cell(row=row_idx, column=col)
                    cell.border = border
                    if row_idx % 2 == 0:
                        cell.fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
            
            # Ajustar anchos
            ws.column_dimensions['A'].width = 40
            ws.column_dimensions['B'].width = 40
            ws.column_dimensions['C'].width = 15
            ws.column_dimensions['D'].width = 12
            ws.column_dimensions['E'].width = 30
            
            wb.save(str(output_path))
            print(f"✅ Reporte de dependencias guardado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando reporte de dependencias: {e}")
    
    def create_comparison_report(self, all_docs_data: List[Dict], output_path: Path):
        """Crea reporte comparativo detallado entre documentos"""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.chart import BarChart, Reference
            
            wb = Workbook()
            
            # Hoja 1: Comparación de Métricas
            ws1 = wb.active
            ws1.title = "Comparación Métricas"
            
            header_font = Font(bold=True, size=12, color="FFFFFF")
            header_fill = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Headers
            headers = ['Documento', 'Líneas', 'Palabras', 'Secciones', 'Calidad', 'Legibilidad', 
                      'Complejidad', 'Enlaces', 'Código', 'Tiempo Lectura']
            for col_idx, header in enumerate(headers, 1):
                cell = ws1.cell(row=1, column=col_idx)
                cell.value = header
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
            
            # Datos
            for row_idx, doc in enumerate(all_docs_data, 2):
                metrics = doc.get('metrics', {})
                ws1.cell(row=row_idx, column=1).value = doc['title'][:40]
                ws1.cell(row=row_idx, column=2).value = metrics.get('total_lines', 0)
                ws1.cell(row=row_idx, column=3).value = metrics.get('total_words', 0)
                ws1.cell(row=row_idx, column=4).value = doc.get('sections_count', 0)
                ws1.cell(row=row_idx, column=5).value = doc.get('quality_score', 0)
                ws1.cell(row=row_idx, column=6).value = doc.get('readability_score', 0)
                ws1.cell(row=row_idx, column=7).value = metrics.get('complexity_level', 'N/A')
                ws1.cell(row=row_idx, column=8).value = metrics.get('total_links', 0)
                ws1.cell(row=row_idx, column=9).value = metrics.get('code_blocks', 0)
                ws1.cell(row=row_idx, column=10).value = metrics.get('reading_time_minutes', 0)
                
                for col in range(1, 11):
                    cell = ws1.cell(row=row_idx, column=col)
                    cell.border = border
                    if row_idx % 2 == 0:
                        cell.fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
            
            # Ajustar anchos
            for col in range(1, 11):
                ws1.column_dimensions[get_column_letter(col)].width = 15
            ws1.column_dimensions['A'].width = 40
            
            # Hoja 2: Palabras Clave Compartidas
            ws2 = wb.create_sheet("Palabras Clave Compartidas")
            
            # Encontrar palabras clave compartidas
            all_keywords = {}
            for doc in all_docs_data:
                keywords = doc.get('top_keywords', [])
                for keyword in keywords:
                    if keyword not in all_keywords:
                        all_keywords[keyword] = []
                    all_keywords[keyword].append(doc['title'])
            
            # Headers
            ws2['A1'] = "Palabra Clave"
            ws2['B1'] = "Documentos"
            ws2['C1'] = "Frecuencia"
            for col in range(1, 4):
                cell = ws2.cell(row=1, column=col)
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
            
            # Datos (solo palabras compartidas)
            shared_keywords = {k: v for k, v in all_keywords.items() if len(v) > 1}
            sorted_shared = sorted(shared_keywords.items(), key=lambda x: len(x[1]), reverse=True)
            
            for row_idx, (keyword, docs) in enumerate(sorted_shared[:50], 2):
                ws2.cell(row=row_idx, column=1).value = keyword
                ws2.cell(row=row_idx, column=2).value = ", ".join(docs[:3])
                ws2.cell(row=row_idx, column=3).value = len(docs)
                
                for col in range(1, 4):
                    cell = ws2.cell(row=row_idx, column=col)
                    cell.border = border
            
            ws2.column_dimensions['A'].width = 25
            ws2.column_dimensions['B'].width = 60
            ws2.column_dimensions['C'].width = 12
            
            wb.save(str(output_path))
            print(f"✅ Reporte comparativo guardado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando reporte comparativo: {e}")
            import traceback
            traceback.print_exc()
    
    def generate_executive_summary(self, doc_data: Dict[str, Any], doc_title: str) -> str:
        """Genera un resumen ejecutivo automático del documento"""
        summary_parts = []
        
        # Información básica
        metrics = doc_data.get('metrics', {})
        summary_parts.append(f"# Resumen Ejecutivo: {doc_title}\n")
        summary_parts.append(f"**Generado el:** {datetime.now().strftime('%d de %B de %Y')}\n")
        
        # Métricas clave
        summary_parts.append("## Métricas Clave\n")
        summary_parts.append(f"- **Líneas totales:** {metrics.get('total_lines', 0):,}")
        summary_parts.append(f"- **Palabras totales:** {metrics.get('total_words', 0):,}")
        summary_parts.append(f"- **Secciones:** {metrics.get('total_sections', 0)}")
        summary_parts.append(f"- **Tiempo de lectura estimado:** ~{metrics.get('reading_time_minutes', 0)} minutos\n")
        
        # Calidad
        quality_score = metrics.get('quality_score', 0)
        quality_level = metrics.get('quality_level', 'N/A')
        summary_parts.append(f"## Calidad del Documento\n")
        summary_parts.append(f"- **Score de calidad:** {quality_score}/100")
        summary_parts.append(f"- **Nivel:** {quality_level}\n")
        
        # Legibilidad
        readability = doc_data.get('readability', {})
        if readability:
            summary_parts.append(f"## Legibilidad\n")
            summary_parts.append(f"- **Flesch Score:** {readability.get('flesch_score', 0):.2f}")
            summary_parts.append(f"- **Nivel de lectura:** {readability.get('reading_level', 'N/A')}\n")
        
        # Secciones principales
        sections = doc_data.get('sections', [])
        main_sections = [s for s in sections if s.get('level', 0) == 1][:5]
        if main_sections:
            summary_parts.append("## Secciones Principales\n")
            for i, section in enumerate(main_sections, 1):
                summary_parts.append(f"{i}. {section.get('title', 'N/A')}")
            summary_parts.append("")
        
        # Palabras clave
        keywords = doc_data.get('top_keywords', [])[:10]
        if keywords:
            summary_parts.append("## Palabras Clave Principales\n")
            keyword_list = ", ".join([kw.get('word', '') for kw in keywords])
            summary_parts.append(f"{keyword_list}\n")
        
        # Recomendaciones
        quality_analysis = doc_data.get('quality_analysis', {})
        recommendations = quality_analysis.get('recommendations', [])
        if recommendations:
            summary_parts.append("## Recomendaciones\n")
            for i, rec in enumerate(recommendations[:5], 1):
                summary_parts.append(f"{i}. {rec}")
            summary_parts.append("")
        
        return "\n".join(summary_parts)
    
    def detect_duplicates(self, all_docs_data: List[Dict]) -> List[Dict]:
        """Detecta documentos duplicados o muy similares"""
        duplicates = []
        seen_hashes = {}
        
        for doc in all_docs_data:
            # Calcular hash del contenido
            content = doc.get('content', '')
            content_hash = hashlib.md5(content[:1000].encode()).hexdigest()
            
            if content_hash in seen_hashes:
                duplicates.append({
                    'doc1': seen_hashes[content_hash],
                    'doc2': doc['title'],
                    'similarity': 1.0,
                    'type': 'exact_duplicate'
                })
            else:
                seen_hashes[content_hash] = doc['title']
        
        # Detectar documentos muy similares (>90% similitud)
        for i, doc1 in enumerate(all_docs_data):
            for j, doc2 in enumerate(all_docs_data[i+1:], i+1):
                content1 = doc1.get('content', '')[:2000]
                content2 = doc2.get('content', '')[:2000]
                similarity = SequenceMatcher(None, content1, content2).ratio()
                
                if similarity >= 0.9:
                    duplicates.append({
                        'doc1': doc1['title'],
                        'doc2': doc2['title'],
                        'similarity': similarity,
                        'type': 'near_duplicate'
                    })
        
        return duplicates
    
    def analyze_document_age(self, all_docs_data: List[Dict]) -> Dict[str, Any]:
        """Analiza la antigüedad y actualización de documentos"""
        age_analysis = {
            'oldest': None,
            'newest': None,
            'avg_age_days': 0,
            'documents_by_age': []
        }
        
        # Extraer fechas de los documentos
        doc_dates = []
        for doc in all_docs_data:
            content = doc.get('content', '')
            # Buscar fechas en formato YYYY-MM-DD
            dates = re.findall(r'\d{4}-\d{2}-\d{2}', content)
            if dates:
                try:
                    latest_date = max(dates)
                    from datetime import datetime as dt
                    doc_date = dt.strptime(latest_date, '%Y-%m-%d')
                    days_old = (datetime.now() - doc_date).days
                    doc_dates.append({
                        'title': doc['title'],
                        'date': latest_date,
                        'days_old': days_old
                    })
                except:
                    pass
        
        if doc_dates:
            doc_dates.sort(key=lambda x: x['days_old'], reverse=True)
            age_analysis['oldest'] = doc_dates[0]
            age_analysis['newest'] = doc_dates[-1]
            age_analysis['avg_age_days'] = sum(d['days_old'] for d in doc_dates) // len(doc_dates)
            age_analysis['documents_by_age'] = doc_dates
        
        return age_analysis
    
    def create_executive_report(self, all_docs_data: List[Dict], output_path: Path):
        """Crea un reporte ejecutivo completo en PDF"""
        try:
            if not PDF_AVAILABLE:
                print("⚠️  reportlab no disponible para crear reporte ejecutivo")
                return
            
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib.enums import TA_CENTER, TA_LEFT
            from reportlab.lib import colors
            
            doc = SimpleDocTemplate(str(output_path), pagesize=A4,
                                   rightMargin=72, leftMargin=72,
                                   topMargin=72, bottomMargin=72)
            story = []
            styles = getSampleStyleSheet()
            
            # Título
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#2E86AB'),
                spaceAfter=30,
                alignment=TA_CENTER
            )
            story.append(Paragraph("Reporte Ejecutivo Consolidado", title_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Fecha
            date_style = ParagraphStyle(
                'DateStyle',
                parent=styles['Normal'],
                fontSize=10,
                textColor=colors.HexColor('#666666'),
                alignment=TA_CENTER
            )
            story.append(Paragraph(f"Generado el {datetime.now().strftime('%d de %B de %Y')}", date_style))
            story.append(Spacer(1, 0.3*inch))
            
            # Resumen general
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=16,
                textColor=colors.HexColor('#A23B72'),
                spaceAfter=12,
                spaceBefore=12
            )
            
            story.append(Paragraph("Resumen General", heading_style))
            story.append(Paragraph(
                f"Este reporte analiza {len(all_docs_data)} documentos importantes, "
                f"proporcionando una visión consolidada de su contenido, calidad y relaciones.",
                styles['Normal']
            ))
            story.append(Spacer(1, 0.2*inch))
            
            # Estadísticas generales
            total_lines = sum(d.get('metrics', {}).get('total_lines', 0) for d in all_docs_data)
            total_words = sum(d.get('metrics', {}).get('total_words', 0) for d in all_docs_data)
            avg_quality = sum(d.get('quality_score', 0) for d in all_docs_data) / len(all_docs_data) if all_docs_data else 0
            
            story.append(Paragraph("Estadísticas Generales", heading_style))
            stats_text = f"""
            <b>Total de líneas:</b> {total_lines:,}<br/>
            <b>Total de palabras:</b> {total_words:,}<br/>
            <b>Calidad promedio:</b> {avg_quality:.1f}/100<br/>
            <b>Documentos analizados:</b> {len(all_docs_data)}
            """
            story.append(Paragraph(stats_text, styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
            
            # Top documentos por calidad
            story.append(Paragraph("Documentos con Mayor Calidad", heading_style))
            sorted_docs = sorted(all_docs_data, key=lambda x: x.get('quality_score', 0), reverse=True)[:5]
            for i, doc in enumerate(sorted_docs, 1):
                doc_text = f"{i}. <b>{doc['title']}</b> - Calidad: {doc.get('quality_score', 0)}/100"
                story.append(Paragraph(doc_text, styles['Normal']))
                story.append(Spacer(1, 0.1*inch))
            
            story.append(PageBreak())
            
            # Resúmenes individuales
            story.append(Paragraph("Resúmenes por Documento", heading_style))
            story.append(Spacer(1, 0.2*inch))
            
            for doc in all_docs_data[:10]:  # Limitar a 10 documentos
                story.append(Paragraph(f"<b>{doc['title']}</b>", styles['Heading3']))
                
                metrics = doc.get('metrics', {})
                doc_summary = f"""
                Líneas: {metrics.get('total_lines', 0):,} | 
                Palabras: {metrics.get('total_words', 0):,} | 
                Secciones: {doc.get('sections_count', 0)} | 
                Calidad: {doc.get('quality_score', 0)}/100
                """
                story.append(Paragraph(doc_summary, styles['Normal']))
                story.append(Spacer(1, 0.15*inch))
            
            # Construir PDF
            doc.build(story)
            print(f"✅ Reporte ejecutivo creado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando reporte ejecutivo: {e}")
            import traceback
            traceback.print_exc()
    
    def create_advanced_analytics(self, all_docs_data: List[Dict], output_path: Path):
        """Crea un archivo Excel con análisis avanzados"""
        try:
            if not EXCEL_AVAILABLE:
                return
            
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.chart import BarChart, Reference
            
            wb = Workbook()
            
            # Hoja 1: Duplicados
            ws1 = wb.active
            ws1.title = "Duplicados Detectados"
            
            duplicates = self.detect_duplicates(all_docs_data)
            
            header_font = Font(bold=True, size=12, color="FFFFFF")
            header_fill = PatternFill(start_color="C73E1D", end_color="C73E1D", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            ws1['A1'] = "Documento 1"
            ws1['B1'] = "Documento 2"
            ws1['C1'] = "Similitud"
            ws1['D1'] = "Tipo"
            
            for col in range(1, 5):
                cell = ws1.cell(row=1, column=col)
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
            
            for row_idx, dup in enumerate(duplicates, 2):
                ws1.cell(row=row_idx, column=1).value = dup['doc1']
                ws1.cell(row=row_idx, column=2).value = dup['doc2']
                ws1.cell(row=row_idx, column=3).value = dup['similarity']
                ws1.cell(row=row_idx, column=4).value = dup['type']
                
                for col in range(1, 5):
                    cell = ws1.cell(row=row_idx, column=col)
                    cell.border = border
            
            ws1.column_dimensions['A'].width = 40
            ws1.column_dimensions['B'].width = 40
            ws1.column_dimensions['C'].width = 15
            ws1.column_dimensions['D'].width = 20
            
            # Hoja 2: Análisis de Antigüedad
            ws2 = wb.create_sheet("Antigüedad de Documentos")
            
            age_analysis = self.analyze_document_age(all_docs_data)
            
            ws2['A1'] = "Documento"
            ws2['B1'] = "Última Fecha"
            ws2['C1'] = "Días desde Actualización"
            
            for col in range(1, 4):
                cell = ws2.cell(row=1, column=col)
                cell.font = header_font
                cell.fill = PatternFill(start_color="F18F01", end_color="F18F01", fill_type="solid")
                cell.border = border
            
            for row_idx, doc_age in enumerate(age_analysis.get('documents_by_age', [])[:20], 2):
                ws2.cell(row=row_idx, column=1).value = doc_age['title'][:50]
                ws2.cell(row=row_idx, column=2).value = doc_age['date']
                ws2.cell(row=row_idx, column=3).value = doc_age['days_old']
                
                for col in range(1, 4):
                    cell = ws2.cell(row=row_idx, column=col)
                    cell.border = border
            
            ws2.column_dimensions['A'].width = 50
            ws2.column_dimensions['B'].width = 15
            ws2.column_dimensions['C'].width = 20
            
            wb.save(str(output_path))
            print(f"✅ Análisis avanzado creado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando análisis avanzado: {e}")
            import traceback
            traceback.print_exc()
    
    def create_concept_network(self, all_docs_data: List[Dict], output_path: Path):
        """Crea análisis de red de conceptos entre documentos"""
        try:
            if not EXCEL_AVAILABLE:
                return
            
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            
            wb = Workbook()
            ws = wb.active
            ws.title = "Red de Conceptos"
            
            # Recopilar todos los conceptos (palabras clave) y sus relaciones
            concept_docs = defaultdict(list)
            for doc in all_docs_data:
                keywords = doc.get('top_keywords', [])
                for kw in keywords[:15]:  # Top 15 por documento
                    concept = kw.get('word', '') if isinstance(kw, dict) else kw
                    if concept:
                        concept_docs[concept].append(doc['title'])
            
            # Filtrar conceptos compartidos (aparecen en al menos 2 documentos)
            shared_concepts = {k: v for k, v in concept_docs.items() if len(v) >= 2}
            sorted_concepts = sorted(shared_concepts.items(), key=lambda x: len(x[1]), reverse=True)
            
            header_font = Font(bold=True, size=12, color="FFFFFF")
            header_fill = PatternFill(start_color="06A77D", end_color="06A77D", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            ws['A1'] = "Concepto"
            ws['B1'] = "Documentos Relacionados"
            ws['C1'] = "Frecuencia"
            ws['D1'] = "Fuerza de Conexión"
            
            for col in range(1, 5):
                cell = ws.cell(row=1, column=col)
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
            
            for row_idx, (concept, docs) in enumerate(sorted_concepts[:50], 2):
                ws.cell(row=row_idx, column=1).value = concept
                ws.cell(row=row_idx, column=2).value = ", ".join(docs[:5])
                ws.cell(row=row_idx, column=3).value = len(docs)
                ws.cell(row=row_idx, column=4).value = len(docs) / len(all_docs_data)  # Fuerza normalizada
                
                for col in range(1, 5):
                    cell = ws.cell(row=row_idx, column=col)
                    cell.border = border
                    if row_idx % 2 == 0:
                        cell.fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
            
            ws.column_dimensions['A'].width = 30
            ws.column_dimensions['B'].width = 60
            ws.column_dimensions['C'].width = 15
            ws.column_dimensions['D'].width = 20
            
            wb.save(str(output_path))
            print(f"✅ Red de conceptos guardada: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando red de conceptos: {e}")
    
    def create_powerpoint_presentation(self, all_docs_data: List[Dict], output_path: Path):
        """Crea una presentación PowerPoint consolidada"""
        try:
            if not PPTX_AVAILABLE:
                print("⚠️  python-pptx no disponible para crear presentación")
                return
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Portada
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = "Análisis Consolidado de Documentos"
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = f"{len(all_docs_data)} Documentos Analizados"
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(162, 59, 114)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
            date_frame = date_box.text_frame
            date_frame.text = f"Generado el {datetime.now().strftime('%d de %B de %Y')}"
            date_frame.paragraphs[0].font.size = Pt(14)
            date_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
            date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Slide 2: Resumen Ejecutivo
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            title2 = slide2.shapes.title
            title2.text = "Resumen Ejecutivo"
            
            content2 = slide2.placeholders[1]
            tf2 = content2.text_frame
            tf2.text = f"Total de Documentos: {len(all_docs_data)}"
            
            total_lines = sum(d.get('metrics', {}).get('total_lines', 0) for d in all_docs_data)
            total_words = sum(d.get('metrics', {}).get('total_words', 0) for d in all_docs_data)
            avg_quality = sum(d.get('quality_score', 0) for d in all_docs_data) / len(all_docs_data) if all_docs_data else 0
            
            p = tf2.add_paragraph()
            p.text = f"Total de Líneas: {total_lines:,}"
            p = tf2.add_paragraph()
            p.text = f"Total de Palabras: {total_words:,}"
            p = tf2.add_paragraph()
            p.text = f"Calidad Promedio: {avg_quality:.1f}/100"
            
            # Slide 3: Top Documentos por Calidad
            slide3 = prs.slides.add_slide(prs.slide_layouts[1])
            title3 = slide3.shapes.title
            title3.text = "Top 5 Documentos por Calidad"
            
            content3 = slide3.placeholders[1]
            tf3 = content3.text_frame
            sorted_docs = sorted(all_docs_data, key=lambda x: x.get('quality_score', 0), reverse=True)[:5]
            
            for i, doc in enumerate(sorted_docs, 1):
                if i == 1:
                    tf3.text = f"{i}. {doc['title'][:60]} - {doc.get('quality_score', 0)}/100"
                else:
                    p = tf3.add_paragraph()
                    p.text = f"{i}. {doc['title'][:60]} - {doc.get('quality_score', 0)}/100"
            
            # Slides individuales para top 3 documentos
            for doc in sorted_docs[:3]:
                slide_doc = prs.slides.add_slide(prs.slide_layouts[1])
                title_doc = slide_doc.shapes.title
                title_doc.text = doc['title'][:80]
                
                content_doc = slide_doc.placeholders[1]
                tf_doc = content_doc.text_frame
                metrics = doc.get('metrics', {})
                
                tf_doc.text = f"Líneas: {metrics.get('total_lines', 0):,}"
                p = tf_doc.add_paragraph()
                p.text = f"Palabras: {metrics.get('total_words', 0):,}"
                p = tf_doc.add_paragraph()
                p.text = f"Secciones: {doc.get('sections_count', 0)}"
                p = tf_doc.add_paragraph()
                p.text = f"Calidad: {doc.get('quality_score', 0)}/100"
                p = tf_doc.add_paragraph()
                p.text = f"Legibilidad: {metrics.get('reading_level', 'N/A')}"
            
            prs.save(str(output_path))
            print(f"✅ Presentación PowerPoint guardada: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando presentación PowerPoint: {e}")
            import traceback
            traceback.print_exc()
    
    def create_interactive_dashboard(self, all_docs_data: List[Dict], output_path: Path):
        """Crea un dashboard HTML interactivo consolidado"""
        try:
            html_content = f"""<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Dashboard Interactivo - Análisis de Documentos</title>
    <script src="https://cdn.jsdelivr.net/npm/chart.js@3.9.1/dist/chart.min.js"></script>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            padding: 20px;
        }}
        .container {{
            max-width: 1400px;
            margin: 0 auto;
        }}
        .header {{
            background: white;
            padding: 30px;
            border-radius: 10px;
            margin-bottom: 20px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        }}
        .header h1 {{
            color: #2E86AB;
            font-size: 2.5em;
            margin-bottom: 10px;
        }}
        .stats-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin-bottom: 20px;
        }}
        .stat-card {{
            background: white;
            padding: 25px;
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            text-align: center;
        }}
        .stat-card h3 {{
            color: #666;
            font-size: 0.9em;
            margin-bottom: 10px;
        }}
        .stat-card .value {{
            font-size: 2.5em;
            font-weight: bold;
            color: #2E86AB;
        }}
        .charts-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(400px, 1fr));
            gap: 20px;
            margin-bottom: 20px;
        }}
        .chart-card {{
            background: white;
            padding: 20px;
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        }}
        .chart-card h3 {{
            color: #2E86AB;
            margin-bottom: 15px;
        }}
        .documents-table {{
            background: white;
            padding: 20px;
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            overflow-x: auto;
        }}
        table {{
            width: 100%;
            border-collapse: collapse;
        }}
        th, td {{
            padding: 12px;
            text-align: left;
            border-bottom: 1px solid #ddd;
        }}
        th {{
            background: #2E86AB;
            color: white;
            font-weight: bold;
        }}
        tr:hover {{
            background: #f5f5f5;
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>📊 Dashboard Interactivo de Documentos</h1>
            <p>Análisis consolidado de {len(all_docs_data)} documentos | Generado el {datetime.now().strftime('%d de %B de %Y')}</p>
        </div>
        
        <div class="stats-grid">
"""
            
            # Calcular estadísticas
            total_lines = sum(d.get('metrics', {}).get('total_lines', 0) for d in all_docs_data)
            total_words = sum(d.get('metrics', {}).get('total_words', 0) for d in all_docs_data)
            total_sections = sum(d.get('sections_count', 0) for d in all_docs_data)
            avg_quality = sum(d.get('quality_score', 0) for d in all_docs_data) / len(all_docs_data) if all_docs_data else 0
            
            html_content += f"""
            <div class="stat-card">
                <h3>Total de Líneas</h3>
                <div class="value">{total_lines:,}</div>
            </div>
            <div class="stat-card">
                <h3>Total de Palabras</h3>
                <div class="value">{total_words:,}</div>
            </div>
            <div class="stat-card">
                <h3>Total de Secciones</h3>
                <div class="value">{total_sections}</div>
            </div>
            <div class="stat-card">
                <h3>Calidad Promedio</h3>
                <div class="value">{avg_quality:.1f}/100</div>
            </div>
        </div>
        
        <div class="charts-grid">
            <div class="chart-card">
                <h3>Calidad por Documento</h3>
                <canvas id="qualityChart"></canvas>
            </div>
            <div class="chart-card">
                <h3>Distribución de Palabras</h3>
                <canvas id="wordsChart"></canvas>
            </div>
        </div>
        
        <div class="documents-table">
            <h3 style="margin-bottom: 15px; color: #2E86AB;">Documentos Analizados</h3>
            <table>
                <thead>
                    <tr>
                        <th>Documento</th>
                        <th>Categoría</th>
                        <th>Líneas</th>
                        <th>Palabras</th>
                        <th>Secciones</th>
                        <th>Calidad</th>
                    </tr>
                </thead>
                <tbody>
"""
            
            for doc in all_docs_data:
                metrics = doc.get('metrics', {})
                html_content += f"""
                    <tr>
                        <td>{doc['title'][:50]}</td>
                        <td>{doc.get('category', 'General')}</td>
                        <td>{metrics.get('total_lines', 0):,}</td>
                        <td>{metrics.get('total_words', 0):,}</td>
                        <td>{doc.get('sections_count', 0)}</td>
                        <td>{doc.get('quality_score', 0)}/100</td>
                    </tr>
"""
            
            html_content += """
                </tbody>
            </table>
        </div>
    </div>
    
    <script>
        // Gráfica de Calidad
        const qualityCtx = document.getElementById('qualityChart').getContext('2d');
        new Chart(qualityCtx, {
            type: 'bar',
            data: {
                labels: """ + json.dumps([d['title'][:30] for d in sorted(all_docs_data, key=lambda x: x.get('quality_score', 0), reverse=True)[:10]]) + """,
                datasets: [{
                    label: 'Calidad',
                    data: """ + json.dumps([d.get('quality_score', 0) for d in sorted(all_docs_data, key=lambda x: x.get('quality_score', 0), reverse=True)[:10]]) + """,
                    backgroundColor: 'rgba(46, 134, 171, 0.8)',
                    borderColor: 'rgba(46, 134, 171, 1)',
                    borderWidth: 1
                }]
            },
            options: {
                responsive: true,
                scales: {
                    y: {
                        beginAtZero: true,
                        max: 100
                    }
                }
            }
        });
        
        // Gráfica de Palabras
        const wordsCtx = document.getElementById('wordsChart').getContext('2d');
        new Chart(wordsCtx, {
            type: 'doughnut',
            data: {
                labels: """ + json.dumps([d['title'][:20] for d in sorted(all_docs_data, key=lambda x: x.get('metrics', {}).get('total_words', 0), reverse=True)[:8]]) + """,
                datasets: [{
                    data: """ + json.dumps([d.get('metrics', {}).get('total_words', 0) for d in sorted(all_docs_data, key=lambda x: x.get('metrics', {}).get('total_words', 0), reverse=True)[:8]]) + """,
                    backgroundColor: [
                        'rgba(46, 134, 171, 0.8)',
                        'rgba(162, 59, 114, 0.8)',
                        'rgba(241, 143, 1, 0.8)',
                        'rgba(6, 167, 125, 0.8)',
                        'rgba(199, 62, 29, 0.8)',
                        'rgba(102, 126, 234, 0.8)',
                        'rgba(118, 75, 162, 0.8)',
                        'rgba(52, 73, 94, 0.8)'
                    ]
                }]
            },
            options: {
                responsive: true
            }
        });
    </script>
</body>
</html>
"""
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            print(f"✅ Dashboard HTML interactivo guardado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando dashboard HTML: {e}")
            import traceback
            traceback.print_exc()
    
    def create_trends_analysis(self, all_docs_data: List[Dict], output_path: Path):
        """Crea análisis de tendencias temporales"""
        try:
            if not EXCEL_AVAILABLE:
                return
            
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.chart import LineChart, Reference
            
            wb = Workbook()
            ws = wb.active
            ws.title = "Tendencias"
            
            header_font = Font(bold=True, size=12, color="FFFFFF")
            header_fill = PatternFill(start_color="F18F01", end_color="F18F01", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Analizar tendencias por categoría
            category_stats = defaultdict(lambda: {'count': 0, 'total_quality': 0, 'total_words': 0})
            for doc in all_docs_data:
                category = doc.get('category', 'General')
                category_stats[category]['count'] += 1
                category_stats[category]['total_quality'] += doc.get('quality_score', 0)
                category_stats[category]['total_words'] += doc.get('metrics', {}).get('total_words', 0)
            
            ws['A1'] = "Categoría"
            ws['B1'] = "Documentos"
            ws['C1'] = "Calidad Promedio"
            ws['D1'] = "Palabras Totales"
            ws['E1'] = "Palabras Promedio"
            
            for col in range(1, 6):
                cell = ws.cell(row=1, column=col)
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
            
            row = 2
            for category, stats in sorted(category_stats.items()):
                avg_quality = stats['total_quality'] / stats['count'] if stats['count'] > 0 else 0
                avg_words = stats['total_words'] / stats['count'] if stats['count'] > 0 else 0
                
                ws.cell(row=row, column=1).value = category
                ws.cell(row=row, column=2).value = stats['count']
                ws.cell(row=row, column=3).value = round(avg_quality, 1)
                ws.cell(row=row, column=4).value = stats['total_words']
                ws.cell(row=row, column=5).value = round(avg_words, 0)
                
                for col in range(1, 6):
                    cell = ws.cell(row=row, column=col)
                    cell.border = border
                    if row % 2 == 0:
                        cell.fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                
                row += 1
            
            ws.column_dimensions['A'].width = 25
            ws.column_dimensions['B'].width = 15
            ws.column_dimensions['C'].width = 18
            ws.column_dimensions['D'].width = 18
            ws.column_dimensions['E'].width = 18
            
            wb.save(str(output_path))
            print(f"✅ Análisis de tendencias guardado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando análisis de tendencias: {e}")
    
    def create_quality_report(self, all_docs_data: List[Dict], output_path: Path):
        """Crea reporte completo de calidad"""
        try:
            if not EXCEL_AVAILABLE:
                return
            
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            
            wb = Workbook()
            
            # Hoja 1: Resumen de Calidad
            ws1 = wb.active
            ws1.title = "Resumen Calidad"
            
            header_font = Font(bold=True, size=12, color="FFFFFF")
            header_fill = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            headers = ['Documento', 'Calidad', 'Legibilidad', 'Coherencia', 'Sentimiento', 'Complejidad']
            for col, header in enumerate(headers, 1):
                cell = ws1.cell(row=1, column=col)
                cell.value = header
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
            
            for row_idx, doc in enumerate(all_docs_data, 2):
                metrics = doc.get('metrics', {})
                ws1.cell(row=row_idx, column=1).value = doc['title'][:50]
                ws1.cell(row=row_idx, column=2).value = doc.get('quality_score', 0)
                ws1.cell(row=row_idx, column=3).value = metrics.get('flesch_score', 0)
                ws1.cell(row=row_idx, column=4).value = metrics.get('coherence_score', 0)
                ws1.cell(row=row_idx, column=5).value = metrics.get('sentiment_tone', 'N/A')
                ws1.cell(row=row_idx, column=6).value = metrics.get('complexity_level', 'N/A')
                
                # Color según calidad
                quality = doc.get('quality_score', 0)
                if quality >= 80:
                    fill_color = "C6E0B4"  # Verde
                elif quality >= 60:
                    fill_color = "FFE699"  # Amarillo
                elif quality >= 40:
                    fill_color = "F8CBAD"  # Naranja
                else:
                    fill_color = "F4B084"  # Rojo
                
                for col in range(1, 7):
                    cell = ws1.cell(row=row_idx, column=col)
                    cell.border = border
                    if col == 2:  # Columna de calidad
                        cell.fill = PatternFill(start_color=fill_color, end_color=fill_color, fill_type="solid")
            
            ws1.column_dimensions['A'].width = 50
            for col in range(2, 7):
                ws1.column_dimensions[get_column_letter(col)].width = 18
            
            # Hoja 2: Distribución de Calidad
            ws2 = wb.create_sheet("Distribución")
            
            quality_ranges = {
                'Excelente (80-100)': 0,
                'Buena (60-79)': 0,
                'Regular (40-59)': 0,
                'Necesita Mejora (<40)': 0
            }
            
            for doc in all_docs_data:
                quality = doc.get('quality_score', 0)
                if quality >= 80:
                    quality_ranges['Excelente (80-100)'] += 1
                elif quality >= 60:
                    quality_ranges['Buena (60-79)'] += 1
                elif quality >= 40:
                    quality_ranges['Regular (40-59)'] += 1
                else:
                    quality_ranges['Necesita Mejora (<40)'] += 1
            
            ws2['A1'] = "Rango de Calidad"
            ws2['B1'] = "Cantidad"
            ws2['C1'] = "Porcentaje"
            
            for col in range(1, 4):
                cell = ws2.cell(row=1, column=col)
                cell.font = header_font
                cell.fill = PatternFill(start_color="A23B72", end_color="A23B72", fill_type="solid")
                cell.border = border
            
            total = len(all_docs_data)
            row = 2
            for range_name, count in quality_ranges.items():
                percentage = (count / total * 100) if total > 0 else 0
                ws2.cell(row=row, column=1).value = range_name
                ws2.cell(row=row, column=2).value = count
                ws2.cell(row=row, column=3).value = f"{percentage:.1f}%"
                
                for col in range(1, 4):
                    cell = ws2.cell(row=row, column=col)
                    cell.border = border
                
                row += 1
            
            ws2.column_dimensions['A'].width = 30
            ws2.column_dimensions['B'].width = 15
            ws2.column_dimensions['C'].width = 15
            
            wb.save(str(output_path))
            print(f"✅ Reporte de calidad guardado: {output_path}")
        except Exception as e:
            print(f"⚠️  Error creando reporte de calidad: {e}")
    
    def export_to_markdown_enhanced(self, all_docs_data: List[Dict], output_path: Path):
        """Exporta todos los documentos a un archivo Markdown mejorado"""
        try:
            md_content = f"""# 📊 Análisis Consolidado de Documentos

**Generado el:** {datetime.now().strftime('%d de %B de %Y %H:%M')}  
**Total de Documentos:** {len(all_docs_data)}

---

## 📈 Resumen Ejecutivo

"""
            
            # Estadísticas generales
            total_lines = sum(d.get('metrics', {}).get('total_lines', 0) for d in all_docs_data)
            total_words = sum(d.get('metrics', {}).get('total_words', 0) for d in all_docs_data)
            total_sections = sum(d.get('sections_count', 0) for d in all_docs_data)
            avg_quality = sum(d.get('quality_score', 0) for d in all_docs_data) / len(all_docs_data) if all_docs_data else 0
            
            md_content += f"""
### Estadísticas Generales

- **Total de Líneas:** {total_lines:,}
- **Total de Palabras:** {total_words:,}
- **Total de Secciones:** {total_sections}
- **Calidad Promedio:** {avg_quality:.1f}/100

---

## 📋 Documentos Analizados

"""
            
            # Lista de documentos
            for i, doc in enumerate(all_docs_data, 1):
                metrics = doc.get('metrics', {})
                md_content += f"""
### {i}. {doc['title']}

- **Categoría:** {doc.get('category', 'General')}
- **Líneas:** {metrics.get('total_lines', 0):,}
- **Palabras:** {metrics.get('total_words', 0):,}
- **Secciones:** {doc.get('sections_count', 0)}
- **Calidad:** {doc.get('quality_score', 0)}/100
- **Legibilidad:** {metrics.get('reading_level', 'N/A')}
- **Complejidad:** {metrics.get('complexity_level', 'N/A')}
- **Tiempo de Lectura:** ~{metrics.get('reading_time_minutes', 0)} minutos

"""
            
            # Top documentos
            md_content += """
---

## 🏆 Top Documentos

### Por Calidad

"""
            sorted_by_quality = sorted(all_docs_data, key=lambda x: x.get('quality_score', 0), reverse=True)[:5]
            for i, doc in enumerate(sorted_by_quality, 1):
                md_content += f"{i}. **{doc['title']}** - {doc.get('quality_score', 0)}/100\n"
            
            md_content += """
### Por Tamaño (Palabras)

"""
            sorted_by_words = sorted(all_docs_data, key=lambda x: x.get('metrics', {}).get('total_words', 0), reverse=True)[:5]
            for i, doc in enumerate(sorted_by_words, 1):
                md_content += f"{i}. **{doc['title']}** - {doc.get('metrics', {}).get('total_words', 0):,} palabras\n"
            
            md_content += f"""

---

## 📊 Distribución por Categoría

"""
            category_count = defaultdict(int)
            for doc in all_docs_data:
                category_count[doc.get('category', 'General')] += 1
            
            for category, count in sorted(category_count.items(), key=lambda x: x[1], reverse=True):
                md_content += f"- **{category}:** {count} documento(s)\n"
            
            md_content += f"""

---

*Este reporte fue generado automáticamente por el sistema de análisis de documentos.*
"""
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            print(f"✅ Exportación Markdown guardada: {output_path}")
        except Exception as e:
            print(f"⚠️  Error exportando a Markdown: {e}")
            import traceback
            traceback.print_exc()
    
    def convert_document(self, doc_info: Dict[str, str]):
        """Convierte un documento a los tres formatos"""
        doc_path = self.base_path / doc_info['path']
        
        if not doc_path.exists():
            print(f"⚠️  Documento no encontrado: {doc_path}")
            return None
        
        print(f"\n📄 Procesando: {doc_info['title']}")
        print(f"   Ruta: {doc_info['path']}")
        
        # Leer contenido
        content = self.read_markdown(doc_path)
        if not content:
            return None
        
        # Parsear
        try:
            doc_data, html_content = self.parse_markdown(content)
        except Exception as e:
            print(f"⚠️  Error parseando markdown: {e}")
            return None
        
        # Crear gráficas
        safe_title = re.sub(r'[^\w\s-]', '', doc_info['title']).strip()
        graphs = self.create_summary_graphs(doc_data, safe_title)
        
        # Crear archivos de salida
        base_name = safe_title.replace(' ', '_').replace('/', '_')
        
        # PDF
        pdf_path = OUTPUT_DIR / f"{base_name}.pdf"
        self.create_pdf(content, doc_data, doc_info['title'], pdf_path, graphs)
        
        # Word
        word_path = OUTPUT_DIR / f"{base_name}.docx"
        self.create_word(content, doc_data, doc_info['title'], word_path, graphs)
        
        # Excel
        excel_path = OUTPUT_DIR / f"{base_name}.xlsx"
        self.create_excel(doc_data, doc_info['title'], excel_path, graphs)
        
        # Exportar a JSON
        json_path = OUTPUT_DIR / f"{base_name}.json"
        self.export_to_json(doc_data, doc_info['title'], json_path)
        
        # Exportar a HTML
        html_path = OUTPUT_DIR / f"{base_name}.html"
        self.export_to_html(content, doc_data, doc_info['title'], html_path, graphs)
        
        print(f"✅ Completado: {doc_info['title']}\n")
        
        # Devolver datos para dashboard consolidado
        return {
            'title': doc_info['title'],
            'category': doc_info.get('category', 'General'),
            'path': doc_info['path'],
            'content': content[:5000],  # Primeros 5000 caracteres para análisis
            'metrics': doc_data.get('metrics', {}),
            'sections_count': len(doc_data.get('sections', [])),
            'graphs_count': len(graphs),
            'quality_score': doc_data.get('metrics', {}).get('quality_score', 0),
            'readability_score': doc_data.get('metrics', {}).get('flesch_score', 0),
            'top_keywords': [kw['word'] for kw in doc_data.get('top_keywords', [])[:10]],
            'sections': [s['title'] for s in doc_data.get('sections', [])[:20]]
        }


def main():
    """Función principal"""
    print("=" * 70)
    print("🚀 Generador de Documentos Profesionales")
    print("   PDF, Word y Excel con gráficas de alta calidad")
    print("=" * 70)
    
    # Verificar dependencias
    missing = []
    if not PDF_AVAILABLE:
        missing.append("reportlab")
    if not WORD_AVAILABLE:
        missing.append("python-docx")
    if not EXCEL_AVAILABLE:
        missing.append("openpyxl")
    if not MATPLOTLIB_AVAILABLE:
        missing.append("matplotlib numpy")
    
    if missing:
        print(f"\n⚠️  Faltan dependencias: {', '.join(missing)}")
        print("   Instala con: pip install " + " ".join(missing))
        print("\n   Continuando con las librerías disponibles...\n")
    
    # Base path
    base_path = Path(__file__).parent
    
    # Crear convertidor
    converter = DocumentConverter(base_path)
    
    # Convertir documentos y recopilar datos
    print(f"\n📚 Documentos a procesar: {len(IMPORTANT_DOCS)}\n")
    
    all_docs_data = []
    for doc_info in IMPORTANT_DOCS:
        doc_result = converter.convert_document(doc_info)
        if doc_result:
            all_docs_data.append(doc_result)
    
    # Crear dashboard consolidado
    if all_docs_data and EXCEL_AVAILABLE:
        try:
            wb = Workbook()
            ws = wb.active
            ws.title = "Dashboard Consolidado"
            
            # Estilos
            title_fill = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
            header_fill = PatternFill(start_color="A23B72", end_color="A23B72", fill_type="solid")
            title_font = Font(bold=True, size=16, color="FFFFFF")
            header_font = Font(bold=True, size=12, color="FFFFFF")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Título
            ws['A1'] = "Dashboard Consolidado - Todos los Documentos"
            ws['A1'].font = title_font
            ws['A1'].fill = title_fill
            ws.merge_cells('A1:I1')
            ws['A1'].alignment = Alignment(horizontal='center', vertical='center')
            
            # Encabezados
            headers = ['Documento', 'Categoría', 'Líneas', 'Palabras', 'Secciones', 'Calidad', 'Legibilidad', 'Gráficas', 'Tiempo Lectura (min)']
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=2, column=col)
                cell.value = header
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
            
            # Datos
            for row_idx, doc_data_item in enumerate(all_docs_data, 3):
                metrics = doc_data_item.get('metrics', {})
                ws.cell(row=row_idx, column=1).value = doc_data_item['title'][:50]
                ws.cell(row=row_idx, column=2).value = doc_data_item.get('category', 'General')
                ws.cell(row=row_idx, column=3).value = metrics.get('total_lines', 0)
                ws.cell(row=row_idx, column=4).value = metrics.get('total_words', 0)
                ws.cell(row=row_idx, column=5).value = doc_data_item.get('sections_count', 0)
                ws.cell(row=row_idx, column=6).value = f"{doc_data_item.get('quality_score', 0)}/100"
                ws.cell(row=row_idx, column=7).value = metrics.get('reading_level', 'N/A')
                ws.cell(row=row_idx, column=8).value = doc_data_item.get('graphs_count', 0)
                ws.cell(row=row_idx, column=9).value = metrics.get('reading_time_minutes', 0)
                
                # Aplicar bordes
                for col in range(1, 10):
                    cell = ws.cell(row=row_idx, column=col)
                    cell.border = border
                    if row_idx % 2 == 0:
                        cell.fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
            
            # Ajustar anchos
            ws.column_dimensions['A'].width = 50
            ws.column_dimensions['B'].width = 20
            for col in range(3, 10):
                ws.column_dimensions[get_column_letter(col)].width = 15
            
            # Guardar dashboard
            dashboard_path = OUTPUT_DIR / "DASHBOARD_CONSOLIDADO.xlsx"
            wb.save(str(dashboard_path))
            print(f"\n📊 Dashboard consolidado creado: {dashboard_path}")
            
            # Análisis avanzado: Similitud y comparación entre documentos
            if len(all_docs_data) > 1:
                print("\n" + "=" * 70)
                print("🔍 Análisis Avanzado: Similitud y Comparación")
                print("=" * 70)
                
                # Análisis de similitud
                similarity_matrix = converter.analyze_document_similarity(all_docs_data)
                similarity_path = OUTPUT_DIR / "ANALISIS_SIMILITUD.xlsx"
                converter.create_similarity_report(similarity_matrix, similarity_path)
                print(f"✅ Análisis de similitud creado: {similarity_path}")
                
                # Análisis de dependencias
                dependencies = converter.analyze_document_dependencies(all_docs_data)
                deps_path = OUTPUT_DIR / "ANALISIS_DEPENDENCIAS.xlsx"
                converter.create_dependencies_report(dependencies, deps_path)
                print(f"✅ Análisis de dependencias creado: {deps_path}")
                
                # Reporte comparativo
                comparison_path = OUTPUT_DIR / "REPORTE_COMPARATIVO.xlsx"
                converter.create_comparison_report(all_docs_data, comparison_path)
                print(f"✅ Reporte comparativo creado: {comparison_path}")
                
                # Análisis avanzado (duplicados, antigüedad)
                advanced_path = OUTPUT_DIR / "ANALISIS_AVANZADO.xlsx"
                converter.create_advanced_analytics(all_docs_data, advanced_path)
                print(f"✅ Análisis avanzado creado: {advanced_path}")
                
                # Reporte ejecutivo en PDF
                exec_path = OUTPUT_DIR / "REPORTE_EJECUTIVO.pdf"
                converter.create_executive_report(all_docs_data, exec_path)
                print(f"✅ Reporte ejecutivo creado: {exec_path}")
                
                # Análisis de red de conceptos
                concepts_path = OUTPUT_DIR / "RED_CONCEPTOS.xlsx"
                converter.create_concept_network(all_docs_data, concepts_path)
                print(f"✅ Red de conceptos creada: {concepts_path}")
                
                # Exportación a PowerPoint
                pptx_path = OUTPUT_DIR / "PRESENTACION_CONSOLIDADA.pptx"
                converter.create_powerpoint_presentation(all_docs_data, pptx_path)
                print(f"✅ Presentación PowerPoint creada: {pptx_path}")
                
                # Dashboard HTML interactivo consolidado
                dashboard_html_path = OUTPUT_DIR / "DASHBOARD_INTERACTIVO.html"
                converter.create_interactive_dashboard(all_docs_data, dashboard_html_path)
                print(f"✅ Dashboard HTML interactivo creado: {dashboard_html_path}")
                
                # Análisis de tendencias temporales
                trends_path = OUTPUT_DIR / "ANALISIS_TENDENCIAS.xlsx"
                converter.create_trends_analysis(all_docs_data, trends_path)
                print(f"✅ Análisis de tendencias creado: {trends_path}")
                
                # Reporte de calidad mejorado
                quality_report_path = OUTPUT_DIR / "REPORTE_CALIDAD_COMPLETO.xlsx"
                converter.create_quality_report(all_docs_data, quality_report_path)
                print(f"✅ Reporte de calidad completo creado: {quality_report_path}")
                
                # Exportación a Markdown mejorado
                markdown_export_path = OUTPUT_DIR / "EXPORTACION_MARKDOWN.md"
                converter.export_to_markdown_enhanced(all_docs_data, markdown_export_path)
                print(f"✅ Exportación Markdown mejorada creada: {markdown_export_path}")
            
        except Exception as e:
            print(f"⚠️  Error creando dashboard: {e}")
            import traceback
            traceback.print_exc()
    
    print("\n" + "=" * 70)
    print("✅ Proceso completado!")
    print(f"📁 Archivos guardados en: {OUTPUT_DIR.absolute()}")
    print(f"📊 Total de documentos procesados: {len(all_docs_data)}")
    print("=" * 70)


if __name__ == "__main__":
    main()
