#!/usr/bin/env python3
"""
Generador Premium V2 de Documentos Importantes
==============================================
Versión mejorada con más funcionalidades:
- Más documentos procesados
- Más tipos de gráficas
- Análisis de texto avanzado
- Tablas de contenido automáticas
- Índices y referencias cruzadas
- Exportación mejorada
"""

import argparse
import os
import sys
import shutil
import zipfile
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Tuple
import re
import json
from collections import Counter
import math
import random
import statistics

# Importar librerías de conversión
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
except ImportError:
    print("Instalando python-docx...")
    os.system("pip install python-docx")
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

try:
    import openpyxl
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference, ScatterChart
    from openpyxl.drawing.image import Image
except ImportError:
    print("Instalando openpyxl...")
    os.system("pip install openpyxl")
    import openpyxl
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image as RLImage, KeepTogether
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
    from reportlab.pdfgen import canvas
except ImportError:
    print("Instalando reportlab...")
    os.system("pip install reportlab")
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image as RLImage
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from matplotlib.backends.backend_pdf import PdfPages
    import seaborn as sns
    sns.set_style("whitegrid")
    sns.set_palette("husl")
    from matplotlib.patches import Rectangle
    import numpy as np
except ImportError:
    print("Instalando matplotlib y seaborn...")
    os.system("pip install matplotlib seaborn numpy")
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import seaborn as sns
    import numpy as np

WEASYPRINT_AVAILABLE = False
try:
    from weasyprint import HTML, CSS
    WEASYPRINT_AVAILABLE = True
except Exception:
    pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Instalando python-pptx...")
    os.system("pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

try:
    import plotly.graph_objects as go
    import plotly.express as px
    PLOTLY_AVAILABLE = True
except ImportError:
    print("Instalando plotly...")
    os.system("pip install plotly")
    try:
        import plotly.graph_objects as go
        import plotly.express as px
        PLOTLY_AVAILABLE = True
    except ImportError:
        PLOTLY_AVAILABLE = False

DEFAULT_IMPORTANT_DOCS = [
    "airflow_automation_prompt.md",
    "truthgpt_collected/integration_code/production_code/ARCHITECTURE_IMPROVEMENTS.md",
    "truthgpt_collected/integration_code/production_code/REFACTORING_PLAN.md",
    "BEST_PRACTICES.md",
    "CHANGELOG.md",
    "CONTRIBUTING.md",
    "README.md",
    "ARCHITECTURE.md",
]


class AdvancedDocumentAnalyzer:
    """Analizador avanzado de documentos con NLP básico"""
    
    @staticmethod
    def analyze_text_complexity(text: str) -> Dict:
        """Analiza la complejidad del texto"""
        sentences = re.split(r'[.!?]+', text)
        words = text.split()
        
        # Calcular estadísticas
        avg_sentence_length = len(words) / len(sentences) if sentences else 0
        avg_word_length = sum(len(w) for w in words) / len(words) if words else 0
        
        # Contar palabras técnicas (mayúsculas, números, guiones)
        technical_words = sum(1 for w in words if any(c.isupper() for c in w) or any(c.isdigit() for c in w) or '-' in w)
        technical_ratio = technical_words / len(words) if words else 0
        
        # Contar palabras largas (>10 caracteres)
        long_words = sum(1 for w in words if len(w) > 10)
        long_words_ratio = long_words / len(words) if words else 0
        
        return {
            'avg_sentence_length': avg_sentence_length,
            'avg_word_length': avg_word_length,
            'technical_ratio': technical_ratio,
            'long_words_ratio': long_words_ratio,
            'complexity_score': (avg_sentence_length * 0.3 + avg_word_length * 0.2 + technical_ratio * 0.3 + long_words_ratio * 0.2) * 10
        }
    
    @staticmethod
    def extract_keywords(text: str, top_n: int = 20) -> List[Tuple[str, int]]:
        """Extrae palabras clave más frecuentes"""
        # Limpiar y tokenizar
        words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
        
        # Filtrar stop words comunes
        stop_words = {'this', 'that', 'with', 'from', 'have', 'will', 'would', 'could', 'should',
                     'been', 'being', 'were', 'what', 'when', 'where', 'which', 'while', 'them',
                     'they', 'their', 'there', 'these', 'those', 'then', 'than', 'more', 'most'}
        
        words = [w for w in words if w not in stop_words and len(w) > 3]
        
        # Contar frecuencia
        word_freq = Counter(words)
        return word_freq.most_common(top_n)
    
    @staticmethod
    def analyze_code_blocks(content: str) -> Dict:
        """Analiza bloques de código"""
        code_blocks = re.findall(r'```[\s\S]*?```', content)
        
        languages = []
        total_lines = 0
        
        for block in code_blocks:
            lines = block.split('\n')
            total_lines += len(lines) - 2  # Excluir delimitadores
            
            # Intentar detectar lenguaje
            first_line = lines[0].strip('`').strip()
            if first_line and not first_line.startswith('```'):
                lang = first_line
            else:
                lang = 'unknown'
            languages.append(lang)
        
        lang_counter = Counter(languages)
        
        return {
            'total_blocks': len(code_blocks),
            'total_code_lines': total_lines,
            'languages': dict(lang_counter),
            'avg_block_size': total_lines / len(code_blocks) if code_blocks else 0
        }

    @staticmethod
    def extract_summary_points(content: str, max_points: int = 5) -> List[str]:
        """Genera puntos clave a partir de la estructura del documento"""
        points: List[str] = []
        sections = re.findall(r'(^#+\s+.+$)', content, re.MULTILINE)
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]

        # Preferir la primera oración de los primeros párrafos informativos
        for paragraph in paragraphs:
            if paragraph.startswith('#') or len(paragraph.split()) < 8:
                continue
            sentence = re.split(r'(?<=[.!?])\s+', paragraph.strip())[0]
            sentence = sentence.replace('\n', ' ').strip()
            if len(sentence.split()) >= 8:
                points.append(sentence)
            if len(points) >= max_points:
                break

        # Si no hay suficientes puntos, usar títulos principales
        if len(points) < max_points and sections:
            for section in sections[:max_points]:
                clean = section.lstrip('#').strip()
                if clean and clean not in points:
                    points.append(clean)
                if len(points) >= max_points:
                    break

        return points[:max_points]

    @staticmethod
    def detect_risks(content: str) -> Dict[str, List[str]]:
        """Detecta riesgos potenciales en el documento"""
        risk_keywords = {
            'Pendientes': [r'\bTODO\b', r'\bPENDING\b', r'\bTBD\b'],
            'Advertencias': [r'\bWARNING\b', r'\bRISK\b', r'\bISSUE\b'],
            'Notas críticas': [r'\bFIXME\b', r'\bCRITICAL\b', r'\bBLOCKER\b']
        }
        findings: Dict[str, List[str]] = {k: [] for k in risk_keywords.keys()}

        lines = content.split('\n')
        for idx, line in enumerate(lines, start=1):
            for category, patterns in risk_keywords.items():
                for pattern in patterns:
                    if re.search(pattern, line, re.IGNORECASE):
                        snippet = line.strip()
                        if len(snippet) > 140:
                            snippet = snippet[:137] + '...'
                        findings[category].append(f"L{idx}: {snippet}")

        return {k: v for k, v in findings.items() if v}


class DocumentGeneratorPremiumV2:
    """Generador premium mejorado de documentos con gráficas"""
    
    def __init__(
        self,
        output_dir: str = "documentos_importantes_premium_v2",
        formats: Optional[List[str]] = None,
    ):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.graphs_dir = self.output_dir / "graficas"
        self.graphs_dir.mkdir(exist_ok=True)
        self.analyzer = AdvancedDocumentAnalyzer()
        self.formats = [fmt.lower() for fmt in (formats or ["pdf", "word", "excel", "pptx"])]
        self.summary_path = self.output_dir / "summary_report.pdf"
        
        # Configurar matplotlib
        plt.rcParams['figure.figsize'] = (12, 8)
        plt.rcParams['font.size'] = 10
        plt.rcParams['axes.labelsize'] = 12
        plt.rcParams['axes.titlesize'] = 14
        plt.rcParams['xtick.labelsize'] = 10
        plt.rcParams['ytick.labelsize'] = 10
        plt.rcParams['legend.fontsize'] = 10
        plt.rcParams['figure.dpi'] = 300
        
    def analyze_document(self, file_path: Path) -> Dict:
        """Analiza un documento y extrae estadísticas avanzadas"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Estadísticas básicas
        lines = content.split('\n')
        words = content.split()
        paragraphs = [p for p in content.split('\n\n') if p.strip()]
        
        # Contar secciones
        sections = len(re.findall(r'^#+\s+', content, re.MULTILINE))
        code_blocks = len(re.findall(r'```', content)) // 2
        tables = len(re.findall(r'\|.*\|', content))
        links = len(re.findall(r'\[.*?\]\(.*?\)', content))
        images = len(re.findall(r'!\[.*?\]\(.*?\)', content))
        
        # Extraer estructura
        structure = []
        for line in lines:
            if line.strip().startswith('#'):
                level = len(line) - len(line.lstrip('#'))
                title = line.lstrip('#').strip()
                structure.append({'level': level, 'title': title})
        
        # Análisis avanzado
        text_complexity = self.analyzer.analyze_text_complexity(content)
        keywords = self.analyzer.extract_keywords(content)
        code_analysis = self.analyzer.analyze_code_blocks(content)
        summary_points = self.analyzer.extract_summary_points(content)
        risks = self.analyzer.detect_risks(content)
        
        return {
            'file_path': str(file_path),
            'name': file_path.stem,
            'total_lines': len(lines),
            'total_words': len(words),
            'total_paragraphs': len(paragraphs),
            'sections': sections,
            'code_blocks': code_blocks,
            'tables': tables,
            'links': links,
            'images': images,
            'structure': structure[:50],  # Primeras 50 secciones
            'content': content,
            'complexity': text_complexity,
            'keywords': keywords,
            'code_analysis': code_analysis,
            'summary_points': summary_points,
            'risks': risks
        }
    
    def create_advanced_statistics_graphs(self, doc_stats: Dict) -> List[str]:
        """Crea gráficas avanzadas de estadísticas"""
        graph_files = []
        
        # Gráfica 1: Distribución de contenido mejorada
        fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))
        fig.suptitle(f'Análisis Completo: {doc_stats["name"]}', fontsize=16, fontweight='bold')
        
        # Subgráfica 1: Métricas principales
        categories = ['Líneas', 'Palabras', 'Párrafos', 'Secciones', 'Código', 'Tablas', 'Enlaces', 'Imágenes']
        values = [
            doc_stats['total_lines'],
            doc_stats['total_words'],
            doc_stats['total_paragraphs'],
            doc_stats['sections'],
            doc_stats['code_blocks'],
            doc_stats['tables'],
            doc_stats['links'],
            doc_stats['images']
        ]
        
        max_val = max(values) if values else 1
        normalized = [v / max_val * 100 if max_val > 0 else 0 for v in values]
        
        bars = ax1.barh(categories, normalized, color=sns.color_palette("husl", len(categories)))
        ax1.set_xlabel('Valor Normalizado (%)', fontsize=11, fontweight='bold')
        ax1.set_title('Métricas de Contenido', fontsize=12, fontweight='bold')
        ax1.grid(axis='x', alpha=0.3)
        
        for i, (bar, val) in enumerate(zip(bars, values)):
            width = bar.get_width()
            ax1.text(width + 1, bar.get_y() + bar.get_height()/2, 
                   f'{val:,}', ha='left', va='center', fontweight='bold', fontsize=9)
        
        # Subgráfica 2: Palabras clave
        if doc_stats['keywords']:
            top_keywords = doc_stats['keywords'][:10]
            words_list = [w[0] for w in top_keywords]
            counts = [w[1] for w in top_keywords]
            
            bars2 = ax2.barh(words_list, counts, color=sns.color_palette("coolwarm", len(words_list)))
            ax2.set_xlabel('Frecuencia', fontsize=11, fontweight='bold')
            ax2.set_title('Top 10 Palabras Clave', fontsize=12, fontweight='bold')
            ax2.grid(axis='x', alpha=0.3)
            
            for bar, count in zip(bars2, counts):
                width = bar.get_width()
                ax2.text(width + max(counts) * 0.01, bar.get_y() + bar.get_height()/2,
                        f'{count}', ha='left', va='center', fontweight='bold', fontsize=9)
        
        # Subgráfica 3: Análisis de complejidad
        complexity = doc_stats['complexity']
        metrics = ['Long. Promedio\nOración', 'Long. Promedio\nPalabra', 'Ratio\nTécnico', 'Ratio\nPalabras Largas']
        complexity_values = [
            complexity['avg_sentence_length'],
            complexity['avg_word_length'],
            complexity['technical_ratio'] * 100,
            complexity['long_words_ratio'] * 100
        ]
        
        bars3 = ax3.bar(metrics, complexity_values, color=sns.color_palette("Set2", len(metrics)), alpha=0.8)
        ax3.set_ylabel('Valor', fontsize=11, fontweight='bold')
        ax3.set_title('Análisis de Complejidad del Texto', fontsize=12, fontweight='bold')
        ax3.grid(axis='y', alpha=0.3)
        ax3.tick_params(axis='x', rotation=45)
        
        for bar, val in zip(bars3, complexity_values):
            height = bar.get_height()
            ax3.text(bar.get_x() + bar.get_width()/2., height + max(complexity_values) * 0.01,
                    f'{val:.1f}', ha='center', va='bottom', fontweight='bold', fontsize=9)
        
        # Subgráfica 4: Análisis de código
        code_analysis = doc_stats['code_analysis']
        if code_analysis['total_blocks'] > 0:
            languages = list(code_analysis['languages'].keys())[:8]
            lang_counts = [code_analysis['languages'].get(lang, 0) for lang in languages]
            
            if languages:
                colors_pie = sns.color_palette("pastel", len(languages))
                wedges, texts, autotexts = ax4.pie(lang_counts, labels=languages, autopct='%1.1f%%',
                                                  colors=colors_pie, startangle=90)
                for autotext in autotexts:
                    autotext.set_color('black')
                    autotext.set_fontweight('bold')
                ax4.set_title(f'Distribución de Lenguajes de Código\n({code_analysis["total_blocks"]} bloques)', 
                            fontsize=12, fontweight='bold')
        else:
            ax4.text(0.5, 0.5, 'Sin bloques\nde código', ha='center', va='center',
                    fontsize=14, transform=ax4.transAxes)
            ax4.set_title('Análisis de Código', fontsize=12, fontweight='bold')
        
        plt.tight_layout()
        graph_file = self.graphs_dir / f"{doc_stats['name']}_analisis_completo.png"
        plt.savefig(graph_file, dpi=300, bbox_inches='tight')
        plt.close()
        graph_files.append(str(graph_file))
        
        # Gráfica 2: Estructura jerárquica mejorada
        if doc_stats['structure']:
            fig, ax = plt.subplots(figsize=(14, 10))
            
            # Contar niveles
            level_counts = {}
            for item in doc_stats['structure']:
                level = item['level']
                level_counts[level] = level_counts.get(level, 0) + 1
            
            if level_counts:
                levels = sorted(level_counts.keys())
                counts = [level_counts[l] for l in levels]
                labels = [f'Nivel {l}\n({counts[i]} secciones)' for i, l in enumerate(levels)]
                
                colors_pie = sns.color_palette("Set3", len(levels))
                wedges, texts, autotexts = ax.pie(counts, labels=labels, autopct='%1.1f%%',
                                                  colors=colors_pie, startangle=90, textprops={'fontsize': 10})
                
                for autotext in autotexts:
                    autotext.set_color('black')
                    autotext.set_fontweight('bold')
                    autotext.set_fontsize(11)
                
                ax.set_title(f'Distribución de Secciones por Nivel: {doc_stats["name"]}', 
                           fontsize=14, fontweight='bold', pad=20)
                
                plt.tight_layout()
                graph_file = self.graphs_dir / f"{doc_stats['name']}_estructura_mejorada.png"
                plt.savefig(graph_file, dpi=300, bbox_inches='tight')
                plt.close()
                graph_files.append(str(graph_file))
        
        # Gráfica 3: Evolución de secciones (si hay suficientes)
        if len(doc_stats['structure']) > 10:
            fig, ax = plt.subplots(figsize=(14, 6))
            
            # Crear histograma de distribución de secciones
            section_positions = list(range(len(doc_stats['structure'])))
            section_levels = [item['level'] for item in doc_stats['structure']]
            
            ax.plot(section_positions, section_levels, marker='o', markersize=4, linewidth=1.5, alpha=0.7)
            ax.fill_between(section_positions, section_levels, alpha=0.3)
            ax.set_xlabel('Posición en el Documento', fontsize=12, fontweight='bold')
            ax.set_ylabel('Nivel de Sección', fontsize=12, fontweight='bold')
            ax.set_title(f'Evolución de la Estructura: {doc_stats["name"]}', fontsize=14, fontweight='bold')
            ax.grid(True, alpha=0.3)
            ax.set_yticks(range(1, max(section_levels) + 2))
            
            plt.tight_layout()
            graph_file = self.graphs_dir / f"{doc_stats['name']}_evolucion_estructura.png"
            plt.savefig(graph_file, dpi=300, bbox_inches='tight')
            plt.close()
            graph_files.append(str(graph_file))
        
        return graph_files
    
    def create_enhanced_comparison_graph(self, all_stats: List[Dict]) -> str:
        """Crea gráfica comparativa mejorada"""
        if len(all_stats) < 2:
            return None
        
        fig = plt.figure(figsize=(18, 12))
        gs = fig.add_gridspec(3, 3, hspace=0.3, wspace=0.3)
        fig.suptitle('Análisis Comparativo de Documentos Importantes', fontsize=18, fontweight='bold', y=0.98)
        
        names = [s['name'] for s in all_stats]
        x = range(len(names))
        width = 0.35
        
        # Gráfica 1: Volumen de contenido
        ax1 = fig.add_subplot(gs[0, 0])
        lines = [s['total_lines'] for s in all_stats]
        words = [s['total_words'] for s in all_stats]
        max_lines = max(lines) if lines else 1
        max_words = max(words) if words else 1
        lines_norm = [l / max_lines * 100 for l in lines]
        words_norm = [w / max_words * 100 for w in words]
        
        ax1.bar([i - width/2 for i in x], lines_norm, width, label='Líneas', alpha=0.8, color='#3498db')
        ax1.bar([i + width/2 for i in x], words_norm, width, label='Palabras', alpha=0.8, color='#e74c3c')
        ax1.set_xlabel('Documentos', fontweight='bold')
        ax1.set_ylabel('Valor Normalizado (%)', fontweight='bold')
        ax1.set_title('Volumen de Contenido', fontweight='bold')
        ax1.set_xticks(x)
        ax1.set_xticklabels(names, rotation=45, ha='right', fontsize=9)
        ax1.legend()
        ax1.grid(axis='y', alpha=0.3)
        
        # Gráfica 2: Estructura
        ax2 = fig.add_subplot(gs[0, 1])
        sections = [s['sections'] for s in all_stats]
        paragraphs = [s['total_paragraphs'] for s in all_stats]
        max_sec = max(sections) if sections else 1
        max_par = max(paragraphs) if paragraphs else 1
        sections_norm = [s / max_sec * 100 for s in sections]
        paragraphs_norm = [p / max_par * 100 for p in paragraphs]
        
        ax2.bar([i - width/2 for i in x], sections_norm, width, label='Secciones', alpha=0.8, color='#2ecc71')
        ax2.bar([i + width/2 for i in x], paragraphs_norm, width, label='Párrafos', alpha=0.8, color='#f39c12')
        ax2.set_xlabel('Documentos', fontweight='bold')
        ax2.set_ylabel('Valor Normalizado (%)', fontweight='bold')
        ax2.set_title('Estructura', fontweight='bold')
        ax2.set_xticks(x)
        ax2.set_xticklabels(names, rotation=45, ha='right', fontsize=9)
        ax2.legend()
        ax2.grid(axis='y', alpha=0.3)
        
        # Gráfica 3: Elementos técnicos
        ax3 = fig.add_subplot(gs[0, 2])
        code_blocks = [s['code_blocks'] for s in all_stats]
        tables = [s['tables'] for s in all_stats]
        max_code = max(code_blocks) if code_blocks else 1
        max_tab = max(tables) if tables else 1
        code_norm = [c / max_code * 100 if max_code > 0 else 0 for c in code_blocks]
        tables_norm = [t / max_tab * 100 if max_tab > 0 else 0 for t in tables]
        
        ax3.bar([i - width/2 for i in x], code_norm, width, label='Bloques Código', alpha=0.8, color='#9b59b6')
        ax3.bar([i + width/2 for i in x], tables_norm, width, label='Tablas', alpha=0.8, color='#1abc9c')
        ax3.set_xlabel('Documentos', fontweight='bold')
        ax3.set_ylabel('Valor Normalizado (%)', fontweight='bold')
        ax3.set_title('Elementos Técnicos', fontweight='bold')
        ax3.set_xticks(x)
        ax3.set_xticklabels(names, rotation=45, ha='right', fontsize=9)
        ax3.legend()
        ax3.grid(axis='y', alpha=0.3)
        
        # Gráfica 4: Complejidad
        ax4 = fig.add_subplot(gs[1, 0])
        complexity_scores = [s['complexity']['complexity_score'] for s in all_stats]
        colors_bar = sns.color_palette("RdYlGn_r", len(names))
        bars = ax4.barh(names, complexity_scores, color=colors_bar, alpha=0.8)
        ax4.set_xlabel('Puntuación de Complejidad', fontweight='bold')
        ax4.set_title('Complejidad del Texto', fontweight='bold')
        ax4.grid(axis='x', alpha=0.3)
        
        for bar, score in zip(bars, complexity_scores):
            width = bar.get_width()
            ax4.text(width + max(complexity_scores) * 0.01, bar.get_y() + bar.get_height()/2,
                    f'{score:.1f}', ha='left', va='center', fontweight='bold', fontsize=9)
        
        # Gráfica 5: Palabras clave compartidas
        ax5 = fig.add_subplot(gs[1, 1])
        all_keywords = {}
        for stats in all_stats:
            for word, count in stats['keywords'][:5]:
                all_keywords[word] = all_keywords.get(word, 0) + count
        
        if all_keywords:
            top_shared = sorted(all_keywords.items(), key=lambda x: x[1], reverse=True)[:10]
            words_list = [w[0] for w in top_shared]
            counts = [w[1] for w in top_shared]
            
            bars5 = ax5.barh(words_list, counts, color=sns.color_palette("viridis", len(words_list)))
            ax5.set_xlabel('Frecuencia Total', fontweight='bold')
            ax5.set_title('Top 10 Palabras Clave Compartidas', fontweight='bold')
            ax5.grid(axis='x', alpha=0.3)
            
            for bar, count in zip(bars5, counts):
                width = bar.get_width()
                ax5.text(width + max(counts) * 0.01, bar.get_y() + bar.get_height()/2,
                        f'{count}', ha='left', va='center', fontweight='bold', fontsize=9)
        
        # Gráfica 6: Resumen radial
        ax6 = fig.add_subplot(gs[1, 2], projection='polar')
        categories = ['Líneas', 'Palabras', 'Secciones', 'Código', 'Complejidad']
        angles = np.linspace(0, 2 * np.pi, len(categories), endpoint=False).tolist()
        angles += angles[:1]  # Cerrar el círculo
        
        for i, stats in enumerate(all_stats):
            values = [
                stats['total_lines'] / max([s['total_lines'] for s in all_stats]) * 100,
                stats['total_words'] / max([s['total_words'] for s in all_stats]) * 100,
                stats['sections'] / max([s['sections'] for s in all_stats]) * 100 if max([s['sections'] for s in all_stats]) > 0 else 0,
                stats['code_blocks'] / max([s['code_blocks'] for s in all_stats]) * 100 if max([s['code_blocks'] for s in all_stats]) > 0 else 0,
                stats['complexity']['complexity_score'] / max([s['complexity']['complexity_score'] for s in all_stats]) * 100
            ]
            values += values[:1]
            
            ax6.plot(angles, values, 'o-', linewidth=2, label=stats['name'][:20], alpha=0.7)
            ax6.fill(angles, values, alpha=0.15)
        
        ax6.set_xticks(angles[:-1])
        ax6.set_xticklabels(categories, fontsize=9)
        ax6.set_ylim(0, 100)
        ax6.set_title('Comparativa Radial', fontweight='bold', pad=20)
        ax6.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1), fontsize=8)
        ax6.grid(True)
        
        # Gráfica 7: Distribución de palabras clave por documento
        ax7 = fig.add_subplot(gs[2, :])
        keyword_matrix = []
        all_unique_keywords = set()
        for stats in all_stats:
            all_unique_keywords.update([w[0] for w in stats['keywords'][:15]])
        
        all_unique_keywords = sorted(list(all_unique_keywords))[:20]
        
        for stats in all_stats:
            keyword_dict = dict(stats['keywords'])
            row = [keyword_dict.get(kw, 0) for kw in all_unique_keywords]
            keyword_matrix.append(row)
        
        if keyword_matrix:
            im = ax7.imshow(keyword_matrix, aspect='auto', cmap='YlOrRd', interpolation='nearest')
            ax7.set_xticks(range(len(all_unique_keywords)))
            ax7.set_xticklabels(all_unique_keywords, rotation=45, ha='right', fontsize=8)
            ax7.set_yticks(range(len(names)))
            ax7.set_yticklabels(names, fontsize=9)
            ax7.set_title('Mapa de Calor: Palabras Clave por Documento', fontweight='bold', pad=15)
            plt.colorbar(im, ax=ax7, label='Frecuencia')
        
        plt.tight_layout()
        graph_file = self.graphs_dir / "comparativa_avanzada_documentos.png"
        plt.savefig(graph_file, dpi=300, bbox_inches='tight')
        plt.close()
        
        return str(graph_file)
    
    def create_pdf(self, doc_stats: Dict, graph_files: List[str]) -> str:
        """Crea un PDF profesional mejorado"""
        pdf_file = self.output_dir / f"{doc_stats['name']}_premium_v2.pdf"
        
        doc = SimpleDocTemplate(str(pdf_file), pagesize=A4,
                               rightMargin=72, leftMargin=72,
                               topMargin=72, bottomMargin=72)
        
        styles = getSampleStyleSheet()
        
        # Estilos personalizados
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a1a1a'),
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        subtitle_style = ParagraphStyle(
            'CustomSubtitle',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#2c3e50'),
            spaceAfter=12,
            spaceBefore=20,
            fontName='Helvetica-Bold'
        )
        
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontSize=11,
            textColor=colors.HexColor('#333333'),
            spaceAfter=12,
            leading=14,
            alignment=TA_JUSTIFY
        )
        
        story = []
        
        # Portada mejorada
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph(doc_stats['name'].replace('_', ' ').title(), title_style))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"Análisis Completo y Documentación", 
                              ParagraphStyle('Subtitle', parent=styles['Normal'], fontSize=14, 
                                           alignment=TA_CENTER, textColor=colors.HexColor('#7f8c8d'))))
        story.append(Spacer(1, 0.3*inch))
        story.append(Paragraph(f"Generado: {datetime.now().strftime('%d de %B de %Y, %H:%M')}", 
                              ParagraphStyle('Date', parent=styles['Normal'], fontSize=10, 
                                           alignment=TA_CENTER, textColor=colors.HexColor('#95a5a6'))))
        story.append(PageBreak())
        
        # Resumen ejecutivo mejorado
        story.append(Paragraph("Resumen Ejecutivo", subtitle_style))
        story.append(Spacer(1, 0.2*inch))
        
        summary_data = [
            ['Métrica', 'Valor', 'Análisis'],
            ['Total de Líneas', f"{doc_stats['total_lines']:,}", 
             'Alto' if doc_stats['total_lines'] > 1000 else 'Medio' if doc_stats['total_lines'] > 500 else 'Bajo'],
            ['Total de Palabras', f"{doc_stats['total_words']:,}", 
             'Extenso' if doc_stats['total_words'] > 5000 else 'Moderado' if doc_stats['total_words'] > 2000 else 'Breve'],
            ['Total de Párrafos', f"{doc_stats['total_paragraphs']:,}", 
             f"Promedio: {doc_stats['total_words']/doc_stats['total_paragraphs']:.1f} palabras/párrafo" if doc_stats['total_paragraphs'] > 0 else 'N/A'],
            ['Secciones', f"{doc_stats['sections']}", 
             f"Promedio: {doc_stats['total_words']/doc_stats['sections']:.0f} palabras/sección" if doc_stats['sections'] > 0 else 'N/A'],
            ['Bloques de Código', f"{doc_stats['code_blocks']}", 
             f"{doc_stats['code_analysis']['total_code_lines']} líneas de código" if doc_stats['code_blocks'] > 0 else 'Sin código'],
            ['Tablas', f"{doc_stats['tables']}", 'Con datos estructurados' if doc_stats['tables'] > 0 else 'Sin tablas'],
            ['Enlaces', f"{doc_stats['links']}", 'Con referencias externas' if doc_stats['links'] > 0 else 'Sin enlaces'],
            ['Complejidad', f"{doc_stats['complexity']['complexity_score']:.1f}", 
             'Alta' if doc_stats['complexity']['complexity_score'] > 50 else 'Media' if doc_stats['complexity']['complexity_score'] > 30 else 'Baja']
        ]
        
        summary_table = Table(summary_data, colWidths=[2.5*inch, 1.5*inch, 2*inch])
        summary_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#34495e')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#ecf0f1')),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#bdc3c7')),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8f9fa')])
        ]))
        story.append(summary_table)
        story.append(Spacer(1, 0.3*inch))

        # Interactive placeholder
        if PLOTLY_AVAILABLE:
            story.append(
                Paragraph(
                    "Explora visualizaciones interactivas en dashboard.html con los gráficos de Plotly.",
                    normal_style,
                )
            )

        # Puntos clave
        if doc_stats['summary_points']:
            story.append(Paragraph("Puntos Clave Detectados", subtitle_style))
            story.append(Spacer(1, 0.1*inch))
            for point in doc_stats['summary_points']:
                point_text = point.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                story.append(Paragraph(f"• {point_text}", normal_style))
            story.append(Spacer(1, 0.2*inch))

        # Riesgos
        if doc_stats['risks']:
            story.append(Paragraph("Alertas y Riesgos Detectados", subtitle_style))
            story.append(Spacer(1, 0.1*inch))
            risk_table_data = [['Categoría', 'Detalle']]
            for category, items in doc_stats['risks'].items():
                for detail in items[:5]:  # limitar cada categoria
                    safe_detail = detail.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    risk_table_data.append([category, safe_detail])
            risk_table = Table(risk_table_data, colWidths=[2*inch, 3.5*inch])
            risk_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#c0392b')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#fdecea')),
                ('TEXTCOLOR', (0, 1), (-1, -1), colors.HexColor('#8e2800')),
                ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e6b0aa')),
                ('VALIGN', (0, 0), (-1, -1), 'TOP')
            ]))
            story.append(risk_table)
            story.append(Spacer(1, 0.3*inch))
        
        # Palabras clave
        if doc_stats['keywords']:
            story.append(Paragraph("Palabras Clave Principales", subtitle_style))
            story.append(Spacer(1, 0.1*inch))
            keywords_text = ", ".join([f"{w[0]} ({w[1]})" for w in doc_stats['keywords'][:15]])
            story.append(Paragraph(keywords_text, normal_style))
            story.append(Spacer(1, 0.2*inch))
        
        # Agregar gráficas
        for graph_file in graph_files:
            if Path(graph_file).exists():
                try:
                    img = RLImage(graph_file, width=6*inch, height=4.5*inch)
                    story.append(KeepTogether([img, Spacer(1, 0.2*inch)]))
                except Exception as e:
                    print(f"  ⚠️  Error agregando gráfica {graph_file}: {e}")
        
        story.append(PageBreak())
        
        # Contenido del documento (mejorado)
        story.append(Paragraph("Contenido del Documento", subtitle_style))
        story.append(Spacer(1, 0.2*inch))
        
        content = doc_stats['content']
        lines = content.split('\n')
        max_lines = min(1500, len(lines))
        
        for i, line in enumerate(lines[:max_lines]):
            line = line.strip()
            if not line:
                if i < max_lines - 1:
                    story.append(Spacer(1, 0.1*inch))
                continue
            
            try:
                if line.startswith('#'):
                    level = len(line) - len(line.lstrip('#'))
                    text = line.lstrip('#').strip()
                    text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    if level == 1:
                        story.append(Paragraph(text, subtitle_style))
                    elif level == 2:
                        story.append(Paragraph(text, styles['Heading2']))
                    else:
                        story.append(Paragraph(text, styles['Heading3']))
                elif line.startswith('-') or line.startswith('*'):
                    text = line.lstrip('-*').strip()
                    text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    story.append(Paragraph(f"• {text}", normal_style))
                elif line.startswith('|') and '|' in line[1:]:
                    cells = [c.strip() for c in line.split('|')[1:-1]]
                    if cells and not all(c.startswith('-') for c in cells):
                        text = ' | '.join(cells)
                        text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        story.append(Paragraph(text, normal_style))
                else:
                    text = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    if text:
                        story.append(Paragraph(text, normal_style))
            except Exception as e:
                continue
        
        doc.build(story)
        return str(pdf_file)
    
    def create_word(self, doc_stats: Dict, graph_files: List[str]) -> str:
        """Crea un documento Word profesional mejorado"""
        doc_file = self.output_dir / f"{doc_stats['name']}_premium_v2.docx"
        doc = Document()
        
        # Configurar estilos
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = Pt(11)
        
        # Portada mejorada
        title = doc.add_heading(doc_stats['name'].replace('_', ' ').title(), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        subtitle = doc.add_paragraph("Análisis Completo y Documentación")
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle.runs[0].font.size = Pt(14)
        subtitle.runs[0].font.italic = True
        
        para = doc.add_paragraph(f"Generado: {datetime.now().strftime('%d de %B de %Y, %H:%M')}")
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        para.runs[0].font.size = Pt(10)
        para.runs[0].font.color.rgb = RGBColor(149, 165, 166)
        
        doc.add_page_break()
        
        # Resumen ejecutivo mejorado
        doc.add_heading('Resumen Ejecutivo', 1)
        
        # Tabla de resumen expandida
        table = doc.add_table(rows=1, cols=3)
        table.style = 'Light Grid Accent 1'
        
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Métrica'
        hdr_cells[1].text = 'Valor'
        hdr_cells[2].text = 'Análisis'
        for cell in hdr_cells:
            cell.paragraphs[0].runs[0].font.bold = True
        
        metrics = [
            ('Total de Líneas', f"{doc_stats['total_lines']:,}", 
             'Alto' if doc_stats['total_lines'] > 1000 else 'Medio' if doc_stats['total_lines'] > 500 else 'Bajo'),
            ('Total de Palabras', f"{doc_stats['total_words']:,}", 
             'Extenso' if doc_stats['total_words'] > 5000 else 'Moderado' if doc_stats['total_words'] > 2000 else 'Breve'),
            ('Total de Párrafos', f"{doc_stats['total_paragraphs']:,}", 
             f"Promedio: {doc_stats['total_words']/doc_stats['total_paragraphs']:.1f} palabras/párrafo" if doc_stats['total_paragraphs'] > 0 else 'N/A'),
            ('Secciones', f"{doc_stats['sections']}", 
             f"Promedio: {doc_stats['total_words']/doc_stats['sections']:.0f} palabras/sección" if doc_stats['sections'] > 0 else 'N/A'),
            ('Bloques de Código', f"{doc_stats['code_blocks']}", 
             f"{doc_stats['code_analysis']['total_code_lines']} líneas de código" if doc_stats['code_blocks'] > 0 else 'Sin código'),
            ('Tablas', f"{doc_stats['tables']}", 'Con datos estructurados' if doc_stats['tables'] > 0 else 'Sin tablas'),
            ('Enlaces', f"{doc_stats['links']}", 'Con referencias externas' if doc_stats['links'] > 0 else 'Sin enlaces'),
            ('Complejidad', f"{doc_stats['complexity']['complexity_score']:.1f}", 
             'Alta' if doc_stats['complexity']['complexity_score'] > 50 else 'Media' if doc_stats['complexity']['complexity_score'] > 30 else 'Baja')
        ]
        
        for metric, value, analysis in metrics:
            row_cells = table.add_row().cells
            row_cells[0].text = metric
            row_cells[1].text = value
            row_cells[2].text = analysis
        
        doc.add_paragraph()

        if doc_stats['summary_points']:
            doc.add_heading('Puntos Clave Detectados', 2)
            for point in doc_stats['summary_points']:
                doc.add_paragraph(point, style='List Bullet')
            doc.add_paragraph()

        if doc_stats['risks']:
            doc.add_heading('Alertas y Riesgos', 2)
            for category, items in doc_stats['risks'].items():
                doc.add_paragraph(category, style='List Number')
                for detail in items[:5]:
                    doc.add_paragraph(detail, style='List Bullet')
            doc.add_paragraph()
        
        # Palabras clave
        if doc_stats['keywords']:
            doc.add_heading('Palabras Clave Principales', 2)
            keywords_text = ", ".join([f"{w[0]} ({w[1]})" for w in doc_stats['keywords'][:20]])
            doc.add_paragraph(keywords_text)
            doc.add_paragraph()
        
        # Agregar gráficas
        for graph_file in graph_files:
            if Path(graph_file).exists():
                try:
                    doc.add_paragraph()
                    doc.add_picture(graph_file, width=Inches(6))
                    doc.add_paragraph()
                except Exception as e:
                    print(f"  ⚠️  Error agregando gráfica {graph_file}: {e}")
        
        doc.add_page_break()
        
        # Contenido
        doc.add_heading('Contenido del Documento', 1)
        
        content = doc_stats['content']
        lines = content.split('\n')
        max_lines = min(2500, len(lines))
        
        for i, line in enumerate(lines[:max_lines]):
            line = line.strip()
            if not line:
                if i < max_lines - 1:
                    doc.add_paragraph()
                continue
            
            try:
                if line.startswith('#'):
                    level = min(len(line) - len(line.lstrip('#')), 9)
                    text = line.lstrip('#').strip()
                    doc.add_heading(text, level)
                elif line.startswith('-') or line.startswith('*'):
                    text = line.lstrip('-*').strip()
                    if text:
                        doc.add_paragraph(text, style='List Bullet')
                elif line.startswith('|') and '|' in line[1:]:
                    cells = [c.strip() for c in line.split('|')[1:-1]]
                    if cells and not all(c.startswith('-') for c in cells):
                        text = ' | '.join(cells)
                        doc.add_paragraph(text)
                else:
                    if line:
                        doc.add_paragraph(line)
            except Exception as e:
                continue
        
        doc.save(str(doc_file))
        return str(doc_file)
    
    def create_ppt(self, doc_stats: Dict, graph_files: List[str]) -> str:
        """Genera una presentación ejecutiva en PPTX"""
        ppt_file = self.output_dir / f"{doc_stats['name']}_premium_v2.pptx"
        prs = Presentation()

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = doc_stats['name'].replace('_', ' ').title()
        subtitle = slide.placeholders[1]
        subtitle.text = f"Resumen ejecutivo generado el {datetime.now().strftime('%d/%m/%Y %H:%M')}"

        # Metrics slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Métricas Principales"
        body = slide.shapes.placeholders[1].text_frame
        metrics = [
            f"Líneas: {doc_stats['total_lines']:,}",
            f"Palabras: {doc_stats['total_words']:,}",
            f"Secciones: {doc_stats['sections']}",
            f"Bloques de código: {doc_stats['code_blocks']}",
            f"Complejidad: {doc_stats['complexity']['complexity_score']:.1f}",
        ]
        body.text = metrics[0]
        for metric in metrics[1:]:
            body.add_paragraph().text = metric

        # Summary slide
        if doc_stats['summary_points']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Puntos Clave"
            body = slide.shapes.placeholders[1].text_frame
            body.text = doc_stats['summary_points'][0]
            for point in doc_stats['summary_points'][1:5]:
                para = body.add_paragraph()
                para.text = point

        # Risks slide
        if doc_stats['risks']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Alertas / Riesgos"
            body = slide.shapes.placeholders[1].text_frame
            rendered = False
            for category, items in doc_stats['risks'].items():
                text = f"{category}: {', '.join(items[:2])}"
                if not rendered:
                    body.text = text
                    rendered = True
                else:
                    body.add_paragraph().text = text

        # Graphs slide
        if graph_files:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Visualizaciones Clave"
            left = Inches(0.5)
            top = Inches(1.5)
            max_width = Inches(4)
            max_height = Inches(3)
            added = 0
            for graph in graph_files[:2]:
                if Path(graph).exists():
                    slide.shapes.add_picture(graph, left + Inches(added * 5), top, width=max_width, height=max_height)
                    added += 1
                    if added >= 2:
                        break

        prs.save(str(ppt_file))
        return str(ppt_file)
    
    def create_excel(self, doc_stats: Dict, all_stats: List[Dict]) -> str:
        """Crea un archivo Excel mejorado con análisis y gráficas"""
        excel_file = self.output_dir / f"{doc_stats['name']}_premium_v2.xlsx"
        wb = Workbook()
        
        # Hoja 1: Resumen mejorado
        ws1 = wb.active
        ws1.title = "Resumen"
        
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=12)
        title_font = Font(bold=True, size=14)
        border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        
        ws1['A1'] = doc_stats['name'].replace('_', ' ').title()
        ws1['A1'].font = title_font
        ws1.merge_cells('A1:C1')
        
        ws1['A2'] = f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}"
        ws1['A2'].font = Font(italic=True)
        
        # Tabla de resumen expandida
        ws1['A4'] = 'Métrica'
        ws1['B4'] = 'Valor'
        ws1['C4'] = 'Análisis'
        for col in ['A4', 'B4', 'C4']:
            ws1[col].fill = header_fill
            ws1[col].font = header_font
        
        metrics = [
            ('Total de Líneas', doc_stats['total_lines'], 
             'Alto' if doc_stats['total_lines'] > 1000 else 'Medio' if doc_stats['total_lines'] > 500 else 'Bajo'),
            ('Total de Palabras', doc_stats['total_words'], 
             'Extenso' if doc_stats['total_words'] > 5000 else 'Moderado' if doc_stats['total_words'] > 2000 else 'Breve'),
            ('Total de Párrafos', doc_stats['total_paragraphs'], 
             f"Promedio: {doc_stats['total_words']/doc_stats['total_paragraphs']:.1f} palabras/párrafo" if doc_stats['total_paragraphs'] > 0 else 'N/A'),
            ('Secciones', doc_stats['sections'], 
             f"Promedio: {doc_stats['total_words']/doc_stats['sections']:.0f} palabras/sección" if doc_stats['sections'] > 0 else 'N/A'),
            ('Bloques de Código', doc_stats['code_blocks'], 
             f"{doc_stats['code_analysis']['total_code_lines']} líneas" if doc_stats['code_blocks'] > 0 else 'Sin código'),
            ('Tablas', doc_stats['tables'], 'Con datos estructurados' if doc_stats['tables'] > 0 else 'Sin tablas'),
            ('Enlaces', doc_stats['links'], 'Con referencias' if doc_stats['links'] > 0 else 'Sin enlaces'),
            ('Complejidad', f"{doc_stats['complexity']['complexity_score']:.1f}", 
             'Alta' if doc_stats['complexity']['complexity_score'] > 50 else 'Media' if doc_stats['complexity']['complexity_score'] > 30 else 'Baja')
        ]
        
        for i, (metric, value, analysis) in enumerate(metrics, start=5):
            ws1[f'A{i}'] = metric
            ws1[f'B{i}'] = value
            ws1[f'C{i}'] = analysis
            for col in ['A', 'B', 'C']:
                ws1[f'{col}{i}'].border = border
            ws1[f'B{i}'].alignment = Alignment(horizontal='right')

        next_row = 5 + len(metrics)
        if doc_stats['summary_points']:
            ws1[f'A{next_row}'] = 'Puntos Clave'
            ws1[f'A{next_row}'].font = header_font
            for offset, point in enumerate(doc_stats['summary_points'], start=1):
                ws1[f'A{next_row + offset}'] = f'• {point}'
                ws1.merge_cells(start_row=next_row + offset, start_column=1,
                                end_row=next_row + offset, end_column=3)
        next_row = 5 + len(metrics) + len(doc_stats['summary_points']) + 2
        if doc_stats['risks']:
            ws1[f'A{next_row}'] = 'Alertas / Riesgos'
            ws1[f'A{next_row}'].font = header_font
            current = next_row + 1
            for category, entries in doc_stats['risks'].items():
                for entry in entries[:5]:
                    ws1[f'A{current}'] = category
                    ws1[f'B{current}'] = entry
                    ws1[f'A{current}'].border = border
                    ws1[f'B{current}'].border = border
                    current += 1
        
        ws1.column_dimensions['A'].width = 25
        ws1.column_dimensions['B'].width = 15
        ws1.column_dimensions['C'].width = 30
        
        # Hoja 2: Palabras clave
        ws2 = wb.create_sheet("Palabras Clave")
        ws2['A1'] = 'Palabra'
        ws2['B1'] = 'Frecuencia'
        ws2['A1'].fill = header_fill
        ws2['A1'].font = header_font
        ws2['B1'].fill = header_fill
        ws2['B1'].font = header_font
        
        for i, (word, count) in enumerate(doc_stats['keywords'], start=2):
            ws2[f'A{i}'] = word
            ws2[f'B{i}'] = count
            ws2[f'A{i}'].border = border
            ws2[f'B{i}'].border = border
            ws2[f'B{i}'].alignment = Alignment(horizontal='right')
        
        ws2.column_dimensions['A'].width = 25
        ws2.column_dimensions['B'].width = 15
        
        # Gráfica de palabras clave
        chart = BarChart()
        chart.type = "col"
        chart.style = 10
        chart.title = "Top 15 Palabras Clave"
        chart.y_axis.title = 'Frecuencia'
        chart.x_axis.title = 'Palabras'
        
        data = Reference(ws2, min_col=2, min_row=1, max_row=min(16, len(doc_stats['keywords'])+1))
        cats = Reference(ws2, min_col=1, min_row=2, max_row=min(16, len(doc_stats['keywords'])+1))
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.height = 10
        chart.width = 15
        
        ws2.add_chart(chart, "D2")
        
        # Hoja 3: Comparativa (si hay múltiples documentos)
        if len(all_stats) > 1:
            ws3 = wb.create_sheet("Comparativa")
            
            headers = ['Documento', 'Líneas', 'Palabras', 'Párrafos', 'Secciones', 'Código', 'Tablas', 'Enlaces', 'Complejidad']
            for col, header in enumerate(headers, start=1):
                cell = ws3.cell(row=1, column=col, value=header)
                cell.fill = header_fill
                cell.font = header_font
                cell.border = border
            
            for row, stats in enumerate(all_stats, start=2):
                ws3.cell(row=row, column=1, value=stats['name'])
                ws3.cell(row=row, column=2, value=stats['total_lines'])
                ws3.cell(row=row, column=3, value=stats['total_words'])
                ws3.cell(row=row, column=4, value=stats['total_paragraphs'])
                ws3.cell(row=row, column=5, value=stats['sections'])
                ws3.cell(row=row, column=6, value=stats['code_blocks'])
                ws3.cell(row=row, column=7, value=stats['tables'])
                ws3.cell(row=row, column=8, value=stats['links'])
                ws3.cell(row=row, column=9, value=f"{stats['complexity']['complexity_score']:.1f}")
                
                for col in range(1, 10):
                    ws3.cell(row=row, column=col).border = border
            
            for col in range(1, 10):
                ws3.column_dimensions[chr(64 + col)].width = 15
        
        # Hoja 4: Estructura
        ws4 = wb.create_sheet("Estructura")
        ws4['A1'] = 'Nivel'
        ws4['B1'] = 'Título'
        ws4['A1'].fill = header_fill
        ws4['A1'].font = header_font
        ws4['B1'].fill = header_fill
        ws4['B1'].font = header_font
        
        for i, item in enumerate(doc_stats['structure'], start=2):
            ws4.cell(row=i, column=1, value=item['level'])
            ws4.cell(row=i, column=2, value=item['title'])
            ws4.cell(row=i, column=1).border = border
            ws4.cell(row=i, column=2).border = border
        
        ws4.column_dimensions['A'].width = 10
        ws4.column_dimensions['B'].width = 60
        
        wb.save(str(excel_file))
        return str(excel_file)
    
    def process_documents(self, doc_files: List[str]) -> Dict[str, List[str]]:
        """Procesa múltiples documentos con análisis avanzado"""
        results = {}
        all_stats = []
        summary_rows = []
        
        for doc_file in doc_files:
            file_path = Path(doc_file)
            if not file_path.exists():
                print(f"⚠️  Archivo no encontrado: {doc_file}")
                continue
            
            try:
                print(f"📄 Procesando: {file_path.name}")
                
                # Analizar documento
                stats = self.analyze_document(file_path)
                all_stats.append(stats)
                
                # Crear gráficas avanzadas
                print("  📊 Generando gráficas avanzadas...")
                try:
                    graph_files = self.create_advanced_statistics_graphs(stats)
                except Exception as e:
                    print(f"  ⚠️  Error generando gráficas: {e}")
                    graph_files = []
                
                # Crear PDF
                if "pdf" in self.formats:
                    print("  📑 Generando PDF mejorado...")
                    try:
                        pdf_file = self.create_pdf(stats, graph_files)
                    except Exception as e:
                        print(f"  ⚠️  Error generando PDF: {e}")
                        pdf_file = None
                else:
                    pdf_file = None
                
                # Crear Word
                if "word" in self.formats:
                    print("  📝 Generando Word mejorado...")
                    try:
                        word_file = self.create_word(stats, graph_files)
                    except Exception as e:
                        print(f"  ⚠️  Error generando Word: {e}")
                        word_file = None
                else:
                    word_file = None
                
                # Crear Excel
                if "excel" in self.formats:
                    print("  📊 Generando Excel mejorado...")
                    try:
                        excel_file = self.create_excel(stats, all_stats)
                    except Exception as e:
                        print(f"  ⚠️  Error generando Excel: {e}")
                        excel_file = None
                else:
                    excel_file = None

                if "pptx" in self.formats:
                    print("  📽️  Generando presentación PPTX...")
                    try:
                        ppt_file = self.create_ppt(stats, graph_files)
                    except Exception as e:
                        print(f"  ⚠️  Error generando PPTX: {e}")
                        ppt_file = None
                else:
                    ppt_file = None
                
                results[stats['name']] = {
                    'pdf': pdf_file,
                    'word': word_file,
                    'excel': excel_file,
                    'pptx': ppt_file,
                    'graphs': graph_files,
                    'stats': stats
                }
                summary_rows.append(stats)
                
                print(f"  ✅ Completado: {stats['name']}\n")
                
            except Exception as e:
                print(f"  ❌ Error procesando {file_path.name}: {e}\n")
                continue
        
        # Crear gráfica comparativa avanzada
        if len(all_stats) > 1:
            print("📊 Generando gráfica comparativa avanzada...")
            try:
                comparison_graph = self.create_enhanced_comparison_graph(all_stats)
                if comparison_graph:
                    print(f"  ✅ Gráfica comparativa avanzada: {comparison_graph}\n")
            except Exception as e:
                print(f"  ⚠️  Error generando gráfica comparativa: {e}\n")
        
        summary_path = self.generate_summary_report(summary_rows)
        if summary_path:
            print(f"🗂️  Resumen global: {summary_path}")
        return results

    def generate_summary_report(self, all_stats: List[Dict]) -> Optional[str]:
        """Genera un resumen ejecutivo PDF del conjunto de documentos"""
        if not all_stats:
            return None
        try:
            doc = SimpleDocTemplate(
                str(self.summary_path),
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72,
            )
            styles = getSampleStyleSheet()
            story = []

            title = ParagraphStyle(
                'SummaryTitle',
                parent=styles['Heading1'],
                fontSize=24,
                alignment=TA_CENTER,
                textColor=colors.HexColor('#1b4f72'),
            )
            story.append(Paragraph("Resumen Global de Documentos", title))
            story.append(Spacer(1, 0.2 * inch))
            story.append(
                Paragraph(
                    f"Generado: {datetime.now().strftime('%d de %B de %Y, %H:%M')}",
                    styles['Normal'],
                )
            )
            story.append(Spacer(1, 0.3 * inch))

            total_docs = len(all_stats)
            total_words = sum(s['total_words'] for s in all_stats)
            total_lines = sum(s['total_lines'] for s in all_stats)
            avg_complexity = statistics.mean(s['complexity']['complexity_score'] for s in all_stats)

            summary_data = [
                ['Indicador', 'Valor'],
                ['Documentos procesados', total_docs],
                ['Total de palabras', f"{total_words:,}"],
                ['Total de líneas', f"{total_lines:,}"],
                ['Complejidad promedio', f"{avg_complexity:.1f}"],
            ]
            table = Table(summary_data, colWidths=[3 * inch, 3 * inch])
            table.setStyle(
                TableStyle(
                    [
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1b4f72')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#ebf5fb')),
                        ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#aed6f1')),
                    ]
                )
            )
            story.append(table)
            story.append(Spacer(1, 0.3 * inch))

            # Destacar documentos con mayor complejidad
            top_complex = sorted(
                all_stats, key=lambda s: s['complexity']['complexity_score'], reverse=True
            )[:3]
            story.append(Paragraph("Top 3 Documentos por Complejidad", styles['Heading2']))
            for doc_stats in top_complex:
                story.append(
                    Paragraph(
                        f"{doc_stats['name']} – {doc_stats['complexity']['complexity_score']:.1f}",
                        styles['Normal'],
                    )
                )

            doc.build(story)
            return str(self.summary_path)
        except Exception as e:
            print(f"⚠️  Error generando resumen global: {e}")
            return None

    def bundle_outputs(self) -> Optional[str]:
        """Genera un ZIP con todos los artefactos construidos"""
        if not self.output_dir.exists():
            return None
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        archive_name = f"{self.output_dir.name}_bundle_{timestamp}.zip"
        archive_path = self.output_dir.parent / archive_name

        with zipfile.ZipFile(archive_path, 'w', compression=zipfile.ZIP_DEFLATED) as bundle:
            for file_path in self.output_dir.rglob('*'):
                if file_path.is_file():
                    bundle.write(file_path, arcname=file_path.relative_to(self.output_dir))
        return str(archive_path)

    def save_metadata_report(self, results: Dict[str, Dict]) -> Tuple[Optional[str], Optional[str]]:
        """Genera reportes JSON y Markdown con estadísticas agregadas."""
        if not results:
            return None, None

        metadata = {
            "generated_at": datetime.now().isoformat(),
            "output_directory": str(self.output_dir),
            "documents": []
        }

        for name, payload in results.items():
            stats = payload.get('stats', {})
            doc_entry = {
                "name": name,
                "files": {
                    "pdf": payload.get('pdf'),
                    "word": payload.get('word'),
                    "excel": payload.get('excel'),
                    "pptx": payload.get('pptx'),
                    "graphs": payload.get('graphs')
                },
                "metrics": {
                    "lines": stats.get('total_lines'),
                    "words": stats.get('total_words'),
                    "sections": stats.get('sections'),
                    "code_blocks": stats.get('code_blocks'),
                    "complexity_score": stats.get('complexity', {}).get('complexity_score')
                },
                "summary_points": stats.get('summary_points', []),
                "risks": stats.get('risks', {})
            }
            metadata["documents"].append(doc_entry)

        json_path = self.output_dir / "metadata_summary.json"
        with open(json_path, 'w', encoding='utf-8') as jf:
            json.dump(metadata, jf, ensure_ascii=False, indent=2)

        # Markdown report
        md_path = self.output_dir / "metadata_summary.md"
        lines = [
            "# Reporte de Documentos Generados",
            f"- Generado: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}",
            f"- Total documentos: {len(metadata['documents'])}",
            ""
        ]
        for doc in metadata['documents']:
            lines.append(f"## {doc['name']}")
            lines.append("- Métricas:")
            lines.append(f"  - Líneas: {doc['metrics']['lines']}")
            lines.append(f"  - Palabras: {doc['metrics']['words']}")
            lines.append(f"  - Secciones: {doc['metrics']['sections']}")
            lines.append(f"  - Bloques de código: {doc['metrics']['code_blocks']}")
            lines.append(f"  - Complejidad: {doc['metrics']['complexity_score']}")
            if doc['summary_points']:
                lines.append("  - Puntos clave:")
                for point in doc['summary_points']:
                    lines.append(f"    - {point}")
            if doc['risks']:
                lines.append("  - Alertas detectadas:")
                for category, items in doc['risks'].items():
                    sample = ", ".join(items[:3])
                    lines.append(f"    - {category}: {sample}")
            lines.append("")

        with open(md_path, 'w', encoding='utf-8') as mf:
            mf.write("\n".join(lines))

        return str(json_path), str(md_path)

    def generate_html_dashboard(self, metadata_json: Optional[str]) -> Optional[str]:
        """Construye un dashboard HTML con la información agregada."""
        if not metadata_json:
            return None
        metadata_path = Path(metadata_json)
        if not metadata_path.exists():
            return None

        with open(metadata_path, 'r', encoding='utf-8') as f:
            metadata = json.load(f)

        documents = metadata.get('documents', [])
        if not documents:
            return None

        total_words = sum(doc.get('metrics', {}).get('words', 0) or 0 for doc in documents)
        total_lines = sum(doc.get('metrics', {}).get('lines', 0) or 0 for doc in documents)
        avg_complexity = (
            sum(doc.get('metrics', {}).get('complexity_score', 0) or 0 for doc in documents) / len(documents)
        )

        rows_html = []
        for doc in documents:
            metrics = doc.get('metrics', {})
            summary_points = doc.get('summary_points', [])
            risks = doc.get('risks', {})
            summary_html = "<br>".join(f"• {point}" for point in summary_points[:3]) or "—"
            risk_entries = []
            for category, items in risks.items():
                if not items:
                    continue
                risk_entries.append(f"<strong>{category}:</strong> {', '.join(items[:2])}")
            risks_html = "<br>".join(risk_entries) or "—"

            rows_html.append(
                f"""
                <tr>
                    <td>{doc.get('name')}</td>
                    <td>{metrics.get('lines', 0):,}</td>
                    <td>{metrics.get('words', 0):,}</td>
                    <td>{metrics.get('sections', 0)}</td>
                    <td>{metrics.get('code_blocks', 0)}</td>
                    <td>{metrics.get('complexity_score', 0):.1f}</td>
                    <td>{summary_html}</td>
                    <td>{risks_html}</td>
                </tr>
                """
            )

        graphs_preview = ""
        interactive_sections = ""
        if PLOTLY_AVAILABLE:
            try:
                gauges = []
                for idx, doc in enumerate(documents[:3]):
                    fig = go.Figure(
                        go.Indicator(
                            mode="gauge+number",
                            value=doc['metrics'].get('complexity_score') or 0,
                            title={'text': doc['name']},
                            gauge={'axis': {'range': [0, 100]}},
                        )
                    )
                    gauge_path = self.graphs_dir / f"interactive_gauge_{idx}.html"
                    fig.write_html(gauge_path, include_plotlyjs='cdn')
                    gauges.append(gauge_path)
                cards = []
                for gauge in gauges:
                    cards.append(
                        f"""
                        <div class="card">
                            <iframe src="{gauge}" sandbox="allow-scripts allow-same-origin"></iframe>
                        </div>
                        """
                    )
                graphs_preview = (
                    "<section><h2>Indicadores Interactivos</h2><div class=\"cards\">"
                    + "".join(cards)
                    + "</div></section>"
                )
            except Exception as e:
                print(f"⚠️  No se pudieron generar los indicadores interactivos: {e}")

        dashboard_html = f"""
<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8" />
    <title>Dashboard Documentos Importantes</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 2rem; background: #f6f8fb; color: #2c3e50; }}
        h1, h2 {{ color: #1a5276; }}
        .stats {{ display: flex; gap: 1.5rem; flex-wrap: wrap; margin-bottom: 2rem; }}
        .stat {{ background: #fff; border-radius: 12px; padding: 1rem 1.5rem; box-shadow: 0 4px 10px rgba(0,0,0,0.08); }}
        table {{ width: 100%; border-collapse: collapse; background: #fff; box-shadow: 0 4px 10px rgba(0,0,0,0.08); }}
        th, td {{ padding: 0.75rem 1rem; text-align: left; vertical-align: top; }}
        th {{ background: #154360; color: #fff; position: sticky; top: 0; }}
        tr:nth-child(even) {{ background: #f2f4f7; }}
        .cards {{ display: flex; gap: 1rem; flex-wrap: wrap; }}
        .card {{ background: #fff; padding: 0.5rem; border-radius: 10px; box-shadow: 0 3px 8px rgba(0,0,0,0.1); width: 240px; }}
        .card img {{ width: 100%; border-radius: 8px; }}
        footer {{ margin-top: 2rem; font-size: 0.9rem; color: #566573; }}
    </style>
</head>
<body>
    <h1>Dashboard de Documentos Importantes</h1>
    <p>Generado el {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}</p>
    <section class="stats">
        <div class="stat">
            <h3>Total documentos</h3>
            <p><strong>{len(documents)}</strong></p>
        </div>
        <div class="stat">
            <h3>Total de palabras</h3>
            <p><strong>{total_words:,}</strong></p>
        </div>
        <div class="stat">
            <h3>Total de líneas</h3>
            <p><strong>{total_lines:,}</strong></p>
        </div>
        <div class="stat">
            <h3>Complejidad promedio</h3>
            <p><strong>{avg_complexity:.1f}</strong></p>
        </div>
    </section>
    {graphs_preview}
    <section>
        <h2>Detalle de documentos</h2>
        <table>
            <thead>
                <tr>
                    <th>Documento</th>
                    <th>Líneas</th>
                    <th>Palabras</th>
                    <th>Secciones</th>
                    <th>Código</th>
                    <th>Complejidad</th>
                    <th>Puntos clave</th>
                    <th>Alertas</th>
                </tr>
            </thead>
            <tbody>
                {''.join(rows_html)}
            </tbody>
        </table>
    </section>
    <footer>
        Dashboard generado automáticamente por DocumentGeneratorPremiumV2.
    </footer>
</body>
</html>
        """

        dashboard_path = self.output_dir / "dashboard.html"
        with open(dashboard_path, 'w', encoding='utf-8') as html_file:
            html_file.write(dashboard_html)

        return str(dashboard_path)


def discover_additional_documents(base_dir: Path,
                                  existing: List[str],
                                  limit: int = 5,
                                  min_size_kb: int = 50) -> List[str]:
    """Descubre automáticamente más documentos importantes según tamaño."""
    existing_set = {Path(p).resolve() for p in existing}
    candidates = []
    for md_file in base_dir.rglob("*.md"):
        if md_file.resolve() in existing_set:
            continue
        try:
            size_kb = md_file.stat().st_size / 1024
        except OSError:
            continue
        if size_kb >= min_size_kb:
            candidates.append((size_kb, md_file))
    candidates.sort(reverse=True, key=lambda x: x[0])
    return [str(path) for _, path in candidates[:limit]]


def resolve_doc_path(doc_path: str, base_dir: Path) -> Optional[str]:
    """Convierte rutas relativas a absolutas y valida su existencia."""
    candidate = Path(doc_path)
    if not candidate.is_absolute():
        candidate = base_dir / candidate
    candidate = candidate.resolve()
    if candidate.exists():
        return str(candidate)
    return None


def load_docs_from_file(file_path: Path, base_dir: Path) -> List[str]:
    docs: List[str] = []
    if not file_path.exists():
        print(f"⚠️  Archivo de documentos no encontrado: {file_path}")
        return docs
    with open(file_path, "r", encoding="utf-8") as f:
        for line in f:
            cleaned = line.strip()
            if not cleaned or cleaned.startswith("#"):
                continue
            resolved = resolve_doc_path(cleaned, base_dir)
            if resolved:
                docs.append(resolved)
            else:
                print(f"⚠️  Ruta inválida en docs-file: {cleaned}")
    return docs


def deduplicate_preserve_order(items: List[str]) -> List[str]:
    seen = set()
    deduped = []
    for item in items:
        if item and item not in seen:
            deduped.append(item)
            seen.add(item)
    return deduped


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generador premium de documentos importantes con análisis avanzado."
    )
    parser.add_argument(
        "--base-dir",
        type=str,
        default=str(Path(__file__).parent),
        help="Directorio base desde el cual resolver documentos (por defecto, carpeta del script).",
    )
    parser.add_argument(
        "--output-dir",
        type=str,
        default="documentos_importantes_premium_v2",
        help="Directorio donde se guardarán los archivos generados.",
    )
    parser.add_argument(
        "--docs-file",
        type=str,
        help="Archivo de texto con rutas (una por línea) de documentos a procesar.",
    )
    parser.add_argument(
        "--include-doc",
        action="append",
        default=[],
        help="Ruta adicional de documento a incluir (se puede usar múltiples veces).",
    )
    parser.add_argument(
        "--max-auto",
        type=int,
        default=5,
        help="Cantidad máxima de documentos descubiertos automáticamente.",
    )
    parser.add_argument(
        "--min-size-kb",
        type=int,
        default=80,
        help="Tamaño mínimo en KB para considerar un documento en la detección automática.",
    )
    parser.add_argument(
        "--no-auto",
        action="store_true",
        help="Desactiva el descubrimiento automático de documentos grandes.",
    )
    parser.add_argument(
        "--formats",
        nargs="+",
        choices=["pdf", "word", "excel", "pptx"],
        default=["pdf", "word", "excel", "pptx"],
        help="Formatos a generar (por defecto todos).",
    )
    return parser.parse_args()


def main():
    """Función principal"""
    args = parse_args()
    base_dir = Path(args.base_dir).resolve()
    output_dir = Path(args.output_dir)
    if not output_dir.is_absolute():
        output_dir = (base_dir / output_dir).resolve()

    important_docs: List[str] = []

    # Documentos predeterminados
    for doc_rel in DEFAULT_IMPORTANT_DOCS:
        resolved = resolve_doc_path(doc_rel, base_dir)
        if resolved:
            important_docs.append(resolved)

    # Desde archivo externo
    if args.docs_file:
        important_docs.extend(load_docs_from_file(Path(args.docs_file), base_dir))

    # Documentos añadidos manualmente
    for doc in args.include_doc:
        resolved = resolve_doc_path(doc, base_dir)
        if resolved:
            important_docs.append(resolved)
        else:
            print(f"⚠️  Ruta inválida en --include-doc: {doc}")

    # Descubrimiento automático
    if not args.no_auto:
        auto_docs = discover_additional_documents(
            base_dir,
            important_docs,
            limit=args.max_auto,
            min_size_kb=args.min_size_kb,
        )
        if auto_docs:
            print(f"🔍 Documentos descubiertos automáticamente: {len(auto_docs)}")
        important_docs.extend(auto_docs)

    important_docs = deduplicate_preserve_order(important_docs)

    if not important_docs:
        print("❌ No se encontraron documentos para procesar. Verifica las rutas.")
        return

    print("=" * 70)
    print("GENERADOR PREMIUM V2 DE DOCUMENTOS IMPORTANTES")
    print("Versión Mejorada con Análisis Avanzado")
    print("=" * 70)
    print(f"Documentos a procesar ({len(important_docs)}):")
    for path in important_docs:
        print(f"  - {path}")
    print()

    generator = DocumentGeneratorPremiumV2(output_dir=str(output_dir), formats=args.formats)
    results = generator.process_documents(important_docs)

    print("=" * 70)
    print("RESUMEN DE ARCHIVOS GENERADOS")
    print("=" * 70)
    print()
    
    for doc_name, payload in results.items():
        print(f"📄 {doc_name}:")
        if payload.get('pdf'):
            print(f"   📑 PDF:  {payload['pdf']}")
        if payload.get('word'):
            print(f"   📝 Word: {payload['word']}")
        if payload.get('excel'):
            print(f"   📊 Excel: {payload['excel']}")
        if payload.get('pptx'):
            print(f"   📽️  PPTX: {payload['pptx']}")
        print(f"   📈 Gráficas: {len(payload.get('graphs', []))} archivos")
        print()
    
    print(f"✅ Todos los archivos se han generado en: {generator.output_dir}")
    print(f"📁 Gráficas disponibles en: {generator.graphs_dir}")
    metadata_json, metadata_md = generator.save_metadata_report(results)
    if metadata_json:
        print(f"🧾 Reporte JSON: {metadata_json}")
    if metadata_md:
        print(f"📝 Reporte Markdown: {metadata_md}")
    dashboard_path = generator.generate_html_dashboard(metadata_json)
    if dashboard_path:
        print(f"📊 Dashboard HTML: {dashboard_path}")
    bundle_path = generator.bundle_outputs()
    if bundle_path:
        print(f"🗜️  Paquete consolidado: {bundle_path}")

    # Historial de ejecuciones
    history_entry = {
        "timestamp": datetime.now().isoformat(),
        "base_dir": str(base_dir),
        "output_dir": str(output_dir),
        "documents_processed": list(results.keys()),
        "total_documents": len(results),
        "formats": args.formats,
        "metadata_json": metadata_json,
        "dashboard": dashboard_path,
        "bundle": bundle_path,
    }
    history_path = output_dir / "run_history.log"
    with open(history_path, "a", encoding="utf-8") as history_file:
        history_file.write(json.dumps(history_entry, ensure_ascii=False) + "\n")
    print(f"🕘 Historial actualizado: {history_path}")
    print()


if __name__ == "__main__":
    main()


