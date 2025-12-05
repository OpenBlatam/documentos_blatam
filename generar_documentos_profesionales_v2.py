#!/usr/bin/env python3
"""
Script MEJORADO para convertir documentos importantes a PDF, Word y Excel
con gráficas profesionales, portadas, índices y formato de alta calidad
"""

import os
import re
import json
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Tuple, Any, Optional
from collections import Counter
import markdown
from markdown.extensions import tables, fenced_code, codehilite

# PDF
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch, cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image, KeepTogether
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("⚠️  reportlab no disponible. Instala con: pip install reportlab")

# Word
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    WORD_AVAILABLE = True
except ImportError:
    WORD_AVAILABLE = False
    print("⚠️  python-docx no disponible. Instala con: pip install python-docx")

# Excel
try:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, PieChart, ScatterChart, Reference
    from openpyxl.chart.series import DataPoint
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("⚠️  openpyxl no disponible. Instala con: pip install openpyxl")

# PowerPoint
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as PPTXRGB
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx no disponible. Instala con: pip install python-pptx")

# Gráficas
try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from matplotlib.backends.backend_pdf import PdfPages
    import numpy as np
    from matplotlib.patches import Rectangle
    import seaborn as sns
    sns.set_style("whitegrid")
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False
    print("⚠️  matplotlib no disponible. Instala con: pip install matplotlib numpy seaborn")

# Configuración
OUTPUT_DIR = Path("documentos_exportados")
OUTPUT_DIR.mkdir(exist_ok=True)

# Documentos importantes - LISTA EXPANDIDA
IMPORTANT_DOCS = [
    {
        "path": "airflow_automation_prompt.md",
        "title": "Guía de Automatización con Airflow",
        "category": "Automatización",
        "description": "Guía completa de automatización de workflows con Apache Airflow"
    },
    {
        "path": "ARCHITECTURE.md",
        "title": "Arquitectura del Proyecto",
        "category": "Arquitectura",
        "description": "Documentación de la arquitectura y estructura del proyecto"
    },
    {
        "path": "README.md",
        "title": "Documentación Principal - Documentos BLATAM",
        "category": "Documentación",
        "description": "Documentación principal del ecosistema de documentos empresariales"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/ARCHITECTURE_IMPROVEMENTS.md",
        "title": "Mejoras Arquitectónicas",
        "category": "Arquitectura",
        "description": "Plan de mejoras arquitectónicas y refactorización"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/REFACTORING_PLAN.md",
        "title": "Plan de Refactorización",
        "category": "Desarrollo",
        "description": "Plan detallado de refactorización del código de producción"
    },
    {
        "path": "truthgpt_collected/integration_code/production_code/README.md",
        "title": "Código de Producción - README",
        "category": "Desarrollo",
        "description": "Documentación del código de producción y modelos implementados"
    },
    {
        "path": "06_documentation/resumen_final_completo.md",
        "title": "Resumen Final Completo - Documentos BLATAM",
        "category": "Documentación",
        "description": "Resumen ejecutivo completo del ecosistema de documentación"
    },
    {
        "path": "04_business_strategy/Other/ecosystem_summary.md",
        "title": "Resumen del Ecosistema Empresarial",
        "category": "Estrategia",
        "description": "Resumen completo del ecosistema de estrategia empresarial"
    },
    {
        "path": "06_strategy/Strategy_other/readme.md",
        "title": "Estrategias y Frameworks",
        "category": "Estrategia",
        "description": "Documentación completa de estrategias y frameworks empresariales"
    },
    {
        "path": "BEST_PRACTICES.md",
        "title": "Mejores Prácticas",
        "category": "Documentación",
        "description": "Mejores prácticas y guías de implementación"
    },
    {
        "path": "CHANGELOG.md",
        "title": "Registro de Cambios",
        "category": "Documentación",
        "description": "Historial completo de cambios y actualizaciones del proyecto"
    },
    {
        "path": "CONTRIBUTING.md",
        "title": "Guía de Contribución",
        "category": "Documentación",
        "description": "Guía para contribuir al proyecto y estándares de desarrollo"
    },
    {
        "path": "06_documentation/Master_documents/indice_maestro_documentacion.md",
        "title": "Índice Maestro de Documentación",
        "category": "Documentación",
        "description": "Índice completo y navegación de toda la documentación"
    },
    {
        "path": "06_documentation/Index_files/docs_index.md",
        "title": "Índice de Documentación Técnica",
        "category": "Documentación",
        "description": "Índice completo de documentación técnica y de negocio"
    },
    {
        "path": "06_strategy/Strategy_other/readme.md",
        "title": "Estrategias y Frameworks Completos",
        "category": "Estrategia",
        "description": "Documentación completa de estrategias y frameworks empresariales"
    },
    {
        "path": "FAQ.md",
        "title": "Preguntas Frecuentes",
        "category": "Documentación",
        "description": "Respuestas a las preguntas más frecuentes del proyecto"
    }
]


class EnhancedDocumentConverter:
    """Convertidor mejorado con funcionalidades avanzadas"""
    
    def __init__(self, base_path: Path):
        self.base_path = base_path
        self.graphs_dir = OUTPUT_DIR / "graficas"
        self.graphs_dir.mkdir(exist_ok=True)
        self.graphs_created = []
        
    def read_markdown(self, file_path: Path) -> str:
        """Lee un archivo markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"❌ Error leyendo {file_path}: {e}")
            return ""
    
    def parse_markdown_enhanced(self, content: str) -> Dict[str, Any]:
        """Parsea markdown con extracción avanzada de datos"""
        md = markdown.Markdown(extensions=['tables', 'fenced_code', 'codehilite', 'nl2br'])
        html_content = md.convert(content)
        
        data = {
            'sections': [],
            'code_blocks': [],
            'tables': [],
            'metrics': {},
            'lists': [],
            'links': [],
            'images': [],
            'keywords': [],
            'phases': [],
            'checklists': []
        }
        
        lines = content.split('\n')
        
        # Extraer secciones con jerarquía
        for i, line in enumerate(lines):
            if line.strip().startswith('#'):
                level = len(line) - len(line.lstrip('#'))
                title = line.lstrip('#').strip()
                data['sections'].append({
                    'title': title,
                    'level': level,
                    'line': i + 1
                })
        
        # Extraer tablas markdown
        in_table = False
        table_lines = []
        for line in lines:
            if '|' in line and line.strip().startswith('|'):
                if not in_table:
                    in_table = True
                    table_lines = []
                table_lines.append(line)
            elif in_table:
                if table_lines:
                    data['tables'].append('\n'.join(table_lines))
                in_table = False
                table_lines = []
        
        # Extraer bloques de código
        in_code = False
        code_block = []
        for line in lines:
            if line.strip().startswith('```'):
                if in_code:
                    if code_block:
                        data['code_blocks'].append('\n'.join(code_block))
                    code_block = []
                    in_code = False
                else:
                    in_code = True
            elif in_code:
                code_block.append(line)
        
        # Extraer listas y checklists
        for line in lines:
            if re.match(r'^[-*+]\s+', line):
                data['lists'].append(line.strip())
            elif re.match(r'^[-*+]\s+\[[ xX]\]', line):
                data['checklists'].append(line.strip())
        
        # Extraer enlaces
        links = re.findall(r'\[([^\]]+)\]\(([^\)]+)\)', content)
        data['links'] = [{'text': l[0], 'url': l[1]} for l in links]
        
        # Extraer palabras clave importantes (palabras en mayúsculas, términos técnicos)
        words = re.findall(r'\b[A-Z][a-z]+\b', content)
        keywords = [w for w in words if len(w) > 3]
        data['keywords'] = list(set(keywords))[:50]  # Top 50 únicos
        
        # Extraer fases/tareas (líneas con ✅, ⏳, etc.)
        for line in lines:
            if re.search(r'[✅⏳🔄❌]', line) or re.search(r'\[[ xX]\]', line):
                data['phases'].append(line.strip())
        
        # Métricas mejoradas
        total_chars = len(content)
        total_words_list = content.split()
        data['metrics'] = {
            'total_lines': len(lines),
            'total_words': len(total_words_list),
            'total_chars': total_chars,
            'total_sections': len(data['sections']),
            'total_code_blocks': len(data['code_blocks']),
            'total_tables': len(data['tables']),
            'total_lists': len(data['lists']),
            'total_links': len(data['links']),
            'total_checklists': len(data['checklists']),
            'avg_section_length': len(content) / max(len(data['sections']), 1),
            'code_ratio': sum(len(cb) for cb in data['code_blocks']) / max(len(content), 1) * 100,
            'avg_words_per_section': len(total_words_list) / max(len(data['sections']), 1),
            'avg_line_length': total_chars / max(len(lines), 1),
            'readability_score': min(100, max(0, 206.835 - (1.015 * (len(total_words_list) / max(len([l for l in lines if l.strip()]), 1))) - (84.6 * (len([w for w in total_words_list if len(w) > 6]) / max(len(total_words_list), 1))))),
            'complexity_score': len(data['sections']) * 0.3 + len(data['code_blocks']) * 0.4 + len(data['tables']) * 0.2 + len(data['links']) * 0.1
        }
        
        # Contar números
        numbers = re.findall(r'\b(\d+)\b', content)
        if numbers:
            data['metrics']['total_numbers'] = len(numbers)
            data['metrics']['max_number'] = max([int(n) for n in numbers if n.isdigit() and int(n) < 1e10])
        
        return data, html_content
    
    def create_enhanced_graphs(self, doc_data: Dict[str, Any], doc_title: str) -> List[str]:
        """Crea gráficas mejoradas y profesionales"""
        graphs = []
        
        if not MATPLOTLIB_AVAILABLE:
            return graphs
        
        try:
            # Configuración de estilo
            plt.style.use('seaborn-v0_8-darkgrid')
            colors_palette = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E', '#7209B7', '#F72585']
            
            # Gráfica 1: Distribución de secciones (mejorada)
            if doc_data.get('sections'):
                levels = Counter(s['level'] for s in doc_data['sections'])
                
                fig, ax = plt.subplots(figsize=(12, 7))
                levels_sorted = sorted(levels.items())
                bars = ax.bar([f'Nivel {l}' for l, _ in levels_sorted], 
                             [c for _, c in levels_sorted],
                             color=colors_palette[:len(levels_sorted)],
                             edgecolor='white', linewidth=2)
                
                ax.set_title(f'Distribución de Secciones por Nivel\n{doc_title}', 
                           fontsize=16, fontweight='bold', pad=20)
                ax.set_xlabel('Nivel de Sección', fontsize=13, fontweight='bold')
                ax.set_ylabel('Cantidad de Secciones', fontsize=13, fontweight='bold')
                ax.grid(axis='y', alpha=0.3, linestyle='--')
                
                # Agregar valores en las barras
                for bar in bars:
                    height = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2., height,
                           f'{int(height)}', ha='center', va='bottom', 
                           fontweight='bold', fontsize=11)
                
                plt.tight_layout()
                graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_').replace('/', '_')}_secciones.png"
                plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                plt.close()
                graphs.append(str(graph_path))
            
            # Gráfica 2: Métricas del documento (mejorada)
            if doc_data.get('metrics'):
                metrics = doc_data['metrics']
                fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6))
                
                # Gráfica de barras horizontales
                metric_names = ['Líneas', 'Palabras', 'Secciones', 'Bloques\nCódigo', 'Tablas', 'Enlaces']
                metric_values = [
                    metrics.get('total_lines', 0),
                    metrics.get('total_words', 0),
                    metrics.get('total_sections', 0),
                    metrics.get('total_code_blocks', 0),
                    metrics.get('total_tables', 0),
                    metrics.get('total_links', 0)
                ]
                
                # Normalizar para visualización
                max_val = max(metric_values) if metric_values else 1
                normalized = [min(v / max_val * 100, 100) if max_val > 0 else 0 for v in metric_values]
                
                bars = ax1.barh(metric_names, normalized, color=colors_palette[:len(metric_names)])
                ax1.set_title('Métricas Normalizadas (%)', fontsize=14, fontweight='bold')
                ax1.set_xlabel('Porcentaje', fontsize=12)
                ax1.grid(axis='x', alpha=0.3, linestyle='--')
                
                # Agregar valores reales
                for i, (bar, val) in enumerate(zip(bars, metric_values)):
                    width = bar.get_width()
                    ax1.text(width + 2, bar.get_y() + bar.get_height()/2,
                           f'{val:,}', ha='left', va='center', 
                           fontweight='bold', fontsize=10)
                
                # Gráfica de pastel para composición
                pie_data = {
                    'Texto': metrics.get('total_words', 0),
                    'Código': sum(len(cb) for cb in doc_data.get('code_blocks', [])) // 10,
                    'Estructura': metrics.get('total_sections', 0) * 50
                }
                pie_data = {k: v for k, v in pie_data.items() if v > 0}
                
                if pie_data:
                    ax2.pie(pie_data.values(), labels=pie_data.keys(), autopct='%1.1f%%',
                           colors=colors_palette[:len(pie_data)], startangle=90,
                           textprops={'fontsize': 11, 'fontweight': 'bold'})
                    ax2.set_title('Composición del Documento', fontsize=14, fontweight='bold')
                
                plt.suptitle(f'Análisis Completo - {doc_title}', 
                           fontsize=16, fontweight='bold', y=1.02)
                plt.tight_layout()
                graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_').replace('/', '_')}_metricas.png"
                plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                plt.close()
                graphs.append(str(graph_path))
            
            # Gráfica 3: Timeline de secciones (si hay suficientes)
            if len(doc_data.get('sections', [])) > 5:
                sections = doc_data['sections'][:20]  # Primeras 20
                fig, ax = plt.subplots(figsize=(14, max(8, len(sections) * 0.4)))
                
                y_pos = np.arange(len(sections))
                colors_timeline = [colors_palette[min(s['level']-1, len(colors_palette)-1)] 
                                  for s in sections]
                
                ax.barh(y_pos, [1] * len(sections), color=colors_timeline, alpha=0.7)
                ax.set_yticks(y_pos)
                ax.set_yticklabels([f"{'  ' * (s['level']-1)}• {s['title'][:60]}" 
                                   for s in sections], fontsize=9)
                ax.set_xlabel('Estructura del Documento', fontsize=12, fontweight='bold')
                ax.set_title(f'Índice Visual - {doc_title}', fontsize=14, fontweight='bold')
                ax.set_xlim(0, 1.2)
                ax.axis('off')
                
                plt.tight_layout()
                graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_').replace('/', '_')}_timeline.png"
                plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                plt.close()
                graphs.append(str(graph_path))
            
            # Gráfica 4: Frecuencia de palabras clave (si hay)
            if doc_data.get('keywords'):
                keywords = doc_data['keywords'][:15]  # Top 15
                keyword_counts = Counter(doc_data['keywords'])
                top_keywords = keyword_counts.most_common(15)
                
                if top_keywords:
                    fig, ax = plt.subplots(figsize=(12, 7))
                    words, counts = zip(*top_keywords)
                    
                    bars = ax.barh(range(len(words)), counts, 
                                 color=colors_palette[:len(words)], 
                                 edgecolor='white', linewidth=1.5)
                    ax.set_yticks(range(len(words)))
                    ax.set_yticklabels(words, fontsize=10)
                    ax.set_xlabel('Frecuencia', fontsize=12, fontweight='bold')
                    ax.set_title(f'Palabras Clave Más Frecuentes\n{doc_title}', 
                               fontsize=14, fontweight='bold', pad=15)
                    ax.grid(axis='x', alpha=0.3, linestyle='--')
                    
                    # Agregar valores
                    for i, (bar, count) in enumerate(zip(bars, counts)):
                        ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2,
                               f'{count}', ha='left', va='center', 
                               fontweight='bold', fontsize=9)
                    
                    plt.tight_layout()
                    graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_').replace('/', '_')}_keywords.png"
                    plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                    plt.close()
                    graphs.append(str(graph_path))
            
            # Gráfica 5: Análisis de Complejidad y Legibilidad
            if doc_data.get('metrics'):
                metrics = doc_data['metrics']
                fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6))
                
                # Radar chart de métricas clave
                categories = ['Líneas', 'Palabras', 'Secciones', 'Código', 'Enlaces']
                values = [
                    min(metrics.get('total_lines', 0) / 1000, 100),
                    min(metrics.get('total_words', 0) / 100, 100),
                    min(metrics.get('total_sections', 0) * 10, 100),
                    min(metrics.get('total_code_blocks', 0) * 20, 100),
                    min(metrics.get('total_links', 0) * 5, 100)
                ]
                
                # Gráfica de barras comparativa
                bars = ax1.bar(categories, values, color=colors_palette[:len(categories)],
                              edgecolor='white', linewidth=2)
                ax1.set_ylim(0, 100)
                ax1.set_ylabel('Índice Normalizado', fontsize=12, fontweight='bold')
                ax1.set_title('Análisis de Complejidad', fontsize=14, fontweight='bold')
                ax1.grid(axis='y', alpha=0.3, linestyle='--')
                
                for bar, val in zip(bars, values):
                    height = bar.get_height()
                    ax1.text(bar.get_x() + bar.get_width()/2., height + 2,
                           f'{val:.1f}%', ha='center', va='bottom',
                           fontweight='bold', fontsize=9)
                
                # Gráfica de legibilidad y complejidad
                readability = metrics.get('readability_score', 50)
                complexity = metrics.get('complexity_score', 0)
                
                ax2.scatter([complexity], [readability], s=500, 
                          color=colors_palette[0], alpha=0.7, edgecolors='black', linewidth=2)
                ax2.set_xlabel('Complejidad', fontsize=12, fontweight='bold')
                ax2.set_ylabel('Legibilidad', fontsize=12, fontweight='bold')
                ax2.set_title('Análisis de Legibilidad vs Complejidad', 
                            fontsize=14, fontweight='bold')
                ax2.set_xlim(0, max(complexity * 1.5, 10))
                ax2.set_ylim(0, 100)
                ax2.grid(alpha=0.3, linestyle='--')
                
                # Agregar etiqueta
                ax2.text(complexity, readability + 5, 
                        f'Leg: {readability:.1f}\nComp: {complexity:.1f}',
                        ha='center', va='bottom', fontweight='bold', fontsize=10,
                        bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))
                
                plt.suptitle(f'Análisis Avanzado - {doc_title}', 
                           fontsize=16, fontweight='bold', y=1.02)
                plt.tight_layout()
                graph_path = self.graphs_dir / f"{doc_title.replace(' ', '_').replace('/', '_')}_complejidad.png"
                plt.savefig(graph_path, dpi=300, bbox_inches='tight', facecolor='white')
                plt.close()
                graphs.append(str(graph_path))
            
        except Exception as e:
            print(f"⚠️  Error creando gráficas mejoradas: {e}")
            import traceback
            traceback.print_exc()
        
        return graphs
    
    def create_cover_page_pdf(self, canvas_obj, doc_title: str, category: str, description: str):
        """Crea una portada profesional para el PDF"""
        width, height = A4
        
        # Fondo con gradiente simulado
        canvas_obj.setFillColor(colors.HexColor('#2E86AB'))
        canvas_obj.rect(0, height - 200, width, 200, fill=1, stroke=0)
        
        # Título principal
        canvas_obj.setFillColor(colors.white)
        canvas_obj.setFont("Helvetica-Bold", 32)
        title_lines = self._wrap_text(doc_title, width - 144, "Helvetica-Bold", 32)
        y_pos = height - 100
        for line in title_lines:
            canvas_obj.drawCentredString(width/2, y_pos, line)
            y_pos -= 40
        
        # Categoría
        canvas_obj.setFont("Helvetica", 18)
        canvas_obj.drawCentredString(width/2, y_pos - 20, category)
        
        # Descripción
        canvas_obj.setFont("Helvetica", 12)
        desc_lines = self._wrap_text(description, width - 144, "Helvetica", 12)
        y_pos = y_pos - 60
        for line in desc_lines:
            canvas_obj.drawCentredString(width/2, y_pos, line)
            y_pos -= 20
        
        # Fecha
        canvas_obj.setFont("Helvetica-Oblique", 10)
        date_str = datetime.now().strftime("%d de %B de %Y")
        canvas_obj.drawCentredString(width/2, 100, f"Generado el {date_str}")
        
        # Línea decorativa
        canvas_obj.setStrokeColor(colors.HexColor('#F18F01'))
        canvas_obj.setLineWidth(3)
        canvas_obj.line(72, 80, width - 72, 80)
    
    def _wrap_text(self, text: str, max_width: float, font_name: str, font_size: int) -> List[str]:
        """Envuelve texto para que quepa en el ancho especificado"""
        # Implementación simplificada
        words = text.split()
        lines = []
        current_line = []
        current_width = 0
        
        for word in words:
            word_width = len(word) * font_size * 0.6  # Aproximación
            if current_width + word_width > max_width and current_line:
                lines.append(' '.join(current_line))
                current_line = [word]
                current_width = word_width
            else:
                current_line.append(word)
                current_width += word_width + font_size * 0.3
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return lines if lines else [text]
    
    def create_pdf_enhanced(self, content: str, doc_data: Dict[str, Any], 
                           doc_title: str, output_path: Path, graphs: List[str],
                           category: str = "", description: str = ""):
        """Crea un PDF mejorado con portada e índice"""
        if not PDF_AVAILABLE:
            print("⚠️  PDF no disponible")
            return
        
        try:
            class NumberedCanvas(canvas.Canvas):
                def __init__(self, *args, **kwargs):
                    canvas.Canvas.__init__(self, *args, **kwargs)
                    self._saved_page_states = []
                
                def showPage(self):
                    self._saved_page_states.append(dict(self.__dict__))
                    self._startPage()
                
                def save(self):
                    num_pages = len(self._saved_page_states)
                    for state in self._saved_page_states:
                        self.__dict__.update(state)
                        self.draw_page_number(num_pages)
                        canvas.Canvas.showPage(self)
                    canvas.Canvas.save(self)
                
                def draw_page_number(self, page_count):
                    self.setFont("Helvetica", 9)
                    self.drawRightString(
                        A4[0] - 72,
                        30,
                        f"Página {self._pageNumber} de {page_count}"
                    )
            
            doc = SimpleDocTemplate(str(output_path), pagesize=A4,
                                  rightMargin=72, leftMargin=72,
                                  topMargin=72, bottomMargin=50)
            
            # Estilos mejorados
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#2E86AB'),
                spaceAfter=30,
                alignment=TA_CENTER,
                fontName='Helvetica-Bold'
            )
            
            heading1_style = ParagraphStyle(
                'CustomHeading1',
                parent=styles['Heading1'],
                fontSize=18,
                textColor=colors.HexColor('#2E86AB'),
                spaceAfter=12,
                spaceBefore=12,
                fontName='Helvetica-Bold'
            )
            
            heading2_style = ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=colors.HexColor('#A23B72'),
                spaceAfter=10,
                spaceBefore=10,
                fontName='Helvetica-Bold'
            )
            
            normal_style = ParagraphStyle(
                'CustomNormal',
                parent=styles['Normal'],
                fontSize=11,
                leading=14,
                spaceAfter=6,
                alignment=TA_JUSTIFY
            )
            
            story = []
            
            # Portada (se agregará después)
            story.append(PageBreak())
            
            # Índice
            if doc_data.get('sections'):
                story.append(Paragraph("📑 Índice de Contenidos", heading1_style))
                story.append(Spacer(1, 0.2*inch))
                
                for section in doc_data['sections'][:30]:  # Primeras 30 secciones
                    indent = (section['level'] - 1) * 0.3
                    style = heading2_style if section['level'] == 1 else normal_style
                    text = f"{'&nbsp;' * int(indent * 10)}• {section['title']}"
                    story.append(Paragraph(text, style))
                    story.append(Spacer(1, 0.05*inch))
                
                story.append(PageBreak())
            
            # Gráficas
            if graphs:
                story.append(Paragraph("📊 Análisis Visual del Documento", heading1_style))
                story.append(Spacer(1, 0.2*inch))
                for graph_path in graphs:
                    if os.path.exists(graph_path):
                        try:
                            img = Image(graph_path, width=6*inch, height=3.6*inch)
                            story.append(img)
                            story.append(Spacer(1, 0.2*inch))
                        except Exception as e:
                            print(f"⚠️  Error agregando gráfica {graph_path}: {e}")
                story.append(PageBreak())
            
            # Contenido
            story.append(Paragraph("📄 Contenido del Documento", heading1_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Procesar contenido
            lines = content.split('\n')
            in_code_block = False
            code_lines = []
            code_lang = ""
            
            for line in lines:
                if line.startswith('# '):
                    story.append(Paragraph(line[2:], heading1_style))
                elif line.startswith('## '):
                    story.append(Paragraph(line[3:], heading2_style))
                elif line.startswith('### '):
                    story.append(Paragraph(line[4:], heading2_style))
                elif line.startswith('```'):
                    if in_code_block:
                        if code_lines:
                            code_text = '\n'.join(code_lines)
                            para = Paragraph(f"<font face='Courier' size=9>{code_text}</font>", normal_style)
                            story.append(para)
                            code_lines = []
                        in_code_block = False
                        code_lang = ""
                    else:
                        in_code_block = True
                        code_lang = line[3:].strip()
                elif in_code_block:
                    code_lines.append(line)
                elif line.strip().startswith('- ') or line.strip().startswith('* '):
                    text = line.strip()[2:].strip()
                    # Limpiar markdown básico
                    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
                    text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', text)
                    story.append(Paragraph(f"• {text}", normal_style))
                elif not line.strip():
                    story.append(Spacer(1, 0.1*inch))
                elif line.strip():
                    clean_line = line
                    # Convertir markdown a HTML de forma segura
                    clean_line = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', clean_line)
                    clean_line = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r'<i>\1</i>', clean_line)
                    # Escapar caracteres especiales
                    parts = re.split(r'(<[bi]>|</[bi]>)', clean_line)
                    result_parts = []
                    for part in parts:
                        if re.match(r'<[bi]>|</[bi]>', part):
                            result_parts.append(part)
                        else:
                            result_parts.append(part.replace('<', '&lt;').replace('>', '&gt;').replace('&', '&amp;'))
                    clean_line = ''.join(result_parts)
                    
                    try:
                        story.append(Paragraph(clean_line, normal_style))
                    except Exception:
                        plain_text = re.sub(r'[*_`]', '', line)
                        story.append(Paragraph(plain_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'), normal_style))
            
            # Construir PDF con portada personalizada
            def on_first_page(canvas_obj, doc):
                self.create_cover_page_pdf(canvas_obj, doc_title, category, description)
            
            doc.build(story, onFirstPage=on_first_page, 
                     canvasmaker=NumberedCanvas)
            print(f"✅ PDF mejorado creado: {output_path}")
            
        except Exception as e:
            print(f"❌ Error creando PDF mejorado: {e}")
            import traceback
            traceback.print_exc()
    
    def create_word_enhanced(self, content: str, doc_data: Dict[str, Any],
                             doc_title: str, output_path: Path, graphs: List[str],
                             category: str = "", description: str = ""):
        """Crea un documento Word mejorado"""
        if not WORD_AVAILABLE:
            print("⚠️  Word no disponible")
            return
        
        try:
            doc = Document()
            
            # Portada
            title = doc.add_heading(doc_title, 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title.runs[0]
            title_run.font.color.rgb = RGBColor(46, 134, 171)
            title_run.font.size = Pt(28)
            
            if category:
                cat_para = doc.add_paragraph(category)
                cat_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                cat_run = cat_para.runs[0]
                cat_run.font.size = Pt(16)
                cat_run.font.color.rgb = RGBColor(162, 59, 114)
            
            doc.add_paragraph()
            if description:
                desc_para = doc.add_paragraph(description)
                desc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                desc_run = desc_para.runs[0]
                desc_run.italic = True
                desc_run.font.size = Pt(12)
            
            date_para = doc.add_paragraph()
            date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            date_run = date_para.add_run(f"Generado el {datetime.now().strftime('%d de %B de %Y')}")
            date_run.italic = True
            date_run.font.size = Pt(10)
            
            doc.add_page_break()
            
            # Índice
            if doc_data.get('sections'):
                doc.add_heading('📑 Índice de Contenidos', 1)
                for section in doc_data['sections'][:30]:
                    para = doc.add_paragraph()
                    para.add_run('  ' * (section['level'] - 1) + '• ' + section['title'])
                    para.style = 'List Bullet' if section['level'] > 1 else 'Normal'
                doc.add_page_break()
            
            # Gráficas
            if graphs:
                doc.add_heading('📊 Análisis Visual', 1)
                for graph_path in graphs:
                    if os.path.exists(graph_path):
                        try:
                            doc.add_picture(graph_path, width=Inches(6))
                            doc.add_paragraph()
                        except Exception as e:
                            print(f"⚠️  Error agregando gráfica: {e}")
                doc.add_page_break()
            
            # Contenido
            doc.add_heading('📄 Contenido', 1)
            
            lines = content.split('\n')
            in_code_block = False
            code_lines = []
            
            for line in lines:
                if line.startswith('# '):
                    doc.add_heading(line[2:], 1)
                elif line.startswith('## '):
                    doc.add_heading(line[3:], 2)
                elif line.startswith('### '):
                    doc.add_heading(line[4:], 3)
                elif line.startswith('```'):
                    if in_code_block:
                        if code_lines:
                            para = doc.add_paragraph('\n'.join(code_lines))
                            para.style = 'No Spacing'
                            for run in para.runs:
                                run.font.name = 'Courier New'
                                run.font.size = Pt(9)
                            code_lines = []
                        in_code_block = False
                    else:
                        in_code_block = True
                elif in_code_block:
                    code_lines.append(line)
                elif line.strip().startswith('- ') or line.strip().startswith('* '):
                    doc.add_paragraph(line.strip()[2:].strip(), style='List Bullet')
                elif line.strip():
                    para = doc.add_paragraph()
                    text = re.sub(r'\*\*(.+?)\*\*', r'\1', line)
                    text = re.sub(r'\*(.+?)\*', r'\1', text)
                    para.add_run(text)
            
            doc.save(str(output_path))
            print(f"✅ Word mejorado creado: {output_path}")
            
        except Exception as e:
            print(f"❌ Error creando Word mejorado: {e}")
            import traceback
            traceback.print_exc()
    
    def create_excel_enhanced(self, doc_data: Dict[str, Any], doc_title: str, 
                             output_path: Path, graphs: List[str]):
        """Crea un archivo Excel mejorado con múltiples hojas y gráficas"""
        if not EXCEL_AVAILABLE:
            print("⚠️  Excel no disponible")
            return
        
        try:
            wb = Workbook()
            
            # Hoja 1: Resumen Ejecutivo
            ws1 = wb.active
            ws1.title = "Resumen Ejecutivo"
            
            # Estilos
            title_fill = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
            header_fill = PatternFill(start_color="A23B72", end_color="A23B72", fill_type="solid")
            title_font = Font(bold=True, size=16, color="FFFFFF")
            header_font = Font(bold=True, size=12, color="FFFFFF")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Título
            ws1['A1'] = doc_title
            ws1['A1'].font = title_font
            ws1['A1'].fill = title_fill
            ws1['A1'].alignment = Alignment(horizontal='center', vertical='center')
            ws1.merge_cells('A1:D1')
            ws1.row_dimensions[1].height = 35
            
            # Fecha
            ws1['A2'] = f"Generado el {datetime.now().strftime('%d/%m/%Y %H:%M')}"
            ws1['A2'].font = Font(italic=True, size=10)
            ws1.merge_cells('A2:D2')
            
            # Métricas principales
            row = 4
            ws1[f'A{row}'] = "Métrica"
            ws1[f'B{row}'] = "Valor"
            ws1[f'A{row}'].font = header_font
            ws1[f'B{row}'].font = header_font
            ws1[f'A{row}'].fill = header_fill
            ws1[f'B{row}'].fill = header_fill
            ws1[f'A{row}'].border = border
            ws1[f'B{row}'].border = border
            
            row += 1
            metrics = doc_data.get('metrics', {})
            metric_data = [
                ("Total de Líneas", metrics.get('total_lines', 0)),
                ("Total de Palabras", metrics.get('total_words', 0)),
                ("Total de Secciones", metrics.get('total_sections', 0)),
                ("Bloques de Código", metrics.get('total_code_blocks', 0)),
                ("Tablas", metrics.get('total_tables', 0)),
                ("Enlaces", metrics.get('total_links', 0)),
                ("Listas", metrics.get('total_lists', 0)),
            ]
            
            for metric_name, metric_value in metric_data:
                ws1[f'A{row}'] = metric_name
                ws1[f'B{row}'] = metric_value
                ws1[f'A{row}'].border = border
                ws1[f'B{row}'].border = border
                ws1[f'B{row}'].alignment = Alignment(horizontal='right')
                row += 1
            
            ws1.column_dimensions['A'].width = 25
            ws1.column_dimensions['B'].width = 15
            
            # Gráfica de barras
            if metric_data and any(m[1] > 0 for m in metric_data):
                chart1 = BarChart()
                chart1.type = "col"
                chart1.style = 10
                chart1.title = "Métricas del Documento"
                chart1.y_axis.title = 'Valor'
                chart1.x_axis.title = 'Métrica'
                
                data_start_row = 5
                data_end_row = data_start_row + len(metric_data) - 1
                data_ref = Reference(ws1, min_col=2, min_row=data_start_row, max_row=data_end_row)
                cats_ref = Reference(ws1, min_col=1, min_row=data_start_row, max_row=data_end_row)
                
                chart1.add_data(data_ref, titles_from_data=False)
                chart1.set_categories(cats_ref)
                chart1.height = 10
                chart1.width = 15
                
                ws1.add_chart(chart1, "D4")
            
            # Hoja 2: Secciones
            if doc_data.get('sections'):
                ws2 = wb.create_sheet("Secciones")
                
                ws2['A1'] = "Nivel"
                ws2['B1'] = "Título"
                ws2['C1'] = "Línea"
                ws2['A1'].font = header_font
                ws2['B1'].font = header_font
                ws2['C1'].font = header_font
                ws2['A1'].fill = header_fill
                ws2['B1'].fill = header_fill
                ws2['C1'].fill = header_fill
                ws2['A1'].border = border
                ws2['B1'].border = border
                ws2['C1'].border = border
                
                row = 2
                for section in doc_data['sections']:
                    ws2[f'A{row}'] = section['level']
                    ws2[f'B{row}'] = section['title']
                    ws2[f'C{row}'] = section.get('line', row)
                    ws2[f'A{row}'].border = border
                    ws2[f'B{row}'].border = border
                    ws2[f'C{row}'].border = border
                    ws2[f'A{row}'].alignment = Alignment(horizontal='center')
                    ws2[f'C{row}'].alignment = Alignment(horizontal='center')
                    row += 1
                
                ws2.column_dimensions['A'].width = 10
                ws2.column_dimensions['B'].width = 60
                ws2.column_dimensions['C'].width = 10
                
                # Gráfica de pastel
                levels = {}
                for section in doc_data['sections']:
                    level = section['level']
                    levels[level] = levels.get(level, 0) + 1
                
                if levels:
                    chart_row = 2
                    for level in sorted(levels.keys()):
                        ws2[f'D{chart_row}'] = f"Nivel {level}"
                        ws2[f'E{chart_row}'] = levels[level]
                        chart_row += 1
                    
                    chart2 = PieChart()
                    chart2.title = "Distribución de Secciones por Nivel"
                    
                    data_end_row = 2 + len(levels) - 1
                    data_ref = Reference(ws2, min_col=5, min_row=2, max_row=data_end_row)
                    cats_ref = Reference(ws2, min_col=4, min_row=2, max_row=data_end_row)
                    
                    chart2.add_data(data_ref, titles_from_data=False)
                    chart2.set_categories(cats_ref)
                    chart2.height = 10
                    chart2.width = 15
                    
                    ws2.add_chart(chart2, "G2")
            
            # Hoja 3: Estadísticas
            ws3 = wb.create_sheet("Estadísticas")
            ws3['A1'] = "Categoría"
            ws3['B1'] = "Cantidad"
            ws3['A1'].font = header_font
            ws3['B1'].font = header_font
            ws3['A1'].fill = header_fill
            ws3['B1'].fill = header_fill
            
            stats_data = [
                ("Bloques de Código", len(doc_data.get('code_blocks', []))),
                ("Tablas", len(doc_data.get('tables', []))),
                ("Listas", len(doc_data.get('lists', []))),
                ("Enlaces", len(doc_data.get('links', []))),
                ("Checklists", len(doc_data.get('checklists', []))),
                ("Palabras Clave", len(doc_data.get('keywords', []))),
            ]
            
            row = 2
            for stat_name, stat_value in stats_data:
                ws3[f'A{row}'] = stat_name
                ws3[f'B{row}'] = stat_value
                ws3[f'A{row}'].border = border
                ws3[f'B{row}'].border = border
                row += 1
            
            ws3.column_dimensions['A'].width = 25
            ws3.column_dimensions['B'].width = 15
            
            # Hoja 4: Análisis de Complejidad
            ws4 = wb.create_sheet("Análisis Avanzado")
            ws4['A1'] = "Métrica"
            ws4['B1'] = "Valor"
            ws4['C1'] = "Interpretación"
            ws4['A1'].font = header_font
            ws4['B1'].font = header_font
            ws4['C1'].font = header_font
            ws4['A1'].fill = header_fill
            ws4['B1'].fill = header_fill
            ws4['C1'].fill = header_fill
            
            metrics = doc_data.get('metrics', {})
            advanced_metrics = [
                ("Legibilidad", metrics.get('readability_score', 0), 
                 "Alto (>70): Fácil | Medio (50-70): Moderado | Bajo (<50): Complejo"),
                ("Complejidad", metrics.get('complexity_score', 0),
                 "Bajo (<5): Simple | Medio (5-15): Moderado | Alto (>15): Complejo"),
                ("Promedio Palabras/Sección", metrics.get('avg_words_per_section', 0),
                 "Indica densidad de contenido por sección"),
                ("Promedio Caracteres/Línea", metrics.get('avg_line_length', 0),
                 "Indica longitud promedio de líneas"),
                ("Ratio de Código (%)", metrics.get('code_ratio', 0),
                 "Porcentaje de contenido que es código"),
            ]
            
            row = 2
            for metric_name, metric_value, interpretation in advanced_metrics:
                ws4[f'A{row}'] = metric_name
                ws4[f'B{row}'] = round(metric_value, 2) if isinstance(metric_value, float) else metric_value
                ws4[f'C{row}'] = interpretation
                ws4[f'A{row}'].border = border
                ws4[f'B{row}'].border = border
                ws4[f'C{row}'].border = border
                ws4[f'B{row}'].alignment = Alignment(horizontal='right')
                row += 1
            
            ws4.column_dimensions['A'].width = 25
            ws4.column_dimensions['B'].width = 15
            ws4.column_dimensions['C'].width = 60
            
            # Hoja 5: Enlaces y Referencias
            if doc_data.get('links'):
                ws5 = wb.create_sheet("Enlaces")
                ws5['A1'] = "Texto"
                ws5['B1'] = "URL"
                ws5['A1'].font = header_font
                ws5['B1'].font = header_font
                ws5['A1'].fill = header_fill
                ws5['B1'].fill = header_fill
                
                row = 2
                for link in doc_data['links'][:100]:  # Primeros 100
                    ws5[f'A{row}'] = link.get('text', '')
                    ws5[f'B{row}'] = link.get('url', '')
                    ws5[f'A{row}'].border = border
                    ws5[f'B{row}'].border = border
                    row += 1
                
                ws5.column_dimensions['A'].width = 40
                ws5.column_dimensions['B'].width = 60
            
            wb.save(str(output_path))
            print(f"✅ Excel mejorado creado: {output_path}")
            
        except Exception as e:
            print(f"❌ Error creando Excel mejorado: {e}")
            import traceback
            traceback.print_exc()
    
    def create_powerpoint_enhanced(self, doc_data: Dict[str, Any], doc_title: str,
                                  output_path: Path, graphs: List[str],
                                  category: str = "", description: str = ""):
        """Crea una presentación PowerPoint profesional"""
        if not PPTX_AVAILABLE:
            print("⚠️  PowerPoint no disponible")
            return
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Portada
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            
            # Fondo con color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = PPTXRGB(46, 134, 171)  # #2E86AB
            
            # Título
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = doc_title
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = PPTXRGB(255, 255, 255)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Categoría
            if category:
                cat_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
                cat_frame = cat_box.text_frame
                cat_frame.text = category
                cat_frame.paragraphs[0].font.size = Pt(24)
                cat_frame.paragraphs[0].font.color.rgb = PPTXRGB(241, 143, 1)  # #F18F01
                cat_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Fecha
            date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
            date_frame = date_box.text_frame
            date_frame.text = f"Generado el {datetime.now().strftime('%d de %B de %Y')}"
            date_frame.paragraphs[0].font.size = Pt(14)
            date_frame.paragraphs[0].font.color.rgb = PPTXRGB(255, 255, 255)
            date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Slide 2: Resumen Ejecutivo
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Resumen Ejecutivo"
            
            metrics = doc_data.get('metrics', {})
            content = slide.placeholders[1].text_frame
            content.text = f"Total de Palabras: {metrics.get('total_words', 0):,}"
            
            p = content.add_paragraph()
            p.text = f"Total de Secciones: {metrics.get('total_sections', 0)}"
            p.level = 0
            
            p = content.add_paragraph()
            p.text = f"Bloques de Código: {len(doc_data.get('code_blocks', []))}"
            p.level = 0
            
            p = content.add_paragraph()
            p.text = f"Legibilidad: {metrics.get('readability_score', 0):.1f}/100"
            p.level = 0
            
            p = content.add_paragraph()
            p.text = f"Complejidad: {metrics.get('complexity_score', 0):.2f}"
            p.level = 0
            
            # Slide 3: Gráficas (si hay)
            if graphs:
                for i, graph_path in enumerate(graphs[:3]):  # Máximo 3 gráficas
                    if os.path.exists(graph_path):
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        try:
                            slide.shapes.add_picture(graph_path, Inches(0.5), Inches(0.5), 
                                                    width=Inches(9), height=Inches(6.5))
                        except Exception as e:
                            print(f"⚠️  Error agregando gráfica a PowerPoint: {e}")
            
            # Slide 4: Estructura del Documento
            if doc_data.get('sections'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Estructura del Documento"
                
                content = slide.placeholders[1].text_frame
                for section in doc_data['sections'][:10]:  # Primeras 10 secciones
                    p = content.add_paragraph()
                    p.text = f"{'  ' * (section['level'] - 1)}• {section['title']}"
                    p.level = min(section['level'] - 1, 2)
            
            prs.save(str(output_path))
            print(f"✅ PowerPoint mejorado creado: {output_path}")
            
        except Exception as e:
            print(f"❌ Error creando PowerPoint mejorado: {e}")
            import traceback
            traceback.print_exc()
    
    def create_document_summary(self, doc_data: Dict[str, Any], doc_title: str, output_path: Path):
        """Crea un resumen automático del documento"""
        try:
            metrics = doc_data.get('metrics', {})
            
            summary = f"""
RESUMEN AUTOMÁTICO - {doc_title}
{'=' * 70}

Fecha: {datetime.now().strftime('%d de %B de %Y %H:%M')}

MÉTRICAS PRINCIPALES
{'-' * 70}
• Total de Palabras: {metrics.get('total_words', 0):,}
• Total de Líneas: {metrics.get('total_lines', 0):,}
• Total de Secciones: {metrics.get('total_sections', 0)}
• Bloques de Código: {len(doc_data.get('code_blocks', []))}
• Tablas: {len(doc_data.get('tables', []))}
• Enlaces: {len(doc_data.get('links', []))}
• Listas: {len(doc_data.get('lists', []))}

ANÁLISIS DE CALIDAD
{'-' * 70}
• Legibilidad: {metrics.get('readability_score', 0):.1f}/100
  {'✓ Excelente' if metrics.get('readability_score', 0) > 70 else '✓ Buena' if metrics.get('readability_score', 0) > 50 else '⚠ Mejorable'}
• Complejidad: {metrics.get('complexity_score', 0):.2f}
  {'✓ Simple' if metrics.get('complexity_score', 0) < 5 else '✓ Moderado' if metrics.get('complexity_score', 0) < 15 else '⚠ Complejo'}
• Promedio Palabras/Sección: {metrics.get('avg_words_per_section', 0):.0f}
• Ratio de Código: {metrics.get('code_ratio', 0):.1f}%

ESTRUCTURA
{'-' * 70}
"""
            
            if doc_data.get('sections'):
                summary += f"Total de secciones: {len(doc_data['sections'])}\n"
                summary += "\nPrincipales secciones:\n"
                for section in doc_data['sections'][:10]:
                    summary += f"  {'  ' * (section['level'] - 1)}• {section['title']}\n"
            
            if doc_data.get('keywords'):
                summary += f"""
PALABRAS CLAVE
{'-' * 70}
"""
                top_keywords = Counter(doc_data['keywords']).most_common(10)
                for keyword, count in top_keywords:
                    summary += f"• {keyword} ({count})\n"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(summary)
            
        except Exception as e:
            print(f"⚠️  Error creando resumen: {e}")
    
    def convert_document(self, doc_info: Dict[str, str]) -> Optional[Dict[str, Any]]:
        """Convierte un documento a los tres formatos mejorados y retorna los datos"""
        doc_path = self.base_path / doc_info['path']
        
        if not doc_path.exists():
            print(f"⚠️  Documento no encontrado: {doc_path}")
            return None
        
        print(f"\n📄 Procesando: {doc_info['title']}")
        print(f"   Ruta: {doc_info['path']}")
        
        content = self.read_markdown(doc_path)
        if not content:
            return None
        
        # Parsear con método mejorado
        doc_data, html_content = self.parse_markdown_enhanced(content)
        
        # Agregar información del documento
        doc_data['title'] = doc_info['title']
        doc_data['category'] = doc_info.get('category', '')
        doc_data['description'] = doc_info.get('description', '')
        doc_data['path'] = doc_info['path']
        
        # Crear gráficas mejoradas
        safe_title = re.sub(r'[^\w\s-]', '', doc_info['title']).strip().replace(' ', '_').replace('/', '_')
        graphs = self.create_enhanced_graphs(doc_data, safe_title)
        
        # Crear archivos
        base_name = safe_title
        
        # PDF mejorado
        pdf_path = OUTPUT_DIR / f"{base_name}.pdf"
        self.create_pdf_enhanced(
            content, doc_data, doc_info['title'], pdf_path, graphs,
            doc_info.get('category', ''), doc_info.get('description', '')
        )
        
        # Word mejorado
        word_path = OUTPUT_DIR / f"{base_name}.docx"
        self.create_word_enhanced(
            content, doc_data, doc_info['title'], word_path, graphs,
            doc_info.get('category', ''), doc_info.get('description', '')
        )
        
        # Excel mejorado
        excel_path = OUTPUT_DIR / f"{base_name}.xlsx"
        self.create_excel_enhanced(doc_data, doc_info['title'], excel_path, graphs)
        
        # PowerPoint mejorado
        if PPTX_AVAILABLE:
            pptx_path = OUTPUT_DIR / f"{base_name}.pptx"
            self.create_powerpoint_enhanced(doc_data, doc_info['title'], pptx_path, graphs,
                                           doc_info.get('category', ''), doc_info.get('description', ''))
        
        # Resumen automático del documento
        summary_path = OUTPUT_DIR / "resumenes" / f"{base_name}_resumen.txt"
        summary_path.parent.mkdir(exist_ok=True)
        self.create_document_summary(doc_data, doc_info['title'], summary_path)
        
        print(f"✅ Completado: {doc_info['title']}\n")
        
        return doc_data


def create_html_dashboard(all_docs_data: List[Dict], output_path: Path):
    """Crea un dashboard HTML interactivo con todos los documentos"""
    html_content = f"""<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Dashboard - Documentos BLATAM</title>
    <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 20px;
            min-height: 100vh;
        }}
        .container {{
            max-width: 1400px;
            margin: 0 auto;
            background: white;
            border-radius: 20px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
            padding: 40px;
        }}
        h1 {{ color: #2E86AB; text-align: center; margin-bottom: 10px; font-size: 2.5em; }}
        .subtitle {{ text-align: center; color: #666; margin-bottom: 40px; font-size: 1.1em; }}
        .stats-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin-bottom: 40px;
        }}
        .stat-card {{
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 25px;
            border-radius: 15px;
            text-align: center;
            box-shadow: 0 10px 30px rgba(0,0,0,0.2);
            transition: transform 0.3s;
        }}
        .stat-card:hover {{ transform: translateY(-5px); }}
        .stat-value {{ font-size: 2.5em; font-weight: bold; margin-bottom: 10px; }}
        .stat-label {{ font-size: 1.1em; opacity: 0.9; }}
        .charts-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(500px, 1fr));
            gap: 30px;
            margin-bottom: 40px;
        }}
        .chart-container {{
            background: #f8f9fa;
            padding: 25px;
            border-radius: 15px;
            box-shadow: 0 5px 15px rgba(0,0,0,0.1);
        }}
        .chart-title {{
            color: #2E86AB;
            margin-bottom: 20px;
            font-size: 1.3em;
            font-weight: bold;
        }}
        .docs-table {{
            width: 100%;
            border-collapse: collapse;
            margin-top: 30px;
            background: white;
            border-radius: 10px;
            overflow: hidden;
            box-shadow: 0 5px 15px rgba(0,0,0,0.1);
        }}
        .docs-table th {{
            background: #2E86AB;
            color: white;
            padding: 15px;
            text-align: left;
            font-weight: bold;
        }}
        .docs-table td {{ padding: 12px 15px; border-bottom: 1px solid #eee; }}
        .docs-table tr:hover {{ background: #f8f9fa; }}
        .badge {{
            display: inline-block;
            padding: 5px 12px;
            border-radius: 20px;
            font-size: 0.85em;
            font-weight: bold;
        }}
        .badge-automation {{ background: #F18F01; color: white; }}
        .badge-architecture {{ background: #2E86AB; color: white; }}
        .badge-documentation {{ background: #6A994E; color: white; }}
        .badge-development {{ background: #A23B72; color: white; }}
        .badge-strategy {{ background: #7209B7; color: white; }}
        .search-box {{
            margin: 20px 0;
            padding: 15px;
            background: #f8f9fa;
            border-radius: 10px;
        }}
        .search-box input {{
            width: 100%;
            padding: 12px;
            border: 2px solid #2E86AB;
            border-radius: 8px;
            font-size: 1em;
        }}
        .filter-buttons {{
            display: flex;
            gap: 10px;
            flex-wrap: wrap;
            margin: 15px 0;
        }}
        .filter-btn {{
            padding: 8px 16px;
            border: 2px solid #2E86AB;
            background: white;
            color: #2E86AB;
            border-radius: 20px;
            cursor: pointer;
            transition: all 0.3s;
        }}
        .filter-btn:hover, .filter-btn.active {{
            background: #2E86AB;
            color: white;
        }}
    </style>
</head>
<body>
    <div class="container">
        <h1>📊 Dashboard de Documentos BLATAM</h1>
        <p class="subtitle">Análisis completo y visualización de todos los documentos procesados</p>
        
        <div class="search-box">
            <input type="text" id="searchInput" placeholder="🔍 Buscar documentos..." onkeyup="filterTable()">
        </div>
        
        <div class="filter-buttons">
            <button class="filter-btn active" onclick="filterByCategory('all')">Todos</button>
"""
    
    # Agregar botones de filtro por categoría
    categories = set(d.get('category', 'Otros') for d in all_docs_data)
    for cat in sorted(categories):
        cat_class = cat.lower().replace(' ', '-')
        html_content += f'            <button class="filter-btn" onclick="filterByCategory(\'{cat_class}\')">{cat}</button>\n'
    
    html_content += """        </div>
        
        <div class="stats-grid">
            <div class="stat-card">
                <div class="stat-value">{len(all_docs_data)}</div>
                <div class="stat-label">Documentos Procesados</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">{sum(d.get('metrics', {}).get('total_words', 0) for d in all_docs_data):,}</div>
                <div class="stat-label">Total de Palabras</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">{sum(d.get('metrics', {}).get('total_sections', 0) for d in all_docs_data)}</div>
                <div class="stat-label">Total de Secciones</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">{sum(len(d.get('code_blocks', [])) for d in all_docs_data)}</div>
                <div class="stat-label">Bloques de Código</div>
            </div>
        </div>
        
        <div class="charts-grid">
            <div class="chart-container">
                <div class="chart-title">Distribución por Categoría</div>
                <canvas id="categoryChart"></canvas>
            </div>
            <div class="chart-container">
                <div class="chart-title">Tamaño de Documentos (Palabras)</div>
                <canvas id="sizeChart"></canvas>
            </div>
        </div>
        
        <h2 style="color: #2E86AB; margin-top: 40px; margin-bottom: 20px;">📄 Lista de Documentos</h2>
        <table class="docs-table" id="docsTable">
            <thead>
                <tr>
                    <th>Documento</th>
                    <th>Categoría</th>
                    <th>Palabras</th>
                    <th>Secciones</th>
                    <th>Legibilidad</th>
                    <th>Complejidad</th>
                </tr>
            </thead>
            <tbody>
"""
    
    for doc in all_docs_data:
        metrics = doc.get('metrics', {})
        category = doc.get('category', 'N/A')
        badge_class = f"badge-{category.lower().replace(' ', '-')}" if category else "badge-documentation"
        category_data = category.lower().replace(' ', '-')
        
        html_content += f"""
                <tr data-category="{category_data}" data-title="{doc.get('title', 'N/A').lower()}">
                    <td><strong>{doc.get('title', 'N/A')}</strong></td>
                    <td><span class="badge {badge_class}">{category}</span></td>
                    <td>{metrics.get('total_words', 0):,}</td>
                    <td>{metrics.get('total_sections', 0)}</td>
                    <td>{metrics.get('readability_score', 0):.1f}</td>
                    <td>{metrics.get('complexity_score', 0):.2f}</td>
                </tr>
"""
    
    # Contar categorías para el gráfico
    categories = Counter(d.get('category', 'Otros') for d in all_docs_data)
    category_js = "{" + ", ".join(f"'{k}': {v}" for k, v in categories.items()) + "}"
    
    # Datos de tamaño
    size_data = {doc.get('title', 'N/A').replace("'", "\\'"): doc.get('metrics', {}).get('total_words', 0) for doc in all_docs_data}
    size_js = "{" + ", ".join(f"'{k}': {v}" for k, v in list(size_data.items())[:10]) + "}"  # Primeros 10
    
    html_content += f"""
            </tbody>
        </table>
    </div>
    
    <script>
        new Chart(document.getElementById('categoryChart'), {{
            type: 'doughnut',
            data: {{
                labels: Object.keys({category_js}),
                datasets: [{{
                    data: Object.values({category_js}),
                    backgroundColor: ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E', '#7209B7', '#F72585']
                }}]
            }},
            options: {{ responsive: true, plugins: {{ legend: {{ position: 'bottom' }} }} }}
        }});
        
        new Chart(document.getElementById('sizeChart'), {{
            type: 'bar',
            data: {{
                labels: Object.keys({size_js}),
                datasets: [{{
                    label: 'Palabras',
                    data: Object.values({size_js}),
                    backgroundColor: '#2E86AB'
                }}]
            }},
            options: {{ responsive: true, scales: {{ y: {{ beginAtZero: true }} }} }}
        }});
        
        // Funciones de búsqueda y filtrado
        function filterTable() {{
            const input = document.getElementById('searchInput');
            const filter = input.value.toLowerCase();
            const table = document.getElementById('docsTable');
            const rows = table.getElementsByTagName('tr');
            
            for (let i = 1; i < rows.length; i++) {{
                const row = rows[i];
                const title = row.getAttribute('data-title') || '';
                if (title.includes(filter)) {{
                    row.style.display = '';
                }} else {{
                    row.style.display = 'none';
                }}
            }}
        }}
        
        function filterByCategory(category) {{
            // Actualizar botones activos
            document.querySelectorAll('.filter-btn').forEach(btn => {{
                btn.classList.remove('active');
            }});
            event.target.classList.add('active');
            
            const table = document.getElementById('docsTable');
            const rows = table.getElementsByTagName('tr');
            
            for (let i = 1; i < rows.length; i++) {{
                const row = rows[i];
                const rowCategory = row.getAttribute('data-category') || '';
                if (category === 'all' || rowCategory === category) {{
                    row.style.display = '';
                }} else {{
                    row.style.display = 'none';
                }}
            }}
        }}
    </script>
</body>
</html>
"""
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"✅ Dashboard HTML creado: {output_path}")


def create_comparison_report(all_docs_data: List[Dict], output_path: Path):
    """Crea un reporte comparativo entre documentos"""
    report = f"""
REPORTE COMPARATIVO - DOCUMENTOS BLATAM
{'=' * 70}

Fecha de Generación: {datetime.now().strftime('%d de %B de %Y %H:%M')}

ANÁLISIS COMPARATIVO
{'-' * 70}

DOCUMENTOS POR TAMAÑO (Palabras)
{'-' * 70}
"""
    
    # Ordenar por tamaño
    sorted_docs = sorted(all_docs_data, 
                        key=lambda x: x.get('metrics', {}).get('total_words', 0), 
                        reverse=True)
    
    for i, doc in enumerate(sorted_docs, 1):
        words = doc.get('metrics', {}).get('total_words', 0)
        report += f"{i:2d}. {doc.get('title', 'N/A'):<50} {words:>10,} palabras\n"
    
    report += f"""
DOCUMENTOS POR LEGIBILIDAD
{'-' * 70}
"""
    
    sorted_readability = sorted(all_docs_data,
                               key=lambda x: x.get('metrics', {}).get('readability_score', 0),
                               reverse=True)
    
    for i, doc in enumerate(sorted_readability, 1):
        readability = doc.get('metrics', {}).get('readability_score', 0)
        report += f"{i:2d}. {doc.get('title', 'N/A'):<50} {readability:>6.1f}/100\n"
    
    report += f"""
DOCUMENTOS POR COMPLEJIDAD
{'-' * 70}
"""
    
    sorted_complexity = sorted(all_docs_data,
                              key=lambda x: x.get('metrics', {}).get('complexity_score', 0),
                              reverse=True)
    
    for i, doc in enumerate(sorted_complexity, 1):
        complexity = doc.get('metrics', {}).get('complexity_score', 0)
        report += f"{i:2d}. {doc.get('title', 'N/A'):<50} {complexity:>6.2f}\n"
    
    report += f"""
RANKING GENERAL (Promedio de métricas normalizadas)
{'-' * 70}
"""
    
    # Calcular score general
    for doc in all_docs_data:
        metrics = doc.get('metrics', {})
        words = metrics.get('total_words', 0)
        sections = metrics.get('total_sections', 0)
        readability = metrics.get('readability_score', 0)
        complexity = metrics.get('complexity_score', 0)
        
        # Normalizar y promediar
        max_words = max(d.get('metrics', {}).get('total_words', 0) for d in all_docs_data) or 1
        max_sections = max(d.get('metrics', {}).get('total_sections', 0) for d in all_docs_data) or 1
        max_complexity = max(d.get('metrics', {}).get('complexity_score', 0) for d in all_docs_data) or 1
        
        score = (
            (words / max_words) * 0.3 +
            (sections / max_sections) * 0.2 +
            (readability / 100) * 0.3 +
            (complexity / max_complexity) * 0.2
        ) * 100
        
        doc['_comparison_score'] = score
    
    sorted_general = sorted(all_docs_data,
                           key=lambda x: x.get('_comparison_score', 0),
                           reverse=True)
    
    for i, doc in enumerate(sorted_general, 1):
        score = doc.get('_comparison_score', 0)
        report += f"{i:2d}. {doc.get('title', 'N/A'):<50} {score:>6.1f} pts\n"
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(report)
    
    print(f"✅ Reporte comparativo creado: {output_path}")


def create_executive_summary(all_docs_data: List[Dict], output_path: Path):
    """Crea un resumen ejecutivo en formato texto y JSON"""
    summary = {
        'generated_date': datetime.now().isoformat(),
        'total_documents': len(all_docs_data),
        'total_metrics': {
            'total_words': sum(d.get('metrics', {}).get('total_words', 0) for d in all_docs_data),
            'total_sections': sum(d.get('metrics', {}).get('total_sections', 0) for d in all_docs_data),
            'total_code_blocks': sum(len(d.get('code_blocks', [])) for d in all_docs_data),
            'total_links': sum(len(d.get('links', [])) for d in all_docs_data),
            'total_tables': sum(len(d.get('tables', [])) for d in all_docs_data),
        },
        'average_metrics': {
            'avg_words_per_doc': sum(d.get('metrics', {}).get('total_words', 0) for d in all_docs_data) / max(len(all_docs_data), 1),
            'avg_sections_per_doc': sum(d.get('metrics', {}).get('total_sections', 0) for d in all_docs_data) / max(len(all_docs_data), 1),
            'avg_readability': sum(d.get('metrics', {}).get('readability_score', 0) for d in all_docs_data) / max(len(all_docs_data), 1),
            'avg_complexity': sum(d.get('metrics', {}).get('complexity_score', 0) for d in all_docs_data) / max(len(all_docs_data), 1),
        },
        'categories': dict(Counter(d.get('category', 'Otros') for d in all_docs_data)),
    }
    
    # Guardar JSON
    json_path = output_path.with_suffix('.json')
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(summary, f, indent=2, ensure_ascii=False)
    
    # Guardar texto
    text_content = f"""
RESUMEN EJECUTIVO - DOCUMENTOS BLATAM
{'=' * 60}

Fecha de Generación: {datetime.now().strftime('%d de %B de %Y %H:%M')}

ESTADÍSTICAS GENERALES
{'-' * 60}
Total de Documentos Procesados: {summary['total_documents']}
Total de Palabras: {summary['total_metrics']['total_words']:,}
Total de Secciones: {summary['total_metrics']['total_sections']}
Total de Bloques de Código: {summary['total_metrics']['total_code_blocks']}
Total de Enlaces: {summary['total_metrics']['total_links']}
Total de Tablas: {summary['total_metrics']['total_tables']}

MÉTRICAS PROMEDIO
{'-' * 60}
Promedio de Palabras por Documento: {summary['average_metrics']['avg_words_per_doc']:,.0f}
Promedio de Secciones por Documento: {summary['average_metrics']['avg_sections_per_doc']:.1f}
Legibilidad Promedio: {summary['average_metrics']['avg_readability']:.1f}/100
Complejidad Promedio: {summary['average_metrics']['avg_complexity']:.2f}

DISTRIBUCIÓN POR CATEGORÍA
{'-' * 60}
"""
    
    for category, count in summary['categories'].items():
        text_content += f"{category}: {count} documentos\n"
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(text_content)
    
    print(f"✅ Resumen ejecutivo creado: {output_path}")
    print(f"✅ Resumen JSON creado: {json_path}")


def main():
    """Función principal"""
    print("=" * 70)
    print("🚀 Generador MEJORADO de Documentos Profesionales")
    print("   PDF, Word y Excel con gráficas avanzadas, portadas e índices")
    print("   + Dashboard HTML interactivo + Resumen Ejecutivo")
    print("=" * 70)
    
    # Verificar dependencias
    missing = []
    if not PDF_AVAILABLE:
        missing.append("reportlab")
    if not WORD_AVAILABLE:
        missing.append("python-docx")
    if not EXCEL_AVAILABLE:
        missing.append("openpyxl")
    if not MATPLOTLIB_AVAILABLE:
        missing.append("matplotlib numpy seaborn")
    
    if missing:
        print(f"\n⚠️  Faltan dependencias: {', '.join(missing)}")
        print("   Instala con: pip install " + " ".join(missing))
        print("\n   Continuando con las librerías disponibles...\n")
    
    base_path = Path(__file__).parent
    converter = EnhancedDocumentConverter(base_path)
    
    print(f"\n📚 Documentos a procesar: {len(IMPORTANT_DOCS)}\n")
    
    # Acumular datos de todos los documentos
    all_docs_data = []
    
    for doc_info in IMPORTANT_DOCS:
        doc_data = converter.convert_document(doc_info)
        if doc_data:
            all_docs_data.append(doc_data)
    
    # Crear dashboard HTML
    if all_docs_data:
        dashboard_path = OUTPUT_DIR / "dashboard_documentos.html"
        create_html_dashboard(all_docs_data, dashboard_path)
        
        # Crear resumen ejecutivo
        summary_path = OUTPUT_DIR / "resumen_ejecutivo.txt"
        create_executive_summary(all_docs_data, summary_path)
        
        # Crear reporte comparativo
        comparison_path = OUTPUT_DIR / "reporte_comparativo.txt"
        create_comparison_report(all_docs_data, comparison_path)
    
    print("\n" + "=" * 70)
    print("✅ Proceso completado!")
    print(f"📁 Archivos guardados en: {OUTPUT_DIR.absolute()}")
    print(f"📊 Dashboard HTML: {OUTPUT_DIR / 'dashboard_documentos.html'}")
    print(f"📄 Resumen Ejecutivo: {OUTPUT_DIR / 'resumen_ejecutivo.txt'}")
    print(f"📊 Reporte Comparativo: {OUTPUT_DIR / 'reporte_comparativo.txt'}")
    print("=" * 70)


if __name__ == "__main__":
    main()

