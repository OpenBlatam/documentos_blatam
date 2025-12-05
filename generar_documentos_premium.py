#!/usr/bin/env python3
"""
Script para generar documentos premium (PDF, Word, Excel) con gráficas
de alta calidad a partir de los documentos más importantes.
"""

import os
import re
from datetime import datetime
from pathlib import Path
import markdown
from collections import Counter, defaultdict
import json
import math

# Importar mejoras avanzadas
try:
    from mejoras_documentos import (
        calculate_readability, analyze_sentiment, analyze_structure,
        analyze_topics, generate_json_report, generate_html_dashboard
    )
    ENHANCEMENTS_AVAILABLE = True
except ImportError:
    ENHANCEMENTS_AVAILABLE = False

# Importar mejoras avanzadas adicionales
try:
    from mejoras_avanzadas import (
        analyze_coherence, analyze_link_structure, detect_patterns,
        analyze_trends, calculate_engagement_score, generate_markdown_report,
        export_to_csv, create_topic_network, analyze_document_quality
    )
    ADVANCED_ENHANCEMENTS_AVAILABLE = True
except ImportError:
    ADVANCED_ENHANCEMENTS_AVAILABLE = False

# Importar funcionalidades avanzadas
try:
    from funcionalidades_avanzadas import (
        analyze_seo, detect_common_errors, analyze_accessibility,
        generate_auto_index, suggest_improvements, calculate_overall_quality_score
    )
    ADVANCED_FEATURES_AVAILABLE = True
except ImportError:
    ADVANCED_FEATURES_AVAILABLE = False

# PDF
PDF_AVAILABLE = False
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    from reportlab.pdfgen import canvas
    PDF_AVAILABLE = True
except ImportError as e:
    PDF_AVAILABLE = False

# Word
try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.shared import RGBColor as DocxRGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    WORD_AVAILABLE = True
except ImportError:
    WORD_AVAILABLE = False
    print("⚠️  python-docx no disponible. Instala con: pip install python-docx")

# Excel
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.drawing.image import Image as XLImage
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("⚠️  openpyxl no disponible. Instala con: pip install openpyxl")

# PowerPoint
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx no disponible. Instala con: pip install python-pptx")

# CSV
CSV_AVAILABLE = True
import csv

# Gráficas
try:
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from matplotlib.backends.backend_pdf import PdfPages
    import numpy as np
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False
    print("⚠️  matplotlib no disponible. Instala con: pip install matplotlib numpy")

# Configuración
BASE_DIR = Path(__file__).parent
OUTPUT_DIR = BASE_DIR / "documentos_premium"
OUTPUT_DIR.mkdir(exist_ok=True)

# Documentos importantes a procesar
IMPORTANT_DOCS = [
    # Documentos principales del proyecto
    "ARCHITECTURE.md",
    "README.md",
    "airflow_automation_prompt.md",
    
    # Arquitectura y código de producción
    "truthgpt_collected/integration_code/production_code/ARCHITECTURE_IMPROVEMENTS.md",
    "truthgpt_collected/integration_code/production_code/REFACTORING_PLAN.md",
    "truthgpt_collected/integration_code/production_code/ARCHITECTURE.md",
    "truthgpt_collected/integration_code/production_code/README.md",
    "truthgpt_collected/integration_code/production_code/INDICE_DOCUMENTACION.md",
    
    # Documentación estratégica
    "06_documentation/resumen_final_completo.md",
    "06_documentation/Master_documents/indice_maestro_documentacion.md",
    "06_documentation/Other/Summaries/project_summary.md",
]

class DocumentProcessor:
    """Procesador de documentos con análisis y generación de gráficas"""
    
    def __init__(self):
        self.stats = {
            'total_words': 0,
            'total_sections': 0,
            'code_blocks': 0,
            'links': 0,
            'images': 0,
            'tables': 0,
            'files_processed': 0,
            'avg_words_per_section': 0,
            'code_density': 0,
            'readability_score': 0,
            'complexity_score': 0,
        }
        self.sections = []
        self.keywords = Counter()
        self.doc_names = []
        self.doc_stats = {}  # Estadísticas por documento
        self.topics = []  # Temas detectados
        self.sentiment_scores = []  # Puntuaciones de sentimiento
        self.code_quality_scores = []  # Puntuaciones de calidad de código
        self.structure_scores = []  # Puntuaciones de estructura
        self.summaries = {}  # Resúmenes automáticos por documento
        self.enhanced_stats = {}  # Estadísticas mejoradas por documento
        self.advanced_stats = {}  # Estadísticas avanzadas por documento
        self.seo_stats = {}  # Estadísticas SEO por documento
        self.error_reports = {}  # Reportes de errores por documento
        self.accessibility_stats = {}  # Estadísticas de accesibilidad
        self.improvement_suggestions = {}  # Sugerencias de mejora
        
    def analyze_document(self, content, doc_name=""):
        """Analiza el contenido del documento con métricas avanzadas"""
        # Contar palabras
        words = re.findall(r'\b\w+\b', content.lower())
        word_count = len(words)
        self.stats['total_words'] += word_count
        self.keywords.update(words)
        
        # Contar secciones
        sections = re.findall(r'^#+\s+(.+)$', content, re.MULTILINE)
        section_count = len(sections)
        self.stats['total_sections'] += section_count
        self.sections.extend(sections)
        
        # Contar bloques de código
        code_blocks = re.findall(r'```', content)
        code_count = len(code_blocks) // 2
        self.stats['code_blocks'] += code_count
        
        # Contar enlaces
        links = re.findall(r'\[([^\]]+)\]\([^\)]+\)', content)
        link_count = len(links)
        self.stats['links'] += link_count
        
        # Contar imágenes
        images = re.findall(r'!\[([^\]]*)\]\([^\)]+\)', content)
        image_count = len(images)
        self.stats['images'] += image_count
        
        # Contar tablas
        tables = re.findall(r'\|.*\|', content)
        table_count = len([t for t in tables if '---' not in t]) // 2
        self.stats['tables'] += table_count
        
        # Calcular métricas avanzadas
        avg_words = word_count / max(section_count, 1)
        code_density = (code_count / max(word_count, 1)) * 1000
        
        # Análisis de legibilidad (Flesch simplificado)
        sentences = re.split(r'[.!?]+', content)
        sentences = [s for s in sentences if len(s.strip()) > 10]
        avg_sentence_length = word_count / max(len(sentences), 1)
        readability = max(0, min(100, 206.835 - (1.015 * avg_sentence_length) - (84.6 * (len(re.findall(r'\b[aeiouáéíóúAEIOUÁÉÍÓÚ]\w*', content)) / max(word_count, 1)))))
        
        # Análisis de complejidad
        complexity = (section_count * 0.3) + (code_count * 0.4) + (link_count * 0.1) + (table_count * 0.2)
        
        # Detección de temas (palabras clave más frecuentes)
        doc_keywords = Counter(re.findall(r'\b\w{4,}\b', content.lower()))
        stop_words = {'este', 'esta', 'estos', 'estas', 'también', 'tambien', 'puede', 'pueden', 'debe', 'deben', 'ser', 'son', 'fue', 'fueron', 'hacer', 'hace', 'hacen', 'tiene', 'tienen', 'como', 'para', 'porque', 'cuando', 'donde', 'mientras', 'aunque', 'pero', 'sin', 'sobre', 'bajo', 'entre', 'hacia', 'hasta', 'desde', 'durante', 'mediante', 'según', 'contra', 'tras', 'ante', 'bajo', 'cabe', 'con', 'de', 'desde', 'en', 'entre', 'hacia', 'hasta', 'mediante', 'para', 'por', 'según', 'sin', 'so', 'sobre', 'tras', 'versus', 'vía'}
        filtered_keywords = Counter({k: v for k, v in doc_keywords.items() if k not in stop_words and len(k) > 3})
        top_topics = [word for word, _ in filtered_keywords.most_common(5)]
        
        # Análisis de sentimiento básico (palabras positivas/negativas)
        positive_words = ['excelente', 'bueno', 'mejor', 'óptimo', 'perfecto', 'éxito', 'logro', 'avance', 'mejora', 'solución', 'eficiente', 'efectivo', 'innovador', 'avanzado', 'completo', 'integral']
        negative_words = ['error', 'problema', 'fallo', 'defecto', 'limitación', 'restricción', 'dificultad', 'complejo', 'complicado', 'lento', 'ineficiente']
        positive_count = sum(1 for word in positive_words if word in content.lower())
        negative_count = sum(1 for word in negative_words if word in content.lower())
        sentiment = (positive_count - negative_count) / max(word_count / 100, 1)  # Normalizado
        
        # Análisis de calidad de código
        code_quality = 0
        if code_count > 0:
            # Verificar comentarios en código
            code_blocks = re.findall(r'```[\w]*\n(.*?)```', content, re.DOTALL)
            total_code_lines = sum(len(block.split('\n')) for block in code_blocks)
            comment_lines = sum(len(re.findall(r'#.*|//.*|/\*.*?\*/', block, re.DOTALL)) for block in code_blocks)
            comment_ratio = comment_lines / max(total_code_lines, 1)
            code_quality = min(100, (comment_ratio * 50) + (min(code_count, 20) * 2.5))  # Puntuación 0-100
        
        # Análisis de estructura del documento
        structure_score = 0
        if section_count > 0:
            # Verificar jerarquía correcta (H1 -> H2 -> H3)
            h1_count = len(re.findall(r'^#\s+', content, re.MULTILINE))
            h2_count = len(re.findall(r'^##\s+', content, re.MULTILINE))
            h3_count = len(re.findall(r'^###\s+', content, re.MULTILINE))
            hierarchy_score = min(100, (h1_count * 10) + (h2_count * 5) + (h3_count * 2))
            structure_score = min(100, (hierarchy_score * 0.4) + (section_count * 2) + (link_count * 0.5))
        
        # Generación de resumen automático (primeras 3 oraciones)
        sentences = re.split(r'[.!?]+', content)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
        summary = '. '.join(sentences[:3]) + '.' if sentences else "Resumen no disponible"
        
        # Detección de duplicación básica (frases repetidas)
        phrases = re.findall(r'\b\w{4,}\s+\w{4,}\s+\w{4,}\b', content.lower())
        phrase_counter = Counter(phrases)
        duplicate_phrases = [phrase for phrase, count in phrase_counter.items() if count > 3]
        duplication_ratio = len(duplicate_phrases) / max(len(phrases), 1)
        
        # Guardar estadísticas por documento
        if doc_name:
            self.doc_names.append(doc_name)
            self.doc_stats[doc_name] = {
                'words': word_count,
                'sections': section_count,
                'code_blocks': code_count,
                'links': link_count,
                'images': image_count,
                'tables': table_count,
                'avg_words_per_section': avg_words,
                'code_density': code_density,
                'readability': readability,
                'complexity': complexity,
                'topics': top_topics,
                'sentiment': sentiment,
                'code_quality': code_quality,
                'structure_score': structure_score,
                'duplication_ratio': duplication_ratio,
                'summary': summary[:200],  # Limitar longitud
            }
        
        # Acumular métricas globales
        self.stats['readability_score'] += readability
        self.stats['complexity_score'] += complexity
        self.topics.extend(top_topics)
        self.sentiment_scores.append(sentiment)
        self.code_quality_scores.append(code_quality)
        self.structure_scores.append(structure_score)
        self.summaries[doc_name] = summary[:200] if doc_name else None
        
        self.stats['files_processed'] += 1
        
        # Análisis mejorado si está disponible
        if ENHANCEMENTS_AVAILABLE and doc_name:
            try:
                readability = calculate_readability(content)
                sentiment = analyze_sentiment(content)
                structure = analyze_structure(content)
                topics = analyze_topics(content)
                
                self.enhanced_stats[doc_name] = {
                    'readability': readability,
                    'sentiment': sentiment,
                    'structure': structure,
                    'topics': topics
                }
            except Exception as e:
                print(f"  ⚠️  Error en análisis mejorado para {doc_name}: {e}")
        
        # Análisis avanzado adicional
        if ADVANCED_ENHANCEMENTS_AVAILABLE and doc_name:
            try:
                coherence = analyze_coherence(content)
                link_structure = analyze_link_structure(content)
                patterns = detect_patterns(content)
                trends = analyze_trends(content.split('\n'))
                
                # Calcular engagement score
                doc_stats_for_engagement = {
                    'total_words': word_count,
                    'total_sections': section_count,
                    'code_blocks': code_count,
                    'links': link_count,
                    'images': image_count,
                    'tables': table_count,
                    'readability_score': readability
                }
                engagement = calculate_engagement_score(content, doc_stats_for_engagement)
                
                # Análisis de calidad
                quality_stats = {
                    'readability_score': readability,
                    'links': link_count,
                    'code_blocks': code_count,
                    'tables': table_count,
                    'engagement_score': engagement
                }
                quality = analyze_document_quality(content, quality_stats)
                
                self.advanced_stats[doc_name] = {
                    'coherence': coherence,
                    'link_structure': link_structure,
                    'patterns': patterns,
                    'trends': trends,
                    'engagement_score': engagement,
                    'quality': quality
                }
            except Exception as e:
                print(f"  ⚠️  Error en análisis avanzado para {doc_name}: {e}")
        
        # Análisis de funcionalidades avanzadas (SEO, errores, accesibilidad)
        if ADVANCED_FEATURES_AVAILABLE and doc_name:
            try:
                # Extraer título del documento
                doc_title = ""
                first_line = content.split('\n')[0].strip()
                if first_line.startswith('#'):
                    doc_title = first_line.lstrip('#').strip()
                
                # Análisis SEO
                seo_analysis = analyze_seo(content, doc_title)
                self.seo_stats[doc_name] = seo_analysis
                
                # Detección de errores
                error_report = detect_common_errors(content)
                self.error_reports[doc_name] = error_report
                
                # Análisis de accesibilidad
                accessibility_analysis = analyze_accessibility(content)
                self.accessibility_stats[doc_name] = accessibility_analysis
                
                # Generar sugerencias de mejora
                doc_stats_for_suggestions = {
                    'total_words': word_count,
                    'total_sections': section_count,
                    'code_blocks': code_count,
                    'links': link_count,
                    'code_quality': code_quality,
                    'structure_score': structure_score
                }
                suggestions = suggest_improvements(
                    doc_stats_for_suggestions,
                    error_report,
                    seo_analysis,
                    accessibility_analysis
                )
                self.improvement_suggestions[doc_name] = suggestions
                
                # Calcular score de calidad general
                overall_quality = calculate_overall_quality_score(
                    doc_stats_for_suggestions,
                    error_report,
                    seo_analysis,
                    accessibility_analysis,
                    code_quality,
                    structure_score
                )
                
                # Agregar al advanced_stats
                if doc_name in self.advanced_stats:
                    self.advanced_stats[doc_name]['seo'] = seo_analysis
                    self.advanced_stats[doc_name]['errors'] = error_report
                    self.advanced_stats[doc_name]['accessibility'] = accessibility_analysis
                    self.advanced_stats[doc_name]['suggestions'] = suggestions
                    self.advanced_stats[doc_name]['overall_quality'] = overall_quality
                
            except Exception as e:
                print(f"  ⚠️  Error en funcionalidades avanzadas para {doc_name}: {e}")
        
    def generate_statistics_charts(self):
        """Genera gráficas de estadísticas mejoradas"""
        if not MATPLOTLIB_AVAILABLE:
            return None
            
        charts = {}
        # Intentar usar estilo mejorado, si no está disponible usar default
        try:
            plt.style.use('seaborn-v0_8-darkgrid')
        except:
            try:
                plt.style.use('seaborn-darkgrid')
            except:
                plt.style.use('default')
        
        # Gráfica 1: Distribución de contenido mejorada
        fig1, ax1 = plt.subplots(figsize=(12, 7))
        categories = ['Palabras\n(x1000)', 'Secciones', 'Bloques\nCódigo', 'Enlaces', 'Imágenes', 'Tablas']
        values = [
            self.stats['total_words'] / 1000,  # En miles
            self.stats['total_sections'],
            self.stats['code_blocks'],
            min(self.stats['links'], 100),
            self.stats['images'],
            self.stats['tables']
        ]
        colors_bar = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E', '#BC4749']
        bars = ax1.bar(categories, values, color=colors_bar, alpha=0.85, edgecolor='black', linewidth=2)
        ax1.set_ylabel('Cantidad', fontsize=13, fontweight='bold')
        ax1.set_title('Análisis Completo de Contenido de Documentos', fontsize=16, fontweight='bold', pad=25)
        ax1.grid(axis='y', alpha=0.4, linestyle='--', linewidth=1)
        
        # Añadir valores en las barras con formato mejorado
        for bar, value in zip(bars, values):
            height = bar.get_height()
            if value >= 1:
                label = f'{value:.1f}' if value < 10 else f'{int(value)}'
            else:
                label = f'{value:.2f}'
            ax1.text(bar.get_x() + bar.get_width()/2., height,
                    label,
                    ha='center', va='bottom', fontweight='bold', fontsize=11)
        
        # Añadir línea de promedio
        avg_value = sum(values) / len(values)
        ax1.axhline(y=avg_value, color='red', linestyle='--', linewidth=2, alpha=0.7, label=f'Promedio: {avg_value:.1f}')
        ax1.legend(loc='upper right', fontsize=10)
        
        plt.tight_layout()
        charts['content_distribution'] = fig1
        
        # Gráfica 2: Top keywords mejorada
        if self.keywords:
            fig2, ax2 = plt.subplots(figsize=(12, 8))
            top_keywords = dict(self.keywords.most_common(20))
            # Filtrar palabras comunes más completo
            stop_words = {'el', 'la', 'de', 'y', 'en', 'un', 'una', 'que', 'es', 'se', 'los', 'las', 'del', 'con', 'para', 'por', 'al', 'lo', 'le', 'da', 'su', 'sus', 'este', 'esta', 'como', 'más', 'muy', 'sin', 'sobre', 'también', 'pero', 'entre', 'hasta', 'desde', 'durante', 'mediante', 'según', 'contra', 'tras', 'ante', 'bajo', 'cabe', 'hacia', 'so', 'versus', 'vía', 'son', 'ser', 'tiene', 'tienen', 'fue', 'fueron', 'está', 'están', 'hacer', 'hace', 'hacen', 'puede', 'pueden', 'debe', 'deben', 'también', 'tambien', 'cada', 'todo', 'todos', 'toda', 'todas', 'otro', 'otra', 'otros', 'otras', 'mismo', 'misma', 'mismos', 'mismas'}
            filtered_keywords = {k: v for k, v in top_keywords.items() if k not in stop_words and len(k) > 3}
            if filtered_keywords:
                top_filtered = dict(list(filtered_keywords.items())[:15])
                words = list(top_filtered.keys())
                counts = list(top_filtered.values())
                # Colores degradados
                colors_gradient = plt.cm.viridis(np.linspace(0.2, 0.8, len(words)))
                bars = ax2.barh(words, counts, color=colors_gradient, alpha=0.85, edgecolor='black', linewidth=1.5)
                ax2.set_xlabel('Frecuencia de Aparición', fontsize=13, fontweight='bold')
                ax2.set_title('Top 15 Palabras Clave Más Relevantes', fontsize=16, fontweight='bold', pad=25)
                ax2.grid(axis='x', alpha=0.4, linestyle='--', linewidth=1)
                
                # Añadir valores
                for i, (bar, count) in enumerate(zip(bars, counts)):
                    ax2.text(count, bar.get_y() + bar.get_height()/2, 
                            f' {count}', va='center', fontweight='bold', fontsize=10)
                
                plt.tight_layout()
                charts['keywords'] = fig2
        
        # Gráfica 3: Distribución de secciones mejorada
        if self.sections:
            fig3, ax3 = plt.subplots(figsize=(14, 10))
            # Procesar secciones correctamente
            section_data = []
            for section in self.sections[:25]:  # Top 25
                clean_section = section.strip()
                if clean_section:
                    # Determinar nivel por número de #
                    level = clean_section.count('#') if clean_section.startswith('#') else 1
                    title = clean_section.lstrip('#').strip()[:60]  # Limitar longitud
                    section_data.append((title, level))
            
            if section_data:
                sections_list = [s[0] for s in section_data[:20]]
                levels_list = [s[1] for s in section_data[:20]]
                y_pos = np.arange(len(sections_list))
                colors_level = ['#FF6B6B', '#4ECDC4', '#45B7D1', '#FFA07A', '#98D8C8', '#F7DC6F']
                bars = ax3.barh(y_pos, levels_list, 
                               color=[colors_level[min(l-1, len(colors_level)-1)] for l in levels_list], 
                               alpha=0.85, edgecolor='black', linewidth=1.5)
                ax3.set_yticks(y_pos)
                ax3.set_yticklabels(sections_list, fontsize=10)
                ax3.set_xlabel('Nivel de Jerarquía', fontsize=13, fontweight='bold')
                ax3.set_title('Estructura Jerárquica de Secciones', fontsize=16, fontweight='bold', pad=25)
                ax3.set_xlim(0, max(levels_list) + 1)
                ax3.grid(axis='x', alpha=0.4, linestyle='--', linewidth=1)
                
                # Añadir valores
                for bar, level in zip(bars, levels_list):
                    ax3.text(level, bar.get_y() + bar.get_height()/2, 
                            f' Nivel {level}', va='center', fontweight='bold', fontsize=9)
                
                plt.tight_layout()
                charts['sections'] = fig3
        
        # Gráfica 4: Análisis de densidad de contenido
        fig4, ax4 = plt.subplots(figsize=(10, 6))
        metrics = ['Palabras', 'Secciones', 'Código', 'Enlaces']
        density_values = [
            self.stats['total_words'] / max(self.stats['total_sections'], 1),
            self.stats['total_sections'],
            self.stats['code_blocks'] * 10,  # Escalado
            self.stats['links'] * 2  # Escalado
        ]
        colors_density = ['#E74C3C', '#3498DB', '#2ECC71', '#F39C12']
        wedges, texts, autotexts = ax4.pie(density_values, labels=metrics, autopct='%1.1f%%',
                                          colors=colors_density, startangle=90,
                                          explode=(0.05, 0.05, 0.05, 0.05),
                                          shadow=True, textprops={'fontsize': 12, 'fontweight': 'bold'})
        ax4.set_title('Distribución de Tipos de Contenido', fontsize=16, fontweight='bold', pad=25)
        plt.tight_layout()
        charts['content_density'] = fig4
        
        # Gráfica 5: Comparativa por documento
        if len(self.doc_stats) > 1:
            fig5, ax5 = plt.subplots(figsize=(14, 8))
            doc_names_short = [name.split('/')[-1][:30] for name in self.doc_names]
            word_counts = [self.doc_stats[name]['words'] for name in self.doc_names]
            section_counts = [self.doc_stats[name]['sections'] for name in self.doc_names]
            
            x = np.arange(len(doc_names_short))
            width = 0.35
            
            bars1 = ax5.bar(x - width/2, [w/1000 for w in word_counts], width, 
                          label='Palabras (x1000)', color='#2E86AB', alpha=0.8)
            bars2 = ax5.bar(x + width/2, section_counts, width,
                          label='Secciones', color='#A23B72', alpha=0.8)
            
            ax5.set_xlabel('Documentos', fontsize=13, fontweight='bold')
            ax5.set_ylabel('Cantidad', fontsize=13, fontweight='bold')
            ax5.set_title('Comparativa de Documentos Procesados', fontsize=16, fontweight='bold', pad=25)
            ax5.set_xticks(x)
            ax5.set_xticklabels(doc_names_short, rotation=45, ha='right', fontsize=9)
            ax5.legend(fontsize=11)
            ax5.grid(axis='y', alpha=0.3, linestyle='--')
            plt.tight_layout()
            charts['doc_comparison'] = fig5
        
        # Gráfica 6: Análisis de complejidad
        if len(self.doc_stats) > 0:
            fig6, (ax6a, ax6b) = plt.subplots(1, 2, figsize=(16, 6))
            
            # Complejidad por densidad de código
            doc_names_short = [name.split('/')[-1][:20] for name in self.doc_names]
            code_densities = [self.doc_stats[name]['code_density'] for name in self.doc_names]
            avg_words = [self.doc_stats[name]['avg_words_per_section'] for name in self.doc_names]
            
            # Scatter plot: densidad de código vs palabras por sección
            scatter = ax6a.scatter(code_densities, avg_words, 
                                  s=[self.doc_stats[name]['words']/100 for name in self.doc_names],
                                  c=range(len(doc_names_short)), cmap='viridis', 
                                  alpha=0.6, edgecolors='black', linewidth=2)
            ax6a.set_xlabel('Densidad de Código', fontsize=12, fontweight='bold')
            ax6a.set_ylabel('Palabras por Sección', fontsize=12, fontweight='bold')
            ax6a.set_title('Análisis de Complejidad de Documentos', fontsize=14, fontweight='bold')
            ax6a.grid(alpha=0.3)
            plt.colorbar(scatter, ax=ax6a, label='Tamaño = Palabras/100')
            
            # Heatmap de métricas
            metrics_data = np.array([
                [self.doc_stats[name]['words']/1000 for name in self.doc_names],
                [self.doc_stats[name]['sections'] for name in self.doc_names],
                [self.doc_stats[name]['code_blocks'] for name in self.doc_names],
                [self.doc_stats[name]['links']/10 for name in self.doc_names],
            ])
            im = ax6b.imshow(metrics_data, cmap='YlOrRd', aspect='auto')
            ax6b.set_xticks(range(len(doc_names_short)))
            ax6b.set_xticklabels(doc_names_short, rotation=45, ha='right', fontsize=8)
            ax6b.set_yticks(range(4))
            ax6b.set_yticklabels(['Palabras (K)', 'Secciones', 'Código', 'Enlaces (x10)'], fontsize=10)
            ax6b.set_title('Heatmap de Métricas por Documento', fontsize=14, fontweight='bold')
            plt.colorbar(im, ax=ax6b)
            
            plt.tight_layout()
            charts['complexity_analysis'] = fig6
        
        # Gráfica 7: Análisis de legibilidad (si hay datos mejorados)
        if self.enhanced_stats and ENHANCEMENTS_AVAILABLE:
            try:
                fig7, (ax7a, ax7b) = plt.subplots(1, 2, figsize=(16, 6))
                
                # Distribución de scores de legibilidad
                flesch_scores = [s['readability']['flesch_score'] for s in self.enhanced_stats.values()]
                reading_levels = [s['readability']['reading_level'] for s in self.enhanced_stats.values()]
                
                if flesch_scores:
                    ax7a.hist(flesch_scores, bins=10, color='#2E86AB', alpha=0.7, edgecolor='black')
                    ax7a.axvline(sum(flesch_scores)/len(flesch_scores), color='red', linestyle='--', 
                                linewidth=2, label=f'Promedio: {sum(flesch_scores)/len(flesch_scores):.1f}')
                    ax7a.set_xlabel('Flesch Reading Ease Score', fontsize=12, fontweight='bold')
                    ax7a.set_ylabel('Frecuencia', fontsize=12, fontweight='bold')
                    ax7a.set_title('Distribución de Legibilidad', fontsize=14, fontweight='bold')
                    ax7a.legend()
                    ax7a.grid(alpha=0.3)
                    
                    # Niveles de lectura
                    level_counts = Counter(reading_levels)
                    ax7b.pie(level_counts.values(), labels=level_counts.keys(), autopct='%1.1f%%',
                            startangle=90, colors=plt.cm.Set3(range(len(level_counts))))
                    ax7b.set_title('Niveles de Lectura', fontsize=14, fontweight='bold')
                    
                    plt.tight_layout()
                    charts['readability_analysis'] = fig7
            except Exception as e:
                print(f"  ⚠️  Error generando gráfica de legibilidad: {e}")
        
        # Gráfica 8: Análisis de sentimiento
        if self.enhanced_stats and ENHANCEMENTS_AVAILABLE:
            try:
                fig8, ax8 = plt.subplots(figsize=(10, 6))
                
                sentiments = [s['sentiment']['sentiment'] for s in self.enhanced_stats.values()]
                sentiment_counts = Counter(sentiments)
                
                colors_sent = {'Positivo': '#4CAF50', 'Neutral': '#FFC107', 'Negativo': '#F44336'}
                bars = ax8.bar(sentiment_counts.keys(), 
                              sentiment_counts.values(),
                              color=[colors_sent.get(s, '#9E9E9E') for s in sentiment_counts.keys()],
                              alpha=0.8, edgecolor='black', linewidth=2)
                
                ax8.set_ylabel('Cantidad de Documentos', fontsize=12, fontweight='bold')
                ax8.set_title('Análisis de Sentimiento de Documentos', fontsize=14, fontweight='bold')
                ax8.grid(axis='y', alpha=0.3)
                
                for bar in bars:
                    height = bar.get_height()
                    ax8.text(bar.get_x() + bar.get_width()/2., height,
                            f'{int(height)}', ha='center', va='bottom', fontweight='bold')
                
                plt.tight_layout()
                charts['sentiment_analysis'] = fig8
            except Exception as e:
                print(f"  ⚠️  Error generando gráfica de sentimiento: {e}")
        
        # Gráfica 9: Análisis de coherencia (si hay datos avanzados)
        if self.advanced_stats and ADVANCED_ENHANCEMENTS_AVAILABLE:
            try:
                fig9, ax9 = plt.subplots(figsize=(12, 6))
                
                coherence_scores = [s['coherence']['coherence_score'] for s in self.advanced_stats.values()]
                doc_names_short = [name.split('/')[-1][:25] for name in self.advanced_stats.keys()]
                
                if coherence_scores:
                    bars = ax9.bar(range(len(coherence_scores)), coherence_scores,
                                  color=plt.cm.viridis(np.linspace(0.2, 0.8, len(coherence_scores))),
                                  alpha=0.8, edgecolor='black', linewidth=2)
                    ax9.set_xticks(range(len(doc_names_short)))
                    ax9.set_xticklabels(doc_names_short, rotation=45, ha='right', fontsize=9)
                    ax9.set_ylabel('Score de Coherencia', fontsize=12, fontweight='bold')
                    ax9.set_title('Análisis de Coherencia por Documento', fontsize=14, fontweight='bold')
                    ax9.set_ylim(0, 100)
                    ax9.axhline(y=50, color='red', linestyle='--', linewidth=2, alpha=0.7, label='Umbral (50)')
                    ax9.legend()
                    ax9.grid(axis='y', alpha=0.3)
                    
                    for bar, score in zip(bars, coherence_scores):
                        ax9.text(bar.get_x() + bar.get_width()/2., bar.get_height(),
                                f'{score:.1f}', ha='center', va='bottom', fontweight='bold')
                    
                    plt.tight_layout()
                    charts['coherence_analysis'] = fig9
            except Exception as e:
                print(f"  ⚠️  Error generando gráfica de coherencia: {e}")
        
        # Gráfica 10: Análisis de calidad
        if self.advanced_stats and ADVANCED_ENHANCEMENTS_AVAILABLE:
            try:
                fig10, (ax10a, ax10b) = plt.subplots(1, 2, figsize=(16, 6))
                
                quality_scores = [s['quality']['total_score'] for s in self.advanced_stats.values()]
                quality_factors = defaultdict(list)
                
                for stats in self.advanced_stats.values():
                    factors = stats['quality']['factors']
                    for factor, value in factors.items():
                        quality_factors[factor].append(value)
                
                # Gráfica de barras: scores de calidad
                doc_names_short = [name.split('/')[-1][:20] for name in self.advanced_stats.keys()]
                bars = ax10a.bar(range(len(quality_scores)), quality_scores,
                               color=plt.cm.RdYlGn(np.linspace(0.3, 0.7, len(quality_scores))),
                               alpha=0.8, edgecolor='black', linewidth=2)
                ax10a.set_xticks(range(len(doc_names_short)))
                ax10a.set_xticklabels(doc_names_short, rotation=45, ha='right', fontsize=8)
                ax10a.set_ylabel('Score de Calidad', fontsize=12, fontweight='bold')
                ax10a.set_title('Score de Calidad por Documento', fontsize=14, fontweight='bold')
                ax10a.set_ylim(0, 100)
                ax10a.grid(axis='y', alpha=0.3)
                
                for bar, score in zip(bars, quality_scores):
                    color = 'green' if score >= 80 else 'orange' if score >= 60 else 'red'
                    ax10a.text(bar.get_x() + bar.get_width()/2., bar.get_height(),
                              f'{score:.1f}', ha='center', va='bottom', fontweight='bold', color=color)
                
                # Gráfica de radar: factores de calidad promedio
                if quality_factors:
                    factors_list = list(quality_factors.keys())
                    avg_values = [sum(values) / len(values) for values in quality_factors.values()]
                    
                    # Crear gráfica de barras horizontal para factores
                    y_pos = np.arange(len(factors_list))
                    bars = ax10b.barh(y_pos, avg_values, color=plt.cm.Set3(range(len(factors_list))),
                                     alpha=0.8, edgecolor='black', linewidth=1.5)
                    ax10b.set_yticks(y_pos)
                    ax10b.set_yticklabels([f.replace('_', ' ').title() for f in factors_list], fontsize=10)
                    ax10b.set_xlabel('Score Promedio', fontsize=12, fontweight='bold')
                    ax10b.set_title('Factores de Calidad Promedio', fontsize=14, fontweight='bold')
                    ax10b.set_xlim(0, 20)
                    ax10b.grid(axis='x', alpha=0.3)
                    
                    for i, (bar, value) in enumerate(zip(bars, avg_values)):
                        ax10b.text(value, bar.get_y() + bar.get_height()/2,
                                  f' {value:.1f}', va='center', fontweight='bold')
                
                plt.tight_layout()
                charts['quality_analysis'] = fig10
            except Exception as e:
                print(f"  ⚠️  Error generando gráfica de calidad: {e}")
        
        return charts


class PDFGenerator:
    """Generador de documentos PDF de alta calidad"""
    
    def __init__(self, output_path):
        self.output_path = output_path
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
    
    def _clean_text(self, text):
        """Limpia y convierte texto markdown a texto plano seguro para ReportLab"""
        if not text:
            return ""
        
        # Limpiar emojis primero
        text = re.sub(r'[^\x00-\x7F]+', '', text)
        
        # Remover markdown básico y convertir a texto plano
        # Remover negrita
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'__(.+?)__', r'\1', text)
        
        # Remover cursiva
        text = re.sub(r'\*(.+?)\*', r'\1', text)
        text = re.sub(r'_(.+?)_', r'\1', text)
        
        # Remover código inline
        text = re.sub(r'`([^`]+)`', r'\1', text)
        
        # Escapar caracteres especiales
        text = text.replace('&', '&amp;')
        text = text.replace('<', '&lt;')
        text = text.replace('>', '&gt;')
        
        return text
        
    def _setup_custom_styles(self):
        """Configura estilos personalizados"""
        # Título principal
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#2E86AB'),
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ))
        
        # Subtítulo
        self.styles.add(ParagraphStyle(
            name='CustomHeading2',
            parent=self.styles['Heading2'],
            fontSize=18,
            textColor=colors.HexColor('#A23B72'),
            spaceAfter=12,
            spaceBefore=20,
            fontName='Helvetica-Bold'
        ))
        
        # Párrafo
        self.styles.add(ParagraphStyle(
            name='CustomBody',
            parent=self.styles['Normal'],
            fontSize=11,
            leading=14,
            alignment=TA_JUSTIFY,
            spaceAfter=12
        ))
    
    def generate(self, title, content, stats, charts=None):
        """Genera el PDF"""
        doc = SimpleDocTemplate(
            str(self.output_path),
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        story = []
        
        # Portada
        story.append(Paragraph(title, self.styles['CustomTitle']))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"Generado el: {datetime.now().strftime('%d/%m/%Y %H:%M')}", 
                              self.styles['Normal']))
        story.append(PageBreak())
        
        # Índice de estadísticas
        story.append(Paragraph("Estadisticas del Documento", self.styles['CustomHeading2']))
        story.append(Spacer(1, 0.2*inch))
        
        stats_data = [
            ['Métrica', 'Valor'],
            ['Total de Palabras', f"{stats['total_words']:,}"],
            ['Secciones', str(stats['total_sections'])],
            ['Bloques de Código', str(stats['code_blocks'])],
            ['Enlaces', str(stats['links'])],
            ['Imágenes', str(stats['images'])],
            ['Tablas', str(stats['tables'])],
        ]
        
        stats_table = Table(stats_data, colWidths=[3*inch, 2*inch])
        stats_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2E86AB')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
        ]))
        story.append(stats_table)
        story.append(PageBreak())
        
        # Gráficas si están disponibles
        if charts:
            story.append(Paragraph("Analisis Visual", self.styles['CustomHeading2']))
            story.append(Spacer(1, 0.2*inch))
            
            for chart_name, fig in charts.items():
                chart_path = OUTPUT_DIR / f"temp_chart_{chart_name}.png"
                fig.savefig(chart_path, dpi=300, bbox_inches='tight')
                img = Image(str(chart_path), width=6*inch, height=4*inch)
                story.append(img)
                story.append(Spacer(1, 0.3*inch))
                plt.close(fig)
        
        story.append(PageBreak())
        
        # Contenido del documento
        story.append(Paragraph("Contenido Completo", self.styles['CustomHeading2']))
        story.append(Spacer(1, 0.2*inch))
        
        # Procesar contenido markdown (simplificado y seguro)
        lines = content.split('\n')
        in_code_block = False
        
        for line in lines:
            original_line = line
            line = line.strip()
            
            # Manejar bloques de código
            if line.startswith('```'):
                in_code_block = not in_code_block
                continue
            if in_code_block:
                continue
            
            if not line:
                story.append(Spacer(1, 0.05*inch))
                continue
                
            # Títulos
            if line.startswith('# '):
                clean_title = self._clean_text(line[2:])
                story.append(Paragraph(clean_title, self.styles['Heading1']))
                story.append(Spacer(1, 0.1*inch))
            elif line.startswith('## '):
                clean_title = self._clean_text(line[3:])
                story.append(Paragraph(clean_title, self.styles['CustomHeading2']))
                story.append(Spacer(1, 0.1*inch))
            elif line.startswith('### '):
                clean_title = self._clean_text(line[4:])
                story.append(Paragraph(clean_title, self.styles['Heading3']))
                story.append(Spacer(1, 0.1*inch))
            # Listas
            elif line.startswith('- ') or line.startswith('* '):
                clean_text = self._clean_text(line[2:])
                story.append(Paragraph(f"• {clean_text}", self.styles['CustomBody']))
            # Párrafos normales
            elif not line.startswith('|') and not line.startswith('---') and not line.startswith('```'):
                clean_text = self._clean_text(line)
                if clean_text:
                    try:
                        story.append(Paragraph(clean_text, self.styles['CustomBody']))
                    except Exception as e:
                        # Si falla el parsing, usar texto plano
                        plain_text = clean_text.replace('<b>', '').replace('</b>', '').replace('<i>', '').replace('</i>', '').replace('<font', '').replace('</font>', '')
                        story.append(Paragraph(plain_text, self.styles['CustomBody']))
    
    def _clean_text(self, text):
        """Limpia y convierte texto markdown a texto plano seguro para ReportLab"""
        if not text:
            return ""
        
        # Limpiar emojis primero
        text = re.sub(r'[^\x00-\x7F]+', '', text)
        
        # Remover markdown básico y convertir a texto plano
        # Remover negrita
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'__(.+?)__', r'\1', text)
        
        # Remover cursiva
        text = re.sub(r'\*(.+?)\*', r'\1', text)
        text = re.sub(r'_(.+?)_', r'\1', text)
        
        # Remover código inline
        text = re.sub(r'`([^`]+)`', r'\1', text)
        
        # Escapar caracteres especiales
        text = text.replace('&', '&amp;')
        text = text.replace('<', '&lt;')
        text = text.replace('>', '&gt;')
        
        return text
    
    def generate(self, title, content, stats, charts=None):
        """Genera el PDF"""
        try:
            doc = SimpleDocTemplate(
                str(self.output_path),
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=18
            )
            
            story = []
            
            # Portada
            story.append(Paragraph(title, self.styles['CustomTitle']))
            story.append(Spacer(1, 0.5*inch))
            story.append(Paragraph(f"Generado el: {datetime.now().strftime('%d/%m/%Y %H:%M')}", 
                                  self.styles['Normal']))
            story.append(PageBreak())
            
            # Índice de estadísticas
            story.append(Paragraph("Estadisticas del Documento", self.styles['CustomHeading2']))
            story.append(Spacer(1, 0.2*inch))
            
            stats_data = [
                ['Métrica', 'Valor', 'Promedio'],
                ['Total de Palabras', f"{stats['total_words']:,}", f"{stats['total_words']/max(stats['files_processed'], 1):,.0f} por doc"],
                ['Secciones', str(stats['total_sections']), f"{stats['total_sections']/max(stats['files_processed'], 1):.1f} por doc"],
                ['Bloques de Código', str(stats['code_blocks']), f"{stats['code_blocks']/max(stats['files_processed'], 1):.1f} por doc"],
                ['Enlaces', str(stats['links']), f"{stats['links']/max(stats['files_processed'], 1):.1f} por doc"],
                ['Imágenes', str(stats['images']), f"{stats['images']/max(stats['files_processed'], 1):.1f} por doc"],
                ['Tablas', str(stats['tables']), f"{stats['tables']/max(stats['files_processed'], 1):.1f} por doc"],
                ['Documentos Procesados', str(stats.get('files_processed', 0)), ''],
                ['Palabras/Sección', f"{stats.get('avg_words_per_section', 0):.1f}", ''],
                ['Densidad de Código', f"{stats.get('code_density', 0):.2f}", 'por 1000 palabras'],
            ]
            
            stats_table = Table(stats_data, colWidths=[2.5*inch, 2*inch, 2*inch])
            stats_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2E86AB')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 11),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 9),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
            ]))
            story.append(stats_table)
            story.append(PageBreak())
            
            # Gráficas si están disponibles
            if charts:
                story.append(Paragraph("Analisis Visual", self.styles['CustomHeading2']))
                story.append(Spacer(1, 0.2*inch))
                
                for chart_name, fig in charts.items():
                    chart_path = OUTPUT_DIR / f"temp_chart_{chart_name}.png"
                    fig.savefig(chart_path, dpi=300, bbox_inches='tight')
                    img = Image(str(chart_path), width=6*inch, height=4*inch)
                    story.append(img)
                    story.append(Spacer(1, 0.3*inch))
                    plt.close(fig)
            
            story.append(PageBreak())
            
            # Contenido del documento (limitado para evitar PDFs muy grandes)
            story.append(Paragraph("Contenido Completo", self.styles['CustomHeading2']))
            story.append(Spacer(1, 0.2*inch))
            
            # Procesar contenido markdown (simplificado y seguro)
            lines = content.split('\n')
            in_code_block = False
            line_count = 0
            max_lines = 5000  # Límite para evitar PDFs muy grandes
            
            for line in lines:
                if line_count >= max_lines:
                    story.append(Paragraph("... (contenido adicional omitido por tamaño)", self.styles['Normal']))
                    break
                    
                original_line = line
                line = line.strip()
                
                # Manejar bloques de código
                if line.startswith('```'):
                    in_code_block = not in_code_block
                    continue
                if in_code_block:
                    continue
                
                if not line:
                    story.append(Spacer(1, 0.05*inch))
                    continue
                    
                # Títulos
                if line.startswith('# '):
                    clean_title = self._clean_text(line[2:])
                    story.append(Paragraph(clean_title, self.styles['Heading1']))
                    story.append(Spacer(1, 0.1*inch))
                elif line.startswith('## '):
                    clean_title = self._clean_text(line[3:])
                    story.append(Paragraph(clean_title, self.styles['CustomHeading2']))
                    story.append(Spacer(1, 0.1*inch))
                elif line.startswith('### '):
                    clean_title = self._clean_text(line[4:])
                    story.append(Paragraph(clean_title, self.styles['Heading3']))
                    story.append(Spacer(1, 0.1*inch))
                # Listas
                elif line.startswith('- ') or line.startswith('* '):
                    clean_text = self._clean_text(line[2:])
                    story.append(Paragraph(f"• {clean_text}", self.styles['CustomBody']))
                # Párrafos normales
                elif not line.startswith('|') and not line.startswith('---') and not line.startswith('```'):
                    clean_text = self._clean_text(line)
                    if clean_text and len(clean_text) > 0:
                        try:
                            story.append(Paragraph(clean_text[:500], self.styles['CustomBody']))  # Limitar longitud
                        except Exception as e:
                            # Si falla el parsing, usar texto plano
                            plain_text = clean_text[:500].replace('<', '').replace('>', '')
                            story.append(Paragraph(plain_text, self.styles['CustomBody']))
                
                line_count += 1
            
            doc.build(story)
            print(f"✅ PDF generado: {self.output_path}")
        except Exception as e:
            print(f"❌ Error generando PDF: {e}")
            import traceback
            traceback.print_exc()


class WordGenerator:
    """Generador de documentos Word de alta calidad"""
    
    def __init__(self, output_path):
        self.output_path = output_path
        self.doc = Document()
        self._setup_styles()
        
    def _setup_styles(self):
        """Configura estilos del documento"""
        # Estilo de título
        self.doc.styles['Title'].font.name = 'Calibri'
        self.doc.styles['Title'].font.size = Pt(24)
        self.doc.styles['Title'].font.color.rgb = DocxRGBColor(46, 134, 171)
        
        # Estilo de encabezado 1
        self.doc.styles['Heading 1'].font.name = 'Calibri'
        self.doc.styles['Heading 1'].font.size = Pt(18)
        self.doc.styles['Heading 1'].font.color.rgb = DocxRGBColor(162, 59, 114)
        
    def generate(self, title, content, stats, charts=None):
        """Genera el documento Word"""
        # Portada
        title_para = self.doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        date_para = self.doc.add_paragraph(f"Generado el: {datetime.now().strftime('%d/%m/%Y %H:%M')}")
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        self.doc.add_page_break()
        
        # Resumen Ejecutivo
        self.doc.add_heading('Resumen Ejecutivo', 1)
        exec_summary = self.doc.add_paragraph()
        exec_summary.add_run(f"Este documento contiene un análisis completo de {stats['total_sections']} secciones, "
                           f"con un total de {stats['total_words']:,} palabras. Incluye {stats['code_blocks']} bloques "
                           f"de código, {stats['links']} enlaces y {stats['tables']} tablas.")
        exec_summary.add_run().add_break()
        exec_summary.add_run(f"Fecha de generación: {datetime.now().strftime('%d de %B de %Y a las %H:%M')}")
        
        self.doc.add_page_break()
        
        # Estadísticas detalladas
        self.doc.add_heading('Estadisticas Detalladas del Documento', 1)
        
        stats_table = self.doc.add_table(rows=1, cols=3)
        stats_table.style = 'Light Grid Accent 1'
        hdr_cells = stats_table.rows[0].cells
        hdr_cells[0].text = 'Métrica'
        hdr_cells[1].text = 'Valor'
        hdr_cells[2].text = 'Promedio/Detalle'
        for cell in hdr_cells:
            cell.paragraphs[0].runs[0].font.bold = True
        
        total_items = sum([
            stats['total_sections'],
            stats['code_blocks'],
            stats['links'],
            stats['images'],
            stats['tables']
        ])
        files = max(stats.get('files_processed', 1), 1)
        
        for metric, value, detail in [
            ('Total de Palabras', stats['total_words'], f"{stats['total_words']/files:,.0f} por doc"),
            ('Secciones', stats['total_sections'], f"{stats['total_sections']/files:.1f} por doc"),
            ('Bloques de Código', stats['code_blocks'], f"{stats['code_blocks']/files:.1f} por doc"),
            ('Enlaces', stats['links'], f"{stats['links']/files:.1f} por doc"),
            ('Imágenes', stats['images'], f"{stats['images']/files:.1f} por doc"),
            ('Tablas', stats['tables'], f"{stats['tables']/files:.1f} por doc"),
            ('Documentos Procesados', stats.get('files_processed', 0), ''),
            ('Palabras por Sección', f"{stats.get('avg_words_per_section', 0):.1f}", ''),
            ('Densidad de Código', f"{stats.get('code_density', 0):.2f}", 'por 1000 palabras'),
        ]:
            row_cells = stats_table.add_row().cells
            row_cells[0].text = metric
            row_cells[1].text = f"{value:,}" if isinstance(value, (int, float)) and not isinstance(value, bool) else str(value)
            row_cells[2].text = detail
        
        self.doc.add_page_break()
        
        # Gráficas
        if charts:
            self.doc.add_heading('Analisis Visual', 1)
            for chart_name, fig in charts.items():
                chart_path = OUTPUT_DIR / f"temp_chart_{chart_name}.png"
                fig.savefig(chart_path, dpi=300, bbox_inches='tight')
                self.doc.add_picture(str(chart_path), width=Inches(6))
                self.doc.add_paragraph()  # Espacio
                plt.close(fig)
            self.doc.add_page_break()
        
        # Contenido
        self.doc.add_heading('Contenido Completo', 1)
        
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('# '):
                self.doc.add_heading(line[2:], 1)
            elif line.startswith('## '):
                self.doc.add_heading(line[3:], 2)
            elif line.startswith('### '):
                self.doc.add_heading(line[4:], 3)
            elif line.startswith('- ') or line.startswith('* '):
                self.doc.add_paragraph(line[2:], style='List Bullet')
            elif not line.startswith('```') and not line.startswith('|') and not line.startswith('---'):
                self.doc.add_paragraph(line)
        
        self.doc.save(self.output_path)
        print(f"✅ Word generado: {self.output_path}")


class HTMLDashboardGenerator:
    """Generador de Dashboard HTML Interactivo"""
    
    def __init__(self, output_path):
        self.output_path = output_path
    
    def generate(self, title, stats, doc_stats, charts=None):
        """Genera dashboard HTML interactivo"""
        # Preparar datos para JavaScript
        stats_json = json.dumps(stats, indent=2)
        doc_stats_json = json.dumps(doc_stats, indent=2)
        
        # Calcular métricas adicionales
        total_docs = stats.get('files_processed', 1)
        avg_words = stats['total_words'] / max(total_docs, 1)
        avg_sections = stats['total_sections'] / max(total_docs, 1)
        
        html_content = f"""<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Dashboard - {title}</title>
    <script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js"></script>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 20px;
            min-height: 100vh;
        }}
        .container {{
            max-width: 1600px;
            margin: 0 auto;
            background: white;
            border-radius: 20px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
            padding: 40px;
        }}
        header {{
            text-align: center;
            margin-bottom: 40px;
            border-bottom: 3px solid #667eea;
            padding-bottom: 20px;
        }}
        h1 {{
            color: #333;
            font-size: 2.5em;
            margin-bottom: 10px;
        }}
        .subtitle {{
            color: #666;
            font-size: 1.1em;
        }}
        .metrics-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin-bottom: 40px;
        }}
        .metric-card {{
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 25px;
            border-radius: 15px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.2);
            text-align: center;
            transition: transform 0.3s;
        }}
        .metric-card:hover {{
            transform: translateY(-5px);
        }}
        .metric-label {{
            font-size: 0.9em;
            opacity: 0.9;
            margin-bottom: 10px;
            text-transform: uppercase;
            letter-spacing: 1px;
        }}
        .metric-value {{
            font-size: 2.5em;
            font-weight: bold;
        }}
        .charts-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(500px, 1fr));
            gap: 30px;
            margin-bottom: 40px;
        }}
        .chart-card {{
            background: #f8f9fa;
            padding: 25px;
            border-radius: 15px;
            box-shadow: 0 5px 15px rgba(0,0,0,0.1);
        }}
        .chart-title {{
            font-size: 1.3em;
            color: #333;
            margin-bottom: 20px;
            font-weight: 600;
            text-align: center;
        }}
        .chart-container {{
            position: relative;
            height: 400px;
        }}
        .insights {{
            background: #e3f2fd;
            padding: 25px;
            border-radius: 15px;
            border-left: 5px solid #2196f3;
            margin-top: 30px;
        }}
        .insights h3 {{
            color: #1976d2;
            margin-bottom: 15px;
        }}
        .insights ul {{
            list-style: none;
            padding-left: 0;
        }}
        .insights li {{
            padding: 10px;
            margin: 5px 0;
            background: white;
            border-radius: 8px;
            border-left: 3px solid #2196f3;
        }}
        .footer {{
            text-align: center;
            color: #666;
            margin-top: 40px;
            padding-top: 20px;
            border-top: 1px solid #ddd;
        }}
    </style>
</head>
<body>
    <div class="container">
        <header>
            <h1>📊 Dashboard de Análisis de Documentos</h1>
            <p class="subtitle">{title}</p>
            <p class="subtitle">Generado el {datetime.now().strftime('%d/%m/%Y %H:%M')}</p>
        </header>
        
        <div class="metrics-grid">
            <div class="metric-card">
                <div class="metric-label">Total de Palabras</div>
                <div class="metric-value">{stats['total_words']:,}</div>
            </div>
            <div class="metric-card">
                <div class="metric-label">Secciones</div>
                <div class="metric-value">{stats['total_sections']}</div>
            </div>
            <div class="metric-card">
                <div class="metric-label">Documentos Procesados</div>
                <div class="metric-value">{total_docs}</div>
            </div>
            <div class="metric-card">
                <div class="metric-label">Promedio Palabras/Doc</div>
                <div class="metric-value">{avg_words:,.0f}</div>
            </div>
            <div class="metric-card">
                <div class="metric-label">Bloques de Código</div>
                <div class="metric-value">{stats['code_blocks']}</div>
            </div>
            <div class="metric-card">
                <div class="metric-label">Enlaces</div>
                <div class="metric-value">{stats['links']}</div>
            </div>
        </div>
        
        <div class="charts-grid">
            <div class="chart-card">
                <div class="chart-title">Distribución de Contenido</div>
                <div class="chart-container">
                    <canvas id="contentChart"></canvas>
                </div>
            </div>
            <div class="chart-card">
                <div class="chart-title">Comparativa de Documentos</div>
                <div class="chart-container">
                    <canvas id="docChart"></canvas>
                </div>
            </div>
        </div>
        
        <div class="insights">
            <h3>💡 Insights Clave</h3>
            <ul>
                <li>Se procesaron <strong>{total_docs}</strong> documentos con un total de <strong>{stats['total_words']:,}</strong> palabras</li>
                <li>Promedio de <strong>{avg_words:,.0f}</strong> palabras por documento</li>
                <li>Promedio de <strong>{avg_sections:.1f}</strong> secciones por documento</li>
                <li>Densidad de código: <strong>{stats.get('code_density', 0):.2f}</strong> bloques por 1000 palabras</li>
                <li>Total de <strong>{stats['links']}</strong> enlaces y <strong>{stats['images']}</strong> imágenes</li>
            </ul>
        </div>
        
        <div class="footer">
            <p>Dashboard generado automáticamente por generar_documentos_premium.py</p>
        </div>
    </div>
    
    <script>
        // Datos
        const stats = {stats_json};
        const docStats = {doc_stats_json};
        
        // Gráfico de distribución de contenido
        const contentCtx = document.getElementById('contentChart').getContext('2d');
        new Chart(contentCtx, {{
            type: 'bar',
            data: {{
                labels: ['Palabras (K)', 'Secciones', 'Código', 'Enlaces', 'Imágenes', 'Tablas'],
                datasets: [{{
                    label: 'Cantidad',
                    data: [
                        Math.round(stats.total_words / 1000),
                        stats.total_sections,
                        stats.code_blocks,
                        Math.min(stats.links, 100),
                        stats.images,
                        stats.tables
                    ],
                    backgroundColor: [
                        '#2E86AB',
                        '#A23B72',
                        '#F18F01',
                        '#C73E1D',
                        '#6A994E',
                        '#BC4749'
                    ],
                    borderColor: 'rgba(0,0,0,0.8)',
                    borderWidth: 2
                }}]
            }},
            options: {{
                responsive: true,
                maintainAspectRatio: false,
                plugins: {{
                    legend: {{ display: false }},
                    tooltip: {{
                        callbacks: {{
                            label: function(context) {{
                                return 'Cantidad: ' + context.parsed.y;
                            }}
                        }}
                    }}
                }},
                scales: {{
                    y: {{
                        beginAtZero: true,
                        ticks: {{
                            stepSize: 1
                        }}
                    }}
                }}
            }}
        }});
        
        // Gráfico comparativo de documentos
        const docCtx = document.getElementById('docChart').getContext('2d');
        const docNames = Object.keys(docStats);
        const docWords = docNames.map(name => docStats[name]?.words || 0);
        const docSections = docNames.map(name => docStats[name]?.sections || 0);
        
        new Chart(docCtx, {{
            type: 'line',
            data: {{
                labels: docNames.map(name => name.split('/').pop().substring(0, 20)),
                datasets: [{{
                    label: 'Palabras (x100)',
                    data: docWords.map(w => Math.round(w / 100)),
                    borderColor: '#2E86AB',
                    backgroundColor: 'rgba(46, 134, 171, 0.1)',
                    tension: 0.4,
                    fill: true
                }}, {{
                    label: 'Secciones',
                    data: docSections,
                    borderColor: '#A23B72',
                    backgroundColor: 'rgba(162, 59, 114, 0.1)',
                    tension: 0.4,
                    fill: true
                }}]
            }},
            options: {{
                responsive: true,
                maintainAspectRatio: false,
                plugins: {{
                    legend: {{ display: true }},
                    tooltip: {{
                        callbacks: {{
                            label: function(context) {{
                                if (context.datasetIndex === 0) {{
                                    return 'Palabras: ' + (context.parsed.y * 100).toLocaleString();
                                }}
                                return 'Secciones: ' + context.parsed.y;
                            }}
                        }}
                    }}
                }},
                scales: {{
                    y: {{
                        beginAtZero: true
                    }},
                    x: {{
                        ticks: {{
                            maxRotation: 45,
                            minRotation: 45
                        }}
                    }}
                }}
            }}
        }});
    </script>
</body>
</html>"""
        
        with open(self.output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        print(f"✅ Dashboard HTML generado: {self.output_path}")


class ExcelGenerator:
    """Generador de documentos Excel con gráficas"""
    
    def __init__(self, output_path):
        self.output_path = output_path
        self.wb = Workbook()
        self.ws = self.wb.active
        self.ws.title = "Resumen"
        
    def generate(self, title, content, stats, charts=None):
        """Genera el documento Excel"""
        # Estilos
        header_fill = PatternFill(start_color="2E86AB", end_color="2E86AB", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF", size=14)
        title_font = Font(bold=True, size=16, color="2E86AB")
        
        # Título
        self.ws['A1'] = title
        self.ws['A1'].font = title_font
        self.ws.merge_cells('A1:B1')
        self.ws['A2'] = f"Generado el: {datetime.now().strftime('%d/%m/%Y %H:%M')}"
        
        # Estadísticas
        row = 4
        self.ws[f'A{row}'] = "Estadisticas del Documento"
        self.ws[f'A{row}'].font = Font(bold=True, size=14)
        row += 2
        
        # Encabezados
        self.ws[f'A{row}'] = "Métrica"
        self.ws[f'B{row}'] = "Valor"
        self.ws[f'A{row}'].fill = header_fill
        self.ws[f'B{row}'].fill = header_fill
        self.ws[f'A{row}'].font = header_font
        self.ws[f'B{row}'].font = header_font
        row += 1
        
        # Datos
        stats_data = [
            ('Total de Palabras', stats['total_words']),
            ('Secciones', stats['total_sections']),
            ('Bloques de Código', stats['code_blocks']),
            ('Enlaces', stats['links']),
            ('Imágenes', stats['images']),
            ('Tablas', stats['tables']),
        ]
        
        for metric, value in stats_data:
            self.ws[f'A{row}'] = metric
            self.ws[f'B{row}'] = value
            row += 1
        
        # Ajustar columnas
        self.ws.column_dimensions['A'].width = 25
        self.ws.column_dimensions['B'].width = 15
        
        # Gráfica de barras mejorada
        chart = BarChart()
        chart.type = "col"
        chart.style = 10
        chart.title = "Estadísticas del Documento - Análisis Completo"
        chart.y_axis.title = 'Cantidad'
        chart.x_axis.title = 'Métricas'
        chart.legend = None
        
        data = Reference(self.ws, min_col=2, min_row=6, max_row=11)
        cats = Reference(self.ws, min_col=1, min_row=6, max_row=11)
        chart.add_data(data, titles_from_data=False)
        chart.set_categories(cats)
        chart.height = 12
        chart.width = 18
        
        # Añadir colores
        s1 = chart.series[0]
        s1.graphicalProperties.solidFill = "2E86AB"
        
        self.ws.add_chart(chart, "D4")
        
        # Añadir gráfica de pastel
        pie_chart = PieChart()
        pie_chart.title = "Distribución de Contenido"
        pie_data = Reference(self.ws, min_col=2, min_row=6, max_row=9)
        pie_labels = Reference(self.ws, min_col=1, min_row=6, max_row=9)
        pie_chart.add_data(pie_data, titles_from_data=False)
        pie_chart.set_categories(pie_labels)
        pie_chart.height = 8
        pie_chart.width = 10
        self.ws.add_chart(pie_chart, "D20")
        
        # Hoja comparativa de documentos (si hay múltiples)
        if hasattr(stats, 'doc_stats') and stats.get('doc_stats'):
            ws_comparison = self.wb.create_sheet("Comparativa Docs")
            ws_comparison['A1'] = "Comparativa de Documentos"
            ws_comparison['A1'].font = title_font
            
            # Encabezados
            headers = ['Documento', 'Palabras', 'Secciones', 'Código', 'Enlaces', 'Palabras/Sección', 'Densidad Código']
            for col, header in enumerate(headers, 1):
                cell = ws_comparison.cell(row=3, column=col)
                cell.value = header
                cell.fill = header_fill
                cell.font = header_font
            
            # Datos (si están disponibles en stats)
            row = 4
            # Nota: Esto requeriría pasar doc_stats, por ahora dejamos la estructura
            
            ws_comparison.column_dimensions['A'].width = 40
            for col in range(2, 8):
                ws_comparison.column_dimensions[chr(64+col)].width = 15
        
        # Hoja de contenido
        ws_content = self.wb.create_sheet("Contenido")
        ws_content['A1'] = "Contenido del Documento"
        ws_content['A1'].font = title_font
        
        # Procesar contenido (primeras 1000 líneas para Excel)
        lines = content.split('\n')[:1000]
        row = 3
        for line in lines:
            line = line.strip()
            if line and not line.startswith('```'):
                ws_content[f'A{row}'] = line[:32767]  # Límite de Excel
                row += 1
                if row > 10000:  # Límite razonable
                    break
        
        ws_content.column_dimensions['A'].width = 100
        
        # Hoja de resumen ejecutivo
        ws_summary = self.wb.create_sheet("Resumen", 0)  # Primera hoja
        ws_summary['A1'] = "RESUMEN EJECUTIVO"
        ws_summary['A1'].font = Font(bold=True, size=18, color="2E86AB")
        ws_summary.merge_cells('A1:C1')
        
        ws_summary['A3'] = f"Documentos Analizados: {stats.get('files_processed', 0)}"
        ws_summary['A3'].font = Font(bold=True, size=12)
        ws_summary['A4'] = f"Total de Palabras: {stats['total_words']:,}"
        ws_summary['A5'] = f"Total de Secciones: {stats['total_sections']}"
        ws_summary['A6'] = f"Bloques de Código: {stats['code_blocks']}"
        ws_summary['A7'] = f"Enlaces: {stats['links']}"
        ws_summary['A8'] = f"Fecha de Generación: {datetime.now().strftime('%d/%m/%Y %H:%M')}"
        
        ws_summary.column_dimensions['A'].width = 35
        
        # Guardar
        self.wb.save(self.output_path)
        print(f"✅ Excel generado: {self.output_path}")


class PowerPointGenerator:
    """Generador de presentaciones PowerPoint"""
    
    def __init__(self, output_path):
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def generate(self, title, stats, doc_stats, charts=None):
        """Genera presentación PowerPoint"""
        # Slide 1: Portada
        slide1 = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        title_shape = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
        
        subtitle_shape = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = f"Generado el {datetime.now().strftime('%d/%m/%Y %H:%M')}"
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        
        # Slide 2: Estadísticas principales
        slide2 = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        title_shape2 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        title_frame2 = title_shape2.text_frame
        title_frame2.text = "Estadísticas Principales"
        title_frame2.paragraphs[0].font.size = Pt(36)
        title_frame2.paragraphs[0].font.bold = True
        title_frame2.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
        
        content_shape = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = content_shape.text_frame
        tf.text = f"Total de Palabras: {stats['total_words']:,}"
        p = tf.add_paragraph()
        p.text = f"Secciones: {stats['total_sections']}"
        p.space_after = Pt(12)
        p = tf.add_paragraph()
        p.text = f"Documentos Procesados: {stats.get('files_processed', 0)}"
        p.space_after = Pt(12)
        p = tf.add_paragraph()
        p.text = f"Bloques de Código: {stats['code_blocks']}"
        p.space_after = Pt(12)
        p = tf.add_paragraph()
        p.text = f"Enlaces: {stats['links']}"
        p.space_after = Pt(12)
        p = tf.add_paragraph()
        p.text = f"Legibilidad: {stats.get('readability_score', 0):.1f}/100"
        p.space_after = Pt(12)
        for para in tf.paragraphs:
            para.font.size = Pt(18)
        
        # Slide 3: Análisis avanzado
        if stats.get('main_topics'):
            slide3 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            title_shape3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
            title_frame3 = title_shape3.text_frame
            title_frame3.text = "Temas Principales"
            title_frame3.paragraphs[0].font.size = Pt(36)
            title_frame3.paragraphs[0].font.bold = True
            title_frame3.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
            
            content_shape3 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf3 = content_shape3.text_frame
            for i, topic in enumerate(stats['main_topics'][:10], 1):
                if i == 1:
                    tf3.text = f"{i}. {topic}"
                else:
                    p = tf3.add_paragraph()
                    p.text = f"{i}. {topic}"
                p.space_after = Pt(10)
            for para in tf3.paragraphs:
                para.font.size = Pt(18)
        
        # Slide 4: Comparativa de documentos
        if len(doc_stats) > 1:
            slide4 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            title_shape4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
            title_frame4 = title_shape4.text_frame
            title_frame4.text = "Comparativa de Documentos"
            title_frame4.paragraphs[0].font.size = Pt(36)
            title_frame4.paragraphs[0].font.bold = True
            title_frame4.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
            
            content_shape4 = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf4 = content_shape4.text_frame
            for i, (doc_name, doc_data) in enumerate(list(doc_stats.items())[:8], 1):
                short_name = doc_name.split('/')[-1][:40]
                if i == 1:
                    tf4.text = f"{short_name}: {doc_data.get('words', 0):,} palabras"
                else:
                    p = tf4.add_paragraph()
                    p.text = f"{short_name}: {doc_data.get('words', 0):,} palabras"
                p.space_after = Pt(10)
            for para in tf4.paragraphs:
                para.font.size = Pt(16)
        
        # Slide 5: Análisis de Calidad
        if stats.get('avg_sentiment') is not None or stats.get('readability_score', 0) > 0:
            slide5 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            title_shape5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
            title_frame5 = title_shape5.text_frame
            title_frame5.text = "Análisis de Calidad"
            title_frame5.paragraphs[0].font.size = Pt(36)
            title_frame5.paragraphs[0].font.bold = True
            title_frame5.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
            
            content_shape5 = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf5 = content_shape5.text_frame
            
            if stats.get('readability_score', 0) > 0:
                readability_level = "Excelente" if stats['readability_score'] >= 70 else "Bueno" if stats['readability_score'] >= 50 else "Mejorable"
                tf5.text = f"Legibilidad: {stats['readability_score']:.1f}/100 ({readability_level})"
                p = tf5.add_paragraph()
                p.space_after = Pt(12)
            
            if stats.get('avg_sentiment') is not None:
                sentiment_text = "Positivo" if stats['avg_sentiment'] > 0.6 else "Neutral" if stats['avg_sentiment'] > 0.4 else "Negativo"
                p = tf5.add_paragraph()
                p.text = f"Sentimiento Promedio: {sentiment_text} ({stats['avg_sentiment']:.2f})"
                p.space_after = Pt(12)
            
            if stats.get('complexity_score', 0) > 0:
                p = tf5.add_paragraph()
                p.text = f"Complejidad: {stats['complexity_score']:.1f}"
                p.space_after = Pt(12)
            
            if stats.get('code_density', 0) > 0:
                p = tf5.add_paragraph()
                p.text = f"Densidad de Código: {stats['code_density']:.2f} bloques/1000 palabras"
                p.space_after = Pt(12)
            
            for para in tf5.paragraphs:
                para.font.size = Pt(18)
        
        # Slide 6: Insights y Recomendaciones
        slide6 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        title_shape6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        title_frame6 = title_shape6.text_frame
        title_frame6.text = "Insights Clave"
        title_frame6.paragraphs[0].font.size = Pt(36)
        title_frame6.paragraphs[0].font.bold = True
        title_frame6.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
        
        content_shape6 = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf6 = content_shape6.text_frame
        
        insights = []
        if stats['total_words'] > 10000:
            insights.append("✓ Documento extenso con contenido completo")
        if stats['total_sections'] > 10:
            insights.append("✓ Bien estructurado con múltiples secciones")
        if stats['code_blocks'] > 5:
            insights.append("✓ Incluye ejemplos de código prácticos")
        if stats['links'] > 10:
            insights.append("✓ Buena interconexión con referencias")
        if stats.get('readability_score', 0) >= 60:
            insights.append("✓ Buena legibilidad para la audiencia")
        if not insights:
            insights.append("Documento en proceso de desarrollo")
            insights.append("Considerar agregar más contenido y estructura")
        
        for i, insight in enumerate(insights[:6], 1):
            if i == 1:
                tf6.text = insight
            else:
                p = tf6.add_paragraph()
                p.text = insight
            p.space_after = Pt(10)
        
        for para in tf6.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(50, 50, 50)
        
        # Slide 7: Métricas por Documento (si hay múltiples)
        if len(doc_stats) > 1:
            slide7 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            title_shape7 = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
            title_frame7 = title_shape7.text_frame
            title_frame7.text = "Métricas Detalladas por Documento"
            title_frame7.paragraphs[0].font.size = Pt(32)
            title_frame7.paragraphs[0].font.bold = True
            title_frame7.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
            
            # Crear tabla con métricas
            from pptx.util import Inches as PPTXInches
            rows = min(len(doc_stats) + 1, 9)  # Header + máximo 8 documentos
            cols = 4
            
            left = PPTXInches(0.5)
            top = PPTXInches(2)
            width = PPTXInches(9)
            height = PPTXInches(4.5)
            
            table = slide7.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Encabezados
            table.cell(0, 0).text = "Documento"
            table.cell(0, 1).text = "Palabras"
            table.cell(0, 2).text = "Secciones"
            table.cell(0, 3).text = "Código"
            
            # Estilo de encabezados
            for col in range(cols):
                cell = table.cell(0, col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(46, 134, 171)
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.size = Pt(12)
            
            # Datos
            for idx, (doc_name, doc_data) in enumerate(list(doc_stats.items())[:rows-1], 1):
                short_name = doc_name.split('/')[-1][:25]
                table.cell(idx, 0).text = short_name
                table.cell(idx, 1).text = f"{doc_data.get('words', 0):,}"
                table.cell(idx, 2).text = str(doc_data.get('sections', 0))
                table.cell(idx, 3).text = str(doc_data.get('code_blocks', 0))
                
                # Estilo de celdas
                for col in range(cols):
                    cell = table.cell(idx, col)
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(10)
                    if idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        # Slide final: Resumen
        slide_final = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        title_shape_final = slide_final.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame_final = title_shape_final.text_frame
        title_frame_final.text = "Gracias"
        title_frame_final.paragraphs[0].font.size = Pt(48)
        title_frame_final.paragraphs[0].font.bold = True
        title_frame_final.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
        title_frame_final.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        subtitle_shape = slide_final.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = f"Análisis de {stats.get('files_processed', 0)} documentos completado"
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        date_shape = slide_final.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
        date_frame = date_shape.text_frame
        date_frame.text = f"Generado el {datetime.now().strftime('%d de %B de %Y')}"
        date_frame.paragraphs[0].font.size = Pt(14)
        date_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        self.prs.save(self.output_path)
        print(f"✅ PowerPoint generado: {self.output_path}")


class MarkdownGenerator:
    """Generador de documentos Markdown estructurados"""
    
    def __init__(self, output_path):
        self.output_path = output_path
    
    def generate(self, title, stats, doc_stats, summaries):
        """Genera documento Markdown estructurado"""
        md_content = f"""# {title}

**Generado el:** {datetime.now().strftime('%d/%m/%Y %H:%M')}

---

## 📊 Resumen Ejecutivo

Este documento contiene un análisis completo de {stats.get('files_processed', 0)} documentos con un total de {stats['total_words']:,} palabras.

### Métricas Principales

- **Total de Palabras:** {stats['total_words']:,}
- **Secciones:** {stats['total_sections']}
- **Documentos Procesados:** {stats.get('files_processed', 0)}
- **Legibilidad Promedio:** {stats.get('readability_score', 0):.1f}/100
- **Complejidad Promedio:** {stats.get('complexity_score', 0):.2f}
- **Calidad de Código:** {stats.get('avg_code_quality', 0):.1f}/100
- **Puntuación de Estructura:** {stats.get('avg_structure_score', 0):.1f}/100

---

## 📈 Análisis por Documento

"""
        
        for doc_name, doc_data in doc_stats.items():
            short_name = doc_name.split('/')[-1]
            summary = summaries.get(doc_name, 'No disponible')
            md_content += f"""### {short_name}

**Resumen:** {summary[:200] if len(summary) > 200 else summary}

**Estadísticas:**
- Palabras: {doc_data.get('words', 0):,}
- Secciones: {doc_data.get('sections', 0)}
- Bloques de Código: {doc_data.get('code_blocks', 0)}
- Enlaces: {doc_data.get('links', 0)}
- Legibilidad: {doc_data.get('readability', 0):.1f}/100
- Complejidad: {doc_data.get('complexity', 0):.2f}
- Calidad de Código: {doc_data.get('code_quality', 0):.1f}/100
- Estructura: {doc_data.get('structure_score', 0):.1f}/100
- Sentimiento: {doc_data.get('sentiment', 0):.4f}

**Temas Principales:** {', '.join(doc_data.get('topics', [])[:5]) if doc_data.get('topics') else 'N/A'}

---

"""
        
        md_content += f"""## 🎯 Temas Principales Globales

"""
        for i, topic in enumerate(stats.get('main_topics', [])[:10], 1):
            md_content += f"{i}. {topic}\n"
        
        md_content += f"""
---

## 📋 Recomendaciones

### Mejoras Sugeridas

1. **Legibilidad:** {'Mejorar' if stats.get('readability_score', 0) < 60 else 'Mantener'} la legibilidad de los documentos
2. **Estructura:** {'Mejorar' if stats.get('avg_structure_score', 0) < 70 else 'Mantener'} la estructura jerárquica
3. **Calidad de Código:** {'Añadir más comentarios' if stats.get('avg_code_quality', 0) < 50 else 'Mantener'} en los bloques de código
4. **Complejidad:** {'Simplificar' if stats.get('complexity_score', 0) > 50 else 'Mantener'} la complejidad de los documentos

---

**Generado automáticamente por:** `generar_documentos_premium.py`
"""
        
        with open(self.output_path, 'w', encoding='utf-8') as f:
            f.write(md_content)
        print(f"✅ Markdown generado: {self.output_path}")


class CSVGenerator:
    """Generador de archivos CSV con datos estructurados"""
    
    def __init__(self, output_path):
        self.output_path = output_path
    
    def generate(self, stats, doc_stats):
        """Genera archivo CSV con estadísticas"""
        with open(self.output_path, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile)
            
            # Encabezado de estadísticas generales
            writer.writerow(['Métrica', 'Valor'])
            writer.writerow(['Total de Palabras', stats['total_words']])
            writer.writerow(['Total de Secciones', stats['total_sections']])
            writer.writerow(['Bloques de Código', stats['code_blocks']])
            writer.writerow(['Enlaces', stats['links']])
            writer.writerow(['Imágenes', stats['images']])
            writer.writerow(['Tablas', stats['tables']])
            writer.writerow(['Documentos Procesados', stats.get('files_processed', 0)])
            writer.writerow(['Palabras por Sección', f"{stats.get('avg_words_per_section', 0):.2f}"])
            writer.writerow(['Densidad de Código', f"{stats.get('code_density', 0):.4f}"])
            writer.writerow(['Legibilidad', f"{stats.get('readability_score', 0):.2f}"])
            writer.writerow(['Complejidad', f"{stats.get('complexity_score', 0):.2f}"])
            writer.writerow(['Sentimiento Promedio', f"{stats.get('avg_sentiment', 0):.4f}"])
            writer.writerow([])
            
            # Estadísticas por documento
            writer.writerow(['Documento', 'Palabras', 'Secciones', 'Código', 'Enlaces', 'Legibilidad', 'Complejidad', 'Sentimiento', 'Temas'])
            for doc_name, doc_data in doc_stats.items():
                short_name = doc_name.split('/')[-1]
                topics_str = ', '.join(doc_data.get('topics', [])[:3])
                writer.writerow([
                    short_name,
                    doc_data.get('words', 0),
                    doc_data.get('sections', 0),
                    doc_data.get('code_blocks', 0),
                    doc_data.get('links', 0),
                    f"{doc_data.get('readability', 0):.2f}",
                    f"{doc_data.get('complexity', 0):.2f}",
                    f"{doc_data.get('sentiment', 0):.4f}",
                    topics_str
                ])
        
        print(f"✅ CSV generado: {self.output_path}")


def main():
    """Función principal"""
    print("🚀 Iniciando generación de documentos premium...")
    print(f"📁 Directorio de salida: {OUTPUT_DIR}")
    
    processor = DocumentProcessor()
    all_content = []
    
    # Procesar documentos
    for doc_path in IMPORTANT_DOCS:
        full_path = BASE_DIR / doc_path
        if full_path.exists():
            print(f"📖 Procesando: {doc_path}")
            try:
                with open(full_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    processor.analyze_document(content, doc_path)
                    all_content.append(f"\n\n# {doc_path}\n\n{content}")
            except Exception as e:
                print(f"⚠️  Error procesando {doc_path}: {e}")
        else:
            print(f"⚠️  No encontrado: {doc_path}")
    
    # Calcular métricas finales
    if processor.stats['total_sections'] > 0:
        processor.stats['avg_words_per_section'] = processor.stats['total_words'] / processor.stats['total_sections']
    if processor.stats['total_words'] > 0:
        processor.stats['code_density'] = (processor.stats['code_blocks'] / processor.stats['total_words']) * 1000
    if processor.stats['files_processed'] > 0:
        processor.stats['readability_score'] = processor.stats['readability_score'] / processor.stats['files_processed']
        processor.stats['complexity_score'] = processor.stats['complexity_score'] / processor.stats['files_processed']
    
    # Calcular temas principales globales
    topic_counter = Counter(processor.topics)
    processor.stats['main_topics'] = [topic for topic, _ in topic_counter.most_common(10)]
    
    # Calcular sentimiento promedio
    if processor.sentiment_scores:
        processor.stats['avg_sentiment'] = sum(processor.sentiment_scores) / len(processor.sentiment_scores)
    else:
        processor.stats['avg_sentiment'] = 0
    
    if not all_content:
        print("❌ No se encontraron documentos para procesar")
        return
    
    combined_content = "\n".join(all_content)
    combined_title = "Documentos BLATAM - Resumen Ejecutivo"
    
    # Generar gráficas
    print("📊 Generando gráficas...")
    charts = processor.generate_statistics_charts()
    
    # Generar PDF
    if PDF_AVAILABLE:
        print("📄 Generando PDF...")
        pdf_path = OUTPUT_DIR / "Documentos_BLATAM_Premium.pdf"
        pdf_gen = PDFGenerator(pdf_path)
        pdf_gen.generate(combined_title, combined_content, processor.stats, charts)
    else:
        print("⏭️  Saltando PDF (reportlab no disponible)")
    
    # Generar Word
    if WORD_AVAILABLE:
        print("📝 Generando Word...")
        word_path = OUTPUT_DIR / "Documentos_BLATAM_Premium.docx"
        word_gen = WordGenerator(word_path)
        word_gen.generate(combined_title, combined_content, processor.stats, charts)
    else:
        print("⏭️  Saltando Word (python-docx no disponible)")
    
    # Generar Excel
    if EXCEL_AVAILABLE:
        print("📊 Generando Excel...")
        excel_path = OUTPUT_DIR / "Documentos_BLATAM_Premium.xlsx"
        excel_gen = ExcelGenerator(excel_path)
        excel_gen.generate(combined_title, combined_content, processor.stats, charts)
    else:
        print("⏭️  Saltando Excel (openpyxl no disponible)")
    
    # Generar reportes JSON y dashboard HTML si hay análisis mejorado
    if ENHANCEMENTS_AVAILABLE and processor.enhanced_stats:
        print("\n📊 Generando reportes avanzados...")
        
        # Generar reportes JSON individuales
        json_reports = []
        for doc_path_str in IMPORTANT_DOCS:
            full_path = BASE_DIR / doc_path_str
            if full_path.exists() and doc_path_str in processor.enhanced_stats:
                try:
                    with open(full_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                    
                    doc_stats = processor.doc_stats.get(doc_path_str, {})
                    enhanced = processor.enhanced_stats[doc_path_str]
                    
                    report = generate_json_report(
                        doc_path_str,
                        content,
                        doc_stats,
                        enhanced['readability'],
                        enhanced['sentiment'],
                        enhanced['structure'],
                        enhanced['topics']
                    )
                    json_reports.append(report)
                    
                    # Guardar JSON individual
                    json_path = OUTPUT_DIR / f"{Path(doc_path_str).stem}_report.json"
                    with open(json_path, 'w', encoding='utf-8') as f:
                        json.dump(report, f, indent=2, ensure_ascii=False)
                    print(f"  ✅ JSON: {json_path.name}")
                except Exception as e:
                    print(f"  ⚠️  Error generando JSON para {doc_path_str}: {e}")
        
        # Generar dashboard HTML
        if json_reports:
            try:
                dashboard_path = OUTPUT_DIR / "dashboard_analisis.html"
                generate_html_dashboard(json_reports, dashboard_path)
                print(f"  ✅ Dashboard HTML: {dashboard_path.name}")
            except Exception as e:
                print(f"  ⚠️  Error generando dashboard: {e}")
        
        # Generar reporte JSON consolidado
        try:
            consolidated_report = {
                'generated_at': datetime.now().isoformat(),
                'total_documents': len(json_reports),
                'summary': {
                    'total_words': processor.stats['total_words'],
                    'total_sections': processor.stats['total_sections'],
                    'total_code_blocks': processor.stats['code_blocks'],
                    'total_links': processor.stats['links'],
                    'avg_quality_score': sum([r['quality_score'] for r in json_reports]) / len(json_reports) if json_reports else 0
                },
                'documents': json_reports
            }
            
            consolidated_path = OUTPUT_DIR / "reporte_consolidado.json"
            with open(consolidated_path, 'w', encoding='utf-8') as f:
                json.dump(consolidated_report, f, indent=2, ensure_ascii=False)
            print(f"  ✅ Reporte consolidado: {consolidated_path.name}")
        except Exception as e:
            print(f"  ⚠️  Error generando reporte consolidado: {e}")
        
        # Generar reportes Markdown y CSV si hay análisis avanzado
        if ADVANCED_ENHANCEMENTS_AVAILABLE and processor.advanced_stats:
            print("\n📝 Generando reportes adicionales...")
            
            for doc_path_str in IMPORTANT_DOCS:
                full_path = BASE_DIR / doc_path_str
                if full_path.exists() and doc_path_str in processor.advanced_stats:
                    try:
                        doc_name = Path(doc_path_str).stem
                        
                        # Preparar todos los análisis
                        all_analyses = {
                            'stats': processor.doc_stats.get(doc_path_str, {}),
                            'readability': processor.enhanced_stats.get(doc_path_str, {}).get('readability', {}),
                            'coherence': processor.advanced_stats[doc_path_str].get('coherence', {}),
                            'link_structure': processor.advanced_stats[doc_path_str].get('link_structure', {}),
                            'patterns': processor.advanced_stats[doc_path_str].get('patterns', {}),
                            'engagement': processor.advanced_stats[doc_path_str].get('engagement_score', 0)
                        }
                        
                        # Generar reporte Markdown
                        markdown_report = generate_markdown_report(doc_name, all_analyses)
                        md_path = OUTPUT_DIR / f"{doc_name}_report.md"
                        with open(md_path, 'w', encoding='utf-8') as f:
                            f.write(markdown_report)
                        print(f"  ✅ Markdown: {md_path.name}")
                        
                        # Generar CSV con estadísticas
                        csv_path = OUTPUT_DIR / f"{doc_name}_stats.csv"
                        export_to_csv(processor.doc_stats.get(doc_path_str, {}), csv_path)
                        print(f"  ✅ CSV: {csv_path.name}")
                        
                    except Exception as e:
                        print(f"  ⚠️  Error generando reportes adicionales para {doc_path_str}: {e}")
    
    # Generar Dashboard HTML Interactivo
    print("🌐 Generando Dashboard HTML...")
    html_dashboard = HTMLDashboardGenerator(OUTPUT_DIR / "Documentos_BLATAM_Dashboard.html")
    html_dashboard.generate(combined_title, processor.stats, processor.doc_stats, charts)
    print(f"✅ Dashboard HTML generado: {OUTPUT_DIR / 'Documentos_BLATAM_Dashboard.html'}")
    
    # Generar PowerPoint
    if PPTX_AVAILABLE:
        print("📊 Generando PowerPoint...")
        pptx_path = OUTPUT_DIR / "Documentos_BLATAM_Premium.pptx"
        pptx_gen = PowerPointGenerator(pptx_path)
        pptx_gen.generate(combined_title, processor.stats, processor.doc_stats, charts)
        print(f"✅ PowerPoint generado: {pptx_path}")
    else:
        print("⏭️  Saltando PowerPoint (python-pptx no disponible)")
    
    # Generar CSV con datos estructurados
    if CSV_AVAILABLE:
        print("📋 Generando CSV...")
        csv_path = OUTPUT_DIR / "Documentos_BLATAM_Data.csv"
        csv_gen = CSVGenerator(csv_path)
        csv_gen.generate(processor.stats, processor.doc_stats)
        print(f"✅ CSV generado: {csv_path}")
    
    # Generar Markdown estructurado mejorado
    print("📝 Generando Markdown estructurado...")
    md_path = OUTPUT_DIR / "Documentos_BLATAM_Resumen.md"
    md_gen = MarkdownGenerator(md_path)
    md_gen.generate(combined_title, processor.stats, processor.doc_stats, processor.summaries)
    print(f"✅ Markdown estructurado generado: {md_path}")
    
    # Generar JSON con datos para análisis
    json_path = OUTPUT_DIR / "Documentos_BLATAM_Data.json"
    json_data = {
        'metadata': {
            'generated': datetime.now().isoformat(),
            'title': combined_title,
            'files_processed': processor.stats.get('files_processed', 0),
        },
        'statistics': processor.stats,
        'documents': processor.doc_stats,
        'top_keywords': dict(processor.keywords.most_common(50)),
        'sections': processor.sections[:100],
        'main_topics': processor.stats.get('main_topics', []),
        'readability': processor.stats.get('readability_score', 0),
        'complexity': processor.stats.get('complexity_score', 0),
        'sentiment': processor.stats.get('avg_sentiment', 0),
        'code_quality': processor.stats.get('avg_code_quality', 0),
        'structure_score': processor.stats.get('avg_structure_score', 0),
        'summaries': processor.summaries,
    }
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(json_data, f, indent=2, ensure_ascii=False)
    print(f"✅ JSON de datos generado: {json_path}")
    
    # Limpiar archivos temporales
    for temp_file in OUTPUT_DIR.glob("temp_chart_*.png"):
        temp_file.unlink()
    
    print("\n✅ ¡Proceso completado!")
    print(f"📁 Archivos generados en: {OUTPUT_DIR}")
    print(f"   - PDF: Documentos_BLATAM_Premium.pdf")
    print(f"   - Word: Documentos_BLATAM_Premium.docx")
    print(f"   - Excel: Documentos_BLATAM_Premium.xlsx")
    print(f"   - PowerPoint: Documentos_BLATAM_Premium.pptx")
    print(f"   - HTML Dashboard: Documentos_BLATAM_Dashboard.html")
    print(f"   - CSV Data: Documentos_BLATAM_Data.csv")
    print(f"   - JSON Data: Documentos_BLATAM_Data.json")
    print(f"   - Markdown: Documentos_BLATAM_Resumen.md")
    
    # Resumen final
    pdfs = list(OUTPUT_DIR.glob("*.pdf"))
    words = list(OUTPUT_DIR.glob("*.docx"))
    excels = list(OUTPUT_DIR.glob("*.xlsx"))
    jsons = list(OUTPUT_DIR.glob("*_report.json"))
    
    print(f"\n📊 Resumen de archivos generados:")
    print(f"   📄 PDFs: {len(pdfs)}")
    print(f"   📝 Words: {len(words)}")
    print(f"   📊 Excels: {len(excels)}")
    if jsons:
        print(f"   📋 Reportes JSON: {len(jsons)}")
    if (OUTPUT_DIR / "dashboard_analisis.html").exists():
        print(f"   🌐 Dashboard HTML: 1")


if __name__ == "__main__":
    main()
