#!/usr/bin/env python3
"""
Script ULTIMATE PLUS con PowerPoint, JSON, CSV, análisis de tendencias y más
"""

import os
import sys
import re
import csv
from pathlib import Path
from datetime import datetime
import json
from collections import Counter, defaultdict
import math

try:
    import numpy as np
except ImportError:
    os.system("pip install numpy")
    import numpy as np

try:
    from docx import Document
    from docx.shared import Inches as DocxInches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    os.system("pip install python-docx")
    from docx import Document
    from docx.shared import Inches as DocxInches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Instalando python-pptx...")
    os.system("pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
    from pptx.enum.text import PP_ALIGN

# Importar funcionalidades base
sys.path.insert(0, str(Path(__file__).parent))
try:
    from convert_docs_ultimate import UltimateDocumentConverter
except ImportError:
    from convert_docs_to_formats_improved import AdvancedDocumentConverter as UltimateDocumentConverter

try:
    import matplotlib.pyplot as plt
    import matplotlib
    matplotlib.use('Agg')
    import seaborn as sns
    sns.set_style("whitegrid")
except ImportError:
    os.system("pip install matplotlib seaborn")
    import matplotlib.pyplot as plt
    import matplotlib
    matplotlib.use('Agg')
    import seaborn as sns

class UltimatePlusConverter(UltimateDocumentConverter):
    def __init__(self, output_dir="converted_docs"):
        super().__init__(output_dir)
        self.export_dir = self.output_dir / "exports"
        self.export_dir.mkdir(exist_ok=True)
        
    def export_to_json(self, doc_data, output_name):
        """Exporta análisis a JSON estructurado"""
        json_data = {
            "metadata": {
                "document_name": doc_data["name"],
                "generated_at": datetime.now().isoformat(),
                "source_file": doc_data.get("path", ""),
            },
            "statistics": doc_data["stats"],
            "structure": {
                "total_sections": len(doc_data["sections"]),
                "sections": [
                    {
                        "title": s["title"],
                        "level": s["level"],
                        "word_count": s["words"],
                        "code_blocks": s["code_blocks"],
                        "links": s["links"],
                        "content_preview": s["content"][:200]
                    }
                    for s in doc_data["sections"][:50]  # Limitar para JSON
                ]
            },
            "analysis": {
                "reading_time": doc_data["stats"].get("reading_time", {}),
                "readability": {
                    "score": doc_data["stats"]["readability_score"],
                    "level": "Alto" if doc_data["stats"]["readability_score"] > 70 else 
                            "Medio" if doc_data["stats"]["readability_score"] > 40 else "Bajo"
                },
                "complexity": {
                    "score": doc_data["stats"]["complexity_score"],
                    "level": "Alto" if doc_data["stats"]["complexity_score"] > 70 else 
                            "Medio" if doc_data["stats"]["complexity_score"] > 40 else "Bajo"
                }
            }
        }
        
        output_path = self.export_dir / f"{output_name}_analysis.json"
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(json_data, f, indent=2, ensure_ascii=False)
        
        print(f"✓ JSON exportado: {output_path}")
        return output_path
    
    def export_to_csv(self, doc_data, output_name):
        """Exporta estadísticas a CSV"""
        output_path = self.export_dir / f"{output_name}_stats.csv"
        
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            
            # Encabezados
            writer.writerow(["Métrica", "Valor"])
            
            # Estadísticas básicas
            stats = doc_data["stats"]
            for key, value in stats.items():
                if key not in ['top_keywords', 'section_analysis', 'structure', 'reading_time']:
                    writer.writerow([key.replace('_', ' ').title(), value])
            
            # Palabras clave
            writer.writerow([])
            writer.writerow(["Palabra Clave", "Frecuencia"])
            if 'top_keywords' in stats:
                for word, count in list(stats['top_keywords'].items())[:30]:
                    writer.writerow([word, count])
            
            # Secciones
            writer.writerow([])
            writer.writerow(["Sección", "Nivel", "Palabras", "Código", "Enlaces"])
            for section in doc_data["sections"][:100]:
                if section["title"]:
                    writer.writerow([
                        section["title"][:50],
                        section["level"],
                        section["words"],
                        section["code_blocks"],
                        section["links"]
                    ])
        
        print(f"✓ CSV exportado: {output_path}")
        return output_path
    
    def create_powerpoint_presentation(self, all_docs_data, output_name="Presentation"):
        """Crea presentación PowerPoint profesional"""
        print(f"Creando presentación PowerPoint...")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Portada
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        left = top = Inches(1)
        width = height = Inches(8.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Análisis de Documentos"
        p = tf.paragraphs[0]
        p.font.size = PptPt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = f"Generado el {datetime.now().strftime('%d de %B de %Y')}"
        p.font.size = PptPt(18)
        p.alignment = PP_ALIGN.CENTER
        
        # Slide 2: Resumen Ejecutivo
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        title = slide.shapes.title
        title.text = "Resumen Ejecutivo"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = f"Total de Documentos: {len(all_docs_data)}"
        
        p = tf.add_paragraph()
        p.text = f"Total de Palabras: {sum(d['stats']['total_words'] for d in all_docs_data):,}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"Total de Secciones: {sum(d['stats']['sections'] for d in all_docs_data)}"
        p.level = 0
        
        avg_read = sum(d['stats']['readability_score'] for d in all_docs_data) / len(all_docs_data)
        p = tf.add_paragraph()
        p.text = f"Legibilidad Promedio: {avg_read:.1f}/100"
        p.level = 0
        
        # Slide 3: Comparación de Documentos
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Comparación de Documentos"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = ""
        
        for doc in all_docs_data:
            p = tf.add_paragraph()
            p.text = f"{doc['name']}: {doc['stats']['total_words']:,} palabras, {doc['stats']['sections']} secciones"
            p.level = 0
            p.font.size = PptPt(14)
        
        # Slide 4-6: Análisis individual de cada documento
        for doc in all_docs_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = doc['name']
            
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = f"Palabras: {doc['stats']['total_words']:,}"
            
            p = tf.add_paragraph()
            p.text = f"Secciones: {doc['stats']['sections']}"
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = f"Legibilidad: {doc['stats']['readability_score']}/100"
            p.level = 0
            
            p = tf.add_paragraph()
            p.text = f"Complejidad: {doc['stats']['complexity_score']}/100"
            p.level = 0
            
            reading_time = doc['stats'].get('reading_time', {})
            if reading_time:
                time_str = f"{reading_time.get('hours', 0)}h {reading_time.get('remaining_minutes', 0)}m"
                p = tf.add_paragraph()
                p.text = f"Tiempo de Lectura: {time_str}"
                p.level = 0
        
        # Slide final: Conclusiones
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Conclusiones"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Análisis completado exitosamente"
        
        p = tf.add_paragraph()
        p.text = f"Se analizaron {len(all_docs_data)} documentos"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "Todos los documentos están disponibles en múltiples formatos"
        p.level = 0
        
        output_path = self.output_dir / f"{output_name}.pptx"
        prs.save(str(output_path))
        print(f"✓ PowerPoint guardado: {output_path}")
        return output_path
    
    def create_executive_summary(self, all_docs_data, output_name="Executive_Summary"):
        """Crea resumen ejecutivo en Word"""
        doc = Document()
        
        # Título
        title = doc.add_heading('Resumen Ejecutivo - Análisis de Documentos', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Fecha
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        para.add_run(f"Generado el {datetime.now().strftime('%d de %B de %Y a las %H:%M')}")
        
        doc.add_page_break()
        
        # Resumen General
        doc.add_heading('Resumen General', level=1)
        total_words = sum(d['stats']['total_words'] for d in all_docs_data)
        total_sections = sum(d['stats']['sections'] for d in all_docs_data)
        
        para = doc.add_paragraph()
        para.add_run(f"Se analizaron {len(all_docs_data)} documentos con un total de ")
        para.add_run(f"{total_words:,} palabras").bold = True
        para.add_run(" y ")
        para.add_run(f"{total_sections} secciones").bold = True
        para.add_run(".")
        
        # Métricas Promedio
        doc.add_heading('Métricas Promedio', level=2)
        avg_read = sum(d['stats']['readability_score'] for d in all_docs_data) / len(all_docs_data)
        avg_comp = sum(d['stats']['complexity_score'] for d in all_docs_data) / len(all_docs_data)
        
        para = doc.add_paragraph()
        para.add_run(f"Legibilidad promedio: {avg_read:.1f}/100")
        
        para = doc.add_paragraph()
        para.add_run(f"Complejidad promedio: {avg_comp:.1f}/100")
        
        # Análisis por Documento
        doc.add_heading('Análisis Individual', level=1)
        
        for doc_data in all_docs_data:
            doc.add_heading(doc_data['name'], level=2)
            
            para = doc.add_paragraph()
            para.add_run(f"Palabras: {doc_data['stats']['total_words']:,} | ")
            para.add_run(f"Secciones: {doc_data['stats']['sections']} | ")
            para.add_run(f"Legibilidad: {doc_data['stats']['readability_score']}/100 | ")
            para.add_run(f"Complejidad: {doc_data['stats']['complexity_score']}/100")
            
            reading_time = doc_data['stats'].get('reading_time', {})
            if reading_time:
                time_str = f"{reading_time.get('hours', 0)}h {reading_time.get('remaining_minutes', 0)}m"
                para = doc.add_paragraph()
                para.add_run(f"Tiempo estimado de lectura: {time_str}")
        
        # Recomendaciones
        doc.add_heading('Recomendaciones', level=1)
        
        para = doc.add_paragraph()
        para.add_run("1. Revisar documentos con baja legibilidad para mejorar la claridad")
        
        para = doc.add_paragraph()
        para.add_run("2. Considerar dividir documentos muy extensos en secciones más manejables")
        
        para = doc.add_paragraph()
        para.add_run("3. Asegurar consistencia en el formato y estructura entre documentos")
        
        output_path = self.output_dir / f"{output_name}.docx"
        doc.save(str(output_path))
        print(f"✓ Resumen ejecutivo guardado: {output_path}")
        return output_path
    
    def create_trend_analysis(self, all_docs_data):
        """Crea análisis de tendencias"""
        if len(all_docs_data) < 2:
            return None
        
        # Crear gráfica de tendencias
        fig, axes = plt.subplots(2, 2, figsize=(16, 12))
        
        doc_names = [d['name'] for d in all_docs_data]
        
        # 1. Tendencia de tamaño
        ax = axes[0, 0]
        word_counts = [d['stats']['total_words'] for d in all_docs_data]
        ax.plot(doc_names, word_counts, marker='o', linewidth=3, markersize=10, color='#3498db')
        ax.fill_between(doc_names, word_counts, alpha=0.3, color='#3498db')
        ax.set_title('Tendencia de Tamaño (Palabras)', fontsize=14, fontweight='bold')
        ax.set_ylabel('Número de Palabras', fontweight='bold')
        ax.grid(True, alpha=0.3)
        plt.setp(ax.xaxis.get_majorticklabels(), rotation=45, ha='right')
        
        # 2. Tendencia de legibilidad
        ax = axes[0, 1]
        readability = [d['stats']['readability_score'] for d in all_docs_data]
        ax.plot(doc_names, readability, marker='s', linewidth=3, markersize=10, color='#2ecc71')
        ax.fill_between(doc_names, readability, alpha=0.3, color='#2ecc71')
        ax.set_title('Tendencia de Legibilidad', fontsize=14, fontweight='bold')
        ax.set_ylabel('Score (0-100)', fontweight='bold')
        ax.set_ylim(0, 100)
        ax.grid(True, alpha=0.3)
        plt.setp(ax.xaxis.get_majorticklabels(), rotation=45, ha='right')
        
        # 3. Tendencia de complejidad
        ax = axes[1, 0]
        complexity = [d['stats']['complexity_score'] for d in all_docs_data]
        ax.plot(doc_names, complexity, marker='^', linewidth=3, markersize=10, color='#e74c3c')
        ax.fill_between(doc_names, complexity, alpha=0.3, color='#e74c3c')
        ax.set_title('Tendencia de Complejidad', fontsize=14, fontweight='bold')
        ax.set_ylabel('Score (0-100)', fontweight='bold')
        ax.set_ylim(0, 100)
        ax.grid(True, alpha=0.3)
        plt.setp(ax.xaxis.get_majorticklabels(), rotation=45, ha='right')
        
        # 4. Comparación Legibilidad vs Complejidad
        ax = axes[1, 1]
        ax.scatter(readability, complexity, s=200, alpha=0.6, c=range(len(doc_names)), 
                  cmap='viridis', edgecolors='black', linewidth=2)
        for i, name in enumerate(doc_names):
            ax.annotate(name, (readability[i], complexity[i]), 
                       xytext=(5, 5), textcoords='offset points', fontsize=9)
        ax.set_xlabel('Legibilidad', fontweight='bold')
        ax.set_ylabel('Complejidad', fontweight='bold')
        ax.set_title('Legibilidad vs Complejidad', fontsize=14, fontweight='bold')
        ax.grid(True, alpha=0.3)
        
        plt.tight_layout()
        chart_path = self.temp_dir / "trend_analysis.png"
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        print(f"✓ Análisis de tendencias guardado: {chart_path}")
        return chart_path
    
    def create_all_exports(self, all_docs_data):
        """Crea todas las exportaciones adicionales"""
        print("\n" + "=" * 70)
        print("CREANDO EXPORTACIONES ADICIONALES")
        print("=" * 70)
        
        exports = {}
        
        # Exportar cada documento a JSON y CSV
        for doc_data in all_docs_data:
            name = doc_data['name']
            exports[f"{name}_json"] = self.export_to_json(doc_data, name)
            exports[f"{name}_csv"] = self.export_to_csv(doc_data, name)
        
        # Crear PowerPoint
        exports['powerpoint'] = self.create_powerpoint_presentation(all_docs_data)
        
        # Crear Resumen Ejecutivo
        exports['executive_summary'] = self.create_executive_summary(all_docs_data)
        
        # Crear análisis de tendencias
        if len(all_docs_data) >= 2:
            exports['trend_analysis'] = self.create_trend_analysis(all_docs_data)
        
        return exports

def main():
    """Función principal"""
    base_dir = Path("/Users/adan/Documents/documentos_blatam")
    production_dir = base_dir / "truthgpt_collected/integration_code/production_code"
    
    documents = [
        {
            "path": base_dir / "airflow_automation_prompt.md",
            "name": "Automation_Expert_Prompt"
        },
        {
            "path": production_dir / "ARCHITECTURE_IMPROVEMENTS.md",
            "name": "Architecture_Improvements"
        },
        {
            "path": production_dir / "REFACTORING_PLAN.md",
            "name": "Refactoring_Plan"
        }
    ]
    
    converter = UltimatePlusConverter()
    
    # Primero procesar todos los documentos
    all_docs_data = []
    for doc in documents:
        if not doc["path"].exists():
            continue
        
        content = converter.read_markdown(doc["path"])
        sections = converter.parse_markdown_sections(content)
        stats = converter.create_advanced_statistics(content, sections)
        reading_time = converter.calculate_reading_time(stats["total_words"])
        stats["reading_time"] = reading_time
        
        all_docs_data.append({
            "name": doc["name"],
            "path": str(doc["path"]),
            "stats": stats,
            "sections": sections
        })
    
    # Crear todas las exportaciones
    exports = converter.create_all_exports(all_docs_data)
    
    print("\n" + "=" * 70)
    print("RESUMEN DE EXPORTACIONES")
    print("=" * 70)
    print("\nArchivos exportados:")
    for key, path in exports.items():
        print(f"  - {key}: {path}")
    
    print("\n✅ Todas las exportaciones completadas!")

if __name__ == "__main__":
    main()

