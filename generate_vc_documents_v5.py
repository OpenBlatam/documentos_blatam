import pandas as pd
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Create output directory
output_dir = "Venture_Capital_Documents_v5"
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

# --- 1. TERM SHEET GENERATOR (WORD) ---
def create_term_sheet():
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(11)

    # Header
    header = doc.add_heading('TERM SHEET', 0)
    header.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph('FOR SERIES A PREFERRED STOCK FINANCING OF\n[COMPANY NAME]', style='Subtitle').alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph(f"Date: November 23, 2025").alignment = WD_ALIGN_PARAGRAPH.RIGHT
    doc.add_paragraph("_" * 80)

    intro = doc.add_paragraph("This Term Sheet summarizes the principal terms of the Series A Preferred Stock financing of [Company Name], a Delaware corporation (the “Company”). In consideration of the time and expense devoted and to be devoted by the Investors with respect to this investment, the No Shop/Confidentiality provisions of this Term Sheet shall be binding obligations of the Company whether or not the financing is consummated. No other legally binding obligations will be created until definitive agreements are executed and delivered by all parties.")
    intro.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    # Terms Table
    table = doc.add_table(rows=1, cols=2)
    table.style = 'Table Grid'
    table.autofit = False
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.5)

    terms = [
        ("Issuer", "[Company Name]"),
        ("Investors", "Ventura Capital (Lead Investor) and others to be agreed."),
        ("Amount of Financing", "$15,000,000"),
        ("Valuation", "Pre-money valuation of $60,000,000, resulting in a post-money valuation of $75,000,000."),
        ("Liquidation Preference", "1x Non-Participating Preference. Series A gets paid back first, then participates pro-rata with Common Stock."),
        ("Dividends", "8% non-cumulative, payable if and when declared by Board."),
        ("Board of Directors", "The Board shall consist of 5 members: 2 founders, 1 Ventura Capital representative, 2 independent directors."),
        ("Vesting", "Standard 4-year vesting for founders/employees with 1-year cliff."),
        ("Anti-Dilution", "Broad-based weighted average anti-dilution protection."),
        ("Exclusivity (No Shop)", "30 days from the signing of this Term Sheet.")
    ]

    for term, desc in terms:
        row = table.add_row()
        row.cells[0].text = term
        row.cells[0].paragraphs[0].runs[0].bold = True
        row.cells[1].text = desc

    doc.add_paragraph("\n")
    doc.add_paragraph("ACCEPTED AND AGREED:", style='Heading 3')
    
    sig_table = doc.add_table(rows=1, cols=2)
    sig_table.width = Inches(6)
    sig_table.rows[0].cells[0].text = "VENTURA CAPITAL\n\nBy: __________________\nName: Adan Munoz\nTitle: General Partner"
    sig_table.rows[0].cells[1].text = "COMPANY\n\nBy: __________________\nName: [Founder Name]\nTitle: CEO"

    filename = os.path.join(output_dir, "Ventura_Capital_Term_Sheet_Series_A.docx")
    doc.save(filename)
    print(f"Docx created: {filename}")

# --- 2. QUARTERLY LP REPORT (WORD) ---
def create_lp_report():
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    
    # Cover
    doc.add_heading('VENTURA CAPITAL', 0)
    doc.add_paragraph('QUARTERLY INVESTOR REPORT | Q3 2025', style='Subtitle')
    doc.add_paragraph('\nCONFIDENTIAL')
    doc.add_page_break()

    # 1. Letter
    doc.add_heading('1. General Partner Letter', level=1)
    doc.add_paragraph("Dear Limited Partners,")
    doc.add_paragraph("We are pleased to report that Q3 2025 has been a transformative quarter for Fund I. Our 'Moat Building' strategy is yielding tangible results, with our portfolio companies seeing a 45% aggregate revenue increase.")
    doc.add_paragraph("Key highlights:")
    doc.add_paragraph("• 2 New Investments deployed ($30M total).", style='List Bullet')
    doc.add_paragraph("• 1 Portfolio Company (Nexus AI) raised Series B at 3x markup.", style='List Bullet')
    doc.add_paragraph("• Fund TVPI currently stands at 1.8x.", style='List Bullet')

    # 2. Fund Performance
    doc.add_heading('2. Fund Performance Overview', level=1)
    table = doc.add_table(rows=2, cols=4)
    table.style = 'Table Grid'
    
    headers = ['Metric', 'Q2 2025', 'Q3 2025', 'Change']
    data = ['TVPI (Multiple)', '1.5x', '1.8x', '+0.3x']
    
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    for i, d in enumerate(data):
        table.rows[1].cells[i].text = d
        
    # 3. Portfolio Updates
    doc.add_heading('3. Portfolio Highlights', level=1)
    
    companies = [
        ("Nexus AI", "Generative Design", "Closed $50M Series B led by Sequoia. Revenue up 200% YoY."),
        ("RoboLogic", "Industrial Robotics", "Secured contract with Tesla. Manufacturing capacity expanded."),
        ("SynthGen", "Bio-AI", "FDA Phase 1 approval granted. Stock up 15% in private markets.")
    ]
    
    for name, sector, update in companies:
        p = doc.add_paragraph()
        p.add_run(f"{name} ({sector})").bold = True
        doc.add_paragraph(update)
        doc.add_paragraph("")

    filename = os.path.join(output_dir, "Ventura_Capital_LP_Report_Q3_2025.docx")
    doc.save(filename)
    print(f"Docx created: {filename}")

# --- 3. ONE-PAGER TEASER TEMPLATE (PPTX) ---
def create_teaser_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # Navy Background strip on left
    left_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.5), Inches(7.5))
    left_strip.fill.solid()
    left_strip.fill.fore_color.rgb = PptRGBColor(0, 32, 96)
    left_strip.line.fill.background()
    
    # Company Logo / Name
    tf = left_strip.text_frame
    p = tf.add_paragraph()
    p.text = "[STARTUP NAME]"
    p.font.color.rgb = PptRGBColor(255, 255, 255)
    p.font.bold = True
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\n\nSECTOR:\nGenerative AI\n\nSTAGE:\nSeries A\n\nASK:\n$15M"
    p.font.color.rgb = PptRGBColor(255, 255, 255)
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    # Right side content
    # 1. One Liner
    tbox = slide.shapes.add_textbox(Inches(2.8), Inches(0.5), Inches(7), Inches(1))
    p = tbox.text_frame.add_paragraph()
    p.text = "The Operating System for Autonomous Manufacturing"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(0, 32, 96)

    # 2. Problem/Solution Grid
    headings = [
        ("THE PROBLEM", "Manual workflows cost factories $50B/year in lost productivity.", Inches(2.8), Inches(2)),
        ("THE SOLUTION", "AI-driven agentic workflows that automate 90% of supply chain decisions.", Inches(2.8), Inches(3.5)),
        ("TRACTION", "$2M ARR (300% YoY). 15 Enterprise Customers including BMW & Ford.", Inches(6.5), Inches(2)),
        ("TEAM", "Ex-SpaceX Engineers. PhDs in Computer Vision from MIT.", Inches(6.5), Inches(3.5))
    ]

    for title, content, x, y in headings:
        # Title
        tb = slide.shapes.add_textbox(x, y, Inches(3.5), Inches(0.5))
        p = tb.text_frame.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.color.rgb = PptRGBColor(191, 144, 0) # Gold
        p.font.size = Pt(12)
        
        # Content
        cb = slide.shapes.add_textbox(x, y + Inches(0.3), Inches(3.2), Inches(1))
        p = cb.text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(11)
        p.word_wrap = True

    # Footer
    footer = slide.shapes.add_textbox(Inches(2.8), Inches(6.8), Inches(7), Inches(0.5))
    p = footer.text_frame.add_paragraph()
    p.text = "CONTACT: partners@venturacapital.com | CONFIDENTIAL"
    p.font.size = Pt(10)
    p.font.color.rgb = PptRGBColor(150, 150, 150)

    filename = os.path.join(output_dir, "Ventura_Capital_Deal_Teaser.pptx")
    prs.save(filename)
    print(f"PPTX created: {filename}")

# --- 4. BOARD MEETING AGENDA (WORD) ---
def create_board_agenda():
    doc = Document()
    doc.add_heading('BOARD OF DIRECTORS MEETING', 0)
    doc.add_paragraph('COMPANY: [Portfolio Company Name]\nDATE: November 2025\nLOCATION: Ventura Capital HQ / Zoom')
    doc.add_paragraph('_' * 80)

    items = [
        ("1. Call to Order & Admin (5 min)", "Approval of previous minutes. Agenda review."),
        ("2. CEO Update (15 min)", "High-level executive summary. Key wins and losses."),
        ("3. KPI Review (15 min)", "Financials (Burn, Runway, Revenue). Sales Pipeline. Product Metrics."),
        ("4. Strategic Deep Dive (45 min)", "Focus Topic: 'Scaling Engineering Team' or 'Q4 Go-To-Market Strategy'."),
        ("5. Governance & Approvals (10 min)", "Option grants. Budget approvals."),
        ("6. Executive Session (15 min)", "Closed session without management (if needed)."),
        ("7. Adjournment", "")
    ]

    table = doc.add_table(rows=1, cols=2)
    table.style = 'Table Grid'
    table.rows[0].cells[0].text = "TIME / ITEM"
    table.rows[0].cells[1].text = "DESCRIPTION / NOTES"
    
    for item, desc in items:
        row = table.add_row()
        row.cells[0].text = item
        row.cells[0].paragraphs[0].runs[0].bold = True
        row.cells[1].text = desc
        row.height = Inches(0.5)

    filename = os.path.join(output_dir, "Ventura_Capital_Board_Agenda_Template.docx")
    doc.save(filename)
    print(f"Docx created: {filename}")

if __name__ == "__main__":
    create_term_sheet()
    create_lp_report()
    create_teaser_slide()
    create_board_agenda()

