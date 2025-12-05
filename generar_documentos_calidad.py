#!/usr/bin/env python3
"""
Script mejorado para generar documentos PDF, Word y Excel de alta calidad
con gráficas profesionales a partir de los documentos más importantes del proyecto.
"""

import os
import re
from datetime import datetime
from pathlib import Path
import markdown
from typing import List, Dict, Any, Tuple, Optional
import json
import unicodedata
from collections import Counter, defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
import time
import hashlib

# PDF generation
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    print("⚠️  reportlab no disponible. Instalando...")
    os.system("pip install reportlab -q")

# Word generation
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("⚠️  python-docx no disponible. Instalando...")
    os.system("pip install python-docx -q")

# Excel generation
try:
    import openpyxl
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.drawing.image import Image as ExcelImage
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    print("⚠️  openpyxl no disponible. Instalando...")
    os.system("pip install openpyxl -q")

# PowerPoint generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    # No instalamos automáticamente porque python-pptx puede ser opcional

# Graph generation
try:
    import matplotlib.pyplot as plt
    import matplotlib
    matplotlib.use('Agg')  # Non-interactive backend
    import seaborn as sns
    import numpy as np
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False
    print("⚠️  matplotlib no disponible. Instalando...")
    os.system("pip install matplotlib seaborn numpy -q")

# Markdown parsing
try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False
    print("⚠️  markdown no disponible. Instalando...")
    os.system("pip install markdown -q")

# Configure matplotlib for better quality
if MATPLOTLIB_AVAILABLE:
    plt.rcParams['figure.dpi'] = 300
    plt.rcParams['savefig.dpi'] = 300
    plt.rcParams['figure.figsize'] = (10, 6)
    sns.set_style("whitegrid")
    sns.set_palette("husl")


class DocumentConverter:
    """Clase principal para convertir documentos markdown a PDF, Word y Excel"""
    
    def __init__(self, output_dir: str = "documentos_generados", parallel: bool = True):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.graphs_dir = self.output_dir / "graficas"
        self.graphs_dir.mkdir(exist_ok=True)
        self.parallel = parallel
        self.stats = {
            'total_docs': 0,
            'total_pages': 0,
            'total_tables': 0,
            'total_sections': 0,
            'total_words': 0,
            'processing_time': 0,
            'documents_data': []  # Almacenar datos de cada documento
        }
        
    def read_markdown(self, file_path: str) -> str:
        """Lee un archivo markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                # Actualizar estadísticas
                self.stats['total_words'] += len(content.split())
                return content
        except UnicodeDecodeError:
            # Intentar con diferentes encodings
            for encoding in ['latin-1', 'iso-8859-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except:
                    continue
            raise
    
    def clean_emoji(self, text: str) -> str:
        """Limpia emojis del texto para evitar problemas de renderizado"""
        # Remover emojis pero mantener el texto
        emoji_pattern = re.compile("["
            u"\U0001F600-\U0001F64F"  # emoticons
            u"\U0001F300-\U0001F5FF"  # symbols & pictographs
            u"\U0001F680-\U0001F6FF"  # transport & map symbols
            u"\U0001F1E0-\U0001F1FF"  # flags
            u"\U00002702-\U000027B0"
            u"\U000024C2-\U0001F251"
            "]+", flags=re.UNICODE)
        return emoji_pattern.sub('', text).strip()
    
    def parse_table(self, table_lines: List[str]) -> List[List[str]]:
        """Parsea una tabla markdown a lista de listas"""
        if not table_lines or len(table_lines) < 2:
            return []
        
        rows = []
        for line in table_lines:
            if '|' in line:
                cells = [cell.strip() for cell in line.split('|')[1:-1]]
                if cells and any(cell.strip() for cell in cells):
                    rows.append(cells)
        
        return rows if len(rows) > 1 else []
    
    def parse_markdown(self, content: str) -> Dict[str, Any]:
        """Parsea el contenido markdown y extrae información estructurada"""
        # Extraer título principal
        title_match = re.search(r'^#\s+(.+)$', content, re.MULTILINE)
        title = title_match.group(1) if title_match else "Documento"
        title = self.clean_emoji(title)
        
        # Extraer secciones
        sections = []
        current_section = None
        current_table = []
        in_table = False
        
        lines = content.split('\n')
        for i, line in enumerate(lines):
            # Títulos de sección
            if line.startswith('#'):
                if in_table and current_table:
                    if current_section:
                        current_section['tables'] = current_section.get('tables', [])
                        current_section['tables'].append(self.parse_table(current_table))
                    current_table = []
                    in_table = False
                
                if current_section:
                    sections.append(current_section)
                
                level = len(line) - len(line.lstrip('#'))
                title_text = self.clean_emoji(line.replace('#', '').strip())
                
                current_section = {
                    'title': title_text,
                    'content': [],
                    'level': level,
                    'tables': []
                }
            elif line.strip().startswith('|'):
                # Tabla
                in_table = True
                current_table.append(line.strip())
            elif in_table:
                # Fin de tabla
                if current_table:
                    if current_section:
                        current_section['tables'] = current_section.get('tables', [])
                        current_section['tables'].append(self.parse_table(current_table))
                    current_table = []
                in_table = False
                
                if current_section and line.strip():
                    if line.strip().startswith('-') or line.strip().startswith('*'):
                        current_section['content'].append({
                            'type': 'list_item',
                            'content': self.clean_emoji(line.strip())
                        })
                    elif line.strip():
                        current_section['content'].append({
                            'type': 'paragraph',
                            'content': self.clean_emoji(line.strip())
                        })
            elif current_section and line.strip():
                if line.strip().startswith('-') or line.strip().startswith('*'):
                    current_section['content'].append({
                        'type': 'list_item',
                        'content': self.clean_emoji(line.strip())
                    })
                elif line.strip() and not line.strip().startswith('```'):
                    current_section['content'].append({
                        'type': 'paragraph',
                        'content': self.clean_emoji(line.strip())
                    })
        
        if in_table and current_table and current_section:
            current_section['tables'] = current_section.get('tables', [])
            current_section['tables'].append(self.parse_table(current_table))
        
        if current_section:
            sections.append(current_section)
        
        # Extraer métricas y datos numéricos
        metrics = self.extract_metrics(content)
        
        # Extraer tablas globales para análisis
        all_tables = []
        for section in sections:
            all_tables.extend(section.get('tables', []))
        
        return {
            'title': title,
            'sections': sections,
            'metrics': metrics,
            'tables': all_tables,
            'raw_content': content
        }
    
    def extract_metrics(self, content: str) -> Dict[str, Any]:
        """Extrae métricas y datos numéricos del contenido"""
        metrics = {}
        
        # Buscar números con contexto más específico
        number_patterns = [
            (r'(\d+)\s*(?:archivos?|files?)', 'archivos'),
            (r'(\d+)\s*(?:líneas?|lines?)', 'lineas'),
            (r'(\d+)\s*(?:fases?|phases?)', 'fases'),
            (r'(\d+)\s*%', 'porcentajes'),
            (r'✅\s*(\d+)/(\d+)', 'completitud'),
            (r'(\d+)\s*(?:documentos?|docs?)', 'documentos'),
            (r'(\d+)\s*(?:mejoras?|improvements?)', 'mejoras'),
            (r'(\d+)\s*(?:tareas?|tasks?)', 'tareas'),
            (r'(\d+)\+', 'mas_de'),
        ]
        
        for pattern, key in number_patterns:
            matches = re.findall(pattern, content, re.IGNORECASE)
            if matches:
                if key not in metrics:
                    metrics[key] = []
                metrics[key].extend(matches)
        
        # Extraer datos de tablas
        table_data = self.extract_table_metrics(content)
        if table_data:
            metrics['tablas'] = table_data
        
        return metrics
    
    def extract_table_metrics(self, content: str) -> List[Dict[str, Any]]:
        """Extrae métricas de las tablas en el contenido"""
        table_metrics = []
        lines = content.split('\n')
        current_table = []
        in_table = False
        
        for line in lines:
            if line.strip().startswith('|') and '---' not in line:
                in_table = True
                current_table.append(line.strip())
            elif in_table and not line.strip().startswith('|'):
                if len(current_table) > 1:
                    parsed = self.parse_table(current_table)
                    if parsed and len(parsed) > 1:
                        # Intentar extraer números de la tabla
                        numbers = []
                        for row in parsed[1:]:  # Saltar header
                            for cell in row:
                                nums = re.findall(r'\d+', cell)
                                numbers.extend([int(n) for n in nums])
                        if numbers:
                            table_metrics.append({
                                'rows': len(parsed),
                                'cols': len(parsed[0]) if parsed else 0,
                                'numbers': numbers
                            })
                current_table = []
                in_table = False
        
        return table_metrics
    
    def generate_graphs(self, doc_name: str, parsed_data: Dict[str, Any]) -> List[str]:
        """Genera gráficas profesionales basadas en los datos extraídos"""
        graph_files = []
        
        if not MATPLOTLIB_AVAILABLE:
            return graph_files
        
        metrics = parsed_data.get('metrics', {})
        sections = parsed_data.get('sections', [])
        tables = parsed_data.get('tables', [])
        
        # Gráfica 1: Dashboard de métricas generales
        if metrics:
            fig = plt.figure(figsize=(16, 10))
            gs = fig.add_gridspec(3, 3, hspace=0.3, wspace=0.3)
            fig.suptitle(f'Dashboard de Métricas: {doc_name}', fontsize=18, fontweight='bold', y=0.98)
            
            plot_count = 0
            
            # Gráfica de barras para archivos/líneas
            if 'archivos' in metrics or 'lineas' in metrics:
                ax1 = fig.add_subplot(gs[0, 0])
                categories = []
                values = []
                colors_list = ['#3498db', '#e74c3c', '#f39c12', '#27ae60']
                
                if 'archivos' in metrics:
                    categories.append('Archivos')
                    values.append(int(metrics['archivos'][0][0]) if metrics['archivos'] else 0)
                
                if 'lineas' in metrics:
                    categories.append('Líneas')
                    values.append(int(metrics['lineas'][0][0]) if metrics['lineas'] else 0)
                
                if 'documentos' in metrics:
                    categories.append('Documentos')
                    values.append(int(metrics['documentos'][0][0]) if metrics['documentos'] else 0)
                
                if 'mejoras' in metrics:
                    categories.append('Mejoras')
                    values.append(int(metrics['mejoras'][0][0]) if metrics['mejoras'] else 0)
                
                if categories:
                    bars = ax1.bar(categories, values, color=colors_list[:len(categories)])
                    ax1.set_title('Métricas Principales', fontweight='bold', fontsize=12)
                    ax1.set_ylabel('Cantidad', fontsize=10)
                    ax1.grid(True, alpha=0.3, axis='y')
                    # Agregar valores en las barras
                    for bar in bars:
                        height = bar.get_height()
                        ax1.text(bar.get_x() + bar.get_width()/2., height,
                                f'{int(height)}', ha='center', va='bottom', fontweight='bold')
                    plot_count += 1
            
            # Gráfica de completitud (pie chart mejorado)
            if 'completitud' in metrics:
                ax2 = fig.add_subplot(gs[0, 1])
                completitud = metrics['completitud']
                if completitud:
                    completed = int(completitud[0][0])
                    total = int(completitud[0][1])
                    remaining = total - completed
                    percentage = (completed / total * 100) if total > 0 else 0
                    
                    colors_pie = ['#2ecc71', '#ecf0f1']
                    wedges, texts, autotexts = ax2.pie([completed, remaining], 
                           labels=[f'Completado\n{completed}', f'Pendiente\n{remaining}'],
                           autopct='%1.1f%%',
                           colors=colors_pie,
                           startangle=90,
                           textprops={'fontsize': 10, 'fontweight': 'bold'})
                    ax2.set_title(f'Estado de Completitud\n{percentage:.1f}%', 
                                fontweight='bold', fontsize=12)
                    plot_count += 1
            
            # Gráfica de fases (line chart mejorado)
            if 'fases' in metrics:
                ax3 = fig.add_subplot(gs[0, 2])
                fases_data = [int(f[0]) for f in metrics['fases']]
                if fases_data:
                    x_vals = range(1, len(fases_data) + 1)
                    ax3.plot(x_vals, fases_data, marker='o', linewidth=3, 
                            markersize=10, color='#9b59b6', markerfacecolor='#ffffff',
                            markeredgewidth=2, markeredgecolor='#9b59b6')
                    ax3.fill_between(x_vals, fases_data, alpha=0.3, color='#9b59b6')
                    ax3.set_title('Progreso de Fases', fontweight='bold', fontsize=12)
                    ax3.set_xlabel('Número de Fase', fontsize=10)
                    ax3.set_ylabel('Valor', fontsize=10)
                    ax3.grid(True, alpha=0.3)
                    plot_count += 1
            
            # Gráfica de porcentajes (horizontal bar mejorado)
            if 'porcentajes' in metrics:
                ax4 = fig.add_subplot(gs[1, 0])
                porcentajes = [int(p[0]) for p in metrics['porcentajes'][:10]]  # Top 10
                if porcentajes:
                    y_pos = range(len(porcentajes))
                    bars = ax4.barh(y_pos, porcentajes, color=plt.cm.viridis(np.linspace(0, 1, len(porcentajes))))
                    ax4.set_yticks(y_pos)
                    ax4.set_yticklabels([f'Ítem {i+1}' for i in y_pos])
                    ax4.set_xlabel('Porcentaje (%)', fontsize=10)
                    ax4.set_title('Distribución de Porcentajes', fontweight='bold', fontsize=12)
                    ax4.grid(True, alpha=0.3, axis='x')
                    # Agregar valores
                    for i, (bar, val) in enumerate(zip(bars, porcentajes)):
                        ax4.text(val, bar.get_y() + bar.get_height()/2,
                                f'{val}%', ha='left', va='center', fontweight='bold')
                    plot_count += 1
            
            # Gráfica de estructura de secciones
            if sections:
                ax5 = fig.add_subplot(gs[1, 1])
                section_levels = {}
                for s in sections:
                    level = s.get('level', 2)
                    section_levels[level] = section_levels.get(level, 0) + 1
                
                if section_levels:
                    levels = sorted(section_levels.keys())
                    counts = [section_levels[l] for l in levels]
                    colors_map = plt.cm.Set3(np.linspace(0, 1, len(levels)))
                    bars = ax5.bar([f'Nivel {l}' for l in levels], counts, color=colors_map)
                    ax5.set_title('Distribución por Nivel', fontweight='bold', fontsize=12)
                    ax5.set_ylabel('Cantidad de Secciones', fontsize=10)
                    ax5.grid(True, alpha=0.3, axis='y')
                    for bar in bars:
                        height = bar.get_height()
                        ax5.text(bar.get_x() + bar.get_width()/2., height,
                                f'{int(height)}', ha='center', va='bottom', fontweight='bold')
                    plot_count += 1
            
            # Gráfica de datos de tablas
            if 'tablas' in metrics and metrics['tablas']:
                ax6 = fig.add_subplot(gs[1, 2])
                table_data = metrics['tablas']
                if table_data:
                    rows_counts = [t['rows'] for t in table_data]
                    cols_counts = [t['cols'] for t in table_data]
                    x = range(len(table_data))
                    width = 0.35
                    ax6.bar([i - width/2 for i in x], rows_counts, width, 
                           label='Filas', color='#3498db', alpha=0.8)
                    ax6.bar([i + width/2 for i in x], cols_counts, width,
                           label='Columnas', color='#e74c3c', alpha=0.8)
                    ax6.set_title('Análisis de Tablas', fontweight='bold', fontsize=12)
                    ax6.set_xlabel('Tabla', fontsize=10)
                    ax6.set_ylabel('Cantidad', fontsize=10)
                    ax6.legend()
                    ax6.grid(True, alpha=0.3, axis='y')
                    plot_count += 1
            
            # Si hay menos gráficas, ajustar layout
            if plot_count < 6 and len(fig.axes) > plot_count:
                # Ocultar subplots vacíos
                for i in range(plot_count, min(6, len(fig.axes))):
                    try:
                        fig.delaxes(fig.axes[i])
                    except (IndexError, AttributeError):
                        pass
            
            plt.savefig(self.graphs_dir / f"{doc_name}_dashboard_metricas.png", 
                       bbox_inches='tight', dpi=300, facecolor='white')
            graph_files.append(str(self.graphs_dir / f"{doc_name}_dashboard_metricas.png"))
            plt.close()
        
        # Gráfica 2: Estructura del documento (mejorada)
        if sections:
            fig, ax = plt.subplots(figsize=(14, max(8, len(sections) * 0.3)))
            
            section_titles = [s['title'][:40] + '...' if len(s['title']) > 40 else s['title'] 
                            for s in sections[:15]]  # Primeras 15 secciones
            section_levels = [s.get('level', 2) for s in sections[:15]]
            
            if section_titles:
                colors_map = plt.cm.viridis(np.linspace(0, 1, len(section_titles)))
                bars = ax.barh(range(len(section_titles)), section_levels, color=colors_map, alpha=0.8)
                ax.set_yticks(range(len(section_titles)))
                ax.set_yticklabels(section_titles, fontsize=9)
                ax.set_xlabel('Nivel de Jerarquía', fontsize=11, fontweight='bold')
                ax.set_title('Estructura del Documento - Jerarquía de Secciones', 
                           fontsize=14, fontweight='bold', pad=20)
                ax.grid(True, alpha=0.3, axis='x')
                
                # Agregar valores
                for i, (bar, level) in enumerate(zip(bars, section_levels)):
                    ax.text(level, bar.get_y() + bar.get_height()/2,
                           f'Nivel {level}', ha='left', va='center', 
                           fontweight='bold', fontsize=9)
                
                plt.tight_layout()
                graph_file = self.graphs_dir / f"{doc_name}_estructura.png"
                plt.savefig(graph_file, bbox_inches='tight', dpi=300, facecolor='white')
                graph_files.append(str(graph_file))
                plt.close()
        
        # Gráfica 3: Análisis de contenido (word cloud style)
        if sections:
            fig, ax = plt.subplots(figsize=(12, 8))
            
            # Contar palabras clave por sección
            word_counts = {}
            for section in sections:
                title_words = section['title'].lower().split()
                for word in title_words:
                    if len(word) > 3:  # Solo palabras significativas
                        word_counts[word] = word_counts.get(word, 0) + 1
            
            if word_counts:
                top_words = sorted(word_counts.items(), key=lambda x: x[1], reverse=True)[:15]
                words, counts = zip(*top_words) if top_words else ([], [])
                
                if words:
                    colors_map = plt.cm.plasma(np.linspace(0, 1, len(words)))
                    bars = ax.barh(range(len(words)), counts, color=colors_map)
                    ax.set_yticks(range(len(words)))
                    ax.set_yticklabels(words, fontsize=10)
                    ax.set_xlabel('Frecuencia', fontsize=11, fontweight='bold')
                    ax.set_title('Palabras Clave Más Frecuentes en Títulos', 
                               fontsize=14, fontweight='bold')
                    ax.grid(True, alpha=0.3, axis='x')
                    
                    for bar in bars:
                        width = bar.get_width()
                        ax.text(width, bar.get_y() + bar.get_height()/2,
                               f'{int(width)}', ha='left', va='center', 
                               fontweight='bold', fontsize=9)
                    
                    plt.tight_layout()
                    graph_file = self.graphs_dir / f"{doc_name}_palabras_clave.png"
                    plt.savefig(graph_file, bbox_inches='tight', dpi=300, facecolor='white')
                    graph_files.append(str(graph_file))
                    plt.close()
        
        # Gráfica 4: Análisis temporal (si hay fechas)
        dates = re.findall(r'\d{4}-\d{2}-\d{2}', parsed_data.get('raw_content', ''))
        if dates:
            fig, ax = plt.subplots(figsize=(10, 6))
            date_counts = Counter(dates)
            sorted_dates = sorted(date_counts.items())
            dates_list, counts_list = zip(*sorted_dates) if sorted_dates else ([], [])
            
            if dates_list:
                ax.plot(range(len(dates_list)), counts_list, marker='o', linewidth=2, markersize=8)
                ax.set_xticks(range(len(dates_list)))
                ax.set_xticklabels(dates_list, rotation=45, ha='right')
                ax.set_title('Actividad Temporal del Documento', fontsize=14, fontweight='bold')
                ax.set_xlabel('Fecha', fontsize=11)
                ax.set_ylabel('Frecuencia', fontsize=11)
                ax.grid(True, alpha=0.3)
                plt.tight_layout()
                graph_file = self.graphs_dir / f"{doc_name}_temporal.png"
                plt.savefig(graph_file, bbox_inches='tight', dpi=300, facecolor='white')
                graph_files.append(str(graph_file))
                plt.close()
        
        return graph_files
    
    def generate_pdf(self, file_path: str, parsed_data: Dict[str, Any], graph_files: List[str]):
        """Genera un documento PDF de alta calidad"""
        if not REPORTLAB_AVAILABLE:
            print("⚠️  reportlab no disponible. Saltando generación de PDF.")
            return
        
        doc_name = Path(file_path).stem
        output_file = self.output_dir / f"{doc_name}.pdf"
        
        doc = SimpleDocTemplate(str(output_file), pagesize=A4,
                              rightMargin=72, leftMargin=72,
                              topMargin=72, bottomMargin=18)
        
        # Estilos
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#2c3e50'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#34495e'),
            spaceAfter=12,
            spaceBefore=12
        )
        
        # Contenido
        story = []
        
        # Portada (ya movida arriba)
        date_str = datetime.now().strftime("%d de %B de %Y")
        
        # Gráficas
        if graph_files:
            story.append(Paragraph("<b>Gráficas y Análisis</b>", heading_style))
            for graph_file in graph_files:
                if os.path.exists(graph_file):
                    img = Image(graph_file, width=6*inch, height=4*inch)
                    story.append(img)
                    story.append(Spacer(1, 0.2*inch))
            story.append(PageBreak())
        
        # Portada mejorada
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph(f"<b>{parsed_data['title']}</b>", title_style))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"<i>Generado el {date_str}</i>", styles['Normal']))
        story.append(PageBreak())
        
        # Índice de contenido
        story.append(Paragraph("<b>Índice de Contenido</b>", heading_style))
        story.append(Spacer(1, 0.2*inch))
        toc_items = []
        for i, section in enumerate(parsed_data.get('sections', [])[:30], 1):
            toc_items.append(f"{i}. {section['title']}")
        for item in toc_items:
            story.append(Paragraph(item, styles['Normal']))
        story.append(PageBreak())
        
        # Secciones
        for section in parsed_data.get('sections', []):
            story.append(Paragraph(section['title'], heading_style))
            
            # Tablas de la sección
            if section.get('tables'):
                for table_data in section['tables']:
                    if table_data and len(table_data) > 1:
                        # Crear tabla PDF
                        table = Table(table_data)
                        table.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#34495e')),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                            ('GRID', (0, 0), (-1, -1), 1, colors.grey),
                            ('FONTSIZE', (0, 1), (-1, -1), 9),
                        ]))
                        story.append(table)
                        story.append(Spacer(1, 0.2*inch))
            
            # Contenido de la sección
            for item in section.get('content', [])[:30]:  # Aumentado límite
                if item.get('type') == 'paragraph' and item.get('content'):
                    story.append(Paragraph(item['content'], styles['Normal']))
                    story.append(Spacer(1, 0.1*inch))
                elif item.get('type') == 'list_item':
                    story.append(Paragraph(f"• {item['content']}", styles['Normal']))
                    story.append(Spacer(1, 0.05*inch))
            
            story.append(Spacer(1, 0.2*inch))
        
        # Construir PDF
        doc.build(story)
        print(f"✅ PDF generado: {output_file}")
    
    def generate_word(self, file_path: str, parsed_data: Dict[str, Any], graph_files: List[str]):
        """Genera un documento Word de alta calidad"""
        if not DOCX_AVAILABLE:
            print("⚠️  python-docx no disponible. Saltando generación de Word.")
            return
        
        doc_name = Path(file_path).stem
        output_file = self.output_dir / f"{doc_name}.docx"
        
        doc = Document()
        
        # Configurar estilos
        style = doc.styles['Normal']
        style.font.name = 'Calibri'
        style.font.size = Pt(11)
        
        # Título
        title = doc.add_heading(parsed_data['title'], 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Fecha
        date_para = doc.add_paragraph()
        date_para.add_run(f"Generado el {datetime.now().strftime('%d de %B de %Y')}").italic = True
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph()  # Espacio
        
        # Gráficas
        if graph_files:
            doc.add_heading('Gráficas y Análisis', 1)
            for graph_file in graph_files:
                if os.path.exists(graph_file):
                    doc.add_picture(graph_file, width=Inches(6))
                    doc.add_paragraph()  # Espacio
        
        # Índice de contenido
        doc.add_heading('Índice de Contenido', 1)
        for i, section in enumerate(parsed_data.get('sections', [])[:30], 1):
            p = doc.add_paragraph(f"{i}. {section['title']}", style='List Number')
        
        doc.add_page_break()
        
        # Secciones
        for section in parsed_data.get('sections', []):
            doc.add_heading(section['title'], level=2)
            
            # Tablas de la sección
            if section.get('tables'):
                for table_data in section['tables']:
                    if table_data and len(table_data) > 1:
                        table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                        table.style = 'Light Grid Accent 1'
                        
                        # Llenar tabla
                        for i, row_data in enumerate(table_data):
                            for j, cell_data in enumerate(row_data):
                                if i < len(table.rows) and j < len(table.rows[i].cells):
                                    cell = table.rows[i].cells[j]
                                    cell.text = str(cell_data)
                                    if i == 0:  # Header
                                        for paragraph in cell.paragraphs:
                                            for run in paragraph.runs:
                                                run.bold = True
                        doc.add_paragraph()  # Espacio después de tabla
            
            # Contenido
            for item in section.get('content', [])[:30]:
                if item.get('type') == 'paragraph' and item.get('content'):
                    doc.add_paragraph(item['content'])
                elif item.get('type') == 'list_item':
                    p = doc.add_paragraph(item['content'], style='List Bullet')
            
            doc.add_paragraph()  # Espacio
        
        # Guardar
        doc.save(str(output_file))
        print(f"✅ Word generado: {output_file}")
    
    def generate_excel(self, file_path: str, parsed_data: Dict[str, Any], graph_files: List[str]):
        """Genera un documento Excel con gráficas y datos estructurados"""
        if not OPENPYXL_AVAILABLE:
            print("⚠️  openpyxl no disponible. Saltando generación de Excel.")
            return
        
        doc_name = Path(file_path).stem
        output_file = self.output_dir / f"{doc_name}.xlsx"
        
        wb = Workbook()
        
        # Hoja 1: Resumen
        ws1 = wb.active
        ws1.title = "Resumen"
        
        # Estilos
        title_font = Font(bold=True, size=16, color="FFFFFF")
        title_fill = PatternFill(start_color="2c3e50", end_color="2c3e50", fill_type="solid")
        header_font = Font(bold=True, size=12)
        header_fill = PatternFill(start_color="34495e", end_color="34495e", fill_type="solid")
        
        # Título
        ws1['A1'] = parsed_data['title']
        ws1['A1'].font = title_font
        ws1['A1'].fill = title_fill
        ws1.merge_cells('A1:D1')
        ws1['A1'].alignment = Alignment(horizontal='center', vertical='center')
        
        # Fecha
        ws1['A2'] = f"Generado el {datetime.now().strftime('%d de %B de %Y')}"
        ws1['A2'].font = Font(italic=True)
        
        # Métricas
        row = 4
        ws1[f'A{row}'] = "Métricas Extraídas"
        ws1[f'A{row}'].font = header_font
        ws1[f'A{row}'].fill = header_fill
        
        row += 1
        metrics = parsed_data.get('metrics', {})
        for key, values in metrics.items():
            ws1[f'A{row}'] = key.capitalize()
            ws1[f'B{row}'] = str(values[0]) if values else "N/A"
            row += 1
        
        # Hoja 2: Secciones
        ws2 = wb.create_sheet("Secciones")
        
        ws2['A1'] = "Sección"
        ws2['A1'].font = header_font
        ws2['A1'].fill = header_fill
        ws2['B1'] = "Nivel"
        ws2['B1'].font = header_font
        ws2['B1'].fill = header_fill
        ws2['C1'] = "Contenido"
        ws2['C1'].font = header_font
        ws2['C1'].fill = header_fill
        ws2['D1'] = "Tablas"
        ws2['D1'].font = header_font
        ws2['D1'].fill = header_fill
        
        row = 2
        for section in parsed_data.get('sections', [])[:100]:  # Aumentado límite
            ws2[f'A{row}'] = section['title']
            ws2[f'B{row}'] = section.get('level', 2)
            content_preview = ' '.join([
                item.get('content', '')[:50] 
                for item in section.get('content', [])[:3]
            ])
            ws2[f'C{row}'] = content_preview[:200]  # Limitar longitud
            ws2[f'D{row}'] = len(section.get('tables', []))
            row += 1
        
        # Hoja 3: Tablas extraídas
        ws3 = None
        if parsed_data.get('tables'):
            ws3 = wb.create_sheet("Tablas")
            table_num = 0
            row = 1
            
            for table_data in parsed_data['tables'][:20]:  # Primeras 20 tablas
                if table_data and len(table_data) > 1:
                    # Título de tabla
                    ws3[f'A{row}'] = f"Tabla {table_num + 1}"
                    ws3[f'A{row}'].font = Font(bold=True, size=14)
                    row += 1
                    
                    # Datos de tabla
                    for i, row_data in enumerate(table_data):
                        for j, cell_data in enumerate(row_data):
                            cell = ws3.cell(row=row, column=j+1, value=str(cell_data))
                            if i == 0:  # Header
                                cell.font = Font(bold=True)
                                cell.fill = header_fill
                        row += 1
                    
                    row += 2  # Espacio entre tablas
                    table_num += 1
        
        # Ajustar ancho de columnas
        ws1.column_dimensions['A'].width = 30
        ws1.column_dimensions['B'].width = 20
        ws2.column_dimensions['A'].width = 40
        ws2.column_dimensions['B'].width = 10
        ws2.column_dimensions['C'].width = 60
        ws2.column_dimensions['D'].width = 10
        
        if ws3:
            ws3.column_dimensions['A'].width = 20
            ws3.column_dimensions['B'].width = 20
            ws3.column_dimensions['C'].width = 20
            ws3.column_dimensions['D'].width = 20
        
        # Hoja 3: Gráficas (si hay datos numéricos)
        if metrics:
            ws3 = wb.create_sheet("Gráficas")
            
            # Preparar datos para gráfica
            chart_data = []
            chart_labels = []
            
            if 'archivos' in metrics:
                chart_labels.append('Archivos')
                chart_data.append(int(metrics['archivos'][0][0]) if metrics['archivos'] else 0)
            
            if 'lineas' in metrics:
                chart_labels.append('Líneas')
                chart_data.append(int(metrics['lineas'][0][0]) if metrics['lineas'] else 0)
            
            if chart_data:
                # Escribir datos
                ws3['A1'] = "Categoría"
                ws3['A1'].font = header_font
                ws3['B1'] = "Valor"
                ws3['B1'].font = header_font
                
                for i, (label, value) in enumerate(zip(chart_labels, chart_data), start=2):
                    ws3[f'A{i}'] = label
                    ws3[f'B{i}'] = value
                
                # Crear gráfica de barras
                chart = BarChart()
                chart.type = "col"
                chart.style = 10
                chart.title = "Métricas del Documento"
                chart.y_axis.title = "Cantidad"
                chart.x_axis.title = "Categoría"
                
                data = Reference(ws3, min_col=2, min_row=1, max_row=len(chart_data) + 1)
                cats = Reference(ws3, min_col=1, min_row=2, max_row=len(chart_labels) + 1)
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                chart.height = 10
                chart.width = 15
                
                ws3.add_chart(chart, "D2")
        
        # Guardar
        wb.save(str(output_file))
        print(f"✅ Excel generado: {output_file}")
    
    def generate_powerpoint(self, file_path: str, parsed_data: Dict[str, Any], graph_files: List[str]):
        """Genera una presentación PowerPoint de alta calidad"""
        if not PPTX_AVAILABLE:
            return
        
        doc_name = Path(file_path).stem
        output_file = self.output_dir / f"{doc_name}.pptx"
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Portada
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        left = top = Inches(1)
        width = height = Inches(8)
        txBox = slide1.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = parsed_data['title']
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        
        p = tf.add_paragraph()
        p.text = f"Generado el {datetime.now().strftime('%d de %B de %Y')}"
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(127, 140, 141)
        
        # Slide 2: Resumen ejecutivo
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        title2 = slide2.shapes.title
        title2.text = "Resumen Ejecutivo"
        content2 = slide2.placeholders[1]
        tf2 = content2.text_frame
        tf2.text = f"Documento: {parsed_data['title']}"
        p2 = tf2.add_paragraph()
        p2.text = f"Secciones: {len(parsed_data.get('sections', []))}"
        p3 = tf2.add_paragraph()
        p3.text = f"Tablas: {len(parsed_data.get('tables', []))}"
        p4 = tf2.add_paragraph()
        words_count = len(parsed_data.get('raw_content', '').split())
        p4.text = f"Palabras: {words_count:,}"
        
        # Slides con gráficas
        for i, graph_file in enumerate(graph_files[:10], 3):  # Máximo 10 gráficas
            if os.path.exists(graph_file):
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(6.5)
                slide.shapes.add_picture(graph_file, left, top, width, height)
        
        # Slides con secciones principales
        for i, section in enumerate(parsed_data.get('sections', [])[:10], 1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = section['title'][:50]  # Limitar longitud
            
            content = slide.placeholders[1]
            tf = content.text_frame
            # Agregar contenido de la sección
            for item in section.get('content', [])[:5]:  # Primeros 5 items
                if item.get('type') == 'paragraph' and item.get('content'):
                    p = tf.add_paragraph()
                    p.text = item['content'][:200]  # Limitar longitud
                    p.level = 0
        
        # Guardar
        prs.save(str(output_file))
        print(f"✅ PowerPoint generado: {output_file}")
    
    def convert_document(self, file_path: str):
        """Convierte un documento markdown a PDF, Word y Excel"""
        start_time = time.time()
        print(f"\n📄 Procesando: {file_path}")
        
        try:
            # Leer y parsear
            content = self.read_markdown(file_path)
            parsed_data = self.parse_markdown(content)
            
            # Generar gráficas
            doc_name = Path(file_path).stem
            graph_files = self.generate_graphs(doc_name, parsed_data)
            
            # Generar documentos
            self.generate_pdf(file_path, parsed_data, graph_files)
            self.generate_word(file_path, parsed_data, graph_files)
            self.generate_excel(file_path, parsed_data, graph_files)
            
            # Generar presentación si está disponible
            if PPTX_AVAILABLE:
                self.generate_powerpoint(file_path, parsed_data, graph_files)
            
            # Actualizar estadísticas
            self.stats['total_docs'] += 1
            self.stats['total_sections'] += len(parsed_data.get('sections', []))
            self.stats['total_tables'] += len(parsed_data.get('tables', []))
            
            elapsed = time.time() - start_time
            self.stats['processing_time'] += elapsed
            
            # Guardar datos del documento para análisis comparativo
            doc_data = {
                'name': doc_name,
                'path': file_path,
                'sections_count': len(parsed_data.get('sections', [])),
                'tables_count': len(parsed_data.get('tables', [])),
                'words_count': len(parsed_data.get('raw_content', '').split()),
                'metrics': parsed_data.get('metrics', {}),
                'processing_time': elapsed,
                'graph_files': len(graph_files)
            }
            self.stats['documents_data'].append(doc_data)
            
            print(f"✅ Documento {doc_name} procesado completamente ({elapsed:.2f}s)\n")
        except Exception as e:
            import traceback
            print(f"❌ Error procesando {file_path}: {e}")
            if len(str(e)) < 200:  # Solo mostrar traceback si el error es corto
                print(traceback.format_exc())
    
    def generate_summary_report(self):
        """Genera un reporte resumen de todos los documentos procesados"""
        summary_file = self.output_dir / "RESUMEN_GENERACION.md"
        
        avg_time = self.stats['processing_time'] / self.stats['total_docs'] if self.stats['total_docs'] > 0 else 0
        
        with open(summary_file, 'w', encoding='utf-8') as f:
            f.write("# 📊 Resumen de Generación de Documentos\n\n")
            f.write(f"**Fecha de generación**: {datetime.now().strftime('%d de %B de %Y, %H:%M:%S')}\n\n")
            f.write("## 📈 Estadísticas Generales\n\n")
            f.write(f"- **Total de documentos procesados**: {self.stats['total_docs']}\n")
            f.write(f"- **Total de secciones**: {self.stats['total_sections']}\n")
            f.write(f"- **Total de tablas extraídas**: {self.stats['total_tables']}\n")
            f.write(f"- **Total de palabras procesadas**: {self.stats['total_words']:,}\n")
            f.write(f"- **Tiempo total de procesamiento**: {self.stats['processing_time']:.2f} segundos\n")
            f.write(f"- **Tiempo promedio por documento**: {avg_time:.2f} segundos\n\n")
            
            # Listar archivos generados
            f.write("## 📁 Archivos Generados\n\n")
            
            pdf_files = sorted(self.output_dir.glob("*.pdf"))
            docx_files = sorted(self.output_dir.glob("*.docx"))
            xlsx_files = sorted(self.output_dir.glob("*.xlsx"))
            graph_files = sorted(self.graphs_dir.glob("*.png"))
            
            f.write(f"### PDF ({len(pdf_files)} archivos)\n\n")
            for pdf in pdf_files:
                size = pdf.stat().st_size / 1024  # KB
                f.write(f"- `{pdf.name}` ({size:.1f} KB)\n")
            
            f.write(f"\n### Word ({len(docx_files)} archivos)\n\n")
            for docx in docx_files:
                size = docx.stat().st_size / 1024  # KB
                f.write(f"- `{docx.name}` ({size:.1f} KB)\n")
            
            f.write(f"\n### Excel ({len(xlsx_files)} archivos)\n\n")
            for xlsx in xlsx_files:
                size = xlsx.stat().st_size / 1024  # KB
                f.write(f"- `{xlsx.name}` ({size:.1f} KB)\n")
            
            f.write(f"\n### Gráficas ({len(graph_files)} archivos)\n\n")
            for graph in graph_files[:20]:  # Primeras 20
                size = graph.stat().st_size / 1024  # KB
                f.write(f"- `{graph.name}` ({size:.1f} KB)\n")
            if len(graph_files) > 20:
                f.write(f"- ... y {len(graph_files) - 20} gráficas más\n")
        
        print(f"📋 Resumen generado: {summary_file}")
    
    def generate_html_dashboard(self):
        """Genera un dashboard HTML interactivo con todas las estadísticas"""
        dashboard_file = self.output_dir / "dashboard.html"
        
        pdf_files = sorted(self.output_dir.glob("*.pdf"))
        docx_files = sorted(self.output_dir.glob("*.docx"))
        xlsx_files = sorted(self.output_dir.glob("*.xlsx"))
        graph_files = sorted(self.graphs_dir.glob("*.png"))
        
        total_size = sum(f.stat().st_size for f in pdf_files + docx_files + xlsx_files + graph_files) / (1024*1024)
        
        html_content = f"""<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Dashboard - Documentos Generados</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 20px;
            min-height: 100vh;
        }}
        .container {{
            max-width: 1400px;
            margin: 0 auto;
            background: white;
            border-radius: 20px;
            padding: 40px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
        }}
        h1 {{
            color: #2c3e50;
            text-align: center;
            margin-bottom: 10px;
            font-size: 2.5em;
        }}
        .subtitle {{
            text-align: center;
            color: #7f8c8d;
            margin-bottom: 40px;
            font-size: 1.1em;
        }}
        .stats-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
            gap: 20px;
            margin-bottom: 40px;
        }}
        .stat-card {{
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 30px;
            border-radius: 15px;
            text-align: center;
            box-shadow: 0 10px 30px rgba(0,0,0,0.2);
            transition: transform 0.3s;
        }}
        .stat-card:hover {{
            transform: translateY(-5px);
        }}
        .stat-value {{
            font-size: 3em;
            font-weight: bold;
            margin-bottom: 10px;
        }}
        .stat-label {{
            font-size: 1.1em;
            opacity: 0.9;
        }}
        .section {{
            margin-bottom: 40px;
        }}
        .section-title {{
            color: #2c3e50;
            font-size: 1.8em;
            margin-bottom: 20px;
            padding-bottom: 10px;
            border-bottom: 3px solid #667eea;
        }}
        .file-list {{
            display: grid;
            grid-template-columns: repeat(auto-fill, minmax(300px, 1fr));
            gap: 15px;
        }}
        .file-item {{
            background: #f8f9fa;
            padding: 15px;
            border-radius: 10px;
            border-left: 4px solid #667eea;
            transition: all 0.3s;
        }}
        .file-item:hover {{
            background: #e9ecef;
            transform: translateX(5px);
        }}
        .file-name {{
            font-weight: bold;
            color: #2c3e50;
            margin-bottom: 5px;
        }}
        .file-size {{
            color: #7f8c8d;
            font-size: 0.9em;
        }}
        .graph-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fill, minmax(400px, 1fr));
            gap: 20px;
        }}
        .graph-item {{
            text-align: center;
        }}
        .graph-item img {{
            max-width: 100%;
            border-radius: 10px;
            box-shadow: 0 5px 15px rgba(0,0,0,0.2);
        }}
        .graph-title {{
            margin-top: 10px;
            color: #2c3e50;
            font-weight: bold;
        }}
        .footer {{
            text-align: center;
            color: #7f8c8d;
            margin-top: 40px;
            padding-top: 20px;
            border-top: 2px solid #ecf0f1;
        }}
    </style>
</head>
<body>
    <div class="container">
        <h1>📊 Dashboard de Documentos Generados</h1>
        <p class="subtitle">Generado el {datetime.now().strftime('%d de %B de %Y, %H:%M:%S')}</p>
        
        <div class="stats-grid">
            <div class="stat-card">
                <div class="stat-value">{self.stats['total_docs']}</div>
                <div class="stat-label">Documentos Procesados</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">{self.stats['total_sections']}</div>
                <div class="stat-label">Secciones Totales</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">{self.stats['total_tables']}</div>
                <div class="stat-label">Tablas Extraídas</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">{self.stats['total_words']:,}</div>
                <div class="stat-label">Palabras Procesadas</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">{total_size:.1f} MB</div>
                <div class="stat-label">Tamaño Total</div>
            </div>
            <div class="stat-card">
                <div class="stat-value">{self.stats['processing_time']:.1f}s</div>
                <div class="stat-label">Tiempo Total</div>
            </div>
        </div>
        
        <div class="section">
            <h2 class="section-title">📄 Archivos PDF ({len(pdf_files)})</h2>
            <div class="file-list">
"""
        
        for pdf in pdf_files[:20]:  # Primeras 20
            size = pdf.stat().st_size / 1024
            html_content += f"""
                <div class="file-item">
                    <div class="file-name">{pdf.name}</div>
                    <div class="file-size">{size:.1f} KB</div>
                </div>
"""
        
        html_content += """
            </div>
        </div>
        
        <div class="section">
            <h2 class="section-title">📝 Archivos Word ({})</h2>
            <div class="file-list">
""".format(len(docx_files))
        
        for docx in docx_files[:20]:
            size = docx.stat().st_size / 1024
            html_content += f"""
                <div class="file-item">
                    <div class="file-name">{docx.name}</div>
                    <div class="file-size">{size:.1f} KB</div>
                </div>
"""
        
        html_content += """
            </div>
        </div>
        
        <div class="section">
            <h2 class="section-title">📊 Gráficas Generadas ({})</h2>
            <div class="graph-grid">
""".format(len(graph_files))
        
        for graph in graph_files[:12]:  # Primeras 12 gráficas
            rel_path = graph.relative_to(self.output_dir)
            html_content += f"""
                <div class="graph-item">
                    <img src="{rel_path}" alt="{graph.stem}">
                    <div class="graph-title">{graph.stem.replace('_', ' ').title()}</div>
                </div>
"""
        
        html_content += """
            </div>
        </div>
        
        <div class="footer">
            <p>Generado automáticamente por el sistema de conversión de documentos</p>
            <p>Versión mejorada con procesamiento paralelo y análisis avanzado</p>
        </div>
    </div>
</body>
</html>
"""
        
        with open(dashboard_file, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        print(f"🌐 Dashboard HTML generado: {dashboard_file}")
    
    def generate_json_export(self):
        """Genera exportación JSON con todos los datos estructurados"""
        json_file = self.output_dir / "datos_estructurados.json"
        
        export_data = {
            'metadata': {
                'fecha_generacion': datetime.now().isoformat(),
                'version': '2.0',
                'total_documentos': self.stats['total_docs']
            },
            'estadisticas': {
                'total_secciones': self.stats['total_sections'],
                'total_tablas': self.stats['total_tables'],
                'total_palabras': self.stats['total_words'],
                'tiempo_total': self.stats['processing_time'],
                'tiempo_promedio': self.stats['processing_time'] / self.stats['total_docs'] if self.stats['total_docs'] > 0 else 0
            },
            'documentos': self.stats['documents_data']
        }
        
        with open(json_file, 'w', encoding='utf-8') as f:
            json.dump(export_data, f, ensure_ascii=False, indent=2)
        
        print(f"📦 Exportación JSON generada: {json_file}")
    
    def generate_comparative_analysis(self):
        """Genera análisis comparativo entre documentos"""
        if len(self.stats['documents_data']) < 2:
            return
        
        fig, axes = plt.subplots(2, 2, figsize=(16, 12))
        fig.suptitle('Análisis Comparativo de Documentos', fontsize=18, fontweight='bold')
        
        docs_data = self.stats['documents_data']
        doc_names = [d['name'][:20] + '...' if len(d['name']) > 20 else d['name'] for d in docs_data]
        
        # Gráfica 1: Comparación de secciones
        ax1 = axes[0, 0]
        sections_counts = [d['sections_count'] for d in docs_data]
        bars1 = ax1.barh(range(len(doc_names)), sections_counts, color=plt.cm.viridis(np.linspace(0, 1, len(doc_names))))
        ax1.set_yticks(range(len(doc_names)))
        ax1.set_yticklabels(doc_names, fontsize=9)
        ax1.set_xlabel('Número de Secciones', fontweight='bold')
        ax1.set_title('Comparación de Secciones', fontweight='bold')
        ax1.grid(True, alpha=0.3, axis='x')
        for i, (bar, val) in enumerate(zip(bars1, sections_counts)):
            ax1.text(val, bar.get_y() + bar.get_height()/2, f'{val}', 
                    ha='left', va='center', fontweight='bold')
        
        # Gráfica 2: Comparación de tablas
        ax2 = axes[0, 1]
        tables_counts = [d['tables_count'] for d in docs_data]
        bars2 = ax2.bar(range(len(doc_names)), tables_counts, color=plt.cm.plasma(np.linspace(0, 1, len(doc_names))))
        ax2.set_xticks(range(len(doc_names)))
        ax2.set_xticklabels(doc_names, rotation=45, ha='right', fontsize=9)
        ax2.set_ylabel('Número de Tablas', fontweight='bold')
        ax2.set_title('Comparación de Tablas', fontweight='bold')
        ax2.grid(True, alpha=0.3, axis='y')
        for bar, val in zip(bars2, tables_counts):
            height = bar.get_height()
            ax2.text(bar.get_x() + bar.get_width()/2., height, f'{val}',
                    ha='center', va='bottom', fontweight='bold')
        
        # Gráfica 3: Comparación de palabras
        ax3 = axes[1, 0]
        words_counts = [d['words_count'] for d in docs_data]
        bars3 = ax3.bar(range(len(doc_names)), words_counts, color=plt.cm.coolwarm(np.linspace(0, 1, len(doc_names))))
        ax3.set_xticks(range(len(doc_names)))
        ax3.set_xticklabels(doc_names, rotation=45, ha='right', fontsize=9)
        ax3.set_ylabel('Número de Palabras', fontweight='bold')
        ax3.set_title('Comparación de Palabras', fontweight='bold')
        ax3.grid(True, alpha=0.3, axis='y')
        for bar, val in zip(bars3, words_counts):
            height = bar.get_height()
            ax3.text(bar.get_x() + bar.get_width()/2., height, f'{val:,}',
                    ha='center', va='bottom', fontweight='bold', fontsize=8)
        
        # Gráfica 4: Tiempo de procesamiento
        ax4 = axes[1, 1]
        processing_times = [d['processing_time'] for d in docs_data]
        bars4 = ax4.bar(range(len(doc_names)), processing_times, color=plt.cm.Set3(np.linspace(0, 1, len(doc_names))))
        ax4.set_xticks(range(len(doc_names)))
        ax4.set_xticklabels(doc_names, rotation=45, ha='right', fontsize=9)
        ax4.set_ylabel('Tiempo (segundos)', fontweight='bold')
        ax4.set_title('Tiempo de Procesamiento', fontweight='bold')
        ax4.grid(True, alpha=0.3, axis='y')
        for bar, val in zip(bars4, processing_times):
            height = bar.get_height()
            ax4.text(bar.get_x() + bar.get_width()/2., height, f'{val:.1f}s',
                    ha='center', va='bottom', fontweight='bold')
        
        plt.tight_layout()
        comparison_file = self.graphs_dir / "analisis_comparativo.png"
        plt.savefig(comparison_file, bbox_inches='tight', dpi=300, facecolor='white')
        plt.close()
        
        print(f"📊 Análisis comparativo generado: {comparison_file}")
    
    def generate_executive_summary(self):
        """Genera un resumen ejecutivo automático"""
        summary_file = self.output_dir / "RESUMEN_EJECUTIVO.md"
        
        docs_data = self.stats['documents_data']
        total_words = sum(d['words_count'] for d in docs_data)
        avg_sections = sum(d['sections_count'] for d in docs_data) / len(docs_data) if docs_data else 0
        avg_tables = sum(d['tables_count'] for d in docs_data) / len(docs_data) if docs_data else 0
        
        # Documento más grande y más pequeño
        largest_doc = max(docs_data, key=lambda x: x['words_count']) if docs_data else None
        smallest_doc = min(docs_data, key=lambda x: x['words_count']) if docs_data else None
        
        # Documento con más secciones
        most_sections = max(docs_data, key=lambda x: x['sections_count']) if docs_data else None
        
        with open(summary_file, 'w', encoding='utf-8') as f:
            f.write("# 📋 Resumen Ejecutivo - Generación de Documentos\n\n")
            f.write(f"**Fecha**: {datetime.now().strftime('%d de %B de %Y, %H:%M:%S')}\n\n")
            f.write("## 🎯 Resumen General\n\n")
            f.write(f"Se procesaron exitosamente **{self.stats['total_docs']} documentos** ")
            f.write(f"generando un total de **{self.stats['total_sections']} secciones**, ")
            f.write(f"**{self.stats['total_tables']} tablas** y procesando **{total_words:,} palabras**.\n\n")
            
            f.write("## 📊 Métricas Clave\n\n")
            f.write(f"- **Documentos procesados**: {self.stats['total_docs']}\n")
            f.write(f"- **Secciones promedio por documento**: {avg_sections:.1f}\n")
            f.write(f"- **Tablas promedio por documento**: {avg_tables:.1f}\n")
            f.write(f"- **Palabras totales**: {total_words:,}\n")
            f.write(f"- **Tiempo total de procesamiento**: {self.stats['processing_time']:.2f} segundos\n")
            f.write(f"- **Tiempo promedio por documento**: {self.stats['processing_time']/self.stats['total_docs']:.2f} segundos\n\n")
            
            if largest_doc:
                f.write("## 📈 Documentos Destacados\n\n")
                f.write(f"### Documento Más Extenso\n")
                f.write(f"- **Nombre**: {largest_doc['name']}\n")
                f.write(f"- **Palabras**: {largest_doc['words_count']:,}\n")
                f.write(f"- **Secciones**: {largest_doc['sections_count']}\n")
                f.write(f"- **Tablas**: {largest_doc['tables_count']}\n\n")
            
            if most_sections:
                f.write(f"### Documento con Más Secciones\n")
                f.write(f"- **Nombre**: {most_sections['name']}\n")
                f.write(f"- **Secciones**: {most_sections['sections_count']}\n")
                f.write(f"- **Palabras**: {most_sections['words_count']:,}\n\n")
            
            f.write("## 📁 Archivos Generados\n\n")
            f.write("Para cada documento se generaron:\n")
            f.write("- ✅ Archivo PDF de alta calidad\n")
            f.write("- ✅ Archivo Word editable\n")
            f.write("- ✅ Archivo Excel con datos estructurados\n")
            f.write("- ✅ Múltiples gráficas de análisis (PNG 300 DPI)\n\n")
            
            f.write("## 🎨 Tipos de Análisis Realizados\n\n")
            f.write("1. **Dashboard de Métricas**: Análisis completo de métricas del documento\n")
            f.write("2. **Estructura del Documento**: Jerarquía y organización de secciones\n")
            f.write("3. **Análisis de Palabras Clave**: Frecuencia de términos importantes\n")
            f.write("4. **Análisis Temporal**: Actividad y fechas en el documento\n")
            f.write("5. **Análisis de Complejidad**: Complejidad por sección\n")
            f.write("6. **Análisis Comparativo**: Comparación entre todos los documentos\n\n")
            
            f.write("## 🚀 Próximos Pasos\n\n")
            f.write("Los documentos están listos para:\n")
            f.write("- 📄 Presentación y distribución\n")
            f.write("- 📊 Análisis y revisión\n")
            f.write("- 📈 Reportes ejecutivos\n")
            f.write("- 💼 Documentación profesional\n")
        
        print(f"📋 Resumen ejecutivo generado: {summary_file}")
    
    def generate_content_analysis(self):
        """Genera análisis de contenido y temas"""
        if not self.stats['documents_data']:
            return
        
        # Análisis de temas comunes
        all_section_titles = []
        for doc_data in self.stats['documents_data']:
            # Simular extracción de títulos (en producción se leerían los archivos)
            all_section_titles.extend([f"Sección {i}" for i in range(doc_data['sections_count'])])
        
        # Análisis de distribución
        fig, axes = plt.subplots(2, 2, figsize=(16, 12))
        fig.suptitle('Análisis de Contenido y Calidad', fontsize=18, fontweight='bold')
        
        docs_data = self.stats['documents_data']
        
        # Gráfica 1: Distribución de tamaño de documentos
        ax1 = axes[0, 0]
        words_counts = [d['words_count'] for d in docs_data]
        ax1.hist(words_counts, bins=min(10, len(words_counts)), color='#3498db', edgecolor='black', alpha=0.7)
        ax1.set_xlabel('Número de Palabras', fontweight='bold')
        ax1.set_ylabel('Frecuencia', fontweight='bold')
        ax1.set_title('Distribución de Tamaño de Documentos', fontweight='bold')
        ax1.grid(True, alpha=0.3)
        
        # Gráfica 2: Relación palabras vs secciones
        ax2 = axes[0, 1]
        words = [d['words_count'] for d in docs_data]
        sections = [d['sections_count'] for d in docs_data]
        scatter = ax2.scatter(words, sections, s=100, alpha=0.6, c=range(len(docs_data)), cmap='viridis')
        ax2.set_xlabel('Número de Palabras', fontweight='bold')
        ax2.set_ylabel('Número de Secciones', fontweight='bold')
        ax2.set_title('Relación Palabras vs Secciones', fontweight='bold')
        ax2.grid(True, alpha=0.3)
        
        # Gráfica 3: Eficiencia de procesamiento
        ax3 = axes[1, 0]
        processing_times = [d['processing_time'] for d in docs_data]
        words_per_sec = [d['words_count'] / d['processing_time'] if d['processing_time'] > 0 else 0 
                        for d in docs_data]
        bars = ax3.bar(range(len(docs_data)), words_per_sec, color=plt.cm.coolwarm(np.linspace(0, 1, len(docs_data))))
        ax3.set_xlabel('Documento', fontweight='bold')
        ax3.set_ylabel('Palabras por Segundo', fontweight='bold')
        ax3.set_title('Eficiencia de Procesamiento', fontweight='bold')
        ax3.set_xticks(range(len(docs_data)))
        ax3.set_xticklabels([d['name'][:10] for d in docs_data], rotation=45, ha='right')
        ax3.grid(True, alpha=0.3, axis='y')
        
        # Gráfica 4: Complejidad del documento (tablas + secciones)
        ax4 = axes[1, 1]
        complexity = [d['sections_count'] * 10 + d['tables_count'] * 50 for d in docs_data]
        doc_names = [d['name'][:15] + '...' if len(d['name']) > 15 else d['name'] for d in docs_data]
        bars = ax4.barh(range(len(doc_names)), complexity, color=plt.cm.plasma(np.linspace(0, 1, len(docs_data))))
        ax4.set_yticks(range(len(doc_names)))
        ax4.set_yticklabels(doc_names, fontsize=9)
        ax4.set_xlabel('Índice de Complejidad', fontweight='bold')
        ax4.set_title('Complejidad de Documentos', fontweight='bold')
        ax4.grid(True, alpha=0.3, axis='x')
        
        plt.tight_layout()
        analysis_file = self.graphs_dir / "analisis_contenido.png"
        plt.savefig(analysis_file, bbox_inches='tight', dpi=300, facecolor='white')
        plt.close()
        
        print(f"📊 Análisis de contenido generado: {analysis_file}")
    
    def generate_quality_report(self):
        """Genera reporte de calidad del contenido"""
        quality_file = self.output_dir / "REPORTE_CALIDAD.md"
        
        docs_data = self.stats['documents_data']
        
        # Calcular métricas de calidad
        avg_words = sum(d['words_count'] for d in docs_data) / len(docs_data) if docs_data else 0
        avg_sections = sum(d['sections_count'] for d in docs_data) / len(docs_data) if docs_data else 0
        avg_tables = sum(d['tables_count'] for d in docs_data) / len(docs_data) if docs_data else 0
        
        # Clasificar documentos por calidad
        high_quality = [d for d in docs_data if d['words_count'] > avg_words and d['sections_count'] > avg_sections]
        medium_quality = [d for d in docs_data if not (d in high_quality)]
        
        with open(quality_file, 'w', encoding='utf-8') as f:
            f.write("# 📊 Reporte de Calidad de Contenido\n\n")
            f.write(f"**Fecha**: {datetime.now().strftime('%d de %B de %Y, %H:%M:%S')}\n\n")
            
            f.write("## 🎯 Métricas de Calidad\n\n")
            f.write(f"- **Promedio de palabras por documento**: {avg_words:.0f}\n")
            f.write(f"- **Promedio de secciones por documento**: {avg_sections:.1f}\n")
            f.write(f"- **Promedio de tablas por documento**: {avg_tables:.1f}\n\n")
            
            f.write("## ⭐ Documentos de Alta Calidad\n\n")
            f.write(f"Documentos que superan el promedio en palabras y secciones:\n\n")
            for doc in high_quality:
                f.write(f"### {doc['name']}\n")
                f.write(f"- Palabras: {doc['words_count']:,} (promedio: {avg_words:.0f})\n")
                f.write(f"- Secciones: {doc['sections_count']} (promedio: {avg_sections:.1f})\n")
                f.write(f"- Tablas: {doc['tables_count']}\n")
                f.write(f"- Tiempo de procesamiento: {doc['processing_time']:.2f}s\n\n")
            
            f.write("## 📈 Recomendaciones\n\n")
            f.write("### Para mejorar la calidad:\n\n")
            f.write("1. **Documentos con pocas secciones**: Considerar dividir el contenido en más secciones\n")
            f.write("2. **Documentos con pocas palabras**: Revisar si el contenido está completo\n")
            f.write("3. **Documentos sin tablas**: Considerar agregar tablas para mejor organización\n")
            f.write("4. **Tiempo de procesamiento alto**: Revisar si hay contenido muy extenso\n\n")
            
            f.write("## 📊 Distribución de Calidad\n\n")
            f.write(f"- **Alta calidad**: {len(high_quality)} documentos ({len(high_quality)/len(docs_data)*100:.1f}%)\n")
            f.write(f"- **Calidad media**: {len(medium_quality)} documentos ({len(medium_quality)/len(docs_data)*100:.1f}%)\n")
        
        print(f"⭐ Reporte de calidad generado: {quality_file}")


def find_important_documents(base_path: Path) -> List[str]:
    """Encuentra documentos importantes en el proyecto"""
    important_docs = []
    
    # Documentos principales conocidos
    known_docs = [
        "airflow_automation_prompt.md",
        "truthgpt_collected/integration_code/production_code/ARCHITECTURE_IMPROVEMENTS.md",
        "truthgpt_collected/integration_code/production_code/REFACTORING_PLAN.md",
        "truthgpt_collected/integration_code/production_code/RESUMEN_FINAL_MEJORAS.md",
        "truthgpt_collected/integration_code/production_code/MEJORAS_ARQUITECTURA_COMPLETAS.md",
        "truthgpt_collected/integration_code/production_code/ARCHITECTURE.md",
        "truthgpt_collected/integration_code/production_code/README.md",
        "truthgpt_collected/integration_code/production_code/INDICE_DOCUMENTACION.md",
        "truthgpt_collected/integration_code/production_code/QUICK_START.md",
        "Analisis_Financiero_Banxico.md",
        "README.md",
    ]
    
    # Buscar documentos adicionales con nombres importantes
    important_patterns = [
        "*SUMMARY*.md",
        "*COMPLETE*.md",
        "*FINAL*.md",
        "*IMPROVEMENTS*.md",
        "*PLAN*.md",
        "*GUIDE*.md",
        "*ARCHITECTURE*.md",
    ]
    
    # Agregar documentos conocidos que existan
    for doc in known_docs:
        doc_path = base_path / doc
        if doc_path.exists() and doc_path.is_file():
            important_docs.append(doc)
    
    # Buscar documentos adicionales (limitado para no sobrecargar)
    try:
        for pattern in important_patterns[:3]:  # Solo primeros 3 patrones
            for doc_path in base_path.rglob(pattern):
                if doc_path.is_file() and doc_path.suffix == '.md':
                    rel_path = str(doc_path.relative_to(base_path))
                    if rel_path not in important_docs and len(important_docs) < 20:
                        important_docs.append(rel_path)
    except Exception:
        pass  # Si hay error, continuar con los documentos conocidos
    
    return important_docs[:15]  # Limitar a 15 documentos


def main():
    """Función principal"""
    base_path = Path("/Users/adan/Documents/documentos_blatam")
    converter = DocumentConverter()
    
    print("🚀 Iniciando generación de documentos de alta calidad...")
    print("=" * 60)
    
    # Encontrar documentos importantes
    important_docs = find_important_documents(base_path)
    
    print(f"📚 Encontrados {len(important_docs)} documentos importantes para procesar\n")
    
    processed = 0
    errors = 0
    start_total = time.time()
    
    # Procesar documentos
    if converter.parallel and len(important_docs) > 3:
        # Procesamiento paralelo para múltiples documentos
        print("⚡ Usando procesamiento paralelo...\n")
        with ThreadPoolExecutor(max_workers=3) as executor:
            future_to_doc = {
                executor.submit(converter.convert_document, str(base_path / doc)): doc 
                for doc in important_docs 
                if (base_path / doc).exists() and (base_path / doc).is_file()
            }
            
            for future in as_completed(future_to_doc):
                doc = future_to_doc[future]
                try:
                    future.result()
                    processed += 1
                except Exception as e:
                    print(f"❌ Error procesando {doc}: {e}")
                    errors += 1
    else:
        # Procesamiento secuencial
        for doc in important_docs:
            doc_path = base_path / doc
            if doc_path.exists() and doc_path.is_file():
                try:
                    converter.convert_document(str(doc_path))
                    processed += 1
                except Exception as e:
                    print(f"❌ Error procesando {doc}: {e}")
                    errors += 1
            else:
                print(f"⚠️  Archivo no encontrado: {doc_path}")
                errors += 1
    
    total_time = time.time() - start_total
    
    print("=" * 60)
    print(f"✅ Proceso completado:")
    print(f"   📄 Documentos procesados: {processed}")
    print(f"   ❌ Errores: {errors}")
    print(f"   ⏱️  Tiempo total: {total_time:.2f} segundos")
    print(f"   📁 Documentos guardados en: {converter.output_dir}")
    print(f"   📊 Gráficas guardadas en: {converter.graphs_dir}")
    
    # Resumen de archivos generados
    try:
        pdf_files = list(converter.output_dir.glob("*.pdf"))
        docx_files = list(converter.output_dir.glob("*.docx"))
        xlsx_files = list(converter.output_dir.glob("*.xlsx"))
        graph_files = list(converter.graphs_dir.glob("*.png"))
        
        total_size = sum(f.stat().st_size for f in pdf_files + docx_files + xlsx_files + graph_files) / (1024*1024)
        
        print(f"\n📊 Resumen de archivos generados:")
        print(f"   PDF: {len(pdf_files)} archivos")
        print(f"   Word: {len(docx_files)} archivos")
        print(f"   Excel: {len(xlsx_files)} archivos")
        print(f"   Gráficas: {len(graph_files)} archivos")
        print(f"   Tamaño total: {total_size:.2f} MB")
        print(f"\n📈 Estadísticas del contenido:")
        print(f"   Secciones totales: {converter.stats['total_sections']}")
        print(f"   Tablas extraídas: {converter.stats['total_tables']}")
        print(f"   Palabras procesadas: {converter.stats['total_words']:,}")
        print(f"   Tiempo promedio por doc: {converter.stats['processing_time']/processed:.2f}s" if processed > 0 else "")
        
        # Generar reportes y análisis adicionales
        converter.generate_summary_report()
        converter.generate_html_dashboard()
        converter.generate_json_export()
        converter.generate_comparative_analysis()
        converter.generate_executive_summary()
        converter.generate_content_analysis()
        converter.generate_quality_report()
    except Exception as e:
        print(f"⚠️  Error generando resumen: {e}")


if __name__ == "__main__":
    main()

