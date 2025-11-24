from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_brand_guidelines(filename):
    prs = Presentation()
    
    # Colors
    NAVY = RGBColor(32, 55, 100)
    WHITE = RGBColor(255, 255, 255)

    def create_slide(title_text, content_list=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Cm(3))
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        header.line.fill.background()
        
        # Title
        tf = header.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        p.margin_left = Cm(1)
        p.margin_top = Cm(0.5)
        
        if content_list:
            top = Cm(4)
            for item in content_list:
                textbox = slide.shapes.add_textbox(Cm(1), top, Cm(22), Cm(2))
                p = textbox.text_frame.paragraphs[0]
                p.text = item
                p.font.size = Pt(18)
                top += Cm(1.5)
        
        return slide

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    
    title = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(20), Cm(5))
    p = title.text_frame.paragraphs[0]
    p.text = "VENTURA CAPITAL\nBRAND IDENTITY GUIDELINES"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Core Values
    create_slide("Our Core Values", [
        "1. Cultural Intelligence: We speak the local language, literally and figuratively.",
        "2. Speed with Soul: Automation shouldn't feel robotic.",
        "3. Latin Excellence: Showcasing world-class tech built in LATAM."
    ])

    # Slide 3: Typography
    slide = create_slide("Typography System")
    
    # Font Spec
    tb = slide.shapes.add_textbox(Cm(1), Cm(4), Cm(10), Cm(10))
    p = tb.text_frame.paragraphs[0]
    p.text = "Primary Font: INTER / HELVETICA NOW"
    p.font.size = Pt(24)
    p.font.bold = True
    
    p2 = tb.text_frame.add_paragraph()
    p2.text = "Used for Headlines and UI elements.\nClean, modern, highly readable."
    p2.font.size = Pt(14)

    # Slide 4: Color Palette
    slide = create_slide("Color Palette")
    
    def add_swatch(x, color, name, hex_code):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Cm(5), Cm(4), Cm(4))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        lbl = slide.shapes.add_textbox(x, Cm(9.5), Cm(4), Cm(2))
        p = lbl.text_frame.paragraphs[0]
        p.text = f"{name}\n{hex_code}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)

    add_swatch(Cm(2), NAVY, "Ventura Navy", "#203764")
    add_swatch(Cm(7), RGBColor(47, 117, 181), "Growth Blue", "#2F75B5")
    add_swatch(Cm(12), RGBColor(255, 192, 0), "Accent Gold", "#FFC000")
    add_swatch(Cm(17), RGBColor(242, 242, 242), "Light Gray", "#F2F2F2")

    # Slide 5: Voice & Tone
    create_slide("Voice & Tone", [
        "We sound:",
        "• Professional but approachable (not stiff).",
        "• Confident (we know the market).",
        "• Local (we use 'vos', 'tú', 'você' correctly).",
        "",
        "We do NOT sound:",
        "• Generic or 'Translated'.",
        "• Overly hype-y (no 'crypto-bro' slang)."
    ])

    prs.save(filename)
    print(f"Brand Guidelines created: {filename}")

if __name__ == "__main__":
    create_brand_guidelines("Ventura_Capital_Brand_Guidelines.pptx")


