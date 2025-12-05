from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_enhanced_pitch_deck(filename):
    prs = Presentation()
    
    # --- Design System ---
    # Colors
    NAVY = RGBColor(32, 55, 100)
    BLUE = RGBColor(47, 117, 181) 
    LIGHT_GREY = RGBColor(242, 242, 242)
    WHITE = RGBColor(255, 255, 255)
    ACCENT = RGBColor(255, 192, 0) # Gold/Yellow for highlights

    def apply_slide_design(slide, title_text):
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Title Bar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Cm(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background() # No border
        
        # Title Text
        title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), prs.slide_width - Cm(2), Cm(1.5))
        tf = title_box.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

        # Footer/Logo Placeholder
        footer = slide.shapes.add_textbox(Cm(1), prs.slide_height - Cm(1), Cm(5), Cm(0.5))
        footer.text_frame.text = "Ventura Capital - Series A Deck"
        footer.text_frame.paragraphs[0].font.size = Pt(10)
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

    def add_content_text(slide, text_list, left=Cm(1.5), top=Cm(3.5), width=Cm(22), font_size=18):
        textbox = slide.shapes.add_textbox(left, top, width, Cm(10))
        tf = textbox.text_frame
        tf.word_wrap = True
        
        for item in text_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.space_after = Pt(14)
            p.level = 0
            
            # Custom bullet simulation if needed, or rely on default
            # p.font.color.rgb = NAVY

    # --- Slide 1: Title Slide (Custom) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # Big Navy Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    
    # Title
    title = slide.shapes.add_textbox(Cm(2), Cm(6), Cm(20), Cm(3))
    p = title.text_frame.add_paragraph()
    p.text = "VENTURA CAPITAL"
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = WHITE
    
    # Subtitle
    sub = slide.shapes.add_textbox(Cm(2), Cm(9), Cm(20), Cm(2))
    p = sub.text_frame.add_paragraph()
    p.text = "The First AI-Native Marketing Platform for LATAM"
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT
    
    # Date/Info
    info = slide.shapes.add_textbox(Cm(2), Cm(15), Cm(10), Cm(2))
    p = info.text_frame.add_paragraph()
    p.text = "Series A Investment Opportunity\nDecember 2024"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # --- Slide 2: The Problem ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_design(slide, "The Problem")
    
    # Two boxes comparison
    # Box 1: Old Way
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(4), Cm(10), Cm(10))
    b1.fill.solid()
    b1.fill.fore_color.rgb = LIGHT_GREY
    b1.line.color.rgb = NAVY
    tf = b1.text_frame
    tf.text = "Current State"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\n• Generic AI content (English translated)\n• Expensive Agencies ($5k+/mo)\n• Slow turnaround (Weeks)"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.LEFT
    
    # Box 2: Pain
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(13), Cm(4), Cm(10), Cm(10))
    b2.fill.solid()
    b2.fill.fore_color.rgb = LIGHT_GREY
    b2.line.color.rgb = NAVY
    tf = b2.text_frame
    tf.text = "The Pain"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\n• 40% Budget Waste\n• Low Conversion Rates\n• Brand Disconnection in LATAM"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.LEFT

    # --- Slide 3: The Solution (Secret Sauce) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_design(slide, "Our Secret Sauce")
    
    add_content_text(slide, [
        "Proprietary 'Cultural-Aware' LLMs",
        "• Trained on 50M+ high-converting LATAM ads",
        "• Dialect-specific fine-tuning (MX, AR, CO, CL)",
        "• Context-aware image generation (Local faces/places)",
        "",
        "Results:",
        "• 3x Higher CTR than generic AI",
        "• 90% Cost reduction vs Agencies"
    ])
    
    # --- Slide 4: Market & Why Now ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_design(slide, "Why Now? The LATAM Opportunity")
    
    # Market Stats
    # Circle 1
    s1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(3), Cm(5), Cm(5), Cm(5))
    s1.fill.solid()
    s1.fill.fore_color.rgb = BLUE
    s1.text_frame.text = "$2.8B\nTAM"
    s1.text_frame.paragraphs[0].font.size = Pt(28)
    s1.text_frame.paragraphs[0].font.bold = True
    
    # Text
    add_content_text(slide, [
        "• Digital Ad Spend in LATAM growing at 15% CAGR",
        "• E-commerce boom: 300M digital buyers",
        "• No dominant local AI player yet"
    ], left=Cm(10), top=Cm(5))

    # --- Slide 5: Traction (Metrics) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_design(slide, "Traction & Metrics")
    
    metrics = [
        ("$139K", "ARR (Current)"),
        ("20%", "MoM Growth"),
        ("9:1", "LTV / CAC"),
        ("500+", "Active Customers")
    ]
    
    x_start = Cm(2)
    for val, lbl in metrics:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, Cm(6), Cm(5), Cm(4))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GREY
        box.line.fill.background()
        
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = lbl
        p2.font.size = Pt(14)
        p2.font.color.rgb = BLUE
        p2.alignment = PP_ALIGN.CENTER
        
        x_start += Cm(5.5)

    # --- Slide 6: Go-To-Market ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_design(slide, "Go-To-Market Strategy")
    
    # Funnel Graphic simulation
    # Shape: Triangle inverted? Or just list
    
    add_content_text(slide, [
        "1. Inbound / PLG (Product-Led Growth)",
        "   • Free tier for Solopreneurs",
        "   • Viral loops via 'Made with Ventura' watermarks",
        "",
        "2. Direct Sales (Mid-Market)",
        "   • Target: Digital Agencies & SME E-commerce",
        "   • Value Prop: 'White-label AI for your agency'",
        "",
        "3. Channel Partnerships",
        "   • Exclusive integration with TiendaNube & MercadoShops"
    ])

    # --- Slide 7: Roadmap ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_design(slide, "Product Roadmap")
    
    # Timeline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(8), Cm(21), Cm(0.2))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    
    milestones = [
        ("Q1 2025", "Series A Close\nBrazil Launch"),
        ("Q2 2025", "Video Gen Alpha\nAPI Release"),
        ("Q4 2025", "Enterprise Shield\n$1.5M ARR")
    ]
    
    x_pos = Cm(3)
    for date, text in milestones:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_pos, Cm(7.6), Cm(1), Cm(1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        
        lbl = slide.shapes.add_textbox(x_pos - Cm(1), Cm(9), Cm(4), Cm(3))
        lbl.text_frame.text = f"{date}\n{text}"
        lbl.text_frame.paragraphs[0].font.bold = True
        lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        x_pos += Cm(7)

    # --- Slide 8: The Ask ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_design(slide, "The Ask: Series A")
    
    # Left: Amount
    left_box = slide.shapes.add_textbox(Cm(2), Cm(5), Cm(8), Cm(10))
    p = left_box.text_frame.add_paragraph()
    p.text = "$2.0M"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    p2 = left_box.text_frame.add_paragraph()
    p2.text = "Target Raise"
    p2.font.size = Pt(24)
    p2.font.color.rgb = NAVY
    
    # Right: Use of Funds
    add_content_text(slide, [
        "Use of Proceeds:",
        "• 40% Product (Gen-Video R&D)",
        "• 35% Growth (Brazil Expansion)",
        "• 25% Key Hires (Head of Sales, VP Eng)",
        "",
        "Goal: Reach $3.6M ARR in 18 months"
    ], left=Cm(11), top=Cm(5))

    prs.save(filename)
    print(f"Enhanced Pitch Deck created: {filename}")

if __name__ == "__main__":
    create_enhanced_pitch_deck("Ventura_Capital_Pitch_Deck_V2.pptx")







